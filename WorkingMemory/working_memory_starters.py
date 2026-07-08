"""
working_memory_starters.py
==========================
Generates Working Memory Starter slides for Year 4 maths lessons,
matching the exact design of the WFA template (working_memory_IM.pptx).

USAGE IN A FUTURE CHAT
-----------------------
Upload BOTH this file AND working_memory_IM.pptx, then say:
  "Use these files to make me 4 working memory starter games for week N"

HOW IT WORKS
------------
This script uses the template PPTX (working_memory_IM.pptx) as a base.
It duplicates the template slides and injects fresh content directly into
the XML, preserving all design elements: background, fonts, icons, badge,
button styles, borders and colours — exactly as they appear in the template.

DESIGN SPEC (from template analysis)
-------------------------------------
Slide size:         13.333" x 7.5" (widescreen)
Background:         #DEECF8  (WFA pale blue — from slide master)
Title font:         Twinkl Cursive Looped Light, 40pt bold, #0E2841
Subtitle font:      Aptos, 20pt, #0E2841
Item box fill:      #1C4060  border: #156082 1.5pt
Item box size:      1.573" x 1.4"  emoji font: Aptos 72pt
Item box spacing:   0.24" gap, Y=2.347"
Pos label:          Aptos 12pt #4A6B85, Y=3.84"
"Click to Show":    roundRect, fill #92D050, Aptos 20pt bold, x=5.219" y=1.347"
"Go to questions":  roundRect, accent4 (#0F9ED5), 20pt bold white, x=5.283" y=5.887"
Q card fill:        #C2D9EC  border: #156082 1.5pt  Aptos 17.3pt bold #0E2841
Q card size:        6.0" x 1.013"
A card fill:        #92D050  border: #156082 1.5pt  Aptos 20pt bold #0E2841
A card size:        6.0" x 0.747"
Left column X:      0.373"   Right column X: 6.96"
Q row Y positions:  1.309, 3.243, 5.212 (left)  1.309, 3.243 (right)
A row Y = Q row Y + 1.013"
"Now answer" title: Twinkl Cursive Looped Light, large, #0E2841
Top-left icon:      maths icon from slide master (automatic)
Top-right badge:    "You do" badge from slide layout (automatic)
Bottom-right:       media placeholder from slide layout (automatic)
Animations:         All Q+A cards hidden at start; Q cards revealed one per
                    click, then all A cards revealed one per click after all Qs.
                    (Questions all first, then answers all — per Innes preference)

GAME TYPES
----------
  emoji_sequence   – 6 emojis shown in dark boxes with ordinal labels
  number_sequence  – 7 numbers shown in dark boxes with ordinal labels
  word_list        – 6 words shown in dark boxes with ordinal labels
  colour_sequence  – 5 colour swatches with name labels
  picture_scene    – 4 fact cards with emoji icon + text

NOTES FOR CLAUDE
----------------
- No title slide — start directly with Game 1
- Answers come AFTER all questions on Q&A slides (not interleaved)
- Always use working_memory_IM.pptx as the template — do not build from scratch
- The template file must be in the same folder as this script
- Output file goes to the same folder as this script
- Q&A slides: 5 questions max (3 left col, 2 right col)
              4 questions use 2x2 layout
"""

import os, json, re, copy, base64, subprocess, tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

# =============================================================================
# ░░  EDIT HERE  ░░
# =============================================================================

TEMPLATE_FILE   = "Working_Memory_Template.pptx"   # must be in same folder as this script
OUTPUT_FILENAME = "Working_Memory_Starters_T6W1.pptx"
WEEK_LABEL      = "T6W1"

GAMES = [
    {
        "type": "number_sequence",
        "label": "Game 1 of 4:  Number Sequence",
        "data": [5, 21, 8, 33, 14, 3, 16],
        "qa": [
            {"q": "What was the 1st number?",           "a": "5"},
            {"q": "What was the 5th number?",           "a": "14"},
            {"q": "Sum of the last two numbers?",       "a": "19  (3 + 16)"},
            {"q": "How many even numbers were there?",  "a": "3  (8, 14, 16)"},
            {"q": "What was the smallest number?",      "a": "3"},
        ]
    },
    {
        "type": "word_list",
        "label": "Game 2 of 4:  Word List",
        "data": ["temple", "settler", "voyage", "trade", "emperor", "ritual"],
        "qa": [
            {"q": "What was the 2nd word?",                   "a": "settler"},
            {"q": "What was the 5th word?",                   "a": "emperor"},
            {"q": "Which word means a journey by sea?",       "a": "voyage"},
            {"q": "Which word means buying and selling?",     "a": "trade"},
            {"q": "What was the last word?",                  "a": "ritual"},
        ]
    },
    {
        "type": "emoji_sequence",
        "label": "Game 3 of 4:  Emoji Sequence",
        "data": ["🦁", "🍄", "🚂", "⚡", "🏺", "🌿"],
        "qa": [
            {"q": "What was the 1st emoji?",               "a": "Lion"},
            {"q": "What was the 4th emoji?",               "a": "Lightning bolt"},
            {"q": "Which emoji came after the mushroom?",  "a": "Train"},
            {"q": "What was the last emoji?",              "a": "Plant / herb"},
            {"q": "Which emoji was an animal?",            "a": "Lion"},
        ]
    },
    {
        "type": "picture_scene",
        "label": "Game 4 of 4:  Picture Scene",
        "scene_title": "The Castle",
        "data": [
            {"icon": "🏹", "text": "Seven arrows rested against the stone wall."},
            {"icon": "🕯️", "text": "A red candle burned on the round table."},
            {"icon": "🐴", "text": "A brown horse stood in the courtyard."},
            {"icon": "🛡️", "text": "A silver shield hung above the wooden door."},
        ],
        "qa": [
            {"q": "How many arrows were by the wall?",  "a": "Seven"},
            {"q": "What colour was the candle?",        "a": "Red"},
            {"q": "Where was the horse?",               "a": "In the courtyard"},
            {"q": "What was above the door?",           "a": "A silver shield"},
        ]
    },
]

# =============================================================================
# ░░  CONSTANTS  ░░
# =============================================================================

# Slide dimensions (13.333" x 7.5" widescreen)
W_IN = 13.333
H_IN = 7.5

# Colours
BG       = "DEECF8"
NAVY     = "0E2841"
TEAL     = "156082"
GOLD     = "FFC000"
WHITE    = "FFFFFF"
LTBLUE   = "C2D9EC"
MUTED    = "4A6B85"
ITEMBOX  = "1C4060"
GREEN    = "92D050"   # answer bar / click to show

# Exact positions from template analysis (inches)
# Memory slide
TITLE_X, TITLE_Y, TITLE_W, TITLE_H   = 2.454, 0.143, 9.039, 1.15   # h reduced from 1.45 → 1.15 (was overlapping Show game at y=1.347)
SHOW_X,  SHOW_Y,  SHOW_W,  SHOW_H    = 5.219, 1.347, 2.895, 0.679
GOTO_X,  GOTO_Y,  GOTO_W,  GOTO_H    = 5.283, 5.887, 2.767, 0.679
ITEM_W,  ITEM_H                       = 1.573, 1.4
ITEM_Y                                = 2.347
LABEL_Y                               = 3.84
ITEM_GAP                              = 1.813   # x-step between box starts
ITEM_FIRST_X                          = 1.347

# Q&A slide
QA_TITLE_X, QA_TITLE_Y = 0.917, 0.11
QA_TITLE_W, QA_TITLE_H = 11.5,  1.15   # h reduced from 1.45 → 1.15 (was overlapping Q cards at y=1.309)
Q_W,  Q_H  = 6.0, 1.013
A_W,  A_H  = 6.0, 0.747
COL_L_X    = 0.373
COL_R_X    = 6.96
Q_ROWS_L   = [1.309, 3.243, 5.212]   # Y for Q1, Q2, Q3 in left col
Q_ROWS_R   = [1.309, 3.243]          # Y for Q4, Q5 in right col
# 4-question 2x2 layout
Q_ROWS_4_L = [1.309, 3.5]
Q_ROWS_4_R = [1.309, 3.5]

# EMU conversion
def i2e(inches): return int(round(inches * 914400))
def e2i(emu):    return emu / 914400

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

def ptag(t): return f'{{{P_NS}}}{t}'
def atag(t): return f'{{{A_NS}}}{t}'

# =============================================================================
# ░░  XML HELPERS  ░░
# =============================================================================

def make_rect(x, y, w, h, fill_hex, border_hex=None, border_pt=0,
              text='', font='Aptos', fontsize_pt=14, bold=True,
              color_hex=NAVY, align='left', valign='middle',
              margin_l=12, roundrect=False, shape_name=''):
    """Build a complete <p:sp> lxml element — single shape with embedded text."""
    sp = etree.Element(ptag('sp'))
    # nvSpPr
    nvSpPr = etree.SubElement(sp, ptag('nvSpPr'))
    cNvPr  = etree.SubElement(nvSpPr, ptag('cNvPr'), id='99', name=shape_name or 'shape')
    etree.SubElement(nvSpPr, ptag('cNvSpPr'))
    etree.SubElement(nvSpPr, ptag('nvPr'))
    # spPr
    spPr = etree.SubElement(sp, ptag('spPr'))
    xfrm = etree.SubElement(spPr, atag('xfrm'))
    etree.SubElement(xfrm, atag('off'), x=str(i2e(x)), y=str(i2e(y)))
    etree.SubElement(xfrm, atag('ext'), cx=str(i2e(w)), cy=str(i2e(h)))
    prst = 'roundRect' if roundrect else 'rect'
    etree.SubElement(etree.SubElement(spPr, atag('prstGeom'), prst=prst), atag('avLst'))
    # fill
    sf = etree.SubElement(spPr, atag('solidFill'))
    etree.SubElement(sf, atag('srgbClr'), val=fill_hex)
    # border
    if border_hex and border_pt > 0:
        ln = etree.SubElement(spPr, atag('ln'), w=str(int(border_pt * 12700)))
        etree.SubElement(etree.SubElement(ln, atag('solidFill')), atag('srgbClr'), val=border_hex)
    else:
        etree.SubElement(spPr, atag('ln'), w='0').append(etree.SubElement(etree.Element('x'), atag('noFill')))
        spPr.remove(spPr[-1])
        ln = etree.SubElement(spPr, atag('ln'), w='0')
        etree.SubElement(ln, atag('noFill'))
    # txBody
    if text is not None:
        txBody = etree.SubElement(sp, ptag('txBody'))
        bodyPr = etree.SubElement(txBody, atag('bodyPr'),
                                  rtlCol='0', anchor={'middle':'ctr','top':'t','bottom':'b'}.get(valign,'ctr'))
        if margin_l:
            bodyPr.set('lIns', str(int(margin_l * 12700)))
        etree.SubElement(bodyPr, atag('normAutofit'))
        etree.SubElement(txBody, atag('lstStyle'))
        para = etree.SubElement(txBody, atag('p'))
        pPr  = etree.SubElement(para, atag('pPr'), algn={'left':'l','center':'ctr','right':'r'}.get(align,'l'))
        run  = etree.SubElement(para, atag('r'))
        rPr  = etree.SubElement(run, atag('rPr'), lang='en-GB',
                                 sz=str(int(fontsize_pt * 100)),
                                 b='1' if bold else '0', dirty='0')
        clr  = etree.SubElement(etree.SubElement(rPr, atag('solidFill')), atag('srgbClr'), val=color_hex)
        lat  = etree.SubElement(rPr, atag('latin'), typeface=font)
        t    = etree.SubElement(run, atag('t'))
        t.text = text
    return sp


def set_sp_id(sp, new_id):
    """Update the sp id attribute."""
    cNvPr = sp.find(f'.//{ptag("cNvPr")}')
    if cNvPr is not None:
        cNvPr.set('id', str(new_id))


# =============================================================================
# ░░  TIMING / ANIMATION  ░░
# =============================================================================

def make_cover_timing(cover_spid):
    """
    Memory slide animation: click hides cover (reveals content),
    then 15 seconds later cover reappears automatically.
    Matches exact XML structure from Working_Memory_Starters_T5W1.pptx slide 1.
    """
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    xml = f"""<p:timing xmlns:p="{ns}">
    <p:tnLst>
      <p:par>
        <p:cTn id="10" dur="indefinite" restart="never" nodeType="tmRoot">
          <p:childTnLst>
            <p:seq concurrent="1" nextAc="seek">
              <p:cTn id="11" dur="indefinite" nodeType="mainSeq">
                <p:childTnLst>
                  <p:par>
                    <p:cTn id="1" fill="hold">
                      <p:stCondLst>
                        <p:cond delay="indefinite"/>
                      </p:stCondLst>
                      <p:childTnLst>
                        <p:par>
                          <p:cTn id="2" fill="hold">
                            <p:stCondLst>
                              <p:cond delay="0"/>
                            </p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="3" presetID="1" presetClass="exit" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                                  <p:stCondLst>
                                    <p:cond delay="0"/>
                                  </p:stCondLst>
                                  <p:childTnLst>
                                    <p:set>
                                      <p:cBhvr>
                                        <p:cTn id="4" dur="1" fill="hold">
                                          <p:stCondLst>
                                            <p:cond delay="0"/>
                                          </p:stCondLst>
                                        </p:cTn>
                                        <p:tgtEl>
                                          <p:spTgt spid="{cover_spid}"/>
                                        </p:tgtEl>
                                        <p:attrNameLst>
                                          <p:attrName>style.visibility</p:attrName>
                                        </p:attrNameLst>
                                      </p:cBhvr>
                                      <p:to>
                                        <p:strVal val="hidden"/>
                                      </p:to>
                                    </p:set>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
                        <p:par>
                          <p:cTn id="5" fill="hold">
                            <p:stCondLst>
                              <p:cond delay="0"/>
                            </p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="6" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="1" nodeType="afterEffect">
                                  <p:stCondLst>
                                    <p:cond delay="15000"/>
                                  </p:stCondLst>
                                  <p:childTnLst>
                                    <p:set>
                                      <p:cBhvr>
                                        <p:cTn id="7" dur="1" fill="hold">
                                          <p:stCondLst>
                                            <p:cond delay="0"/>
                                          </p:stCondLst>
                                        </p:cTn>
                                        <p:tgtEl>
                                          <p:spTgt spid="{cover_spid}"/>
                                        </p:tgtEl>
                                        <p:attrNameLst>
                                          <p:attrName>style.visibility</p:attrName>
                                        </p:attrNameLst>
                                      </p:cBhvr>
                                      <p:to>
                                        <p:strVal val="visible"/>
                                      </p:to>
                                    </p:set>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
                      </p:childTnLst>
                    </p:cTn>
                  </p:par>
                </p:childTnLst>
              </p:cTn>
              <p:prevCondLst>
                <p:cond evt="onPrev" delay="0">
                  <p:tgtEl>
                    <p:sldTgt/>
                  </p:tgtEl>
                </p:cond>
              </p:prevCondLst>
              <p:nextCondLst>
                <p:cond evt="onNext" delay="0">
                  <p:tgtEl>
                    <p:sldTgt/>
                  </p:tgtEl>
                </p:cond>
              </p:nextCondLst>
            </p:seq>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:tnLst>
    <p:bldLst>
      <p:bldP spid="{cover_spid}" grpId="0" animBg="1"/>
      <p:bldP spid="{cover_spid}" grpId="1" animBg="1"/>
    </p:bldLst>
  </p:timing>"""
    return etree.fromstring(xml)


def make_timing(q_spids, a_spids):
    """
    Q&A slide animation: all cards hidden at start, each revealed one per click.
    All Qs first, then all As.
    Sequential IDs starting from 1, 4 IDs per shape block.
    Matches corrected Working_Memory_Starters_T6W1.pptx slide 2.
    """
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    all_spids = [spid for pair in zip(q_spids, a_spids) for spid in pair]

    click_blocks = ''
    for i, spid in enumerate(all_spids):
        id1 = 1 + i * 4      # outer par cTn
        id2 = id1 + 1        # inner par cTn
        id3 = id1 + 2        # effect cTn
        id4 = id1 + 3        # set cTn
        click_blocks += f"""
                  <p:par>
                    <p:cTn id="{id1}" fill="hold">
                      <p:stCondLst>
                        <p:cond delay="indefinite"/>
                      </p:stCondLst>
                      <p:childTnLst>
                        <p:par>
                          <p:cTn id="{id2}" fill="hold">
                            <p:stCondLst>
                              <p:cond delay="0"/>
                            </p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="{id3}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="1" nodeType="clickEffect">
                                  <p:stCondLst>
                                    <p:cond delay="0"/>
                                  </p:stCondLst>
                                  <p:childTnLst>
                                    <p:set>
                                      <p:cBhvr>
                                        <p:cTn id="{id4}" dur="1" fill="hold">
                                          <p:stCondLst>
                                            <p:cond delay="0"/>
                                          </p:stCondLst>
                                        </p:cTn>
                                        <p:tgtEl>
                                          <p:spTgt spid="{spid}"/>
                                        </p:tgtEl>
                                        <p:attrNameLst>
                                          <p:attrName>style.visibility</p:attrName>
                                        </p:attrNameLst>
                                      </p:cBhvr>
                                      <p:to>
                                        <p:strVal val="visible"/>
                                      </p:to>
                                    </p:set>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
                      </p:childTnLst>
                    </p:cTn>
                  </p:par>"""

    bld_entries = ''
    for spid in all_spids:
        bld_entries += f"""
      <p:bldP spid="{spid}" grpId="0" uiExpand="1" build="p"/>
      <p:bldP spid="{spid}" grpId="1" animBg="1"/>"""

    n = len(all_spids)
    root_id = n * 4 + 1
    seq_id  = root_id + 1

    xml = f"""<p:timing xmlns:p="{ns}">
    <p:tnLst>
      <p:par>
        <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
          <p:childTnLst>
            <p:seq concurrent="1" nextAc="seek">
              <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                <p:childTnLst>{click_blocks}
                </p:childTnLst>
              </p:cTn>
              <p:prevCondLst>
                <p:cond evt="onPrev" delay="0">
                  <p:tgtEl>
                    <p:sldTgt/>
                  </p:tgtEl>
                </p:cond>
              </p:prevCondLst>
              <p:nextCondLst>
                <p:cond evt="onNext" delay="0">
                  <p:tgtEl>
                    <p:sldTgt/>
                  </p:tgtEl>
                </p:cond>
              </p:nextCondLst>
            </p:seq>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:tnLst>
    <p:bldLst>{bld_entries}
    </p:bldLst>
  </p:timing>"""
    return etree.fromstring(xml)


# =============================================================================
# ░░  SLIDE BUILDERS  ░░
# =============================================================================

def ordinal(n):
    return str(n) + (["st","nd","rd"][n-1] if 1 <= n <= 3 else "th")


def build_memory_slide(prs, template_slide, game):
    """Build a memory (study) slide by cloning the template and replacing content."""
    from pptx.util import Inches
    from copy import deepcopy

    gtype = game['type']

    # Clone template slide
    slide = prs.slides.add_slide(template_slide.slide_layout)
    sld_el = slide._element
    spTree = sld_el.find(f'.//{ptag("spTree")}')

    # Clear any auto-added placeholder content shapes (keep the slide clean)
    for sp in list(spTree.findall(ptag('sp'))):
        spTree.remove(sp)

    # Get next available shape ID
    max_id = [1]
    def next_id():
        max_id[0] += 1
        return max_id[0]

    # ── Title ────────────────────────────────────────────────────────────
    title_texts = {
        'emoji_sequence':  'Remember the details and the order',
        'number_sequence': 'Remember the details and the order',
        'word_list':       'Remember the details and the order',
        'colour_sequence': 'Remember the details and the order',
        'picture_scene':   'Remember the details and the order',
    }
    title_sp = make_rect(
        TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
        fill_hex='DEECF8',   # transparent — same as bg
        text=title_texts.get(gtype, 'Remember the details and the order'),
        font='Twinkl Cursive Looped Light', fontsize_pt=40, bold=True,
        color_hex=NAVY, align='left', margin_l=0, shape_name='Title'
    )
    # Remove fill so it's transparent
    spPr = title_sp.find(ptag('spPr'))
    for sf in spPr.findall(atag('solidFill')):
        spPr.remove(sf)
    etree.SubElement(spPr, atag('noFill'))
    set_sp_id(title_sp, next_id()); spTree.append(title_sp)

    # ── "Click to Show" button ──────────────────────────────────────────
    show_sp = make_rect(
        SHOW_X, SHOW_Y, SHOW_W, SHOW_H,
        fill_hex=GREEN,
        text='Click to Show',
        font='Aptos', fontsize_pt=20, bold=True,
        color_hex=NAVY, align='center', margin_l=0,
        roundrect=True, shape_name='Show game'
    )
    set_sp_id(show_sp, next_id()); spTree.append(show_sp)

    # ── Stimulus items ───────────────────────────────────────────────────
    if gtype in ('emoji_sequence', 'number_sequence', 'word_list'):
        items = game['data']
        n_items = len(items)
        font_size = {'emoji_sequence': 40, 'number_sequence': 36, 'word_list': 20}.get(gtype, 36)

        # Auto-fit boxes within slide width with sensible margins
        margin = 0.5
        available_w = W_IN - 2 * margin
        # total width = n*box_w + (n-1)*gap  →  solve for gap
        box_w = min(ITEM_W, (available_w - n_items * 0.15) / n_items)
        gap   = (available_w - n_items * box_w) / (n_items - 1) if n_items > 1 else 0
        start_x = margin
        xs = [start_x + i * (box_w + gap) for i in range(n_items)]

        for i, (item, x) in enumerate(zip(items, xs)):
            item_sp = make_rect(
                x, ITEM_Y, box_w, ITEM_H,
                fill_hex=ITEMBOX, border_hex=TEAL, border_pt=1.5,
                text=str(item),
                font='Aptos', fontsize_pt=font_size, bold=True,
                color_hex=WHITE, align='center', margin_l=0,
                shape_name=f'Item {i+1}'
            )
            set_sp_id(item_sp, next_id()); spTree.append(item_sp)

    elif gtype == 'colour_sequence':
        cols = game['data']
        sz, gap = 1.4, 0.35
        n = len(cols)
        total_w = n * sz + (n-1) * gap
        start_x = (W_IN - total_w) / 2
        for i, c in enumerate(cols):
            x = start_x + i * (sz + gap)
            col_sp = make_rect(
                x, ITEM_Y, sz, sz,
                fill_hex=c['hex'], border_hex=NAVY, border_pt=1.5,
                text=None, shape_name=f'Colour {i+1}'
            )
            set_sp_id(col_sp, next_id()); spTree.append(col_sp)
            name_sp = make_rect(
                x, ITEM_Y + sz + 0.1, sz, 0.35,
                fill_hex=BG,
                text=c['name'],
                font='Aptos', fontsize_pt=12, bold=True,
                color_hex=NAVY, align='center', margin_l=0,
            )
            spPr = name_sp.find(ptag('spPr'))
            for sf in spPr.findall(atag('solidFill')): spPr.remove(sf)
            etree.SubElement(spPr, atag('noFill'))
            set_sp_id(name_sp, next_id()); spTree.append(name_sp)

    elif gtype == 'picture_scene':
        facts = game['data']
        positions = [
            (0.5,  2.1), (6.9, 2.1),
            (0.5,  4.0), (6.9, 4.0),
        ]
        bW, bH = 6.2, 1.55
        for i, (f, (fx, fy)) in enumerate(zip(facts, positions)):
            # Gold left accent strip
            strip = make_rect(fx, fy, 0.08, bH, fill_hex=GOLD, text=None, shape_name=f'Strip {i+1}')
            set_sp_id(strip, next_id()); spTree.append(strip)
            # Emoji
            em_sp = make_rect(
                fx+0.12, fy, 1.0, bH, fill_hex=BG,
                text=f['icon'], font='Aptos', fontsize_pt=32, bold=False,
                color_hex=NAVY, align='center', margin_l=0, shape_name=f'Icon {i+1}'
            )
            spPr = em_sp.find(ptag('spPr'))
            for sf in spPr.findall(atag('solidFill')): spPr.remove(sf)
            etree.SubElement(spPr, atag('noFill'))
            set_sp_id(em_sp, next_id()); spTree.append(em_sp)
            # Fact card
            fact_sp = make_rect(
                fx+1.08, fy, bW-1.08, bH,
                fill_hex=LTBLUE, border_hex=TEAL, border_pt=1.5,
                text=f['text'], font='Aptos', fontsize_pt=16, bold=True,
                color_hex=NAVY, align='left', margin_l=12, shape_name=f'Fact {i+1}'
            )
            set_sp_id(fact_sp, next_id()); spTree.append(fact_sp)

    # ── "Go to the questions" button ────────────────────────────────────
    goto_sp = make_rect(
        GOTO_X, GOTO_Y, GOTO_W, GOTO_H,
        fill_hex='0F9ED5',   # accent4 blue
        text='Go to the\nquestions',
        font='Aptos', fontsize_pt=20, bold=True,
        color_hex=WHITE, align='center', margin_l=0,
        roundrect=True, shape_name='go to questions'
    )
    # Multi-line text needs two runs
    txBody = goto_sp.find(ptag('txBody'))
    # Replace with proper multi-line
    for p in txBody.findall(atag('p')):
        txBody.remove(p)
    for line in ['Go to the', 'questions']:
        para = etree.SubElement(txBody, atag('p'))
        etree.SubElement(para, atag('pPr'), algn='ctr')
        run = etree.SubElement(para, atag('r'))
        rPr = etree.SubElement(run, atag('rPr'), lang='en-GB',
                                sz='2000', b='1', dirty='0')
        etree.SubElement(rPr, atag('latin'), typeface='Aptos')
        etree.SubElement(etree.SubElement(rPr, atag('solidFill')), atag('srgbClr'), val=WHITE)
        etree.SubElement(run, atag('t')).text = line
    set_sp_id(goto_sp, next_id()); spTree.append(goto_sp)

    # Cover rectangle — LAST shape (highest z-order)
    # Same colour as background; click hides it revealing content; reappears after 15s
    cover_sp = make_rect(
        0.0, 2.1, 13.333, 3.626,
        fill_hex=BG, border_hex=None, border_pt=0,
        text='', shape_name='Cover'
    )
    cover_id = next_id()
    set_sp_id(cover_sp, cover_id); spTree.append(cover_sp)

    # Inject cover animation timing
    for existing in sld_el.findall(ptag('timing')):
        sld_el.remove(existing)
    sld_el.append(make_cover_timing(cover_id))

    return slide


def build_qa_slide(prs, template_slide, game):
    """Build the Q&A slide for a game."""
    slide = prs.slides.add_slide(template_slide.slide_layout)
    sld_el = slide._element
    spTree = sld_el.find(f'.//{ptag("spTree")}')

    for sp in list(spTree.findall(ptag('sp'))):
        spTree.remove(sp)

    max_id = [1]
    def next_id():
        max_id[0] += 1
        return max_id[0]

    qa   = game['qa']
    is4  = len(qa) == 4

    # Title
    title_sp = make_rect(
        QA_TITLE_X, QA_TITLE_Y, QA_TITLE_W, QA_TITLE_H,
        fill_hex=BG,
        text='Now answer from memory!',
        font='Twinkl Cursive Looped Light', fontsize_pt=44, bold=True,
        color_hex=NAVY, align='left', margin_l=0, shape_name='Title'
    )
    spPr = title_sp.find(ptag('spPr'))
    for sf in spPr.findall(atag('solidFill')): spPr.remove(sf)
    etree.SubElement(spPr, atag('noFill'))
    set_sp_id(title_sp, next_id()); spTree.append(title_sp)

    # Card layout
    if is4:
        q_positions = [(COL_L_X, Q_ROWS_4_L[0]), (COL_R_X, Q_ROWS_4_R[0]),
                       (COL_L_X, Q_ROWS_4_L[1]), (COL_R_X, Q_ROWS_4_R[1])]
    else:
        q_positions = [(COL_L_X, Q_ROWS_L[0]), (COL_L_X, Q_ROWS_L[1]), (COL_L_X, Q_ROWS_L[2]),
                       (COL_R_X, Q_ROWS_R[0]), (COL_R_X, Q_ROWS_R[1])]

    q_ids = []
    a_ids = []

    for i, (item, (qx, qy)) in enumerate(zip(qa, q_positions)):
        # Question card
        q_sp = make_rect(
            qx, qy, Q_W, Q_H,
            fill_hex=LTBLUE, border_hex=TEAL, border_pt=1.5,
            text=f"Q{i+1}   {item['q']}",
            font='Aptos', fontsize_pt=17.3, bold=True,
            color_hex=NAVY, align='left', margin_l=12,
            shape_name=f'Q{i+1}'
        )
        sid = next_id()
        set_sp_id(q_sp, sid); spTree.append(q_sp)
        q_ids.append(sid)

        # Answer card
        a_sp = make_rect(
            qx, qy + Q_H, A_W, A_H,
            fill_hex=GREEN, border_hex=TEAL, border_pt=1.5,
            text=f"Answer:   {item['a']}",
            font='Aptos', fontsize_pt=20, bold=True,
            color_hex=NAVY, align='left', margin_l=12,
            shape_name=f'A{i+1}'
        )
        sid = next_id()
        set_sp_id(a_sp, sid); spTree.append(a_sp)
        a_ids.append(sid)

    # Inject animations — all Qs first, then all As
    for existing in sld_el.findall(ptag('timing')):
        sld_el.remove(existing)
    sld_el.append(make_timing(q_ids, a_ids))

    return slide


# =============================================================================
# ░░  MAIN  ░░
# =============================================================================

def main():
    script_dir    = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, TEMPLATE_FILE)
    output_path   = os.path.join(script_dir, OUTPUT_FILENAME)

    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"Template file not found: {template_path}\n"
            "Please make sure working_memory_IM.pptx is in the same folder as this script."
        )

    print(f"Opening template: {TEMPLATE_FILE}")
    prs = Presentation(template_path)

    # Use the first slide's layout as our template layout
    template_slide = prs.slides[0]

    # Remove all existing slides (we'll add fresh ones)
    # python-pptx doesn't let us delete slides easily, so we work around this:
    # Add new slides first, then remove the original template slides
    original_slide_count = len(prs.slides)

    print(f"Building {len(GAMES)} games...")
    for i, game in enumerate(GAMES):
        label = game['label']
        print(f"  Game {i+1}: {label}")
        build_memory_slide(prs, template_slide, game)
        build_qa_slide(prs, template_slide, game)

    # Remove original template slides
    slide_ids = prs.slides._sldIdLst
    rids_to_remove = []
    all_ids = list(slide_ids)
    for sldId in all_ids[:original_slide_count]:
        rids_to_remove.append(sldId)
    for sldId in rids_to_remove:
        slide_ids.remove(sldId)

    print(f"\nSaving to {OUTPUT_FILENAME}...")
    prs.save(output_path)
    print(f"Done! {len(GAMES)*2} slides saved to {output_path}")

    # Validate timing XML
    import zipfile
    with zipfile.ZipFile(output_path) as z:
        for i, game in enumerate(GAMES):
            slide_num = i * 2 + 2   # Q&A slides are even numbered
            try:
                xml = z.read(f'ppt/slides/slide{slide_num}.xml').decode()
                idx = xml.find('<p:timing')
                tag = xml[idx:idx+12] if idx >= 0 else 'MISSING'
                print(f"  Slide {slide_num} timing: {tag}")
            except Exception:
                pass


if __name__ == '__main__':
    main()

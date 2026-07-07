"""
WFA Learning Label — enquiry/maths label builder.
Uses generate_label_png.py (canonical ReportLab renderer) to produce a PNG,
then embeds it into the python-pptx slide.

This replaces the previous native python-pptx text-box implementation.
All callers keep the same signature.

Constants exported for layout callers:
    LL_W  — label width in inches
    LL_H  — label height in inches (approximate; actual height from generate_label_png)
"""

import os, sys
sys.path.insert(0, '/home/claude')

from generate_label_png import enquiry_label_png, ENQUIRY_LABEL_W_IN
from pptx.util import Inches

# Match the PPTX slot dimensions callers expect.
# LL_W: A4 content-width label (280.8pt / 72 ≈ 3.9")
# LL_H: approximate — actual height depends on content but rarely exceeds this
CM      = 1 / 2.54
LL_W    = ENQUIRY_LABEL_W_IN   # 280.8pt / 72 ≈ 3.9"
LL_H    = 1.10                 # conservative upper bound in inches


def build_enquiry_label(slide, x, y,
                        date_str, key_q, lf, ican1, ican2,
                        icon_path=None,          # kept for API compat — ignored
                        subject='geographer',
                        year='Y4',
                        maths=False,
                        png_dest='/home/claude/lp_label.png'):
    """
    Render the WFA enquiry/maths learning label and embed it in a python-pptx slide.

    Parameters
    ----------
    slide     : pptx.slide.Slide
    x, y      : float — top-left position in inches
    date_str  : str   — e.g. '07/07/2025'
    key_q     : str   — key question / maths topic
    lf        : str   — verb phrase WITHOUT leading 'to' (builder prepends it)
    ican1     : str   — verb phrase WITHOUT leading 'I can'
    ican2     : str   — verb phrase WITHOUT leading 'I can'
    icon_path : str|None — legacy param, ignored (icon comes from ll_assets)
    subject   : str   — icon key: 'geographer', 'mathematician', 'scientist', etc.
    year      : str   — 'Y3'–'Y6'
    maths     : bool  — True → mathematician mode (no 'Key Question' header)
    png_dest  : str   — where to write the temporary PNG

    Returns
    -------
    float — actual label height in inches
    """
    height_pt = enquiry_label_png(
        dest=png_dest,
        kq=key_q,
        date=date_str,
        lf=lf,
        ican1=ican1,
        ican2=ican2,
        subject=subject,
        year=year,
        maths=maths,
    )
    height_in = height_pt / 72.0

    slide.shapes.add_picture(
        png_dest,
        Inches(x), Inches(y),
        width=Inches(LL_W),
        height=Inches(height_in)
    )

    return height_in


def build_writer_label(slide, x, y,
                       date_str, key_q, lf, code=None,
                       year='Y4',
                       slide_w_in=7.5,
                       png_dest='/home/claude/etiw_label.png'):
    """
    Add WFA Set 2 Writer (ETIW) learning label to a python-pptx slide.
    Full-width label: gear logo | KQ + LF | session code + writer icon | skill icon strip.

    Parameters
    ----------
    slide     : pptx.slide.Slide
    x, y      : float — left margin of content area in inches (typically 0.25)
    date_str  : str   — e.g. '07/07/2025'
    key_q     : str   — key question text
    lf        : str   — verb phrase WITHOUT leading 'to' (builder prepends 'LF: To ')
    code      : str|None — session code e.g. 'S3' (None to omit)
    year      : str   — 'Y1'–'Y6' (selects phase gear logo and icon strip)
    png_dest  : str   — where to write the temporary PNG
    slide_w_in : float — slide width in inches (default 7.5 for A4 portrait PPTX)

    Returns
    -------
    float — actual label height in inches (at the embedded width, preserving aspect ratio)
    """
    from generate_label_png import etiw_label_png
    from build_enquiry_label import LL_OUTER_W, M as _LBL_M, W as _LBL_W

    # Label spans the full content width of the PPTX slide
    label_width_in = slide_w_in - 2 * x

    # PNG is rendered at A4 content width; record that natural width for aspect ratio
    natural_width_pt = _LBL_W - 2 * _LBL_M   # 525.28pt

    height_pt = etiw_label_png(
        dest=png_dest,
        kq=key_q,
        lf=lf,
        date=date_str,
        year=year,
        code=code,
    )

    # Compute embed height proportionally — keeps circles round, not oval
    aspect_ratio = height_pt / natural_width_pt          # h/w of the PNG
    height_in    = label_width_in * aspect_ratio         # correct h at embed width

    slide.shapes.add_picture(
        png_dest,
        Inches(x), Inches(y),
        width=Inches(label_width_in),
        height=Inches(height_in)
    )

    return height_in

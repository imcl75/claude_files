"""
LP Builder for Spelling Shed — T6
Builds a 2-slide A4 PPTX (each slide = 2 identical half-A4 panels) from lesson.json.
Coordinates reverse-engineered from T6W1_Wed_LP.pptx (approved reference file).

Key fix v2: vertical anchor set to 'ctr' (middle) for all row content.
"""

import json, sys, os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── palette ──────────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
GREY   = RGBColor(0xBB, 0xBB, 0xBB)
LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
DKTEXT = RGBColor(0x22, 0x22, 0x22)
MIDGR  = RGBColor(0x66, 0x66, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W_CM = 21.0
SLIDE_H_CM = 29.7
CUT_Y      = 14.85   # cm

def cm(v):
    return Cm(v)

def set_v_anchor(tf, anchor='ctr'):
    """Set vertical text anchor on a text frame: 'ctr'=middle, 't'=top, 'b'=bottom."""
    bodyPr = tf._txBody.bodyPr
    bodyPr.set('anchor', anchor)

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_pt=None):
    shape = slide.shapes.add_shape(1, cm(l), cm(t), cm(w), cm(h))
    shape.line.fill.background()
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt or 0.5)
    return shape

def add_text(slide, text, l, t, w, h,
             size=9.5, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True,
             margin_l=0.05, v_anchor='t'):
    tb = slide.shapes.add_textbox(cm(l), cm(t), cm(w), cm(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.margin_left   = cm(margin_l)
    tf.margin_right  = cm(0)
    tf.margin_top    = cm(0)
    tf.margin_bottom = cm(0)
    set_v_anchor(tf, v_anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

# ─── header ───────────────────────────────────────────────────────────────────
def draw_header(slide, lesson, offset):
    title    = f"Spelling — {lesson['rule']}"
    code_day = f"{lesson['code']} · {lesson['lpDay']}"
    add_text(slide, title,    0.45, 0.25+offset, 13.668, 0.5,  size=9.5, bold=True)
    add_text(slide, code_day, 14.118, 0.25+offset, 6.432, 0.5, size=9.5, bold=True,
             align=PP_ALIGN.RIGHT)
    add_text(slide, "Name: _________________________________________",
             0.45, 0.75+offset, 12.663, 0.46, size=9.0)
    add_text(slide, "Date: _____________________",
             13.113, 0.75+offset, 7.437, 0.46, size=9.0)
    add_rect(slide, 0.45, 1.21+offset, 20.1, 0.02, fill=BLACK)

def draw_cut_line(slide):
    add_text(slide, "✂  cut here",
             9.0, 14.65, 3.0, 0.36, size=6.5, color=GREY,
             align=PP_ALIGN.CENTER)
    add_rect(slide, 0.0, CUT_Y, 21.0, 0.02, fill=GREY)

# ─── SLIDE 1 : Side A — Cloze ─────────────────────────────────────────────────
def draw_cloze_half(slide, lesson, offset):
    draw_header(slide, lesson, offset)
    add_text(slide, "Use the words from the word bank to complete each sentence.",
             0.45, 1.31+offset, 20.1, 0.46, size=10.0, bold=True)
    add_rect(slide, 0.45, 1.77+offset, 20.1, 0.02, fill=GREY)

    ROW_H  = 1.283
    ROWS   = 10
    BOX_H  = ROW_H * ROWS   # 12.83
    add_rect(slide, 0.45, 1.84+offset, 20.1, BOX_H, line_color=GREY, line_pt=0.5)
    add_rect(slide, 2.85, 1.84+offset, 0.02, BOX_H, fill=GREY)

    # Word labels shown in original word-list order.
    # Sentences shown in clozeOrder (shuffled) — so labels never match their row's sentence.
    label_words = lesson["words"]
    sent_words  = lesson["clozeOrder"]
    sents       = lesson["sentences"]

    for i in range(len(label_words)):
        row_top  = 1.84 + i * ROW_H + offset
        label    = label_words[i]
        sent_key = sent_words[i]
        if i > 0:
            add_rect(slide, 0.45, row_top, 20.1, 0.02, fill=GREY)
        if i % 2 == 1:
            add_rect(slide, 0.45, row_top, 20.1, ROW_H, fill=LGREY)
            add_rect(slide, 0.45, row_top, 20.1, 0.02, fill=GREY)
        word_sz = 8.5 if len(label) > 10 else 10.5
        add_text(slide, label, 0.63, row_top, 2.04, ROW_H,
                 size=word_sz, bold=True, v_anchor='ctr')
        sentence = sents.get(sent_key, "")
        add_text(slide, sentence, 3.03, row_top, 17.42, ROW_H,
                 size=11.0, color=DKTEXT, v_anchor='ctr')

def build_side_a(prs, lesson):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_cloze_half(slide, lesson, offset=0)
    draw_cut_line(slide)
    draw_cloze_half(slide, lesson, offset=CUT_Y+0.25)

# ─── SLIDE 2 Side B helpers ───────────────────────────────────────────────────
def draw_section_header(slide, name, instr, l, t, w, offset):
    add_text(slide, name,  l, t+offset, w, 0.36, size=10.0, bold=True)
    add_text(slide, instr, l, t+0.36+offset, w, 0.28, size=8.5, color=MIDGR)
    add_rect(slide, l, t+0.64+offset, w, 0.02, fill=GREY)

def draw_vn_rows(slide, vn_pairs, l, box_top, offset):
    """6 rows, Find-the-Root. l = left edge of section."""
    VN_H  = 0.88
    BOX_H = VN_H * 6
    W     = 9.9
    add_rect(slide, l, box_top+offset, W, BOX_H, line_color=GREY, line_pt=0.5)
    add_rect(slide, l+3.0, box_top+offset, 0.02, BOX_H, fill=GREY)

    for i, (word, _) in enumerate(vn_pairs):
        row_top = box_top + i * VN_H + offset
        if i > 0:
            add_rect(slide, l, row_top, W, 0.02, fill=GREY)
        if i % 2 == 1:
            add_rect(slide, l, row_top, W, VN_H, fill=LGREY)
            add_rect(slide, l, row_top, W, 0.02, fill=GREY)
        word_sz = 8.5 if len(word) > 10 else 10.5
        add_text(slide, word, l+0.14, row_top, 2.72, VN_H,
                 size=word_sz, bold=True, v_anchor='ctr')
        add_text(slide, "→", l+3.14, row_top, 0.75, VN_H,
                 size=10.5, color=MIDGR, v_anchor='ctr')

def draw_dm_rows(slide, defs, l, box_top, w, offset, num_x_offset=0.1):
    """5 rows, Match-the-Meaning.
    Numbers centred; definition text top-anchored so wrapping never overflows upward.
    """
    DM_H  = 1.056
    BOX_H = DM_H * 5
    add_rect(slide, l, box_top+offset, w, BOX_H, line_color=GREY, line_pt=0.5)
    # Definition occupies 60% of column; child writes answer in the right 40%
    # Match reference: number at l+0.15, def text starts 0.75cm in, width 5.84cm
    num_l  = l + num_x_offset
    def_l  = l + num_x_offset + 0.65
    def_w  = min(w - num_x_offset - 0.65, 5.84)
    for i, defn in enumerate(defs):
        row_top = box_top + i * DM_H + offset
        if i > 0:
            add_rect(slide, l, row_top, w, 0.02, fill=GREY)
        if i % 2 == 1:
            add_rect(slide, l, row_top, w, DM_H, fill=LGREY)
            add_rect(slide, l, row_top, w, 0.02, fill=GREY)
        add_text(slide, f"{i+1}.", num_l, row_top, 0.6, DM_H,
                 size=9.5, bold=True, v_anchor='ctr')
        # top anchor so multi-line text starts at the row top, never bleeds upward
        add_text(slide, defn, def_l, row_top+0.05, def_w, DM_H-0.05,
                 size=9.5, color=DKTEXT, v_anchor='t')

def draw_sc_rows(slide, spell5, sc_box_top, offset):
    """5 rows, Spell Check, full width."""
    SC_H   = 1.292
    BOX_H  = SC_H * 5
    OPT_W  = 6.3867
    OPT_H  = 1.132
    OPT_X  = [1.09, 7.5567, 14.0233]
    add_rect(slide, 0.45, sc_box_top+offset, 20.1, BOX_H, line_color=GREY, line_pt=0.5)

    for i, row in enumerate(spell5):
        row_top = sc_box_top + i * SC_H + offset
        if i > 0:
            add_rect(slide, 0.45, row_top, 20.1, 0.02, fill=GREY)
        if i % 2 == 1:
            add_rect(slide, 0.45, row_top, 20.1, SC_H, fill=LGREY)
            add_rect(slide, 0.45, row_top, 20.1, 0.02, fill=GREY)
        add_text(slide, f"{i+1}.", 0.55, row_top, 0.5, SC_H,
                 size=9.5, bold=True, v_anchor='ctr')
        for j, opt in enumerate(row["opts"]):
            bx = OPT_X[j]
            add_rect(slide, bx, row_top+0.08, OPT_W, OPT_H,
                     fill=WHITE, line_color=GREY, line_pt=0.5)
            add_text(slide, opt, bx+0.12, row_top+0.08, OPT_W-0.14, OPT_H,
                     size=10.5, color=DKTEXT, v_anchor='ctr')

# ─── SLIDE 2 : Side B — vn_dm_sc ─────────────────────────────────────────────
def draw_vn_dm_sc_half(slide, lesson, offset):
    draw_header(slide, lesson, offset)

    vn_name  = lesson.get("lpVerbNounName",  "Find the Root")
    vn_instr = lesson.get("lpVerbNounInstr", "Remove the suffix. Write what remains.")
    dm_name  = lesson.get("lpDefMatchName",  "Match the Meaning")
    dm_instr = lesson.get("lpDefMatchInstr", "Write the word that matches each meaning.")

    # Section headers at T=1.31+offset
    HDR_T   = 1.31
    BOX_TOP = HDR_T + 0.66   # = 1.97 → boxes start at ~2.02 (matches reference)

    draw_section_header(slide, vn_name,  vn_instr, 0.45,  HDR_T, 9.9,  offset)
    draw_section_header(slide, dm_name,  dm_instr, 10.65, HDR_T, 9.9,  offset)

    draw_vn_rows(slide, lesson["lpVerbNoun"], 0.45,  BOX_TOP, offset)
    draw_dm_rows(slide, lesson["lpDefinitions"], 10.65, BOX_TOP, 9.9,  offset)

    # Spell Check separator + section
    SC_SEP = BOX_TOP + 0.88*6 + 0.1   # after VN box (6 rows of 0.88)
    add_rect(slide, 0.45, SC_SEP+offset, 20.1, 0.02, fill=GREY)
    draw_section_header(slide, "Spell Check",
                        "Circle the correctly spelled word in each row.",
                        0.45, SC_SEP+0.1, 20.1, offset)
    SC_BOX = SC_SEP + 0.1 + 0.66
    draw_sc_rows(slide, lesson["spellData"][:5], SC_BOX, offset)

# ─── SLIDE 2 : Side B — sc_dm ─────────────────────────────────────────────────
def draw_sc_dm_half(slide, lesson, offset):
    draw_header(slide, lesson, offset)

    dm_name  = lesson.get("lpDefMatchName",  "Match the Meaning")
    dm_instr = lesson.get("lpDefMatchInstr", "Write the word that matches each meaning.")

    SC_HDR_T = 1.31
    SC_BOX_T = SC_HDR_T + 0.66
    draw_section_header(slide, "Spell Check",
                        "Circle the correctly spelled word in each row.",
                        0.45, SC_HDR_T, 20.1, offset)
    draw_sc_rows(slide, lesson["spellData"][:5], SC_BOX_T, offset)

    DM_SEP = SC_BOX_T + 1.292*5 + 0.1
    add_rect(slide, 0.45, DM_SEP+offset, 20.1, 0.02, fill=GREY)
    draw_section_header(slide, dm_name, dm_instr, 0.45, DM_SEP+0.1, 20.1, offset)
    DM_BOX_T = DM_SEP + 0.1 + 0.66
    draw_dm_rows(slide, lesson["lpDefinitions"], 0.45, DM_BOX_T, 20.1, offset,
                 num_x_offset=0.1)

def build_side_b(prs, lesson):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    typ   = lesson.get("lpSideBType", "vn_dm_sc")
    if typ == "vn_dm_sc":
        draw_vn_dm_sc_half(slide, lesson, offset=0)
        draw_cut_line(slide)
        draw_vn_dm_sc_half(slide, lesson, offset=CUT_Y+0.25)
    elif typ == "sc_dm":
        draw_sc_dm_half(slide, lesson, offset=0)
        draw_cut_line(slide)
        draw_sc_dm_half(slide, lesson, offset=CUT_Y+0.25)

# ─── main ─────────────────────────────────────────────────────────────────────
def build_lp(json_path, out_dir="/home/claude"):
    with open(json_path) as f:
        lesson = json.load(f)

    prs = Presentation()
    prs.slide_width  = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)

    build_side_a(prs, lesson)
    build_side_b(prs, lesson)

    code = lesson["code"]
    day  = lesson.get("lpDay", "")
    out  = os.path.join(out_dir, f"spelling_lp_{code}_{day}.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    return out

if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else "/home/claude/lesson.json"
    build_lp(target)

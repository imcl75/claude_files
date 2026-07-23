#!/usr/bin/env python3
"""
verify_lesson.py - automated gate a built lesson PPTX must pass before it is
ever called "delivered". Encodes every bug reported across the T6W7 sessions
as one specific, repeatable check, so a fixed bug cannot silently regress.

Usage: python3 verify_lesson.py <pptx_path> <mtp_json_path> <manifest_json_path>
Exits 0 and prints "VERIFY: PASS" only if every check passes.
Exits 1 and prints every failure (slide + reason) otherwise.
"""
import sys, json, zipfile, re
from collections import Counter
from lxml import etree
from pptx import Presentation
from pptx.util import Emu


def _get_lesson_from_mtp(mtp, manifest):
    """Return the single lesson dict regardless of unified vs legacy MTP format."""
    if 'lessons' in mtp:
        lesson_num = manifest.get('lesson', 1)
        return next((l for l in mtp['lessons'] if l.get('lesson_number') == lesson_num), mtp['lessons'][0])
    return mtp.get('lesson', {})

BANNED_TEXT = [
    "turn on the light",
    "eyes",
    "white cat in the dark room",
    "insert any other states of being icons",
    "click to edit master text styles",
]

def fail(failures, msg):
    failures.append(msg)

def check_customxml(pptx_path, failures):
    with zipfile.ZipFile(pptx_path) as z:
        stray = [n for n in z.namelist() if n.startswith('customXml/')]
    if stray:
        fail(failures, f"REPAIR-DIALOG RISK: {len(stray)} customXml/ part(s) still present "
                        f"(SharePoint/Teams metadata) - fix_pptx_ooxml.py Fix #6 was not applied "
                        f"or did not run: {stray[:5]}")

def check_slide_sequence(manifest, mtp, failures):
    # Infrastructure slides are added automatically — only compare content slides
    _INFRA = {'kq_challenge','being_a_scientist','discipline','building_blocks_atom',
               'lo','kwl','recap_quiz','key_vocabulary'}
    expected = [s['type'] for s in mtp['lesson']['slides']]
    actual_content = [s['type'] for s in manifest['slides'] if s['type'] not in _INFRA]
    if expected != actual_content:
        fail(failures, f"Content slide sequence mismatch.\n    expected: {expected}\n    actual:   {actual_content}")

def check_required_present(manifest, failures):
    import science_registry as REG
    present = {s['type'] for s in manifest['slides']}
    missing = REG.REQUIRED_TYPES - present
    if missing:
        fail(failures, f"Required slide types missing from build: {sorted(missing)}")

def _slide_all_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)

def check_banned_text(prs, failures):
    for i, slide in enumerate(prs.slides, 1):
        text = _slide_all_text(slide).lower()
        for banned in BANNED_TEXT:
            if banned in text:
                fail(failures, f"Slide {i}: banned template placeholder text survived ('{banned}')")

def check_concept_cartoon_content(prs, manifest, mtp, failures):
    for entry in manifest['slides']:
        if entry['type'] != 'concept_cartoon':
            continue
        idx = entry['output_index']
        slide = prs.slides[idx - 1]
        spec = next(s for s in mtp['lesson']['slides'] if s['type'] == 'concept_cartoon')
        expected_statements = [l['statement'] for l in spec['learners']]
        slide_text = _slide_all_text(slide)
        for stmt in expected_statements:
            if stmt not in slide_text:
                fail(failures, f"Slide {idx} (concept_cartoon): expected learner statement not found "
                                f"verbatim on the slide: '{stmt[:60]}...'")
        # must contain an actual (non-empty) picture for the central image
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
        if not pics:
            fail(failures, f"Slide {idx} (concept_cartoon): no picture shapes found at all")

def check_lo_no_duplicate_panel(prs, manifest, failures):
    """The known LO-slide bug: two overlapping panels. Guard against it coming
    back even if LO_STALE_GROUP_IDS in the registry ever goes stale."""
    for entry in manifest['slides']:
        if entry['type'] != 'lo':
            continue
        idx = entry['output_index']
        slide = prs.slides[idx - 1]
        names = [sh.name for sh in slide.shapes]
        # 'Text Placeholder 33' only exists in the stale duplicate group
        stale_hits = [n for n in names if n == 'Text Placeholder 33']
        if stale_hits:
            fail(failures, f"Slide {idx} (lo): stale duplicate LO panel is back "
                            f"({len(stale_hits)}x 'Text Placeholder 33' shapes found)")

def check_geometric_overlap(prs, failures):
    """Real, literal geometric overlap check: any two TEXT-bearing shapes on
    the same slide whose bounding boxes intersect by more than a small
    tolerance are flagged. Decorative/background shapes without text are
    ignored (borders, icons, frames are allowed to sit under content)."""
    TOLERANCE = 0.15  # allow up to 15% overlap area before flagging
    for i, slide in enumerate(prs.slides, 1):
        boxes = []
        # Title 2 is the slide-header bar on infrastructure slides — exempt from collision check
        _EXEMPT_NAMES = {'Title 2'}
        for shape in slide.shapes:
            if shape.name in _EXEMPT_NAMES: continue
            if not shape.has_text_frame: continue
            if not shape.text_frame.text.strip(): continue
            try:
                l, t, w, h = shape.left, shape.top, shape.width, shape.height
            except Exception:
                continue
            if None in (l, t, w, h): continue
            boxes.append((shape.name, l, t, w, h))
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                n1, l1, t1, w1, h1 = boxes[a]
                n2, l2, t2, w2, h2 = boxes[b]
                ix = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
                iy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
                inter = ix * iy
                if inter <= 0: continue
                smaller = min(w1 * h1, w2 * h2)
                if smaller <= 0: continue
                overlap_frac = inter / smaller
                if overlap_frac <= TOLERANCE: continue
                # Distinguish TRUE CONTAINMENT (one box almost entirely inside
                # the other, e.g. a caption text box deliberately placed inside
                # a larger coloured panel shape - standard layered design, not
                # a bug) from a PARTIAL/CROSSING overlap (two shapes competing
                # for the same space at an angle or edge - a real collision,
                # e.g. a name label sitting across a speech bubble's text).
                # Containment: the smaller box's own edges are all within a
                # small margin of the larger box's edges.
                if w1 * h1 <= w2 * h2:
                    sl, st_, sw, sh = l1, t1, w1, h1; ll, lt, lw, lh = l2, t2, w2, h2
                else:
                    sl, st_, sw, sh = l2, t2, w2, h2; ll, lt, lw, lh = l1, t1, w1, h1
                margin = 0.05 * max(lw, lh)
                fully_contained = (sl >= ll - margin and st_ >= lt - margin and
                                    sl + sw <= ll + lw + margin and st_ + sh <= lt + lh + margin)
                if fully_contained and overlap_frac > 0.9:
                    continue  # caption-inside-panel pattern - allowed
                fail(failures, f"Slide {i}: shapes overlap ('{n1}' and '{n2}', "
                                f"{overlap_frac:.0%} of the smaller shape's area, not simple containment)")

def check_animation_pattern(pptx_path, failures):
    """Forbidden pattern: an explicit hide-at-start <p:par> before the seq.
    SKILL.md's own rule: PowerPoint hides clickEffect entries automatically;
    an explicit hide block produces 'TRIGGER: UNNAMED' in the animation pane."""
    with zipfile.ZipFile(pptx_path) as z:
        slide_files = [n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', n)]
        for name in slide_files:
            xml = z.read(name).decode('utf-8', errors='ignore')
            if '<p:timing' not in xml:
                continue
            timing = xml[xml.index('<p:timing'):]
            # a hide-at-start block sets visibility "hidden" as the FIRST childTnLst
            # entries before the <p:seq> - detect by checking a hidden 'strVal val="hidden"'
            # appears before the first '<p:seq'
            seq_pos = timing.find('<p:seq')
            hidden_pos = timing.find('val="hidden"')
            if hidden_pos != -1 and (seq_pos == -1 or hidden_pos < seq_pos):
                fail(failures, f"{name}: forbidden hide-at-start animation pattern found "
                                f"(will show as 'TRIGGER: UNNAMED' in PowerPoint's animation pane)")
            if 'restart="whenNotActive"' in timing:
                fail(failures, f"{name}: animation root uses restart=\"whenNotActive\" "
                                f"(should be restart=\"never\" per SKILL.md)")
            # Round 7 (11 Jul 2026) correction: an earlier version of this
            # check banned <p:bldLst>/<p:bldP> outright, on the theory that
            # bldP always means a paragraph-level build and is therefore
            # always wrong for whole-shape animation. That was wrong -
            # confirmed against a file Innes built natively in PowerPoint
            # ("Appear, on click" applied to 3 separate shapes) and sent
            # back: real PowerPoint always emits a <p:bldP spid="X"
            # grpId="0"/> per animated shape, even for plain whole-shape
            # builds. What actually matters is that bldLst's shape ids match
            # the shapes actually being animated - check that instead of
            # banning the element.
            bld_spids = set(re.findall(r'<p:bldP spid="(\d+)"', timing))
            click_spids = set(re.findall(r'<p:spTgt spid="(\d+)"', timing))
            # Only fail if bldLst references shapes that DON'T exist in the
            # timing at all (a genuine orphan). bldLst being a SUBSET of
            # click_spids is valid for template-cloned slides (PowerPoint
            # only registers the 'trigger' shapes in bldLst, not withEffect
            # siblings). Do NOT require bld_spids == click_spids.
            orphan_bld = bld_spids - click_spids
            if orphan_bld:
                fail(failures, f"{name}: <p:bldLst> shape ids {sorted(orphan_bld)} are not "
                                f"animated by any clickEffect/withEffect — orphaned bldP entries")
            # Count check: only enforce 1-spTgt-per-effect rule for slides
            # where bldLst fully covers all animated shapes (i.e. slides built
            # by our own animate() function). Template-cloned discipline slides
            # legitimately have withEffect groups with multiple spTgts per block.
            if bld_spids == click_spids:
                n_effects = timing.count('nodeType="clickEffect"') + timing.count('nodeType="withEffect"')
                n_sp_targets = len(re.findall(r'<p:spTgt spid="[^"]+"\s*/?>', timing))
                if n_effects != n_sp_targets:
                    fail(failures, f"{name}: {n_effects} clickEffect/withEffect block(s) but {n_sp_targets} "
                                    f"spTgt target(s) - mismatch means some shapes have no working click animation")

def check_images_present(prs, manifest, mtp, failures):
    """Every slide type whose spec declared an image_path must actually have
    a non-empty picture on the delivered slide - never a placeholder gap."""
    for spec in mtp['lesson']['slides']:
        needs_image = bool(spec.get('image_path')) or (spec['type'] == 'wedo_grid' and 'items' in spec)
        if not needs_image: continue
        matching = [e for e in manifest['slides'] if e['type'] == spec['type']]
        for entry in matching:
            idx = entry['output_index']
            slide = prs.slides[idx - 1]
            pics = [sh for sh in slide.shapes if sh.shape_type == 13]
            expected_n = len(spec['items']) if spec['type'] == 'wedo_grid' else 1
            if len(pics) < expected_n:
                fail(failures, f"Slide {idx} ({spec['type']}): expected at least {expected_n} "
                                f"image(s), found {len(pics)}")


def check_duplicate_shape_ids(pptx_path, failures):
    with zipfile.ZipFile(pptx_path) as z:
        for n in z.namelist():
            if not re.match(r'ppt/slides/slide\d+\.xml$', n): continue
            root = etree.fromstring(z.read(n))
            ids = [el.get('id') for el in root.xpath('.//*[local-name()="cNvPr"]')]
            dupes = [k for k, v in Counter(ids).items() if v > 1]
            if dupes:
                fail(failures, f"{n}: duplicate shape id(s) within the slide: {dupes}")

def check_sldidlst_matches_slide_count(pptx_path, failures):
    with zipfile.ZipFile(pptx_path) as z:
        names = z.namelist()
        slides = [n for n in names if re.match(r'ppt/slides/slide\d+\.xml$', n)]
        prs_root = etree.fromstring(z.read('ppt/presentation.xml'))
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        sld_id_count = len(prs_root.findall('.//p:sldIdLst/p:sldId', ns))
    if sld_id_count != len(slides):
        fail(failures, f"sldIdLst count ({sld_id_count}) does not match actual slide files ({len(slides)})")

def check_animation_targets_exist(pptx_path, failures):
    with zipfile.ZipFile(pptx_path) as z:
        for n in z.namelist():
            if not re.match(r'ppt/slides/slide\d+\.xml$', n): continue
            root = etree.fromstring(z.read(n))
            shape_ids = {el.get('id') for el in root.xpath('.//*[local-name()="cNvPr"]')}
            for spid in root.xpath('.//*[local-name()="spTgt"]/@spid'):
                if spid not in shape_ids:
                    fail(failures, f"{n}: animation targets shape id={spid} which does not exist on this slide")

def check_non_numeric_rids(pptx_path, failures):
    with zipfile.ZipFile(pptx_path) as z:
        if 'ppt/_rels/presentation.xml.rels' not in z.namelist(): return
        content = z.read('ppt/_rels/presentation.xml.rels').decode('utf-8')
    non_numeric = re.findall(r'Id="(rId[A-Za-z][A-Za-z0-9]*)"', content)
    if non_numeric:
        fail(failures, f"presentation.xml.rels has non-numeric rId(s) (known repair-dialog cause): {non_numeric}")

def check_orphaned_media(pptx_path, failures):
    """A media file present in the archive but referenced by no relationship
    anywhere. Found via the T6W7 investigation: replace_image() left the
    original (replaced) image physically in the file even though nothing
    pointed to it any more - in one case that was the exact banned concept-
    cartoon image. Not itself a known repair-dialog trigger, but content that
    was supposed to be fully replaced should not still be sitting in the
    package, referenced or not."""
    with zipfile.ZipFile(pptx_path) as z:
        names = z.namelist()
        referenced = set()
        for n in names:
            if not n.endswith('.rels'): continue
            content = z.read(n).decode('utf-8', errors='ignore')
            part_dir = '/'.join(n.replace('_rels/', '').split('/')[:-1])
            for tgt in re.findall(r'Target="([^"]+)"', content):
                if tgt.startswith('http') or tgt.startswith('#'): continue
                segs = (part_dir.split('/') if part_dir else []) + tgt.split('/')
                out = []
                for s in segs:
                    if s == '..':
                        if out: out.pop()
                    elif s and s != '.':
                        out.append(s)
                referenced.add('/'.join(out))
        media = [n for n in names if n.startswith('ppt/media/')]
        orphaned = [m for m in media if m not in referenced]
    if orphaned:
        fail(failures, f"Orphaned media file(s) present but unreferenced by any relationship: {orphaned}")

def main():
    pptx_path, mtp_path, manifest_path = sys.argv[1:4]
    with open(mtp_path) as f: mtp = json.load(f)
    with open(manifest_path) as f: manifest = json.load(f)
    # Normalise unified format → legacy format so checks work unchanged
    if 'lessons' in mtp and 'lesson' not in mtp:
        mtp['lesson'] = _get_lesson_from_mtp(mtp, manifest)
    prs = Presentation(pptx_path)

    failures = []
    check_customxml(pptx_path, failures)
    check_slide_sequence(manifest, mtp, failures)
    check_required_present(manifest, failures)
    check_banned_text(prs, failures)
    check_concept_cartoon_content(prs, manifest, mtp, failures)
    check_lo_no_duplicate_panel(prs, manifest, failures)
    check_geometric_overlap(prs, failures)
    check_animation_pattern(pptx_path, failures)
    check_images_present(prs, manifest, mtp, failures)
    check_duplicate_shape_ids(pptx_path, failures)
    check_sldidlst_matches_slide_count(pptx_path, failures)
    check_animation_targets_exist(pptx_path, failures)
    check_non_numeric_rids(pptx_path, failures)
    check_orphaned_media(pptx_path, failures)

    if failures:
        print(f"VERIFY: FAIL ({len(failures)} issue(s))\n")
        for f_ in failures:
            print(f"  - {f_}")
        sys.exit(1)
    else:
        n_slides = len(list(prs.slides))
        print(f"VERIFY: PASS ({n_slides} slides checked)")
        sys.exit(0)

if __name__ == '__main__':
    main()

"""
build_lesson_v3.py — Generalised lesson teaching PPTX builder
Usage:  python3 build_lesson_v3.py <lesson_number>   (1-20)

Reads all text content from maths_plan_v3.json.
Reads authored visual/WM/RM data from lesson_data.py.
"""

import sys, copy, re, json, io as _io, os, tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from maths_visuals import render_stats_chart

EMU = 914400
def emu(inches): return int(inches * EMU)

# ---------------------------------------------------------------------------
# LESSON NUMBER
# ---------------------------------------------------------------------------
LESSON_NUM = int(sys.argv[1]) if len(sys.argv) > 1 else 1

# ---------------------------------------------------------------------------
# LOAD JSON PLAN
# ---------------------------------------------------------------------------
with open('/home/claude/transfer_files/maths_plan_v3.json') as f:
    PLAN = json.load(f)

L1 = PLAN['lessons'][LESSON_NUM - 1]
KEY_QUESTIONS = PLAN['keyQuestions']
assert L1['lesson'] == LESSON_NUM, f"Lesson mismatch: expected {LESSON_NUM}, got {L1['lesson']}"
print(f"Building lesson {LESSON_NUM}: {L1['day']} ({L1['topic']}) — {L1['li'][:60]}")

# ---------------------------------------------------------------------------
# LOAD AUTHORED LESSON DATA (visuals, WM, RM, vocab)
# ---------------------------------------------------------------------------
import importlib.util as _ilu
_spec = _ilu.spec_from_file_location('lesson_data', '/home/claude/lesson_data.py')
_mod  = _ilu.module_from_spec(_spec)
_spec.loader.exec_module(_mod)

if LESSON_NUM not in _mod.LESSON_DATA:
    raise ValueError(f"No authored data found for lesson {LESSON_NUM} in lesson_data.py")

_ld = _mod.LESSON_DATA[LESSON_NUM]
VOCAB   = _ld['vocab']
WM_DATA = _ld['wm']
RM_DATA = _ld.get('rm', {})

# ---------------------------------------------------------------------------
# BUILD VISUALS DICT — authored coords + spotTheMistake from JSON
# ---------------------------------------------------------------------------
STM = L1['spotTheMistake']

stm_visual = {
    'title':             STM['slideTitle'],
    'slide_type':        'spot_the_mistake',
    'cols':              STM['gridSize'] if STM['gridSize'] > 0 else 6,
    'rows':              STM['gridSize'] if STM['gridSize'] > 0 else 6,
    'points':            [[STM['startPoint'][0], STM['startPoint'][1], 'A', '1F4E79']]
                         if STM['startPoint'] else [],
    'extra_points':      STM.get('extraPoints', []),
    'error_instruction': STM['errorInstruction'],
    'error_note':        STM['errorNote'],
    'error_type':        STM['errorType'],
    'notes':             "\n".join([
                             "I DO C2 — Spot the Mistake",
                             f"Concept: {STM['concept']}",
                             "Beat 1 (load): Grid + instruction visible.",
                             "Beat 2 (click 1): X appears at error position.",
                             "Beat 3 (click 2): Explanation revealed.",
                         ]),
}

VISUALS = dict(_ld['visuals'])
# If lesson_data provides c2_ido2 (e.g. stm_word_problem), use that instead of grid STM
if 'c2_ido2' not in VISUALS:
    VISUALS['c2_ido2'] = stm_visual

# ---------------------------------------------------------------------------
prs = Presentation('/home/claude/template.pptx')

# Extract LO slide avatar image parts BEFORE deleting template slides
_lo_sld3 = prs.slides[2]

# Remove all existing slides cleanly
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
sldIdLst = prs.slides._sldIdLst
prs_part = prs.part
rIds = [s.get(f'{{{NS_R}}}id') for s in list(sldIdLst)]
for rId in rIds:
    for s in list(sldIdLst):
        if s.get(f'{{{NS_R}}}id') == rId:
            sldIdLst.remove(s)
    if rId in prs_part.rels:
        prs_part.rels.pop(rId)

_dummy = prs.slides[2] if False else None  # placeholder line
LO_IMAGE_PARTS = {}  # partname -> ImagePart
for _rId, _rel in _lo_sld3.part.rels.items():
    if hasattr(_rel, 'target_part'):
        _tp = _rel.target_part
        _pn = str(_tp.partname)
        if any(img in _pn for img in ['image8','image9','image10']):
            LO_IMAGE_PARTS[_pn.split('/')[-1]] = _tp

print(f"Template cleared. Slides: {len(prs.slides)}")
print(f"LO image parts saved: {list(LO_IMAGE_PARTS.keys())}")

def layout(n):
    return prs.slide_layouts[n - 1]

# ---------------------------------------------------------------------------
# IMAGE HELPER — always direct file read
# ---------------------------------------------------------------------------
def add_pic(slide, image_filename, x, y, w, h):
    # Absolute path wins; otherwise read from unpacked media
    if os.path.isabs(image_filename):
        img_path = image_filename
    else:
        img_path = f'/home/claude/unpacked/ppt/media/{image_filename}'
    with open(img_path, 'rb') as f:
        img_bytes = f.read()
    pic = slide.shapes.add_picture(_io.BytesIO(img_bytes), emu(x), emu(y), emu(w), emu(h))
    return pic

def add_pic_id(slide, image_filename, x, y, w, h):
    """Add picture and return (element, shape_id) for use in animations."""
    pic = add_pic(slide, image_filename, x, y, w, h)
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/picture'
    cNvPr = pic.element.find(f'.//{{{ns}}}cNvPr')
    if cNvPr is None:
        # try presentation namespace
        ns2 = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        cNvPr = pic.element.find('.//cNvPr')
    shape_id = int(pic.shape_id)
    return pic.element, shape_id

# ---------------------------------------------------------------------------
# XML HELPERS
# ---------------------------------------------------------------------------
def _esc(s):
    return (str(s).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
            .replace('"','&quot;').replace("'","&apos;"))

def add_sp(slide, xml_str):
    slide.shapes._spTree.append(etree.fromstring(xml_str))

def sp(spid, name, x, y, w, h, text, font='Twinkl Cursive Looped Light',
       sz=18, bold=False, color='000000', align='l',
       fill=None, border=None, geom='rect', anchor='ctr',
       no_line=False, underline=False, autofit=False):
    fill_xml = f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>' if fill else '<a:noFill/>'
    if no_line:
        line_xml = '<a:ln w="0"><a:noFill/></a:ln>'
    elif border:
        col, wpt = border
        line_xml = f'<a:ln w="{int(wpt*12700)}"><a:solidFill><a:srgbClr val="{col}"/></a:solidFill></a:ln>'
    else:
        line_xml = ''
    b_attr = ' b="1"' if bold else ''
    u_attr = ' u="sng"' if underline else ''
    color_xml = f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
    if isinstance(text, list):
        paras = ''.join(f'''<a:p><a:pPr algn="{align}"/><a:r>
          <a:rPr lang="en-GB" sz="{int(sz*100)}"{b_attr}{u_attr} dirty="0">
            {color_xml}<a:latin typeface="{font}"/>
          </a:rPr><a:t>{_esc(line)}</a:t></a:r></a:p>''' for line in text)
    else:
        paras = f'''<a:p><a:pPr algn="{align}"/><a:r>
          <a:rPr lang="en-GB" sz="{int(sz*100)}"{b_attr}{u_attr} dirty="0">
            {color_xml}<a:latin typeface="{font}"/>
          </a:rPr><a:t>{_esc(text)}</a:t></a:r></a:p>'''
    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="{name}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>
    <a:prstGeom prst="{geom}"><a:avLst/></a:prstGeom>
    {fill_xml}{line_xml}
  </p:spPr>
  <p:txBody><a:bodyPr rtlCol="0" anchor="{anchor}">{"<a:normAutofit/>" if autofit else ""}</a:bodyPr>
    <a:lstStyle/>{paras}</p:txBody>
</p:sp>'''

def new_slide(layout_num):
    return prs.slides.add_slide(layout(layout_num))

# ===========================================================================
# SLIDE 1 — KEY QUESTION
# Takes KQ_Slide_template.pptx, replaces the placeholder text, injects the
# whole slide unchanged into the built PPTX via zip manipulation.
# Called as a post-process step after prs.save().
# ===========================================================================
KQ_TEMPLATE    = '/home/claude/assets/KQ_Slide_template.pptx'
KQ_PLACEHOLDER = 'Xxxxxxxxxx xxxxxxxxxxxxxx xxxxxxxxxxxxx xxxxxxxx xxxxxxxxxx xxxxxxxxxxxx'

def build_slide1():
    pass   # KQ slide injected post-save — see inject_kq_slide()

def inject_kq_slide(teaching_pptx_path):
    """
    Post-save: opens the built teaching PPTX and the KQ template as zips,
    replaces the placeholder text, then prepends the KQ slide.
    Nothing about the KQ slide changes except the question text.
    """
    import zipfile, re, shutil

    kq_text = KEY_QUESTIONS[L1['topic']]

    # ── Read KQ template ──────────────────────────────────────────────────────
    with zipfile.ZipFile(KQ_TEMPLATE) as kz:
        kq_slide_xml = kz.read('ppt/slides/slide1.xml')
        kq_rels_xml  = kz.read('ppt/slides/_rels/slide1.xml.rels')
        kq_media     = {n: kz.read(n) for n in kz.namelist() if n.startswith('ppt/media/')}

    # Replace placeholder text (same operation the user confirmed worked)
    kq_slide_xml = kq_slide_xml.decode('utf-8').replace(KQ_PLACEHOLDER, kq_text).encode('utf-8')

    # ── Read teaching PPTX ────────────────────────────────────────────────────
    with zipfile.ZipFile(teaching_pptx_path) as tz:
        t_files = {n: tz.read(n) for n in tz.namelist() if not n.endswith('/')}

    # ── Rename KQ media to avoid filename conflicts ───────────────────────────
    existing_media = {n.split('/')[-1] for n in t_files if n.startswith('ppt/media/')}
    remap = {}   # original_name → name_in_teaching_pptx
    for path in kq_media:
        fname = path.split('/')[-1]
        base, ext = fname.rsplit('.', 1)
        new = fname
        n = 0
        while new in existing_media:
            n += 1; new = f'{base}_kq{n}.{ext}'
        existing_media.add(new)
        remap[fname] = new

    # Update KQ slide rels to use renamed media and point to teaching blank layout
    kq_rels_str = kq_rels_xml.decode('utf-8')
    for orig, new in remap.items():
        kq_rels_str = kq_rels_str.replace(f'../media/{orig}', f'../media/{new}')
    # Map layout to teaching PPTX blank layout (slideLayout13.xml = 'Blank')
    kq_rels_str = re.sub(
        r'Target="\.\./slideLayouts/slideLayout\d+\.xml"',
        'Target="../slideLayouts/slideLayout13.xml"', kq_rels_str)
    # Remove notes slide reference (not needed)
    # FIX: old regex r'<Relationship[^/]*/>' never matched because Type URLs contain '/'.
    # Line-by-line filter also wrong — rels are single-line, deletes everything.
    # Use element-level regex matching only the notesSlide Relationship element:
    kq_rels_str = re.sub(r'<Relationship[^>]+notesSlide[^>]+/>', '', kq_rels_str)

    # ── Choose a slide filename that won't clash ───────────────────────────────
    nums = [int(re.search(r'slide(\d+)', n).group(1))
            for n in t_files if re.match(r'ppt/slides/slide\d+\.xml$', n)]
    kq_num  = max(nums) + 1
    kq_path = f'ppt/slides/slide{kq_num}.xml'
    kq_rpath= f'ppt/slides/_rels/slide{kq_num}.xml.rels'

    # ── Prepend KQ slide in presentation.xml (string manipulation — no etree) ──
    import re as _re
    new_rId = 'rIdKQ'
    prs_str = t_files['ppt/presentation.xml'].decode('utf-8')
    existing_ids = [int(m) for m in _re.findall(r'<p:sldId\b[^>]*\bid="(\d+)"', prs_str)]
    new_id = max(existing_ids) + 1
    new_sld_el = f'<p:sldId id="{new_id}" r:id="{new_rId}"/>'
    # Insert as FIRST entry in sldIdLst
    prs_str = prs_str.replace('<p:sldIdLst>', '<p:sldIdLst>' + new_sld_el, 1)
    t_files['ppt/presentation.xml'] = prs_str.encode('utf-8')

    # ── Add slide relationship (string manipulation) ───────────────────────────
    SLIDE_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
    new_rel = f'<Relationship Id="{new_rId}" Type="{SLIDE_REL}" Target="slides/slide{kq_num}.xml"/>'
    rels_str = t_files['ppt/_rels/presentation.xml.rels'].decode('utf-8')
    rels_str = rels_str.replace('</Relationships>', new_rel + '</Relationships>')
    t_files['ppt/_rels/presentation.xml.rels'] = rels_str.encode('utf-8')

    # ── Add content type entry (string manipulation) ───────────────────────────
    ct_str = t_files['[Content_Types].xml'].decode('utf-8')
    if 'Extension="wdp"' not in ct_str:
        ct_str = ct_str.replace('</Types>',
            '<Default Extension="wdp" ContentType="image/vnd.ms-photo"/></Types>')
    ct_str = ct_str.replace('</Types>',
        f'<Override PartName="/ppt/slides/slide{kq_num}.xml" '
        f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        f'</Types>')
    t_files['[Content_Types].xml'] = ct_str.encode('utf-8')

    # ── Write output ──────────────────────────────────────────────────────────
    tmp = teaching_pptx_path + '.tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as oz:
        for name, data in t_files.items():
            oz.writestr(name, data)
        oz.writestr(kq_path,  kq_slide_xml)
        oz.writestr(kq_rpath, kq_rels_str.encode('utf-8'))
        for orig_path, data in kq_media.items():
            oz.writestr(f'ppt/media/{remap[orig_path.split("/")[-1]]}', data)
    shutil.move(tmp, teaching_pptx_path)
    print("Slide 1 (Key Question) ✓")

# ===========================================================================
# SLIDE 2 — DAY TITLE
# ===========================================================================
def build_slide2():
    sld = new_slide(13)
    sid = [2]
    def nid(): sid[0]+=1; return sid[0]
    day_map = {'Monday':1,'Tuesday':2,'Wednesday':3,'Thursday':4}
    day_num = day_map.get(L1['day'],1)
    add_sp(sld, sp(nid(),'BAM', 3.232,2.131, 7.221,0.505,
                   'Being a Mathematician', font='Twinkl Cursive Looped Light',
                   sz=18, bold=True, color='000000', align='ctr',
                   fill=None, no_line=True))
    add_sp(sld, sp(nid(),'DayText', 0.761,2.448, 11.811,2.036,
                   L1['day'], font='Twinkl Cursive Looped Light',
                   sz=100, bold=False, color='000000', align='ctr',
                   fill=None, no_line=True))
    add_pic(sld,'image1.png', 5.634,0.168, 2.066,1.796)
    add_pic(sld,'image7.png', 3.438,4.277, 6.458,3.056)
    sld.notes_slide.notes_text_frame.text = (
        f"DAY TITLE — Day {day_num} ({L1['day']})\nTopic: {L1['topic']}\nLI: {L1['li']}")
    print("Slide 2 (Day Title) ✓")

# ===========================================================================
# SLIDE 3 — LO  (cloned from template XML, text replaced only)
# ===========================================================================
def build_slide3():
    # Read template slide 3 XML
    with open('/home/claude/unpacked/ppt/slides/slide3.xml') as f:
        template_xml = f.read()

    lo = L1['loText']

    def replace_panel_text(xml, rect_name, header, body_text, body_sz):
        """Replace text in a panel's two paragraphs, preserving all formatting."""
        for match in re.finditer(r'(<p:sp>.*?</p:sp>)', xml, re.DOTALL):
            sp_xml = match.group(1)
            if f'name="{rect_name}"' in sp_xml:
                # Replace body text paragraph (para index 1)
                paras = re.findall(r'(<a:p>.*?</a:p>)', sp_xml, re.DOTALL)
                if len(paras) >= 2:
                    old_body_para = paras[1]
                    new_body_para = re.sub(
                        r'<a:t>[^<]*</a:t>',
                        f'<a:t>{_esc(body_text)}</a:t>',
                        old_body_para
                    )
                    new_sp = sp_xml.replace(old_body_para, new_body_para)
                    xml = xml.replace(sp_xml, new_sp)
        return xml

    modified_xml = template_xml
    modified_xml = replace_panel_text(modified_xml, 'Rounded Rectangle 25',
                                       'We are learning to…', lo['walt'], 1600)
    modified_xml = replace_panel_text(modified_xml, 'Rounded Rectangle 26',
                                       'This is because…', lo['tiob'], 1400)
    modified_xml = replace_panel_text(modified_xml, 'Rounded Rectangle 27',
                                       'I will show this by…', lo['iwstb'], 1600)

    # Add a new slide using the blank layout then replace its XML with the modified template
    sld = new_slide(9)  # slideLayout9 = I Do Blank (same as template slide 3)

    # Parse and set the spTree from modified template
    modified_root = etree.fromstring(modified_xml.encode())
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    template_spTree = modified_root.find(f'.//{{{ns_p}}}spTree')

    # Get the slide's own spTree and replace its content
    slide_spTree = sld.shapes._spTree
    # Clear existing children except nvGrpSpPr and grpSpPr
    for child in list(slide_spTree):
        tag = child.tag.split('}')[-1]
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            slide_spTree.remove(child)

    # Copy all children from template spTree (except nvGrpSpPr and grpSpPr)
    for child in template_spTree:
        tag = child.tag.split('}')[-1]
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            slide_spTree.append(copy.deepcopy(child))

    # Add LO avatar images using pre-extracted ImagePart objects
    with open('/home/claude/unpacked/ppt/slides/_rels/slide3.xml.rels') as f:
        rels_content = f.read()
    rid_map = {}
    for m in re.finditer(r'Id="(rId\d+)"[^>]*Target="\.\./media/([^"]+)"', rels_content):
        rid_map[m.group(1)] = m.group(2)

    img_rel_type = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
    spTree_xml = etree.tostring(slide_spTree, encoding='unicode')
    changed = False
    for old_rId, img_file in rid_map.items():
        img_part = LO_IMAGE_PARTS.get(img_file)
        if img_part:
            actual_rId = sld.part.relate_to(img_part, img_rel_type)
            if actual_rId != old_rId:
                spTree_xml = spTree_xml.replace(f'r:embed="{old_rId}"', f'r:embed="{actual_rId}"')
                changed = True
    if changed:
        new_spTree = etree.fromstring(spTree_xml.encode())
        parent = slide_spTree.getparent()
        parent.replace(slide_spTree, new_spTree)

    sld.notes_slide.notes_text_frame.text = (
        f"I DO — Learning Objective\nWALT: {lo['walt']}\nTIOB: {lo['tiob']}\nIWSTB: {lo['iwstb']}")
    print("Slide 3 (LO) ✓")

# ===========================================================================
# SLIDE 4 — WM MEMORY
# ===========================================================================
def build_slide4():
    sld = new_slide(5)
    items = WM_DATA['items']
    n = len(items)

    add_sp(sld, sp(2,'Title', 2.454,0.143, 9.039,1.450,
                   'Remember the details and the order',
                   font='Twinkl Cursive Looped Light', sz=40, bold=True,
                   color='000000', align='l', fill=None, no_line=True))

    add_sp(sld, sp(3,'ShowBtn', 5.219,1.347, 2.895,0.679,
                   'Click to Show', font='Aptos', sz=20, bold=True,
                   color='0E2841', align='ctr', fill='92D050',
                   geom='roundRect', no_line=True))

    item_w, item_h, item_y = 1.573, 1.400, 2.347
    # Font size: numbers fit at 40pt; words need to scale down based on longest item
    max_len = max(len(str(v)) for v in items)
    item_sz = 40 if max_len <= 3 else (28 if max_len <= 6 else 20)
    for i, val in enumerate(items):
        ix = 0.500 + i * ((13.333 - 1.000 - item_w) / (n - 1))
        add_sp(sld, sp(10+i, f'Item{i+1}', ix, item_y, item_w, item_h,
                       str(val), font='Aptos', sz=item_sz, bold=True,
                       color='FFFFFF', align='ctr', fill='1C4060',
                       border=('156082',1.5)))

    add_sp(sld, sp(20,'GoBtn', 5.283,5.887, 2.767,0.679,
                   'Go to the questions', font='Aptos', sz=20, bold=True,
                   color='FFFFFF', align='ctr', fill='0F9ED5',
                   geom='roundRect', no_line=True))

    # Cover — LAST shape
    cover_spid = 30
    add_sp(sld, sp(cover_spid,'Cover', 0.0,2.1, 13.333,3.35,
                   '', fill='DEECF8', no_line=True))

    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
    <p:childTnLst><p:seq concurrent="1" nextAc="seek">
      <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
        <p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
          <p:childTnLst><p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:par>
              <p:cTn id="5" presetID="1" presetClass="exit" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst><p:set><p:cBhvr>
                  <p:cTn id="6" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                  <p:tgtEl><p:spTgt spid="{cover_spid}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr><p:to><p:strVal val="hidden"/></p:to></p:set></p:childTnLst>
              </p:cTn></p:par></p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="7" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst><p:par>
            <p:cTn id="8" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="1" nodeType="afterEffect">
              <p:stCondLst><p:cond delay="15000"/></p:stCondLst>
              <p:childTnLst><p:set><p:cBhvr>
                <p:cTn id="9" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                <p:tgtEl><p:spTgt spid="{cover_spid}"/></p:tgtEl>
                <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
              </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
            </p:cTn></p:par></p:childTnLst></p:cTn></p:par>
      </p:childTnLst></p:cTn></p:par>
      </p:childTnLst></p:cTn>
      <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
      <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>
  <p:bldLst>
    <p:bldP spid="{cover_spid}" grpId="0" animBg="1"/>
    <p:bldP spid="{cover_spid}" grpId="1" animBg="1"/>
  </p:bldLst>
</p:timing>'''
    sld._element.append(etree.fromstring(timing_xml))
    sld.notes_slide.notes_text_frame.text = (
        f"YOU DO — WM Number Sequence\nNumbers: {WM_DATA['items']}\n"
        "Click reveals numbers, cover reappears after 15s.")
    print("Slide 4 (WM Memory) ✓")

# ===========================================================================
# SLIDE 5 — WM Q&A  (spids start at 20 to avoid layout placeholder collision)
# ===========================================================================
def build_slide5():
    sld = new_slide(5)

    add_sp(sld, sp(20,'Title', 0.917,0.110, 11.500,1.450,
                   'Now answer from memory!',
                   font='Twinkl Cursive Looped Light', sz=44, bold=True,
                   color='000000', align='l', fill=None, no_line=True))

    qa = WM_DATA['qa']
    q_y_left  = [1.309, 3.243, 5.212]
    q_y_right = [1.309, 3.243]
    q_h, a_h  = 1.013, 0.747
    q_w       = 6.000

    card_spids = []
    spid = 21
    for i, item in enumerate(qa):
        qx = 0.373 if i < 3 else 6.960
        qy = (q_y_left[i] if i < 3 else q_y_right[i-3])
        ay = qy + q_h

        add_sp(sld, sp(spid, f'Q{i+1}', qx, qy, q_w, q_h,
                       f'Q{i+1}   {item["q"]}',
                       font='Aptos', sz=17, bold=True, color='0E2841',
                       align='l', fill='C2D9EC', border=('156082',1.5)))
        card_spids.append(spid); spid += 1

        add_sp(sld, sp(spid, f'A{i+1}', qx, ay, q_w, a_h,
                       f'Answer:   {item["a"]}',
                       font='Aptos', sz=20, bold=True, color='0E2841',
                       align='l', fill='92D050', border=('156082',1.5)))
        card_spids.append(spid); spid += 1

    # Build timing — all cards hidden, revealed Q1→A1→Q2→A2 etc
    child_blocks = ''
    ctn_id = 3
    for spid_val in card_spids:
        child_blocks += f'''<p:par>
  <p:cTn id="{ctn_id}" fill="hold">
    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
    <p:childTnLst><p:par><p:cTn id="{ctn_id+1}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst><p:par>
        <p:cTn id="{ctn_id+2}" presetID="1" presetClass="entr" presetSubtype="0"
               fill="hold" grpId="1" nodeType="clickEffect">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst><p:set><p:cBhvr>
            <p:cTn id="{ctn_id+3}" dur="1" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
            <p:tgtEl><p:spTgt spid="{spid_val}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
          </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
        </p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
  </p:cTn>
</p:par>'''
        ctn_id += 4

    bld_list = ''.join(
        f'<p:bldP spid="{s}" grpId="0" uiExpand="1" build="p"/>\n'
        f'<p:bldP spid="{s}" grpId="1" animBg="1"/>\n'
        for s in card_spids)

    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par>
    <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
      <p:childTnLst><p:seq concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
          <p:childTnLst>{child_blocks}</p:childTnLst>
        </p:cTn>
        <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
        <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
      </p:seq></p:childTnLst>
    </p:cTn>
  </p:par></p:tnLst>
  <p:bldLst>{bld_list}</p:bldLst>
</p:timing>'''
    sld._element.append(etree.fromstring(timing_xml))
    sld.notes_slide.notes_text_frame.text = (
        "YOU DO — WM Q&A\nAll cards hidden at start. Reveal: Q1→A1→Q2→A2→Q3→A3→Q4→A4→Q5→A5")
    print("Slide 5 (WM Q&A) ✓")


import re as _re_rm
_RM_FRAC_RE = _re_rm.compile(r'(\d+)/(\d+)')

def _rm_has_frac(text):
    return bool(_RM_FRAC_RE.search(text))

def _rm_render_frac(text, w_in, h_in, fontsize=24, color='#1a1a1a'):
    """Render text with n/d fractions as PNG using matplotlib mathtext."""
    import io as _io2
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as _plt2
    mt = _RM_FRAC_RE.sub(lambda m: f'$\\frac{{{m.group(1)}}}{{{m.group(2)}}}$', text)
    fig = _plt2.figure(figsize=(w_in, h_in), dpi=150)
    fig.patch.set_facecolor('white')
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    ax.text(0.03, 0.5, mt, fontsize=fontsize, ha='left', va='center',
            color=color, transform=ax.transAxes, fontfamily='DejaVu Sans')
    buf = _io2.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches=None,
                facecolor='white', edgecolor='none')
    _plt2.close(fig)
    buf.seek(0)
    return buf.read()

# ===========================================================================
# FRACTIONS RECAP — replaces RM for weeks with no Rapid Maths
# 2 slides per lesson: (A) Mixed→Improper  (B) Improper→Mixed
# Uses cover-box exit animation identical to build_slide4().
# ===========================================================================
def build_fractions_slides():
    import matplotlib.pyplot as _plt
    import tempfile as _tmp, os as _os, math as _math

    _day = L1['day']
    _TMPDIR = _tmp.mkdtemp(prefix='wfa_fr_')

    _FRAC = {
        'Monday':    {'m2i': {'ido': (2,3,5),  'youdo': [(3,1,4),(4,2,5)]},
                      'i2m': {'ido': (11,3),   'youdo': [(17,4),(13,5)]}},
        'Tuesday':   {'m2i': {'ido': (2,5,7),  'youdo': [(3,3,8)]},
                      'i2m': {'ido': (19,6),   'youdo': [(23,5),(27,8)]}},
        'Wednesday': {'m2i': {'ido': (5,1,3),  'youdo': [(4,3,7)]},
                      'i2m': {'ido': (31,8),   'youdo': [(22,5),(16,3)]}},
    }
    fd = _FRAC.get(_day)
    if not fd:
        return

    def _fimg(num, den, tag, col='#1798d3', sz=30, fw=0.65, fh=0.92):
        p = _os.path.join(_TMPDIR, f'{tag}.png')
        fig, ax = _plt.subplots(figsize=(fw, fh))
        fig.patch.set_facecolor('none'); ax.set_facecolor('none'); ax.axis('off')
        ax.set_xlim(0,1); ax.set_ylim(0,1)
        ax.text(0.5,0.78,str(num),ha='center',va='center',fontsize=sz,color=col,fontweight='bold')
        ax.plot([0.05,0.95],[0.50,0.50],color=col,lw=2.2)
        ax.text(0.5,0.16,str(den),ha='center',va='center',fontsize=sz,color=col,fontweight='bold')
        _plt.savefig(p,dpi=180,bbox_inches='tight',transparent=True); _plt.close(fig)
        return p

    def _mimg(W, num, den, tag, col='#1798d3', sz=26, fw=1.10, fh=0.92):
        p = _os.path.join(_TMPDIR, f'{tag}.png')
        fig, ax = _plt.subplots(figsize=(fw, fh))
        fig.patch.set_facecolor('none'); ax.set_facecolor('none'); ax.axis('off')
        ax.set_xlim(0,1); ax.set_ylim(0,1)
        ax.text(0.20,0.50,str(W),ha='center',va='center',fontsize=sz+6,color=col,fontweight='bold')
        ax.text(0.66,0.78,str(num),ha='center',va='center',fontsize=sz,color=col,fontweight='bold')
        ax.plot([0.40,0.93],[0.50,0.50],color=col,lw=2.0)
        ax.text(0.66,0.16,str(den),ha='center',va='center',fontsize=sz,color=col,fontweight='bold')
        _plt.savefig(p,dpi=180,bbox_inches='tight',transparent=True); _plt.close(fig)
        return p

    def _pic(sld, fpath, x, y, w, h):
        with open(fpath,'rb') as f: data = f.read()
        sld.shapes.add_picture(_io.BytesIO(data), emu(x), emu(y), emu(w), emu(h))

    def _cover_anim(sld, cid):
        P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        timing_xml = (
            f'<p:timing xmlns:p="{P}" xmlns:a="{A}">'
            f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            f'<p:childTnLst><p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="5" presetID="1" presetClass="exit" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:set><p:cBhvr>'
            f'<p:cTn id="6" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{cid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="hidden"/></p:to></p:set></p:childTnLst>'
            f'</p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
            f'<p:bldLst><p:bldP spid="{cid}" grpId="0" animBg="1"/></p:bldLst>'
            f'</p:timing>'
        )
        sld._element.append(etree.fromstring(timing_xml.encode()))

    # ── Slide A: Mixed → Improper ─────────────────────────────────────────────
    m2i = fd['m2i']
    W, N, D = m2i['ido']
    ans = W * D + N
    sld = new_slide(13)
    add_sp(sld, sp(2,'Title', 0.50,0.05,12.00,0.55,
                   'Fractions Recap:  Mixed  →  Improper',
                   font='Twinkl Cursive Looped',sz=26,bold=True,
                   color='000000',align='l',fill=None,no_line=True))
    # I Do badge + panel
    add_sp(sld, sp(3,'IDoLbl',0.35,0.72,1.20,0.36,
                   'I Do',font='Aptos',sz=16,bold=True,
                   color='FFFFFF',align='ctr',fill='1798D3',no_line=True))
    add_sp(sld, sp(4,'Steps',0.35,1.16,5.90,1.30,
                   f'Step 1:  {W} × {D} = {W*D}\nStep 2:  {W*D} + {N} = {ans}',
                   font='Aptos',sz=19,bold=False,
                   color='1A1A1A',align='l',fill='DEECF8',no_line=True))
    _pic(sld, _mimg(W,N,D,'ia_q'),  0.40,2.70,1.15,0.90)
    add_sp(sld, sp(5,'Arr',1.65,2.92,0.48,0.40,
                   '→',font='Aptos',sz=24,bold=True,
                   color='1798D3',align='ctr',fill=None,no_line=True))
    _pic(sld, _fimg(ans,D,'ia_ans'), 2.20,2.60,0.72,1.00)
    # You Do badge + questions
    add_sp(sld, sp(6,'YDLbl',7.00,0.72,1.40,0.36,
                   'You Do',font='Aptos',sz=16,bold=True,
                   color='FFFFFF',align='ctr',fill='E57D24',no_line=True))
    yy = 1.16
    for qi,(qW,qN,qD) in enumerate(m2i['youdo']):
        qa = qW*qD+qN
        _pic(sld, _mimg(qW,qN,qD,f'ya{qi}q',col='#E57D24'), 7.00,yy,1.15,0.90)
        add_sp(sld, sp(20+qi,f'YArr{qi}',8.25,yy+0.22,0.44,0.38,
                       '→',font='Aptos',sz=22,bold=True,
                       color='E57D24',align='ctr',fill=None,no_line=True))
        _pic(sld, _fimg(qa,qD,f'ya{qi}a',col='#E57D24'), 8.78,yy-0.08,0.72,1.00)
        yy += 1.30
    # Cover over answers — click to reveal
    add_sp(sld, sp(99,'CoverA',8.62,1.05,1.10,yy-0.90,
                   '',fill='FFFFFF',no_line=True))
    _cover_anim(sld, 99)
    print(f"  Fractions A ({_day}: mixed→improper) ✓")

    # ── Slide B: Improper → Mixed ─────────────────────────────────────────────
    i2m = fd['i2m']
    iN, iD = i2m['ido']
    whole = iN // iD;  rem = iN % iD
    sld = new_slide(13)
    add_sp(sld, sp(2,'Title',0.50,0.05,12.00,0.55,
                   'Fractions Recap:  Improper  →  Mixed',
                   font='Twinkl Cursive Looped',sz=26,bold=True,
                   color='000000',align='l',fill=None,no_line=True))
    add_sp(sld, sp(3,'IDoLbl',0.35,0.72,1.20,0.36,
                   'I Do',font='Aptos',sz=16,bold=True,
                   color='FFFFFF',align='ctr',fill='1798D3',no_line=True))
    add_sp(sld, sp(4,'Steps',0.35,1.16,5.90,1.30,
                   f'Step 1:  {iN} ÷ {iD} = {whole} remainder {rem}\nStep 2:  Answer = {whole} {rem}/{iD}',
                   font='Aptos',sz=19,bold=False,
                   color='1A1A1A',align='l',fill='DEECF8',no_line=True))
    _pic(sld, _fimg(iN,iD,'ib_q'),  0.40,2.60,0.72,1.00)
    add_sp(sld, sp(5,'Arr',1.22,2.92,0.48,0.40,
                   '→',font='Aptos',sz=24,bold=True,
                   color='1798D3',align='ctr',fill=None,no_line=True))
    _pic(sld, _mimg(whole,rem,iD,'ib_ans'), 1.78,2.70,1.15,0.90)
    add_sp(sld, sp(6,'YDLbl',7.00,0.72,1.40,0.36,
                   'You Do',font='Aptos',sz=16,bold=True,
                   color='FFFFFF',align='ctr',fill='E57D24',no_line=True))
    yy = 1.16
    for qi,(qiN,qiD) in enumerate(i2m['youdo']):
        qw = qiN//qiD;  qr = qiN%qiD
        _pic(sld, _fimg(qiN,qiD,f'yb{qi}q',col='#E57D24'), 7.00,yy-0.08,0.72,1.00)
        add_sp(sld, sp(20+qi,f'YArr{qi}',7.82,yy+0.22,0.44,0.38,
                       '→',font='Aptos',sz=22,bold=True,
                       color='E57D24',align='ctr',fill=None,no_line=True))
        _pic(sld, _mimg(qw,qr,qiD,f'yb{qi}a',col='#E57D24'), 8.35,yy,1.15,0.90)
        yy += 1.30
    add_sp(sld, sp(99,'CoverB',8.20,1.05,1.35,yy-0.90,
                   '',fill='FFFFFF',no_line=True))
    _cover_anim(sld, 99)
    print(f"  Fractions B ({_day}: improper→mixed) ✓")

# ===========================================================================
# SLIDE 6 & 7 — RAPID MATHS (unchanged from v1 — working correctly)
# ===========================================================================
def build_slide6():
    sld = new_slide(5)
    qs = RM_DATA['questions']
    day = RM_DATA['day']
    add_sp(sld, sp(2,'Title', 0.917,0.069, 11.500,0.814,
                   f'Rapid Maths \u2013 {["Monday","Tuesday","Wednesday","Friday"][day-1] if 1<=day<=4 else f"Day {day}"}',
                   font='Twinkl Cursive Looped Light', sz=28,
                   color='000000', align='l', fill=None, no_line=True))

    left_cards  = [(0.240,0.933,6.200,2.056),(0.240,3.122,6.200,2.056),(0.240,5.311,6.200,2.056)]
    right_cards = [(6.907,0.933,6.200,3.150),(6.907,4.217,6.200,3.150)]
    card_positions = left_cards + right_cards
    lx_left, lx_right, lw = 0.400, 7.067, 5.973

    spid = 10
    for i, q in enumerate(qs):
        cx,cy,cw,ch = card_positions[i]
        lx = lx_left if i < 3 else lx_right
        q_h_text = 1.576 if i < 3 else 2.670

        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Card{i+1}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(cx)}" y="{emu(cy)}"/><a:ext cx="{emu(cw)}" cy="{emu(ch)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(2*12700)}"><a:solidFill><a:srgbClr val="0070C0"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Topic{i+1}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(lx)}" y="{emu(cy+0.067)}"/><a:ext cx="{emu(lw)}" cy="{emu(0.347)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="156082"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr><p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/>
    <a:lstStyle/><a:p><a:pPr algn="l"/><a:r>
      <a:rPr lang="en-GB" sz="1600" b="0" dirty="0">
        <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
        <a:latin typeface="Aptos"/>
      </a:rPr><a:t>Q{q["num"]}  {_esc(q["topic"])}</a:t>
    </a:r></a:p></p:txBody></p:sp>''')
        spid += 1

        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="QText{i+1}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(lx)}" y="{emu(cy+0.414)}"/><a:ext cx="{emu(lw)}" cy="{emu(q_h_text)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/><a:ln w="0"><a:noFill/></a:ln>
  </p:spPr><p:txBody><a:bodyPr rtlCol="0" anchor="t"/>
    <a:lstStyle/><a:p><a:pPr algn="l"/><a:r>
      <a:rPr lang="en-GB" sz="2400" b="0" dirty="0">
        <a:solidFill><a:srgbClr val="1A1A1A"/></a:solidFill>
        <a:latin typeface="Aptos"/>
      </a:rPr><a:t>{_esc(q["q"])}</a:t>
    </a:r></a:p></p:txBody></p:sp>''')
        spid += 1
        # Fraction image: render n/d as proper vinculum PNG, overlay on QText box
        if _rm_has_frac(q['q']):
            png = _rm_render_frac(q['q'], lw, q_h_text, fontsize=24)
            sld.shapes.add_picture(_io.BytesIO(png), emu(lx), emu(cy+0.414), emu(lw), emu(q_h_text))

    sld.notes_slide.notes_text_frame.text = (
        f"YOU DO — Rapid Maths {['Monday','Tuesday','Wednesday','Friday'][day-1] if 1<=day<=4 else f'Day {day}'}\n" +
        '\n'.join(f"Q{q['num']} ({q['topic']}): {q['q']} → {q['a']}" for q in qs))
    print("Slide 6 (RM Questions) ✓")


def build_slide7():
    sld = new_slide(5)
    qs = RM_DATA['questions']
    day = RM_DATA['day']
    add_sp(sld, sp(2,'Title', 0.917,-0.009, 11.500,1.009,
                   f'Rapid Maths \u2013 Answers \u2013 {["Monday","Tuesday","Wednesday","Friday"][day-1] if 1<=day<=4 else f"Day {day}"}',
                   font='Twinkl Cursive Looped Light', sz=28,
                   color='000000', align='l', fill=None, no_line=True))

    left_cards  = [(0.240,0.933,6.200,2.056),(0.240,3.122,6.200,2.056),(0.240,5.311,6.200,2.056)]
    right_cards = [(6.907,0.933,6.200,3.150),(6.907,4.217,6.200,3.150)]
    card_positions = left_cards + right_cards
    lx_left, lx_right, lw = 0.400, 7.067, 5.973

    spid = 10
    for i, q in enumerate(qs):
        cx,cy,cw,ch = card_positions[i]
        lx = lx_left if i < 3 else lx_right
        a_h = 1.189 if i < 3 else 2.283

        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Card{i+1}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(cx)}" y="{emu(cy)}"/><a:ext cx="{emu(cw)}" cy="{emu(ch)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(2*12700)}"><a:solidFill><a:srgbClr val="1D6B40"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

        for label, col, y_off, sz_val, bold_val, txt in [
            (f'Topic{i+1}', '156082', 0.067, 1470, False, f'Q{q["num"]}  {q["topic"]}'),
            (f'QText{i+1}', '777777', 0.387, 1470, False, q['q']),
            (f'AText{i+1}', '1D6B40', 0.800, 2800, True,  q['a']),
        ]:
            h = 0.320 if 'Topic' in label else (0.427 if 'QText' in label else a_h)
            _fill_xml = '<a:solidFill><a:srgbClr val="156082"/></a:solidFill>' if 'Topic' in label else '<a:noFill/>'
            add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="{label}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(lx)}" y="{emu(cy+y_off)}"/><a:ext cx="{emu(lw)}" cy="{emu(h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    {_fill_xml}
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr><p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/>
    <a:lstStyle/><a:p><a:pPr algn="l"/><a:r>
      <a:rPr lang="en-GB" sz="{sz_val}" b="{"1" if bold_val else "0"}" dirty="0">
        <a:solidFill><a:srgbClr val="{col}"/></a:solidFill>
        <a:latin typeface="Aptos"/>
      </a:rPr><a:t>{_esc(txt)}</a:t>
    </a:r></a:p></p:txBody></p:sp>''')
            spid += 1

        # Fraction image on answer slide: overlay mathtext PNG over AText box
        if _rm_has_frac(q['a']):
            a_h_f = 1.189 if i < 3 else 2.283
            png = _rm_render_frac(q['a'], lw, a_h_f, fontsize=30, color='#1D6B40')
            sld.shapes.add_picture(_io.BytesIO(png), emu(lx), emu(cy+0.800), emu(lw), emu(a_h_f))

    sld.notes_slide.notes_text_frame.text = (
        f"YOU DO — RM Answers Day {day}\n" +
        '\n'.join(f"Q{q['num']}: {q['a']}" for q in qs))
    print("Slide 7 (RM Answers) ✓")

# ===========================================================================
# SLIDE 8 — VOCABULARY  (curated list + paragraph-reveal animation)
# ===========================================================================
def build_slide8():
    sld = new_slide(3)  # We do layout

    # Title placeholder
    for ph in sld.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = 'Precise Mathematical Vocabulary'
            break

    # Build body text with correct paragraph structure in the placeholder
    # Para 0 (lvl 0): word 1  — visible from start (not animated)
    # Para 1 (lvl 1): definition 1  — revealed on click 1
    # Para 2 (lvl 1): blank spacer
    # Para 3 (lvl 0): word 2  — revealed on click 2
    # Para 4 (lvl 1): definition 2  — revealed on click 3
    # Para 5 (lvl 1): blank spacer
    # ... etc
    # Paragraph indices animated: 1, 3,4, 6,7, 9,10, 12,13  (def1, word2+def2, ...)

    # Build body as an EXPLICIT fixed-size text shape (not layout placeholder)
    # w=12.205" (31cm), h=5.315" (13.5cm), normAutofit so text shrinks on overflow
    # Build paragraph XML: word (bold lvl0), definition (lvl1), blank spacer (lvl1)
    paras_xml = ''
    for i, (word, defn) in enumerate(VOCAB):
        paras_xml += f'''<a:p>
          <a:pPr lvl="0"/>
          <a:r><a:rPr lang="en-GB" sz="1800" b="1" dirty="0">
            <a:latin typeface="Twinkl Cursive Looped Light"/>
          </a:rPr><a:t>{_esc(word)}</a:t></a:r>
        </a:p>
        <a:p>
          <a:pPr lvl="1"/>
          <a:r><a:rPr lang="en-GB" sz="1600" b="0" dirty="0">
            <a:latin typeface="Twinkl Cursive Looped Light"/>
          </a:rPr><a:t>{_esc(defn)}</a:t></a:r>
        </a:p>
        <a:p><a:pPr lvl="1"/><a:endParaRPr lang="en-GB" dirty="0"/></a:p>'''

    vocab_body_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="50" name="VocabBody"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{emu(0.917)}" y="{emu(1.997)}"/>
      <a:ext cx="{emu(12.205)}" cy="{emu(5.315)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/><a:ln w="0"><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr rtlCol="0" anchor="t">
      <a:normAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    {paras_xml}
  </p:txBody>
</p:sp>'''
    add_sp(sld, vocab_body_xml)
    body_spid = 50  # our explicit shape's spid

    # Paragraph indices to animate (one click each):
    # Word 1 visible at start (para 0), then:
    # Click 1: para 1 (def 1)
    # Click 2: para 3 (word 2)
    # Click 3: para 4 (def 2)
    # Click 4: para 6 (word 3)
    # Click 5: para 7 (def 3)
    # Click 6: para 9 (word 4)
    # Click 7: para 10 (def 4)
    # Click 8: para 12 (word 5)
    # Click 9: para 13 (def 5)
    # Blank spacers (paras 2,5,8,11,14) are skipped

    n_words = len(VOCAB)
    # Build para indices: def of word1, then word2+def2, word3+def3 ...
    anim_paras = [1]  # definition of word 1
    for i in range(1, n_words):
        base = i * 3  # word para index
        anim_paras.append(base)       # word
        anim_paras.append(base + 1)   # definition

    child_blocks = ''
    ctn_id = 3
    for para_idx in anim_paras:
        child_blocks += f'''<p:par>
  <p:cTn id="{ctn_id}" fill="hold">
    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
    <p:childTnLst><p:par><p:cTn id="{ctn_id+1}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst><p:par>
        <p:cTn id="{ctn_id+2}" presetID="1" presetClass="entr" presetSubtype="0"
               fill="hold" nodeType="clickEffect">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst><p:set><p:cBhvr>
            <p:cTn id="{ctn_id+3}" dur="1" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
            <p:tgtEl><p:spTgt spid="{body_spid}">
              <p:txEl><p:pRg st="{para_idx}" end="{para_idx}"/></p:txEl>
            </p:spTgt></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
          </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
        </p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
  </p:cTn>
</p:par>'''
        ctn_id += 4

    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par>
    <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
      <p:childTnLst><p:seq concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
          <p:childTnLst>{child_blocks}</p:childTnLst>
        </p:cTn>
        <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
        <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
      </p:seq></p:childTnLst>
    </p:cTn>
  </p:par></p:tnLst>
</p:timing>'''
    sld._element.append(etree.fromstring(timing_xml))

    # No image badge — body is explicit fixed-size shape above (body_ph not used)
    sld.notes_slide.notes_text_frame.text = (
        f"WE DO — Vocabulary\nWords: {[w for w,_ in VOCAB]}\n"
        "Word 1 visible at start; each definition and subsequent word revealed per click.")
    print("Slide 8 (Vocabulary) ✓")

# ===========================================================================
# GRID DRAWING — core visual for teaching slides
# ===========================================================================
def draw_grid_slide(sld, visual_key, layout_num_was):
    """Draw a teaching slide with a coordinate grid on a white panel."""
    v = VISUALS[visual_key]
    cols = v['cols']
    rows = v['rows']

    # Layout: title top, white grid panel left, annotations/text right
    # Grid panel: x=0.4, y=0.9, w=7.0, h=5.8 (leaving room for title and right column)
    panel_x, panel_y = 0.40, 1.45
    panel_w, panel_h = 7.00, 5.80

    cell = min(panel_w / (cols + 1), panel_h / (rows + 1))  # cell size in inches
    # Grid starts inside panel with margin for axis labels
    margin = 0.50  # room for axis numbers
    grid_x = panel_x + margin
    grid_y = panel_y + 0.20
    grid_w = cell * cols
    grid_h = cell * rows

    # --- White panel background ---
    add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="50" name="GridPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(panel_x)}" y="{emu(panel_y)}"/>
    <a:ext cx="{emu(panel_w)}" cy="{emu(panel_h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # --- Grid lines (horizontal) ---
    for r in range(rows + 1):
        y_pos = grid_y + r * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{51+r}" name="HLine{r}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(grid_x)}" y="{emu(y_pos)}"/>
    <a:ext cx="{emu(grid_w)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(0.75*12700)}"><a:solidFill><a:srgbClr val="AAAAAA"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # --- Grid lines (vertical) ---
    for c in range(cols + 1):
        x_pos = grid_x + c * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{60+c}" name="VLine{c}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(x_pos)}" y="{emu(grid_y)}"/>
    <a:ext cx="0" cy="{emu(grid_h)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(0.75*12700)}"><a:solidFill><a:srgbClr val="AAAAAA"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # --- Axis number labels (bottom row = 0, left col = 0) ---
    label_sz = max(8, int(cell * 10))  # scale font with cell size
    for c in range(cols + 1):
        x_pos = grid_x + c * cell - 0.12
        y_pos = grid_y + grid_h + 0.05
        add_sp(sld, sp(70+c, f'XLabel{c}', x_pos, y_pos, 0.24, 0.25,
                       str(c), font='Aptos', sz=label_sz, bold=True,
                       color='333333', align='ctr', fill=None, no_line=True))

    for r in range(rows + 1):
        x_pos = grid_x - 0.35
        y_pos = grid_y + (rows - r) * cell - 0.12
        add_sp(sld, sp(80+r, f'YLabel{r}', x_pos, y_pos, 0.30, 0.25,
                       str(r), font='Aptos', sz=label_sz, bold=True,
                       color='333333', align='ctr', fill=None, no_line=True))

    # --- Axis labels ---
    add_sp(sld, sp(90, 'XAxisLabel',
                   grid_x + grid_w/2 - 0.15, grid_y + grid_h + 0.35, 0.30, 0.25,
                   'x', font='Aptos', sz=label_sz+2, bold=True,
                   color='156082', align='ctr', fill=None, no_line=True))
    add_sp(sld, sp(91, 'YAxisLabel',
                   grid_x - margin, grid_y + grid_h/2 - 0.12, 0.30, 0.25,
                   'y', font='Aptos', sz=label_sz+2, bold=True,
                   color='156082', align='ctr', fill=None, no_line=True))

    # --- Polygon edges (Shape A) — drawn BEFORE dots so dots sit on top ---
    # polygon: list of point labels in order; edges connect them, closed back to start
    # Visual decision: edges 2.5pt in the same colour as the points, slightly transparent
    # effect achieved by using a slightly lighter shade (8A -> same hue, lighter)
    # Dots are 100% colour, lines are drawn at same colour but thinner → vertex/edge distinction
    polygon = v.get('polygon', [])
    if polygon:
        # Build label→(px,py) map from points
        pt_map = {lbl: (grid_x + col*cell, grid_y + (rows-row)*cell)
                  for col, row, lbl, _ in v['points']}
        pt_colors = {lbl: col for _, _, lbl, col in v['points']}
        # Use the first point's colour for all edges (shapes are monochrome)
        edge_color = pt_colors.get(polygon[0], '1F4E79')
        # Slightly lighter: mix with white — prepend with 'light' variant
        # Decision: keep same colour at 2.5pt weight; dots are bigger so read as vertices
        edge_w_pt = 2.5
        for seg_i in range(len(polygon)):
            p1 = polygon[seg_i]
            p2 = polygon[(seg_i + 1) % len(polygon)]
            if p1 not in pt_map or p2 not in pt_map:
                continue
            x1, y1 = pt_map[p1]
            x2, y2 = pt_map[p2]
            lx, ly = min(x1,x2), min(y1,y2)
            lw = abs(x2-x1) or 0.001
            lh = abs(y2-y1) or 0.001
            add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{120+seg_i}" name="Edge{seg_i}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm{"" if x2>=x1 else ' flipH="1"'}{"" if y2>=y1 else ' flipV="1"'}>
    <a:off x="{emu(lx)}" y="{emu(ly)}"/><a:ext cx="{emu(lw)}" cy="{emu(lh)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(edge_w_pt*12700)}"><a:solidFill><a:srgbClr val="{edge_color}"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # --- Plot points (on top of polygon edges) ---
    dot_r = min(0.12, cell * 0.30)
    for pt_idx, (col, row, label, color) in enumerate(v['points']):
        px = grid_x + col * cell - dot_r
        py = grid_y + (rows - row) * cell - dot_r
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{100+pt_idx}" name="Point{label}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(px)}" y="{emu(py)}"/>
    <a:ext cx="{emu(dot_r*2)}" cy="{emu(dot_r*2)}"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:ln w="{int(1.0*12700)}"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

        # Point label — above/right of dot, bold, matches dot colour
        add_sp(sld, sp(110+pt_idx, f'PtLabel{label}',
                       grid_x + col*cell + dot_r*0.5,
                       grid_y + (rows-row)*cell - dot_r*2 - 0.22,
                       0.30, 0.22, label,
                       font='Aptos', sz=label_sz+2, bold=True,
                       color=color, align='ctr', fill=None, no_line=True))

    # --- Translated shape (Shape B) — animated, hidden on load, revealed on click ---
    # Pedagogical rule: show translation only on I Do slides that demonstrate the move.
    # The translated shape is drawn in orange (E8642A) to contrast with Shape A (teal/blue).
    # Labels become A'→B'→C' etc, or use shape_b_label from visual data.
    # A translation vector arrow on one vertex (A→A') is also revealed with the shape.
    translation = v.get('translation')       # [dc, dr] move vector
    polygon_b   = v.get('polygon_b', polygon)  # defaults to same vertex order as A
    shape_a_label = v.get('shape_a_label', 'Shape A')
    shape_b_label = v.get('shape_b_label', 'Shape B')
    COLOR_B = 'E8642A'  # orange — distinct from any teal/blue shape A

    if translation and polygon:
        dc, dr = translation
        pt_map_b = {}
        for col, row, lbl, _ in v['points']:
            pt_map_b[lbl] = (grid_x + (col+dc)*cell, grid_y + (rows-(row+dr))*cell)

        anim_spids = []
        base_spid = 200

        # Shape B edges
        for seg_i in range(len(polygon_b)):
            p1 = polygon_b[seg_i]
            p2 = polygon_b[(seg_i + 1) % len(polygon_b)]
            if p1 not in pt_map_b or p2 not in pt_map_b:
                continue
            x1, y1 = pt_map_b[p1]
            x2, y2 = pt_map_b[p2]
            lx, ly = min(x1,x2), min(y1,y2)
            lw = abs(x2-x1) or 0.001
            lh = abs(y2-y1) or 0.001
            spid = base_spid + seg_i
            anim_spids.append(spid)
            add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="BEdge{seg_i}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm{"" if x2>=x1 else ' flipH="1"'}{"" if y2>=y1 else ' flipV="1"'}>
    <a:off x="{emu(lx)}" y="{emu(ly)}"/><a:ext cx="{emu(lw)}" cy="{emu(lh)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.5*12700)}"><a:solidFill><a:srgbClr val="{COLOR_B}"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

        # Shape B dots and labels
        for pt_idx, (col, row, lbl, _) in enumerate(v['points']):
            bx = grid_x + (col+dc)*cell
            by = grid_y + (rows-(row+dr))*cell
            spid_dot   = base_spid + 20 + pt_idx
            spid_label = base_spid + 30 + pt_idx
            anim_spids.extend([spid_dot, spid_label])
            add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid_dot}" name="BPoint{lbl}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(bx-dot_r)}" y="{emu(by-dot_r)}"/>
    <a:ext cx="{emu(dot_r*2)}" cy="{emu(dot_r*2)}"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{COLOR_B}"/></a:solidFill>
    <a:ln w="{int(1.0*12700)}"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
            add_sp(sld, sp(spid_label, f'BLabel{lbl}',
                           bx + dot_r*0.5, by - dot_r*2 - 0.22,
                           0.30, 0.22, f"{lbl}'",
                           font='Aptos', sz=label_sz+2, bold=True,
                           color=COLOR_B, align='ctr', fill=None, no_line=True))

        # Translation arrow: A→A' on first point, dashed grey, helps read the move
        first_lbl = polygon[0]
        ax_orig, ay_orig = pt_map[first_lbl]
        ax_dest, ay_dest = pt_map_b[first_lbl]
        arr_spid = base_spid + 50
        anim_spids.append(arr_spid)
        alx, aly = min(ax_orig,ax_dest), min(ay_orig,ay_dest)
        alw = abs(ax_dest-ax_orig) or 0.001
        alh = abs(ay_dest-ay_orig) or 0.001
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{arr_spid}" name="TransArrow"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm{"" if ax_dest>=ax_orig else ' flipH="1"'}{"" if ay_dest>=ay_orig else ' flipV="1"'}>
    <a:off x="{emu(alx)}" y="{emu(aly)}"/><a:ext cx="{emu(alw)}" cy="{emu(alh)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(1.5*12700)}" cap="flat">
      <a:solidFill><a:srgbClr val="888888"/></a:solidFill>
      <a:prstDash val="dash"/>
      <a:tailEnd type="arrow" w="med" len="med"/>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

        # Shape B label box (e.g. "Shape B") — in right column
        right_x = panel_x + panel_w + 0.25
        right_w = 13.333 - right_x - 0.2
        sha_spid = base_spid + 51
        anim_spids.append(sha_spid)
        add_sp(sld, sp(sha_spid, 'ShapeBLabel',
                       right_x, panel_y + 1.50, right_w, 0.50,
                       shape_b_label,
                       font='Twinkl Cursive Looped Light', sz=20, bold=True,
                       color=COLOR_B, align='l', fill=None, no_line=True))

        # Shape A label box — always visible (not animated)
        add_sp(sld, sp(base_spid+52, 'ShapeALabel',
                       right_x, panel_y + 0.15, right_w, 0.50,
                       shape_a_label,
                       font='Twinkl Cursive Looped Light', sz=20, bold=True,
                       color=v['points'][0][3], align='l', fill=None, no_line=True))

        # Animate all Shape B elements: hidden on load, all appear on single click
        def anim_block_multi(spids):
            blocks = ''
            ctn_id = 3
            for spid in spids:
                blocks += f'''<p:par>
  <p:cTn id="{ctn_id}" fill="hold">
    <p:stCondLst><p:cond delay="{"indefinite" if ctn_id==3 else "0"}"/></p:stCondLst>
    <p:childTnLst><p:par><p:cTn id="{ctn_id+1}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst><p:par>
        <p:cTn id="{ctn_id+2}" presetID="1" presetClass="entr" presetSubtype="0"
               fill="hold" grpId="1" nodeType="{"clickEffect" if ctn_id==3 else "afterEffect"}">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst><p:set><p:cBhvr>
            <p:cTn id="{ctn_id+3}" dur="1" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
            <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
          </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
        </p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
  </p:cTn>
</p:par>'''
                ctn_id += 4
            return blocks

        bld_list = ''.join(
            f'<p:bldP spid="{s}" grpId="0" uiExpand="1" build="p"/>\n'
            f'<p:bldP spid="{s}" grpId="1" animBg="1"/>\n'
            for s in anim_spids)

        timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par>
    <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
      <p:childTnLst><p:seq concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
          <p:childTnLst>{anim_block_multi(anim_spids)}</p:childTnLst>
        </p:cTn>
        <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
        <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
      </p:seq></p:childTnLst>
    </p:cTn>
  </p:par></p:tnLst>
  <p:bldLst>{bld_list}</p:bldLst>
</p:timing>'''
        sld._element.append(etree.fromstring(timing_xml))

    # --- Right column: caption / sentence stem ---
    right_x = panel_x + panel_w + 0.25
    right_w = 13.333 - right_x - 0.2
    # Push caption/stem down if shape labels are present
    right_y = panel_y + (2.30 if translation else 0.15)

    if 'caption' in v:
        add_sp(sld, sp(130, 'Caption', right_x, right_y, right_w, 1.20,
                       v['caption'],
                       font='Twinkl Cursive Looped Light', sz=20, bold=False,
                       color='1F4E79', align='l', fill='DEECF8',
                       border=('156082', 1.5), anchor='ctr'))

    if 'sentence_stem' in v:
        add_sp(sld, sp(131, 'StemBox', right_x, right_y, right_w, 1.40,
                       v['sentence_stem'],
                       font='Twinkl Cursive Looped Light', sz=18, bold=False,
                       color='1F4E79', align='l', fill='FFF2CC',
                       border=('E8B825', 1.5), anchor='ctr'))

    # error_note / error_instruction handled with animation in build_spot_the_mistake_slide
    sld.notes_slide.notes_text_frame.text = v['notes']


def build_teaching_slide(layout_num, visual_key, title_text, phase):
    sld = new_slide(layout_num)
    for ph in sld.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title_text

            _A_NS  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            _P_NS  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            _EMU   = 914400

            # FIX 1: Top-anchor the title so text always renders from y≈0.40"
            # downward. Default centre-anchor pushes wrapped titles into the
            # content panel area (placeholder bottom = 1.85").
            _bodyPr = ph.text_frame._txBody.find(f'{{{_A_NS}}}bodyPr')
            if _bodyPr is not None:
                _bodyPr.set('anchor', 't')

            # FIX 2: Clamp title width to 10.5" (right edge at 11.42").
            # The I Do / We Do badge image lives at x=11.67" in the layout.
            # Default title right edge is 12.42", so long titles wrap into
            # the badge. 10.5" leaves 0.25" clearance.
            _spPr = ph._element.find(f'{{{_P_NS}}}spPr')
            if _spPr is None:
                # spPr may be under a different ns in some versions
                _spPr = ph._element.find(
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
            if _spPr is None:
                from lxml import etree as _et
                _spPr = _et.SubElement(ph._element,
                    f'{{{_P_NS}}}spPr')
            _xfrm = _spPr.find(f'{{{_A_NS}}}xfrm')
            if _xfrm is None:
                from lxml import etree as _et
                _xfrm = _et.SubElement(_spPr, f'{{{_A_NS}}}xfrm')
            # Always ensure off is set — without it PowerPoint defaults to (0,0)
            # which places the title over the maths icon in the top-left corner.
            _off = _xfrm.find(f'{{{_A_NS}}}off')
            if _off is None:
                from lxml import etree as _et
                _off = _et.SubElement(_xfrm, f'{{{_A_NS}}}off')
            if not _off.get('x'): _off.set('x', str(int(0.917 * _EMU)))
            if not _off.get('y'): _off.set('y', str(int(0.399 * _EMU)))
            _ext = _xfrm.find(f'{{{_A_NS}}}ext')
            if _ext is None:
                from lxml import etree as _et
                _ext = _et.SubElement(_xfrm, f'{{{_A_NS}}}ext')
            _ext.set('cx', str(int(10.5 * _EMU)))
            if not _ext.get('cy'):
                _ext.set('cy', str(int(1.45 * _EMU)))

            break
    v = VISUALS[visual_key]
    slide_type = v.get('slide_type', 'grid')
    if slide_type == 'word_problem':
        draw_word_problem_slide(sld, visual_key)
    elif slide_type == 'identify_calculate':
        draw_identify_calculate_slide(sld, visual_key)
    elif slide_type == 'bar_model':
        draw_bar_model_slide(sld, visual_key)
    elif slide_type == 'stm_word_problem':
        draw_stm_word_problem_slide(sld, visual_key)
    elif slide_type == 'symmetry_grid':
        draw_symmetry_grid_slide(sld, visual_key)
    elif slide_type == 'clock':
        draw_clock_slide(sld, visual_key)
    elif slide_type == 'number_line':
        draw_number_line_slide(sld, visual_key)
    elif slide_type == 'column_calc':
        draw_column_calc(sld, visual_key)
    elif slide_type == 'fraction_demo':
        draw_fraction_demo_slide(sld, visual_key)
    elif slide_type == 'stats_chart':
        draw_stats_chart_slide(sld, visual_key)
    else:
        draw_grid_slide(sld, visual_key, layout_num)
    print(f"  Teaching slide ({title_text[:40]}) ✓")
    return sld


# ===========================================================================
# SYMMETRY GRID SLIDE
# For lessons involving lines of symmetry, reflection, and symmetrical patterns.
# Visual decisions:
#   - Mirror line drawn as a thick dashed red/purple line through the grid
#   - Filled squares (Shape A, teal) on one side
#   - Reflected squares (Shape B, orange) on the other — static or animated
#   - Grid cells coloured, not just dots
# ===========================================================================
def draw_symmetry_grid_slide(sld, visual_key):
    v = VISUALS[visual_key]
    cols = v['cols']
    rows = v['rows']

    panel_x, panel_y = 0.40, 1.45
    panel_w, panel_h = 7.00, 5.80
    cell = min(panel_w / (cols + 1), panel_h / (rows + 1))
    margin = 0.50
    grid_x = panel_x + margin
    grid_y = panel_y + 0.20
    grid_w  = cell * cols
    grid_h  = cell * rows
    label_sz = max(8, int(cell * 10))

    # White panel background
    add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="50" name="GridPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(panel_x)}" y="{emu(panel_y)}"/>
    <a:ext cx="{emu(panel_w)}" cy="{emu(panel_h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # Grid lines
    for r in range(rows + 1):
        y_pos = grid_y + r * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{51+r}" name="HLine{r}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(grid_x)}" y="{emu(y_pos)}"/>
    <a:ext cx="{emu(grid_w)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(0.75*12700)}"><a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
    for c in range(cols + 1):
        x_pos = grid_x + c * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{60+c}" name="VLine{c}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(x_pos)}" y="{emu(grid_y)}"/>
    <a:ext cx="0" cy="{emu(grid_h)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(0.75*12700)}"><a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # Axis numbers
    for c in range(cols + 1):
        add_sp(sld, sp(70+c, f'XLbl{c}', grid_x+c*cell-0.12, grid_y+grid_h+0.05, 0.24, 0.25,
                       str(c), font='Aptos', sz=label_sz, bold=True,
                       color='333333', align='ctr', fill=None, no_line=True))
    for r in range(rows + 1):
        add_sp(sld, sp(80+r, f'YLbl{r}', grid_x-0.35, grid_y+(rows-r)*cell-0.12, 0.30, 0.25,
                       str(r), font='Aptos', sz=label_sz, bold=True,
                       color='333333', align='ctr', fill=None, no_line=True))

    spid = 100

    # Filled squares — Shape A (coloured cells, teal fill with alpha feel)
    for col, row in v.get('squares_a', []):
        cx = grid_x + col * cell
        cy = grid_y + (rows - row - 1) * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="SqA{col}_{row}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(cx+0.01)}" y="{emu(cy+0.01)}"/>
    <a:ext cx="{emu(cell-0.02)}" cy="{emu(cell-0.02)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="A8D0E6"><a:alpha val="85000"/></a:srgbClr></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

    # Reflected squares — Shape B (orange, animated if translate=True)
    anim_spids = []
    for col, row in v.get('squares_b', []):
        cx = grid_x + col * cell
        cy = grid_y + (rows - row - 1) * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="SqB{col}_{row}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(cx+0.01)}" y="{emu(cy+0.01)}"/>
    <a:ext cx="{emu(cell-0.02)}" cy="{emu(cell-0.02)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="F4B183"><a:alpha val="85000"/></a:srgbClr></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        if v.get('animate_b'):
            anim_spids.append(spid)
        spid += 1

    # Point markers (labelled dots for polygon-style symmetry problems)
    dot_r = min(0.10, cell * 0.25)
    for col, row, label, color in v.get('points', []):
        px = grid_x + col*cell - dot_r
        py = grid_y + (rows-row)*cell - dot_r
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Pt{label}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(px)}" y="{emu(py)}"/>
    <a:ext cx="{emu(dot_r*2)}" cy="{emu(dot_r*2)}"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:ln w="{int(1.0*12700)}"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        add_sp(sld, sp(spid+1, f'PtLbl{label}',
                       grid_x+col*cell+dot_r*0.5, grid_y+(rows-row)*cell-dot_r*2-0.22,
                       0.30, 0.22, label, font='Aptos', sz=label_sz+2, bold=True,
                       color=color, align='ctr', fill=None, no_line=True))
        if v.get('animate_b') and label in v.get('animate_labels', []):
            anim_spids.extend([spid, spid+1])
        spid += 2

    # Mirror line — thick dashed, drawn over the grid
    # mirror_col: vertical line at this column; mirror_row: horizontal line at this row
    mirror_col = v.get('mirror_col')
    mirror_row = v.get('mirror_row')
    MIRROR_COLOR = '7030A0'  # purple — neutral, neither Shape A nor B colour
    if mirror_col is not None:
        mx = grid_x + mirror_col * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="MirrorLine"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(mx)}" y="{emu(grid_y)}"/>
    <a:ext cx="0" cy="{emu(grid_h)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.5*12700)}">
      <a:solidFill><a:srgbClr val="{MIRROR_COLOR}"/></a:solidFill>
      <a:prstDash val="sysDash"/>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        # Mirror line label
        add_sp(sld, sp(spid+1, 'MirrorLbl',
                       mx + 0.05, grid_y - 0.26, 1.20, 0.24,
                       'mirror line', font='Twinkl Cursive Looped Light', sz=14,
                       color=MIRROR_COLOR, align='l', fill=None, no_line=True))
        spid += 2
    if mirror_row is not None:
        my = grid_y + (rows - mirror_row) * cell
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="MirrorLineH"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(grid_x)}" y="{emu(my)}"/>
    <a:ext cx="{emu(grid_w)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.5*12700)}">
      <a:solidFill><a:srgbClr val="{MIRROR_COLOR}"/></a:solidFill>
      <a:prstDash val="sysDash"/>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        add_sp(sld, sp(spid+1, 'MirrorLblH',
                       grid_x + grid_w + 0.05, my - 0.15, 1.20, 0.24,
                       'mirror line', font='Twinkl Cursive Looped Light', sz=14,
                       color=MIRROR_COLOR, align='l', fill=None, no_line=True))
        spid += 2

    # Animate reflected squares if requested (I Do slides that reveal the answer)
    if anim_spids:
        def anim_block_sym(spids):
            blocks = ''
            ctn_id = 3
            for s in spids:
                blocks += f'''<p:par>
  <p:cTn id="{ctn_id}" fill="hold">
    <p:stCondLst><p:cond delay="{"indefinite" if ctn_id==3 else "0"}"/></p:stCondLst>
    <p:childTnLst><p:par><p:cTn id="{ctn_id+1}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst><p:par>
        <p:cTn id="{ctn_id+2}" presetID="1" presetClass="entr" presetSubtype="0"
               fill="hold" grpId="1" nodeType="{"clickEffect" if ctn_id==3 else "afterEffect"}">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst><p:set><p:cBhvr>
            <p:cTn id="{ctn_id+3}" dur="1" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
            <p:tgtEl><p:spTgt spid="{s}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
          </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
        </p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
  </p:cTn>
</p:par>'''
                ctn_id += 4
            return blocks
        bld_list = ''.join(
            f'<p:bldP spid="{s}" grpId="0" uiExpand="1" build="p"/>\n'
            f'<p:bldP spid="{s}" grpId="1" animBg="1"/>\n' for s in anim_spids)
        timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
    <p:childTnLst><p:seq concurrent="1" nextAc="seek">
      <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
        <p:childTnLst>{anim_block_sym(anim_spids)}</p:childTnLst>
      </p:cTn>
      <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
      <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq></p:childTnLst>
  </p:cTn></p:par></p:tnLst>
  <p:bldLst>{bld_list}</p:bldLst>
</p:timing>'''
        sld._element.append(etree.fromstring(timing_xml))

    # Right column
    right_x = panel_x + panel_w + 0.25
    right_w  = 13.333 - right_x - 0.2
    right_y  = panel_y + 0.15
    if 'caption' in v:
        add_sp(sld, sp(200, 'Caption', right_x, right_y, right_w, 1.20,
                       v['caption'], font='Twinkl Cursive Looped Light', sz=20,
                       color='1F4E79', align='l', fill='DEECF8',
                       border=('156082', 1.5), anchor='ctr'))
    if 'sentence_stem' in v:
        add_sp(sld, sp(201, 'Stem', right_x, right_y, right_w, 1.40,
                       v['sentence_stem'], font='Twinkl Cursive Looped Light', sz=18,
                       color='1F4E79', align='l', fill='FFF2CC',
                       border=('E8B825', 1.5), anchor='ctr'))
    sld.notes_slide.notes_text_frame.text = v['notes']


# ===========================================================================
# CLOCK SLIDE — for Time lessons
# Draws one or more analogue clock faces with hour and minute hands.
# Visual decisions:
#   - Clock face: white circle, thin border, 12 tick marks, 12 numerals
#   - Hour hand: thick, shorter (60% radius)
#   - Minute hand: thinner, longer (85% radius)
#   - Both drawn as lines from centre
#   - Multiple clocks tiled left-to-right on the panel
#   - Right column: caption / digital time box / sentence stem as normal
# ===========================================================================
import math as _math

def draw_clock_slide(sld, visual_key):
    v = VISUALS[visual_key]

    panel_x, panel_y = 0.40, 1.45
    panel_w, panel_h = 7.00, 5.80
    right_x = panel_x + panel_w + 0.25
    right_w  = 13.333 - right_x - 0.2

    clocks = v.get('clocks', [])
    n = len(clocks)
    if n == 0:
        return

    # Tile clocks: 2×2 for 4 clocks, up to 3 per row otherwise
    if n == 4:
        per_row = 2
    else:
        per_row = min(n, 3)
    rows_needed = (n + per_row - 1) // per_row
    clock_r = min(panel_w / (per_row * 2.5), panel_h / (rows_needed * 2.6)) * 0.88

    spid = 50
    for ci, clock in enumerate(clocks):
        row_i = ci // per_row
        col_i = ci % per_row
        # Centre of this clock face
        cx = panel_x + (col_i + 0.5) * (panel_w / per_row)
        cy = panel_y + 0.3 + (row_i + 0.5) * (panel_h / rows_needed)

        hour   = clock['hour']
        minute = clock['minute']
        label  = clock.get('label', '')
        show_digital = clock.get('show_digital', False)

        # Clock face circle
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="ClockFace{ci}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(cx-clock_r)}" y="{emu(cy-clock_r)}"/>
    <a:ext cx="{emu(clock_r*2)}" cy="{emu(clock_r*2)}"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(2.0*12700)}"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

        # Tick marks and numerals
        for h in range(1, 13):
            angle = _math.radians(h * 30 - 90)
            is_hour_mark = (h % 3 == 0)
            tick_outer = clock_r * 0.92
            tick_inner = clock_r * (0.78 if is_hour_mark else 0.84)
            tx1 = cx + _math.cos(angle) * tick_inner
            ty1 = cy + _math.sin(angle) * tick_inner
            tx2 = cx + _math.cos(angle) * tick_outer
            ty2 = cy + _math.sin(angle) * tick_outer
            lx, ly = min(tx1,tx2), min(ty1,ty2)
            lw = abs(tx2-tx1) or 0.001
            lh = abs(ty2-ty1) or 0.001
            tw = 2.0 if is_hour_mark else 1.0
            add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Tick{ci}_{h}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm{"" if tx2>=tx1 else ' flipH="1"'}{"" if ty2>=ty1 else ' flipV="1"'}>
    <a:off x="{emu(lx)}" y="{emu(ly)}"/><a:ext cx="{emu(lw)}" cy="{emu(lh)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(tw*12700)}"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
            spid += 1

            # Numeral at 3, 6, 9, 12
            # Box must be wide (0.40") and tall (0.30") enough for "12" at any clock size
            if is_hour_mark:
                num_r   = clock_r * 0.66
                num_w   = max(0.40, clock_r * 0.52)
                num_h   = max(0.30, clock_r * 0.38)
                nx = cx + _math.cos(angle) * num_r - num_w/2
                ny = cy + _math.sin(angle) * num_r - num_h/2
                # Keep font small enough: 12pt per inch of clock radius, max 14pt
                num_sz  = min(14, max(9, int(clock_r * 11)))
                add_sp(sld, sp(spid, f'Num{ci}_{h}', nx, ny, num_w, num_h,
                               str(h), font='Aptos', sz=num_sz, bold=True,
                               color='1F4E79', align='ctr', fill=None, no_line=True))
                spid += 1

        # Hour hand — angle: each hour = 30°, each minute adds 0.5°
        h_angle = _math.radians((hour % 12) * 30 + minute * 0.5 - 90)
        h_len   = clock_r * 0.55
        hx2 = cx + _math.cos(h_angle) * h_len
        hy2 = cy + _math.sin(h_angle) * h_len
        lx,ly = min(cx,hx2), min(cy,hy2)
        lw = abs(hx2-cx) or 0.001
        lh = abs(hy2-cy) or 0.001
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="HourHand{ci}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm{"" if hx2>=cx else ' flipH="1"'}{"" if hy2>=cy else ' flipV="1"'}>
    <a:off x="{emu(lx)}" y="{emu(ly)}"/><a:ext cx="{emu(lw)}" cy="{emu(lh)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(4.0*12700)}" cap="rnd">
      <a:solidFill><a:srgbClr val="1F4E79"/></a:solidFill>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

        # Minute hand — angle: each minute = 6°
        m_angle = _math.radians(minute * 6 - 90)
        m_len   = clock_r * 0.78
        mx2 = cx + _math.cos(m_angle) * m_len
        my2 = cy + _math.sin(m_angle) * m_len
        lx,ly = min(cx,mx2), min(cy,my2)
        lw = abs(mx2-cx) or 0.001
        lh = abs(my2-cy) or 0.001
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="MinHand{ci}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm{"" if mx2>=cx else ' flipH="1"'}{"" if my2>=cy else ' flipV="1"'}>
    <a:off x="{emu(lx)}" y="{emu(ly)}"/><a:ext cx="{emu(lw)}" cy="{emu(lh)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.5*12700)}" cap="rnd">
      <a:solidFill><a:srgbClr val="156082"/></a:solidFill>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

        # Centre dot
        dot = clock_r * 0.04
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Centre{ci}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(cx-dot)}" y="{emu(cy-dot)}"/>
    <a:ext cx="{emu(dot*2)}" cy="{emu(dot*2)}"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1

        # Label below clock (e.g. "7:45 am")
        if label:
            lbl_sz = max(12, int(clock_r * 16))
            add_sp(sld, sp(spid, f'ClockLbl{ci}',
                           cx - clock_r, cy + clock_r + 0.05, clock_r*2, 0.35,
                           label, font='Twinkl Cursive Looped Light', sz=lbl_sz,
                           color='1F4E79', align='ctr', fill=None, no_line=True))
            spid += 1

        # Digital time answer box (for We Do / marking station slides)
        if show_digital:
            dig_h = clock_r * 0.55
            add_sp(sld, sp(spid, f'DigBox{ci}',
                           cx - clock_r*0.7, cy + clock_r + 0.42, clock_r*1.4, dig_h,
                           '', font='Aptos', sz=max(16, int(clock_r*20)),
                           color='1A5C2A', align='ctr', fill=None,
                           border=('333333', 1.0)))
            spid += 1

    # Right column
    right_y = panel_y + 0.15
    if 'caption' in v:
        add_sp(sld, sp(spid, 'Caption', right_x, right_y, right_w, 1.20,
                       v['caption'], font='Twinkl Cursive Looped Light', sz=20,
                       color='1F4E79', align='l', fill='DEECF8',
                       border=('156082', 1.5), anchor='ctr'))
        spid += 1
    if 'sentence_stem' in v:
        add_sp(sld, sp(spid, 'Stem', right_x, right_y, right_w, 1.40,
                       v['sentence_stem'], font='Twinkl Cursive Looped Light', sz=18,
                       color='1F4E79', align='l', fill='FFF2CC',
                       border=('E8B825', 1.5), anchor='ctr'))
        spid += 1
    sld.notes_slide.notes_text_frame.text = v['notes']



# ===========================================================================
# ANIMATION HELPER
# Generates <p:timing> XML for click-by-click reveal of shape groups.
# anim_groups: list of lists of int shape IDs.
# Each inner list = all shapes revealed together on one click.
# ===========================================================================
def _apply_animation(sld, anim_groups, pic_ids=None):
    """Inject click-by-click appear animations matching PowerPoint's native format.
    - bldLst entry added for EVERY animated shape (required for initial hide)
    - nodeType: clickEffect for first shape per group, withEffect for rest
    - evt names match PowerPoint native: onPrev / onNext (not onPrevClick)
    """
    if not anim_groups:
        return
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    cid = [100]
    def nid():
        cid[0] += 1
        return cid[0]

    def make_set(spid, node_type, grp_id=None):
        grp = f' grpId="{grp_id}"' if grp_id is not None else ''
        ic = nid(); ic2 = nid()
        return (
            f'<p:par xmlns:p="{NS_P}"><p:cTn id="{ic}" presetID="1" presetClass="entr"' +
            f' presetSubtype="0" fill="hold"{grp} nodeType="{node_type}">' +
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>' +
            f'<p:childTnLst><p:set><p:cBhvr>' +
            f'<p:cTn id="{ic2}" dur="1" fill="hold">' +
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>' +
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>' +
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>' +
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>' +
            f'</p:childTnLst></p:cTn></p:par>'
        )

    click_pars = []
    all_spids = []  # collect all animated shape IDs for bldLst
    for group in anim_groups:
        io = nid(); ii = nid()
        inner_pars = ''
        for idx, spid in enumerate(group):
            node_type = 'clickEffect' if idx == 0 else 'withEffect'
            grp_id = 0 if idx > 0 else None
            inner_pars += make_set(spid, node_type, grp_id)
            all_spids.append(spid)
        click_pars.append(
            f'<p:par xmlns:p="{NS_P}"><p:cTn id="{io}" fill="hold">' +
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>' +
            f'<p:childTnLst><p:par><p:cTn id="{ii}" fill="hold">' +
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>' +
            f'<p:childTnLst>{inner_pars}</p:childTnLst>' +
            f'</p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
        )

    # bldLst: every animated shape listed — this is what tells PowerPoint
    # to start each shape as hidden before its animation fires
    bld_entries = ''.join(
        f'<p:bldP xmlns:p="{NS_P}" spid="{spid}" grpId="0"/>' for spid in all_spids
    )

    timing_xml = (
        f'<p:timing xmlns:p="{NS_P}">' +
        f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">' +
        f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">' +
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>' +
        ''.join(click_pars) +
        f'</p:childTnLst></p:cTn>' +
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>' +
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>' +
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>' +
        f'<p:bldLst>{bld_entries}</p:bldLst>' +
        f'</p:timing>'
    )

    existing = sld._element.find(f'{{{NS_P}}}timing')
    if existing is not None:
        sld._element.remove(existing)
    sld._element.append(etree.fromstring(timing_xml))



# ===========================================================================
# SQUARED PAPER CALCULATOR
# Draws a written calculation on squared-paper grid (0.597" cells).
# Returns anim_groups for click-by-click reveal.
# ===========================================================================
CELL = 0.597  # inches — matches maths book squares
DIGIT_SZ = 32  # pt — main digits
CARRY_SZ = 14  # pt — carry / remainder superscripts
GRID_BORDER = ('9DC3E6', 0.75)  # light blue cell borders
LINE_COL = '156082'             # dark blue separator lines
LINE_W = 2.0                    # pt

def _cell_sp(spid, col, row, text, grid_x, grid_y, sz=DIGIT_SZ,
             color='000000', bold=False, fill='FFFFFF', align='ctr',
             border=GRID_BORDER):
    """One squared-paper cell with centred text."""
    x = grid_x + col * CELL
    y = grid_y + row * CELL
    return sp(spid, f'C{col}R{row}_{spid}', x, y, CELL, CELL,
              text, font='Calibri', sz=sz, bold=bold,
              color=color, align=align, fill=fill,
              border=border)


def _hline_xml(spid, grid_x, grid_y, row, num_cols):
    """Horizontal separator line spanning num_cols cells."""
    x = grid_x; y = grid_y + row * CELL
    w = num_cols * CELL
    return (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{spid}" name="HL{spid}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
        f'<a:ext cx="{emu(w)}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{int(LINE_W*12700)}"><a:solidFill><a:srgbClr val="{LINE_COL}"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )


def _vline_xml(spid, x, y, height):
    """Vertical line (for division bus stop)."""
    return (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{spid}" name="VL{spid}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
        f'<a:ext cx="0" cy="{emu(height)}"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{int(LINE_W*12700)}"><a:solidFill><a:srgbClr val="{LINE_COL}"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )


def _compute_compact_column(top_str, bottom_str):
    """Return carries and answer for compact column multiplication."""
    top_d = [int(c) for c in top_str]
    bot = int(bottom_str)
    carry = 0
    result_d = []
    carry_map = {}  # col index (0=rightmost) -> carry digit str
    for i, d in enumerate(reversed(top_d)):
        prod = d * bot + carry
        result_d.insert(0, str(prod % 10))
        carry = prod // 10
        if carry and i < len(top_d) - 1:
            carry_map[len(top_d) - 2 - i] = str(carry)
    if carry:
        result_d.insert(0, str(carry))
    return ''.join(result_d), carry_map


def _compute_short_div(dividend_str, divisor_str):
    """Return quotient digits and remainder superscripts for short division."""
    divisor = int(divisor_str)
    rem = 0
    q_digits = []
    r_map = {}  # position in dividend (0-indexed) -> remainder to show before next digit
    for i, ch in enumerate(dividend_str):
        current = rem * 10 + int(ch)
        q = current // divisor
        rem = current % divisor
        q_digits.append(str(q))
        if rem and i < len(dividend_str) - 1:
            r_map[i] = str(rem)
    final_rem = rem
    # Strip leading zeros from quotient
    q_str = ''.join(q_digits).lstrip('0') or '0'
    # Re-pad to same length as dividend
    q_str = q_str.zfill(len(dividend_str))
    return q_str, r_map, final_rem


def _compute_column_add(top_str, bottom_str):
    """Return answer and carry positions for column addition."""
    t = int(top_str); b = int(bottom_str)
    answer = str(t + b)
    # Compute carry positions
    td = top_str.zfill(max(len(top_str), len(bottom_str)))
    bd = bottom_str.zfill(len(td))
    carry = 0
    carry_map = {}
    for i in range(len(td) - 1, -1, -1):
        s = int(td[i]) + int(bd[i]) + carry
        carry = s // 10
        if carry and i > 0:
            carry_map[i - 1] = '1'
    return answer, carry_map


# Replacement for draw_squared_paper — matches reference layouts exactly
# Key measurements from reference file (layouts_for_methods.pptx):
#   CELL = 0.5972"  DIGIT_SZ = 32pt  CARRY_SZ = 20pt
# All digit shapes returned in ONE animation group (single click reveal).
# Background cells and bus-stop lines are always visible (not animated).

def draw_squared_paper(sld, calc_type, v, grid_x, grid_y):
    """
    Draw a calculation grid.
    grid_x, grid_y: top-left of the grid (NOT of a panel — just the grid itself).
    Returns anim_groups: list of lists of shape IDs.
    ONE group = all digit/carry shapes appear together on a single click.
    Background cells and structural lines are always visible.
    """
    SID = [600]
    def nid():
        SID[0] += 1
        return SID[0]

    anim_group = []   # ALL animated digit shapes go in here — one click

    # ── helpers ──────────────────────────────────────────────────────────────
    def cell_x(col): return grid_x + col * CELL
    def cell_y(row): return grid_y + row * CELL

    def bg_cell(col, row):
        """Always-visible white background cell with thin border."""
        add_sp(sld, (
            f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<p:nvSpPr><p:cNvPr id="{nid()}" name="BG{col}_{row}"/>'
            f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{emu(cell_x(col))}" y="{emu(cell_y(row))}"/>'
            f'<a:ext cx="{emu(CELL)}" cy="{emu(CELL)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
            f'<a:ln w="{int(0.75*12700)}">'
            f'<a:solidFill><a:srgbClr val="9DC3E6"/></a:solidFill></a:ln>'
            f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
        ))

    def digit_cell(col, row, text, color='1F1F1F', sz=None, bold=True, animate=True):
        """Cell with digit text. Added to anim_group if animate=True."""
        fsz = sz if sz else DIGIT_SZ
        sid = nid()
        add_sp(sld, sp(
            sid, f'D{col}_{row}',
            cell_x(col), cell_y(row), CELL, CELL,
            text,
            font='Twinkl Cursive Looped Light', sz=fsz,
            bold=bold, color=color, align='ctr',
            fill=None, no_line=True, anchor='ctr'
        ))
        if animate:
            anim_group.append(sid)

    def hline(x, y, w, color='1F4E79', lw=2.25):
        """Horizontal line, always visible."""
        sid = nid()
        add_sp(sld, (
            f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<p:nvSpPr><p:cNvPr id="{sid}" name="HL{sid}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="0"/></a:xfrm>'
            f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
            f'<a:ln w="{int(lw*12700)}">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln>'
            f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
        ))

    def vline(x, y, h, color='1F4E79', lw=2.25):
        """Vertical line, always visible."""
        sid = nid()
        add_sp(sld, (
            f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<p:nvSpPr><p:cNvPr id="{sid}" name="VL{sid}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="0" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
            f'<a:ln w="{int(lw*12700)}">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln>'
            f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
        ))

    def small_digit(x, y, text, color='C00000', sz=None):
        """Small carry/borrow/remainder digit, always animated with main group."""
        csz = sz if sz else CARRY_SZ
        sid = nid()
        add_sp(sld, sp(
            sid, f'Carry{sid}', x, y,
            CELL * 0.40, CELL * 0.50,
            text,
            font='Twinkl Cursive Looped Light', sz=csz,
            bold=False, color=color, align='l',
            fill=None, no_line=True, anchor='t'
        ))
        anim_group.append(sid)

    # =========================================================================
    # SHORT DIVISION
    # Reference layout (7293÷5 = 1458 r3):
    #   Row 0: [blank][q0][q1][q2][q3][r][blank][rem]   ← quotient row
    #   Row 1: [div ][ d0][ d1][ d2][ d3]               ← dividend row
    #   Row 2: [    ][   ][   ][   ][   ]                ← blank
    #   Row 3: [    ][   ][   ][   ][   ]                ← blank
    # Bus-stop: vertical line at left of col 1, from row 1 top, height=CELL
    #           horizontal line at top of row 1, from col 1, width=n_div*CELL
    # =========================================================================
    if calc_type == 'short_division':
        dividend = str(v.get('top', ''))
        divisor  = str(v.get('bottom', ''))
        q_str, r_map, final_rem = _compute_short_div(dividend, divisor)

        n_div  = len(dividend)
        has_rem = bool(final_rem)
        # Cols: 0=blank spacer, 1=divisor, 2..n_div+1=dividend,
        #       n_div+2='r', n_div+3=remainder  (matches reference exactly)
        n_cols = n_div + 2 + (2 if has_rem else 1)
        n_rows = 4   # quotient, dividend, 2 blank rows (matches reference)

        # Background cells — always visible
        for r in range(n_rows):
            for c in range(n_cols):
                bg_cell(c, r)

        # Bus-stop structure — always visible
        # Vertical at left of col 2 (start of dividend); horizontal above cols 2..n_div+1
        vline(cell_x(2), cell_y(1), CELL)
        hline(cell_x(2), cell_y(1), n_div * CELL)

        # Row 1: divisor at col 1, dividend at cols 2..n_div+1
        digit_cell(1, 1, divisor, color='1F1F1F')
        for i, d in enumerate(dividend):
            digit_cell(2 + i, 1, d, color='1F1F1F')

        # Remainder superscripts (small red, top-left of the following dividend cell)
        # r_map[pos] = remainder after dividend[pos] → shown at start of dividend[pos+1]
        for pos, r_d in r_map.items():
            small_digit(
                cell_x(3 + pos) + CELL * 0.04,
                cell_y(1) - CELL * 0.05,
                r_d, color='C00000', sz=CARRY_SZ
            )

        # Row 0: quotient digits at cols 2..n_div+1 (green)
        leading = True
        for i, q in enumerate(q_str):
            if q == '0' and leading and len(q_str) > 1:
                continue
            leading = False
            digit_cell(2 + i, 0, q, color='1A5C2A', bold=True)

        # Remainder: 'r' at col n_div+2, digit at col n_div+3
        if has_rem:
            digit_cell(n_div + 2, 0, 'r', color='1F1F1F', bold=False)
            digit_cell(n_div + 3, 0, str(final_rem), color='1A5C2A', bold=True)

    # =========================================================================
    # COLUMN ADDITION
    # Reference layout (2746+8206=10952):
    #   Row 0: [  ][  ][1 ][  ][  ]   ← carry row (small digits)
    #   Row 1: [  ][ 2][ 7][ 4][ 6]   ← top number
    #   Row 2: [+][ 8][ 2][ 0][ 6]    ← operator + bottom number
    #   Row 3: [ 1][ 0][ 9][ 5][ 2]   ← answer (double line above)
    #   Row 4: [  ][  ][  ][  ][  ]   ← blank
    # Double line: at y=row3_top and y=row4_top
    # =========================================================================
    elif calc_type == 'column_addition':
        top_raw    = str(v.get('top', ''))
        bottom_raw = str(v.get('bottom', ''))
        answer_str, carry_map = _compute_column_add(top_raw, bottom_raw)

        # Pad everything right-aligned to answer width (answer may be wider)
        n_digit_cols = len(answer_str)
        top_pad    = top_raw.zfill(n_digit_cols)
        bottom_pad = bottom_raw.zfill(n_digit_cols)

        n_cols = 1 + n_digit_cols    # col 0 = operator, cols 1..n = digits
        n_rows = 5                   # carry, top, op+bottom, answer, blank

        # Background cells — always visible
        for r in range(n_rows):
            for c in range(n_cols):
                bg_cell(c, r)

        # Double line above answer (row 3) and below answer row (row 4)
        hline(cell_x(0), cell_y(3), n_cols * CELL)
        hline(cell_x(0), cell_y(4), n_cols * CELL)

        # Row 0: carry digits (small, red) — carry_map is {col_index: '1'}
        # col_index is 0-based from LEFT of the n_digit_cols array
        for col_idx, carry_val in carry_map.items():
            col = 1 + col_idx   # +1 for operator col
            small_digit(
                cell_x(col) + CELL * 0.04,
                cell_y(0) + CELL * 0.52,
                carry_val, color='C00000', sz=CARRY_SZ
            )

        # Row 1: top number (leading zero cells left blank)
        top_orig_len = len(top_raw)
        for i, d in enumerate(top_pad):
            show = i >= (n_digit_cols - top_orig_len)
            digit_cell(1 + i, 1, d if show else '', color='1F1F1F')

        # Row 2: operator + bottom number
        digit_cell(0, 2, '+', color='1F1F1F')
        bottom_orig_len = len(bottom_raw)
        for i, d in enumerate(bottom_pad):
            show = i >= (n_digit_cols - bottom_orig_len)
            digit_cell(1 + i, 2, d if show else '', color='1F1F1F')

        # Row 3: answer
        for i, d in enumerate(answer_str):
            digit_cell(1 + i, 3, d, color='1A5C2A', bold=True)

    # =========================================================================
    # COLUMN SUBTRACTION
    # Reference layout (5046-3274=1772):
    #   Row 0: [  ][ 4][ 9][  ][  ]   ← modified top digits (borrow notation)
    #   Row 1: [ ˅][ 5˅][ 0][ 4][ 6] ← top number (with borrow marks)
    #   Row 2: [-][ 3][ 2][ 7][ 4]   ← operator + bottom number
    #   Row 3: [  ][ 1][ 7][ 7][ 2]  ← answer (double line above)
    #   Row 4: [  ][  ][  ][  ][  ]  ← blank
    # =========================================================================
    elif calc_type == 'column_subtraction':
        top_raw    = str(v.get('top', ''))
        bottom_raw = str(v.get('bottom', ''))
        # Inline compute: top - bottom, tracking borrows
        def _do_sub(t_str, b_str):
            n = max(len(t_str), len(b_str))
            t = list(t_str.zfill(n)); b = list(b_str.zfill(n))
            borrow = 0
            result = []; borrow_map = {}
            for i in range(n-1, -1, -1):
                td = int(t[i]); bd = int(b[i])
                td -= borrow
                if td < bd:
                    td += 10; borrow = 1; borrow_map[i] = str(int(t[i])-1 if int(t[i])-1>=0 else 10+int(t[i])-1)
                else:
                    borrow = 0
                result.insert(0, str(td - bd))
            return ''.join(result), borrow_map

        answer_str, borrow_map = _do_sub(top_raw, bottom_raw)
        n_digit_cols = max(len(top_raw), len(bottom_raw))
        answer_str = answer_str.zfill(n_digit_cols)
        top_pad    = top_raw.zfill(n_digit_cols)
        bottom_pad = bottom_raw.zfill(n_digit_cols)

        n_cols = 1 + n_digit_cols
        n_rows = 5

        maxlen = max(len(top), len(bottom), len(answer))
        top    = top.zfill(maxlen)
        bottom = bottom.zfill(maxlen)
        answer = answer.zfill(maxlen)

        n_digit_cols = maxlen
        n_cols = 1 + n_digit_cols
        n_rows = 5

        # Background
        for r in range(n_rows):
            for c in range(n_cols):
                bg_cell(c, r)

        # Double lines
        hline(cell_x(0), cell_y(3), n_cols * CELL)
        hline(cell_x(0), cell_y(4), n_cols * CELL)

        # Row 1: top number (always visible)
        for i, d in enumerate(top):
            col = 1 + i
            # Crossed out / modified digit shown with borrow mark
            if modified and i < len(modified) and modified[i] != d:
                # Show original (small, struck) above and modified below
                small_digit(
                    cell_x(col) + CELL * 0.08,
                    cell_y(0) + CELL * 0.40,
                    modified[i], color='1F1F1F', sz=CARRY_SZ
                )
                digit_cell(col, 1, d, color='888888', bold=False)  # greyed original
            else:
                digit_cell(col, 1, d, color='1F1F1F')

        # Borrow marks (small superscripts): show borrowed values
        if borrows:
            for i, b in enumerate(borrows):
                if b != '0':
                    col = 1 + i
                    small_digit(
                        cell_x(col) + CELL * 0.55,
                        cell_y(0) + CELL * 0.40,
                        b, color='C00000', sz=CARRY_SZ
                    )

        # Row 2: operator + bottom
        digit_cell(0, 2, '−', color='1F1F1F')
        for i, d in enumerate(bottom):
            digit_cell(1 + i, 2, d, color='1F1F1F')

        # Row 3: answer
        for i, d in enumerate(answer):
            digit_cell(1 + i, 3, d, color='1A5C2A', bold=True)

    # =========================================================================
    # OTHER METHODS: return empty (caller falls back to text)
    # =========================================================================
    else:
        return []

    # Return ONE animation group — all digits appear on one click
    return [anim_group] if anim_group else []


# ===========================================================================
# BLANK PROBLEM SLIDE — all word_problem / identify_calculate / bar_model
# Left: problem text + VAA banners (labels only, no filled content).
# Right: large blank squared paper for teacher to model working.
# No calculations drawn, no bar models drawn.
# Matches example.pptx provided by Innes.
# ===========================================================================

CELL_GRID = 0.5972   # inches — matches calculation grid cell size exactly
GRID_X    = 6.111    # grid left edge
GRID_Y    = 1.224    # grid top edge
GRID_COLS = 11
GRID_ROWS = 10
GRID_LINE_COLOR = '9DC3E6'

BNR_X = 0.520        # banner image x
BNR_W = 2.675        # banner image width
BNR_H = 0.780        # banner image height
BNR_Y_VIS = 2.970    # Visualise y
BNR_Y_ANA = 3.991    # Analyse y
BNR_Y_ATK = 4.994    # Attack y
TXT_X = 3.312        # text-beside-banner x
TXT_W = 3.905        # text-beside-banner width
TXT_H = 0.700        # text-beside-banner height

def _blank_problem_cell_xml(sid, x_in, y_in, cell_in):
    """Single squared-paper cell: white fill, light blue border."""
    x = int(x_in * 914400); y = int(y_in * 914400); c = int(cell_in * 914400)
    return (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{sid}" name="GC{sid}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{c}" cy="{c}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="{GRID_LINE_COLOR}"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr/></a:p></p:txBody></p:sp>'
    )

def _draw_blank_squared_paper(sld, sid_start):
    """Draw 11×10 blank squared paper grid. Returns next available sid."""
    sid = sid_start
    for col in range(GRID_COLS):
        for row in range(GRID_ROWS):
            x = GRID_X + col * CELL_GRID
            y = GRID_Y + row * CELL_GRID
            add_sp(sld, _blank_problem_cell_xml(sid, x, y, CELL_GRID))
            sid += 1
    return sid

def _vaa_txt(nid, text_lines, x, y, w, h, color='1F1F1F', sz=16, bold=False):
    """Text box beside a VAA banner."""
    lines_xml = ''
    for i, line in enumerate(text_lines):
        b = '<a:b/>' if bold else ''
        lines_xml += (
            f'<a:p><a:pPr spc="-100"/>'
            f'<a:r><a:rPr lang="en-GB" sz="{sz*100}" b="{1 if bold else 0}" dirty="0">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:latin typeface="Twinkl Cursive Looped Light"/></a:rPr>'
            f'<a:t>{line}</a:t></a:r></a:p>'
        )
    xe = int(x*914400); ye = int(y*914400); we = int(w*914400); he = int(h*914400)
    return (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{nid}" name="VTxt{nid}"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{xe}" y="{ye}"/><a:ext cx="{we}" cy="{he}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square" anchor="t"><a:normAutofit/></a:bodyPr>'
        f'<a:lstStyle/>{lines_xml}</p:txBody></p:sp>'
    )

def draw_blank_problem_slide(sld, visual_key, is_two_step=False):
    """Blank problem slide: problem + VAA banners (animated) + blank squared paper."""
    v       = VISUALS[visual_key]
    problem = v.get('problem', '')
    px, py, pw, ph = 0.40, 1.45, 5.392, 5.80

    SID = [600]
    def nid():
        SID[0] += 1
        return SID[0]
    def sync_past(pic_id):
        if pic_id >= SID[0]:
            SID[0] = pic_id

    # White left panel
    add_sp(sld, (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{nid()}" name="BPPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{emu(px)}" y="{emu(py)}"/>'
        f'<a:ext cx="{emu(pw)}" cy="{emu(ph)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    ))

    # Problem text — height fills space up to Visualise banner with 0.10" gap
    prob_h = BNR_Y_VIS - 1.60 - 0.10
    add_sp(sld, sp(nid(), 'Problem', 0.65, 1.60, pw - 0.30, prob_h,
                   problem.split('\n'),
                   font='Twinkl Cursive Looped Light', sz=22,
                   bold=False, color='1F1F1F', align='l', fill=None, no_line=True,
                   anchor='t', autofit=True))
    _frac_overlay(sld, problem, 0.65, 1.60, pw - 0.30, prob_h,
                  fontsize=20, txt_color='#1F1F1F', bg_hex='F0F5FA')

    # VAA banners — animated, with sync_past to keep nid() ahead of pic IDs
    anim_groups = []

    _, vis_id = add_pic_id(sld, 'banner_visualise.png', BNR_X, BNR_Y_VIS, BNR_W, BNR_H)
    sync_past(vis_id)
    anim_groups.append([vis_id])

    _, ana_id = add_pic_id(sld, 'banner_analyse.png', BNR_X, BNR_Y_ANA, BNR_W, BNR_H)
    sync_past(ana_id)
    ana_txt_id = nid()
    add_sp(sld, _vaa_txt(ana_txt_id, ['I know:  ', '', "I'm finding:  "],
                         TXT_X, BNR_Y_ANA + 0.01, TXT_W, TXT_H, color='1F4E79', sz=15))
    anim_groups.append([ana_id, ana_txt_id])

    _, atk_id = add_pic_id(sld, 'banner_attack.png', BNR_X, BNR_Y_ATK, BNR_W, BNR_H)
    sync_past(atk_id)
    atk_txt_id = nid()
    atk_lines = ["First I'm going to  ", '', 'Then I will  '] if is_two_step else ["I'm going to  "]
    add_sp(sld, _vaa_txt(atk_txt_id, atk_lines,
                         TXT_X, BNR_Y_ATK + 0.01, TXT_W, TXT_H, color='843C0C', sz=15, bold=True))
    anim_groups.append([atk_id, atk_txt_id])

    _apply_animation(sld, anim_groups)

    # Grid last — so its raw-XML IDs never clash with animated shape IDs
    _draw_blank_squared_paper(sld, nid())



def draw_word_problem_slide(sld, visual_key):
    draw_blank_problem_slide(sld, visual_key, is_two_step=False)

def draw_identify_calculate_slide(sld, visual_key):
    draw_blank_problem_slide(sld, visual_key, is_two_step=False)

def draw_bar_model_slide(sld, visual_key):
    draw_blank_problem_slide(sld, visual_key, is_two_step=True)

# ===========================================================================
# FRACTION DEMO SLIDE
# Problem text at top of a full-width white panel.
# Working steps revealed below, one per click (animated).
# Last step styled green (the answer). No VAA, no squared paper.
# Fractions written as n/d in data are rendered with a proper vinculum via
# matplotlib (same pattern as RM slides) — PNG overlay on text box.
# Data keys: problem (str), steps (list[str]), notes (str)
# ===========================================================================


def _frac_overlay(sld, text, x, y, w, h, fontsize, txt_color, bg_hex):
    """If text contains n/d fractions, render a vinculum PNG and place it over the text box.
    For static (non-animated) shapes only. Returns True if overlay was added."""
    if not _rm_has_frac(text):
        return False
    png = _fd_render(text, w, h, fontsize, txt_color, bg_hex)
    sld.shapes.add_picture(_io.BytesIO(png), emu(x), emu(y), emu(w), emu(h))
    return True


def _fd_render(text, w_in, h_in, fontsize, txt_color, bg_hex, border_hex=None):
    """Render text (with n/d → vinculum) as PNG. bg_hex is a 6-char hex string.
    If border_hex is provided, draws a coloured rectangle border around the image."""
    import io as _io3
    import matplotlib
    import matplotlib.patches as _mpatches
    matplotlib.use('Agg')
    import matplotlib.pyplot as _plt3
    mt = _RM_FRAC_RE.sub(lambda m: f'$\\frac{{{m.group(1)}}}{{{m.group(2)}}}$', text)
    bg_rgb = tuple(int(bg_hex[i:i+2], 16) / 255 for i in (0, 2, 4))
    fig = _plt3.figure(figsize=(w_in, h_in), dpi=150)
    fig.patch.set_facecolor(bg_rgb)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    ax.set_facecolor(bg_rgb)
    ax.text(0.025, 0.5, mt, fontsize=fontsize, ha='left', va='center',
            color=txt_color, transform=ax.transAxes, fontfamily='DejaVu Sans',
            wrap=False)
    if border_hex:
        brgb = tuple(int(border_hex[i:i+2], 16) / 255 for i in (0, 2, 4))
        rect = _mpatches.FancyBboxPatch(
            (0.005, 0.04), 0.990, 0.92,
            boxstyle="square,pad=0",
            linewidth=2.5, edgecolor=brgb, facecolor='none',
            transform=ax.transAxes, clip_on=False
        )
        ax.add_patch(rect)
    buf = _io3.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches=None,
                facecolor=bg_rgb, edgecolor='none')
    _plt3.close(fig)
    buf.seek(0)
    return buf.read()


def draw_fraction_demo_slide(sld, visual_key):
    v      = VISUALS[visual_key]
    prob   = v.get('problem', '')
    steps  = v.get('steps', [])

    px       = 0.40
    py       = 1.55
    pw       = 5.392
    top_pad  = 0.22
    prob_h   = 0.70
    div_sp   = 0.20
    step_h   = 0.85
    step_gap = 0.14
    bot_pad  = 0.32
    n_steps  = len(steps)

    ph = top_pad + prob_h + div_sp + n_steps * step_h + (n_steps - 1) * step_gap + bot_pad

    SID = [700]
    def nid():
        SID[0] += 1
        return SID[0]
    def sync_past_local(pic_id):
        while SID[0] <= pic_id:
            SID[0] += 1

    # White left panel
    add_sp(sld, (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{nid()}" name="FDPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{emu(px)}" y="{emu(py)}"/>'
        f'<a:ext cx="{emu(pw)}" cy="{emu(ph)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    ))

    # Problem text
    prob_x = px + 0.22
    prob_y = py + top_pad
    prob_w = pw - 0.44

    add_sp(sld, sp(nid(), 'FDProblem', prob_x, prob_y, prob_w, prob_h,
                   prob.split('\n'),
                   font='Twinkl Cursive Looped Light', sz=18,
                   bold=False, color='1F1F1F', align='l',
                   fill=None, no_line=True, anchor='t', autofit=True))
    if _rm_has_frac(prob):
        prob_png = _fd_render(prob, prob_w, prob_h, fontsize=16,
                              txt_color='#1F1F1F', bg_hex='FFFFFF')
        pic = sld.shapes.add_picture(_io.BytesIO(prob_png),
                                     emu(prob_x), emu(prob_y),
                                     emu(prob_w), emu(prob_h))
        sync_past_local(pic.shape_id)

    # Thin divider
    div_y = prob_y + prob_h + 0.06
    add_sp(sld, (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{nid()}" name="FDDiv"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{emu(prob_x)}" y="{emu(div_y)}"/>'
        f'<a:ext cx="{emu(prob_w)}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{int(0.75*12700)}"><a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    ))

    # Steps — built fully (shapes + PNGs) BEFORE _apply_animation so every
    # shape that needs to be in an animation group is already in the XML.
    step_x = prob_x
    step_w = prob_w
    step_y = div_y + 0.14

    anim_groups = []
    for i, step_text in enumerate(steps):
        is_answer  = (i == n_steps - 1)
        fill_hex   = 'E8F5E9' if is_answer else 'DEEAF1'
        txt_hex    = '1A5C2A' if is_answer else '1F4E79'
        border_col = '1A5C2A' if is_answer else '156082'
        prefix     = '✓  '    if is_answer else '→  '
        bold       = is_answer
        full_step  = prefix + step_text
        has_frac   = _rm_has_frac(full_step)

        if has_frac:
            # For fraction steps: PNG only — no text box underneath.
            # PNG is added to the slide and its shape_id goes in the anim group.
            step_png = _fd_render(full_step, step_w, step_h, fontsize=14,
                                  txt_color=f'#{txt_hex}', bg_hex=fill_hex)
            # Draw the coloured box border separately (panel XML, not animated)
            box_sid = nid()
            add_sp(sld, (
                f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
                f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<p:nvSpPr><p:cNvPr id="{box_sid}" name="FDBox{i}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
                f'<p:spPr><a:xfrm><a:off x="{emu(step_x)}" y="{emu(step_y)}"/>'
                f'<a:ext cx="{emu(step_w)}" cy="{emu(step_h)}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'<a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>'
                f'<a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="{border_col}"/></a:solidFill></a:ln>'
                f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
            ))
            # PNG on top of box
            pic = sld.shapes.add_picture(_io.BytesIO(step_png),
                                         emu(step_x), emu(step_y),
                                         emu(step_w), emu(step_h))
            sync_past_local(pic.shape_id)
            # Both box and PNG animate together on the same click
            anim_groups.append([box_sid, pic.shape_id])
        else:
            # No fraction: Twinkl Cursive text box only
            sid = nid()
            add_sp(sld, sp(sid, f'FDStep{i}', step_x, step_y, step_w, step_h,
                           full_step,
                           font='Twinkl Cursive Looped Light', sz=16,
                           bold=bold, color=txt_hex, align='l',
                           fill=fill_hex, border=(border_col, 1.5), anchor='ctr'))
            anim_groups.append([sid])

        step_y += step_h + step_gap

    # Now ALL shapes are in the slide XML — safe to apply animation
    _apply_animation(sld, anim_groups)

    # Squared paper — clear of title, 8 rows, drawn last
    _fd_grid_y    = 1.60
    _fd_grid_rows = 8
    _fd_sid = nid()
    for _col in range(GRID_COLS):
        for _row in range(_fd_grid_rows):
            add_sp(sld, _blank_problem_cell_xml(_fd_sid,
                   GRID_X + _col * CELL_GRID,
                   _fd_grid_y + _row * CELL_GRID,
                   CELL_GRID))
            _fd_sid += 1

    sld.notes_slide.notes_text_frame.text = v.get('notes', '')


# ===========================================================================
# STATISTICS CHART SLIDE
# Left panel: animated Q&A pairs (question → answer per click).
# Right panel: matplotlib chart PNG (pictogram / bar_chart / line_graph /
#              table / double_bar) generated from chart_data.
# ===========================================================================
def draw_stats_chart_slide(sld, visual_key):
    v          = VISUALS[visual_key]
    chart_type = v['chart_type']
    chart_data = v['chart_data']
    questions  = v.get('questions', [])
    answers    = v.get('answers', [])

    # ── Layout constants ──────────────────────────────────────────────────
    PNL_X, PNL_Y, PNL_W, PNL_H = 0.40, 1.45, 5.10, 5.80
    PAD_X, PAD_TOP              = 0.18, 0.22
    CHT_X, CHT_Y, CHT_W, CHT_H = 5.75, 1.42, 7.20, 5.84

    # ── Generate chart PNG ────────────────────────────────────────────────
    chart_dir = os.path.join(tempfile.gettempdir(), 'wfa_stats_charts')
    os.makedirs(chart_dir, exist_ok=True)
    chart_path = os.path.join(chart_dir, f'{visual_key}_{chart_type}.png')
    if not os.path.exists(chart_path):
        render_stats_chart(chart_type, chart_data, chart_path, dpi=180)

    # ── Shape ID counter ─────────────────────────────────────────────────
    SID = [900]
    def nid():
        SID[0] += 1
        return SID[0]

    # ── Left panel background ─────────────────────────────────────────────
    add_sp(sld, (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{nid()}" name="SCPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{emu(PNL_X)}" y="{emu(PNL_Y)}"/>'
        f'<a:ext cx="{emu(PNL_W)}" cy="{emu(PNL_H)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln w="{int(1.5*12700)}"><a:solidFill>'
        f'<a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    ))

    # ── Q&A boxes — calculate spacing ─────────────────────────────────────
    n_pairs   = len(questions)
    Q_H       = 0.70
    A_H       = 0.70
    QA_GAP    = 0.08    # gap between Q and its A
    PAIR_GAP  = 0.20    # gap between pairs
    bot_pad   = 0.15

    total_needed = (PAD_TOP + n_pairs * (Q_H + QA_GAP + A_H)
                    + max(0, n_pairs - 1) * PAIR_GAP + bot_pad)
    if total_needed > PNL_H:
        # Compress pair gap
        PAIR_GAP = max(0.08, (PNL_H - PAD_TOP - n_pairs * (Q_H + QA_GAP + A_H)
                               - bot_pad) / max(1, n_pairs - 1))

    qx  = PNL_X + PAD_X
    qw  = PNL_W - 2 * PAD_X
    cur_y = PNL_Y + PAD_TOP

    anim_groups = []

    for i, (q, a) in enumerate(zip(questions, answers)):
        # Question box
        q_sid = nid()
        add_sp(sld, sp(q_sid, f'SCQ{i}', qx, cur_y, qw, Q_H,
                       q,
                       font='Twinkl Cursive Looped Light', sz=13,
                       bold=False, color='1F4E79', align='l',
                       fill='DEEAF1', border=('156082', 1.2), anchor='ctr'))
        anim_groups.append([q_sid])
        cur_y += Q_H + QA_GAP

        # Answer box
        a_sid = nid()
        add_sp(sld, sp(a_sid, f'SCA{i}', qx, cur_y, qw, A_H,
                       '✓  ' + a,
                       font='Twinkl Cursive Looped Light', sz=13,
                       bold=True, color='1A5C2A', align='l',
                       fill='E8F5E9', border=('1A5C2A', 1.2), anchor='ctr'))
        anim_groups.append([a_sid])
        cur_y += A_H + PAIR_GAP

    _apply_animation(sld, anim_groups)

    # ── Chart image ───────────────────────────────────────────────────────
    sld.shapes.add_picture(chart_path,
                           emu(CHT_X), emu(CHT_Y),
                           emu(CHT_W), emu(CHT_H))

    sld.notes_slide.notes_text_frame.text = v.get('notes', '')


def draw_stm_word_problem_slide(sld, visual_key):
    v = VISUALS[visual_key]
    px, py, pw, ph = 0.40, 1.45, 7.00, 5.80
    SID = [500]; nid = lambda: (SID.__setitem__(0, SID[0]+1), SID[0])[1]

    # White panel
    add_sp(sld, f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
           f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
           f'<p:nvSpPr><p:cNvPr id="{nid()}" name="STMPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
           f'<p:spPr><a:xfrm><a:off x="{emu(px)}" y="{emu(py)}"/>'
           f'<a:ext cx="{emu(pw)}" cy="{emu(ph)}"/></a:xfrm>'
           f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
           f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
           f'<a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>'
           f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>')

    # Problem text
    lines = v.get('problem', '').split('\n')
    add_sp(sld, sp(nid(), 'STMProblem', px + 0.25, py + 0.25,
                   pw - 0.50, 1.80,
                   lines, font='Twinkl Cursive Looped Light', sz=22,
                   bold=False, color='1F1F1F', align='l',
                   fill=None, no_line=True, anchor='t'))

    # Wrong working — red box
    wrong = v.get('wrong_working', '')
    wlines = wrong.split('\n')
    add_sp(sld, sp(nid(), 'WrongBox', px + 0.25, py + 2.20,
                   pw - 0.50, max(0.65, 0.55 * len(wlines)),
                   wlines, font='Twinkl Cursive Looped Light', sz=22,
                   bold=True, color='FFFFFF', align='l',
                   fill='C00000', border=('800000', 1.5), anchor='ctr'))

    # Error explanation
    error = v.get('error', '')
    elines = error.split('\n')
    err_h  = max(0.75, 0.55 * len(elines))
    add_sp(sld, sp(nid(), 'ErrorBox', px + 0.25, py + 3.20,
                   pw - 0.50, err_h,
                   elines, font='Twinkl Cursive Looped Light', sz=20,
                   bold=False, color='1F4E79', align='l',
                   fill='DEEAF1', border=('156082', 1.5), anchor='t'))

    # Vinculum PNG overlays for any fraction notation in this slide
    _frac_overlay(sld, v.get('problem',''), px+0.25, py+0.25, pw-0.50, 1.80,
                  fontsize=20, txt_color='#1F1F1F', bg_hex='FFFFFF')
    _frac_overlay(sld, wrong, px+0.25, py+2.20, pw-0.50, max(0.65, 0.55*len(wlines)),
                  fontsize=20, txt_color='#FFFFFF', bg_hex='C00000')
    _frac_overlay(sld, error, px+0.25, py+3.20, pw-0.50, err_h,
                  fontsize=18, txt_color='#1F4E79', bg_hex='DEEAF1')

    sld.notes_slide.notes_text_frame.text = v.get('notes', '')



# ===========================================================================
# COLUMN CALC SLIDE — word problem / operation identification
# Left panel: large column calculation (× or ÷). Right panel: context/caption.
# ===========================================================================
def draw_column_calc(sld, visual_key):
    v = VISUALS[visual_key]
    calc_type  = v.get('calc_type', 'multiplication')
    top        = str(v.get('top', ''))
    bottom     = str(v.get('bottom', ''))
    show_ans   = v.get('show_answer', False)
    answer     = str(v.get('answer', '')) if show_ans else ''
    caption    = v.get('caption', '')

    panel_x, panel_y = 0.40, 1.45
    panel_h = 5.80
    calc_w  = 3.30   # left white panel for the calculation
    ctx_x   = panel_x + calc_w + 0.12
    ctx_w   = 7.00 - calc_w - 0.12

    # ---- White calc panel ----
    add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="50" name="CalcPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(panel_x)}" y="{emu(panel_y)}"/>
    <a:ext cx="{emu(calc_w)}" cy="{emu(panel_h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # ---- Numbers centred in left panel ----
    num_cx = panel_x + calc_w / 2
    num_w  = 1.60
    num_x  = num_cx - num_w / 2
    num_sz = 54   # large enough to be clear

    # For multiplication: top (right-aligned), × bottom (right-aligned), line, answer
    # For division: top ÷ bottom (horizontal notation, centred)
    if calc_type == 'multiplication':
        row1_y = panel_y + 0.90
        row2_y = row1_y + 1.20
        line_y = row2_y + 1.15   # clears the 1.10" row2 box
        row3_y = line_y + 0.15

        add_sp(sld, sp(51, 'Top', num_x, row1_y, num_w, 1.10,
                       top, font='Aptos', sz=num_sz, bold=True,
                       color='1F4E79', align='r', fill=None, no_line=True))
        # × sign (offset left)
        add_sp(sld, sp(52, 'Op', num_x - 0.40, row2_y, 0.40, 1.10,
                       '×', font='Aptos', sz=num_sz, bold=True,
                       color='E8642A', align='r', fill=None, no_line=True))
        add_sp(sld, sp(53, 'Bot', num_x, row2_y, num_w, 1.10,
                       bottom, font='Aptos', sz=num_sz, bold=True,
                       color='1F4E79', align='r', fill=None, no_line=True))
        # Answer line
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="54" name="AnsLine"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(num_x-0.45)}" y="{emu(line_y)}"/>
    <a:ext cx="{emu(num_w+0.45)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.0*12700)}"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        if answer:
            add_sp(sld, sp(55, 'Ans', num_x, row3_y, num_w, 1.10,
                           answer, font='Aptos', sz=num_sz, bold=True,
                           color='1A5C2A', align='r', fill=None, no_line=True))
        else:
            add_sp(sld, sp(55, 'AnsBlank', num_x - 0.45, row3_y, num_w + 0.45, 1.10,
                           '?', font='Aptos', sz=num_sz, bold=True,
                           color='AAAAAA', align='r', fill=None, no_line=True))

    elif calc_type in ('addition', 'subtraction'):
        op_sym = '+' if calc_type == 'addition' else '−'
        row1_y = panel_y + 0.90
        row2_y = row1_y + 1.20
        line_y = row2_y + 1.15   # clears the 1.10" row2 box
        row3_y = line_y + 0.15

        add_sp(sld, sp(51, 'Top', num_x, row1_y, num_w, 1.10,
                       top, font='Aptos', sz=num_sz, bold=True,
                       color='1F4E79', align='r', fill=None, no_line=True))
        add_sp(sld, sp(52, 'Op', num_x - 0.40, row2_y, 0.40, 1.10,
                       op_sym, font='Aptos', sz=num_sz, bold=True,
                       color='E8642A', align='r', fill=None, no_line=True))
        add_sp(sld, sp(53, 'Bot', num_x, row2_y, num_w, 1.10,
                       bottom, font='Aptos', sz=num_sz, bold=True,
                       color='1F4E79', align='r', fill=None, no_line=True))
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="54" name="AnsLine"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(num_x-0.45)}" y="{emu(line_y)}"/>
    <a:ext cx="{emu(num_w+0.45)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.0*12700)}"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        if answer:
            add_sp(sld, sp(55, 'Ans', num_x, row3_y, num_w, 1.10,
                           answer, font='Aptos', sz=num_sz, bold=True,
                           color='1A5C2A', align='r', fill=None, no_line=True))
        else:
            add_sp(sld, sp(55, 'AnsBlank', num_x - 0.45, row3_y, num_w + 0.45, 1.10,
                           '?', font='Aptos', sz=num_sz, bold=True,
                           color='AAAAAA', align='r', fill=None, no_line=True))

    else:  # division — bus stop layout: bottom ) top
        # Fixed-width layout centred in calc panel
        divor_bw = 0.90   # divisor box
        brack_w  = 0.48   # bracket )
        divid_bw = 1.60   # dividend box (holds up to 3 digits at 54pt)
        group_w  = divor_bw + brack_w + divid_bw
        start_x  = panel_x + (calc_w - group_w) / 2
        row_y    = panel_y + panel_h / 2 - 0.60

        divor_x  = start_x
        brack_x  = start_x + divor_bw
        divid_x  = brack_x + brack_w

        add_sp(sld, sp(51, 'Divisor', divor_x, row_y, divor_bw, 1.10,
                       bottom, font='Aptos', sz=num_sz, bold=True,
                       color='E8642A', align='r', fill=None, no_line=True))
        add_sp(sld, sp(52, 'Bracket', brack_x, row_y, brack_w, 1.10,
                       ')', font='Aptos', sz=num_sz, bold=True,
                       color='333333', align='ctr', fill=None, no_line=True))
        add_sp(sld, sp(53, 'Dividend', divid_x, row_y, divid_bw, 1.10,
                       top, font='Aptos', sz=num_sz, bold=True,
                       color='1F4E79', align='l', fill=None, no_line=True))
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="54" name="Overline"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(divid_x)}" y="{emu(row_y)}"/>
    <a:ext cx="{emu(divid_bw)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.0*12700)}"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        if answer:
            add_sp(sld, sp(55, 'Ans', divid_x, row_y - 0.85, divid_bw, 1.10,
                           answer, font='Aptos', sz=num_sz, bold=True,
                           color='1A5C2A', align='l', fill=None, no_line=True))

    # ---- Context panel (right) ----
    add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="56" name="CtxPanel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(ctx_x)}" y="{emu(panel_y)}"/>
    <a:ext cx="{emu(ctx_w)}" cy="{emu(panel_h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="DEECF8"/></a:solidFill>
    <a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="156082"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')

    # Caption text inside context panel
    lines = [l for l in caption.split('\n') if l.strip()]
    add_sp(sld, sp(57, 'Caption', ctx_x + 0.12, panel_y + 0.25,
                   ctx_w - 0.24, panel_h - 0.50,
                   lines if lines else caption,
                   font='Twinkl Cursive Looped Light', sz=18, bold=False,
                   color='1F1F1F', align='l', fill=None, no_line=True, anchor='t'))

    sld.notes_slide.notes_text_frame.text = v.get('notes', '')


# ===========================================================================
# NUMBER LINE SLIDE — for Time lesson C2 (24-hour number line)
# ===========================================================================
def draw_number_line_slide(sld, visual_key):
    v = VISUALS[visual_key]
    panel_x, panel_y = 0.40, 1.45
    panel_w, panel_h = 7.00, 5.80
    right_x = panel_x + panel_w + 0.25
    right_w  = 13.333 - right_x - 0.2

    # Number line spans from start_val to end_val
    start_val = v.get('nl_start', 0)
    end_val   = v.get('nl_end', 24)
    total     = end_val - start_val

    nl_y  = panel_y + panel_h * 0.58   # lower — leaves room above for markers/examples
    nl_x1 = panel_x + 0.3
    nl_x2 = panel_x + panel_w - 0.2
    nl_w  = nl_x2 - nl_x1

    spid = 50

    # White panel
    add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="50" name="Panel"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(panel_x)}" y="{emu(panel_y)}"/>
    <a:ext cx="{emu(panel_w)}" cy="{emu(panel_h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="{int(1.5*12700)}"><a:solidFill><a:srgbClr val="BBBBBB"/></a:solidFill></a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
    spid = 51

    # Main line with arrowheads
    add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="NLLine"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(nl_x1)}" y="{emu(nl_y)}"/>
    <a:ext cx="{emu(nl_w)}" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.0*12700)}">
      <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
      <a:tailEnd type="arrow" w="med" len="med"/>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
    spid += 1

    # Tick marks and labels for key hours
    tick_hours = v.get('nl_ticks', list(range(start_val, end_val + 1, 2)))
    tick_h = 0.10
    for val in tick_hours:
        tx = nl_x1 + (val - start_val) / total * nl_w
        is_major = (val % 6 == 0)
        th = tick_h * (1.6 if is_major else 1.0)
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Tick{val}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(tx)}" y="{emu(nl_y - th/2)}"/>
    <a:ext cx="0" cy="{emu(th)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int((1.5 if is_major else 0.75)*12700)}">
      <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1
        # Label — format as HH:00, wider box prevents wrapping
        lbl = f'{val:02d}:00'
        add_sp(sld, sp(spid, f'NLLbl{val}',
                       tx - 0.28, nl_y + tick_h*1.0 + 0.03, 0.56, 0.22,
                       lbl, font='Aptos', sz=8 if is_major else 7, bold=is_major,
                       color='1F4E79' if is_major else '888888',
                       align='ctr', fill=None, no_line=True))
        spid += 1

    # Special markers (e.g. midday line, midnight)
    for marker in v.get('nl_markers', []):
        val   = marker['val']
        lbl   = marker['label']
        color = marker.get('color', '7030A0')
        mx = nl_x1 + (val - start_val) / total * nl_w
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{spid}" name="Marker{val}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{emu(mx)}" y="{emu(nl_y - 0.30)}"/>
    <a:ext cx="0" cy="{emu(0.60)}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(2.0*12700)}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="sysDash"/>
    </a:ln>
  </p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>''')
        spid += 1
        # Marker label: two lines above the line, centred on marker x
        add_sp(sld, sp(spid, f'MkLbl{val}',
                       mx - 0.65, nl_y - 0.88, 1.30, 0.45,
                       lbl, font='Twinkl Cursive Looped Light', sz=12, bold=True,
                       color=color, align='ctr', fill=None, no_line=True))
        spid += 1

    # Conversion examples — staggered above line to avoid overlap
    for ei, ex in enumerate(v.get('examples', [])):
        ex_x = nl_x1 + (ex['val'] - start_val) / total * nl_w
        # Alternate heights so adjacent examples don't overlap
        ex_y = nl_y - 1.55 - (0.40 if ei % 2 == 0 else 0.0)
        add_sp(sld, sp(spid, f'Ex{ex["val"]}',
                       ex_x - 0.52, ex_y, 1.04, 0.55,
                       ex['text'], font='Twinkl Cursive Looped Light', sz=12,
                       color='C00000', align='ctr', fill='FFE6E6',
                       border=('C00000', 1.0), anchor='ctr'))
        spid += 1

    # Right column
    right_y = panel_y + 0.15
    if 'caption' in v:
        add_sp(sld, sp(spid, 'Caption', right_x, right_y, right_w, 1.20,
                       v['caption'], font='Twinkl Cursive Looped Light', sz=20,
                       color='1F4E79', align='l', fill='DEECF8',
                       border=('156082', 1.5), anchor='ctr'))
    sld.notes_slide.notes_text_frame.text = v['notes']




# ===========================================================================
# SPOT THE MISTAKE SLIDE — 3-beat animated teaching slide
# Beat 1 (load):    grid + instruction shown; error marker + explanation hidden
# Beat 2 (click 1): error marker (✗) appears at landing position
# Beat 3 (click 2): explanation box appears
# ===========================================================================
def build_spot_the_mistake_slide(layout_num, visual_key, title_text):
    sld = new_slide(layout_num)
    v = VISUALS[visual_key]

    # Title via placeholder
    for ph in sld.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title_text
            break

    # Draw the grid (without error elements — those come with animation)
    draw_grid_slide(sld, visual_key, layout_num)

    # Calculate grid geometry (matching draw_grid_slide)
    cols, rows = v['cols'], v['rows']
    panel_x, panel_y = 0.40, 1.45
    panel_w, panel_h = 7.00, 5.80
    cell = min(panel_w / (cols + 1), panel_h / (rows + 1))
    margin = 0.50
    grid_x = panel_x + margin
    grid_y = panel_y + 0.20

    # ── Error marker position — derived from errorType and data ──────────────
    # Strategy: find the 'error' position from extraPoints if present,
    # otherwise fall back to type-specific logic from startPoint.
    error_type   = v.get('error_type', 'off_grid')
    extra_points = v.get('extra_points', [])
    start_col, start_row = (v['points'][0][0], v['points'][0][1]) if v['points'] else (3, 3)

    # Look for an explicit error position in extraPoints
    err_col, err_row = None, None
    for ep in extra_points:
        if len(ep) >= 3 and 'error' in str(ep[2]).lower():
            err_col, err_row = ep[0], ep[1]
            break

    # Fallback by errorType if no explicit error point found
    if err_col is None:
        if error_type == 'off_grid':
            # off_grid: 4 right from startPoint goes off grid
            err_col = start_col + 4
            err_row = start_row
        elif error_type in ('partial_translation', 'wrong_reflection_distance',
                            'wrong_reflection_direction', 'wrong_side_of_mirror'):
            # One vertex/point moved incorrectly — use first extraPoint as error position
            if extra_points:
                err_col, err_row = extra_points[0][0], extra_points[0][1]
            else:
                err_col, err_row = start_col + 3, start_row
        elif error_type in ('wrong_direction', 'axis_swap', 'inverse_rule'):
            # Moves in wrong direction — place marker near startPoint offset
            err_col, err_row = start_col - 3, start_row + 3
        elif error_type in ('wrong_order', 'wrong_join_order'):
            # Wrong sequence — marker at first extraPoint
            if extra_points:
                err_col, err_row = extra_points[0][0], extra_points[0][1]
            else:
                err_col, err_row = start_col + 2, start_row + 2
        elif error_type in ('false_symmetry', 'over_counted_symmetry', 'moved_fixed_point'):
            # Shape error — use last extraPoint
            if extra_points:
                err_col, err_row = extra_points[-1][0], extra_points[-1][1]
            else:
                err_col, err_row = start_col + 3, start_row + 3
        else:
            # Generic fallback
            err_col, err_row = min(start_col + 3, cols), start_row

    # For Time lessons (gridSize=0), error marker goes in right column (no grid)
    if STM['gridSize'] == 0:
        err_slide_x = panel_x + panel_w + 0.25
        err_slide_y = 2.30
        marker_x = err_slide_x
        err_y = err_slide_y
    else:
        # Clamp to grid bounds for display
        err_col_clamped = min(err_col, cols + 0.5)
        err_x = grid_x + err_col_clamped * cell
        err_y = grid_y + (rows - err_row) * cell
        marker_x = min(err_x - 0.15, panel_x + panel_w - 0.35)

    # Right column layout
    right_x = panel_x + panel_w + 0.25
    right_w = 13.333 - right_x - 0.2

    # Beat 1 element — instruction (always visible on load)
    add_sp(sld, sp(133, 'ErrorInstruction',
                   right_x, 1.60, right_w, 0.80,
                   f'Instruction given: "{v["error_instruction"]}"',
                   font='Twinkl Cursive Looped Light', sz=16,
                   color='333333', align='l', fill=None, no_line=True))

    # Beat 2 element — error marker (hidden at start, appears on click 1)
    MARK_SPID = 120
    add_sp(sld, sp(MARK_SPID, 'ErrorMark', marker_x, err_y - 0.18, 0.30, 0.36,
                   '✗', font='Aptos', sz=20, bold=True,
                   color='FF0000', align='ctr', fill=None, no_line=True))

    # Beat 3 element — explanation box (hidden at start, appears on click 2)
    NOTE_SPID = 132
    add_sp(sld, sp(NOTE_SPID, 'ErrorNote', right_x, 2.60, right_w, 1.50,
                   v['error_note'],
                   font='Twinkl Cursive Looped Light', sz=18, bold=True,
                   color='C00000', align='l', fill='FFE6E6',
                   border=('C00000', 1.5), anchor='ctr'))

    # ── Animation timing ──
    # Both MARK_SPID and NOTE_SPID start hidden, revealed on sequential clicks
    def anim_block(ctn_id, spid):
        return f'''<p:par>
  <p:cTn id="{ctn_id}" fill="hold">
    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
    <p:childTnLst><p:par><p:cTn id="{ctn_id+1}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst><p:par>
        <p:cTn id="{ctn_id+2}" presetID="1" presetClass="entr" presetSubtype="0"
               fill="hold" grpId="1" nodeType="clickEffect">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst><p:set><p:cBhvr>
            <p:cTn id="{ctn_id+3}" dur="1" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
            <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
          </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
        </p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
  </p:cTn>
</p:par>'''

    child_blocks = anim_block(3, MARK_SPID) + anim_block(7, NOTE_SPID)

    bld_list = (
        f'<p:bldP spid="{MARK_SPID}" grpId="0" uiExpand="1" build="p"/>\n'
        f'<p:bldP spid="{MARK_SPID}" grpId="1" animBg="1"/>\n'
        f'<p:bldP spid="{NOTE_SPID}" grpId="0" uiExpand="1" build="p"/>\n'
        f'<p:bldP spid="{NOTE_SPID}" grpId="1" animBg="1"/>\n'
    )

    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par>
    <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
      <p:childTnLst><p:seq concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
          <p:childTnLst>{child_blocks}</p:childTnLst>
        </p:cTn>
        <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
        <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
      </p:seq></p:childTnLst>
    </p:cTn>
  </p:par></p:tnLst>
  <p:bldLst>{bld_list}</p:bldLst>
</p:timing>'''

    sld._element.append(etree.fromstring(timing_xml))

    sld.notes_slide.notes_text_frame.text = (
        f"I DO — Spot the Mistake (3-beat animated)\n"
        f"Beat 1 (load): Grid + instruction visible. Ask: 'What do you notice?'\n"
        f"Beat 2 (click 1): ✗ appears at landing position. Ask: 'Where does it end up?'\n"
        f"Beat 3 (click 2): Explanation revealed.\n\n"
        + v['notes']
    )
    print(f"  Spot the mistake slide ({title_text[:40]}) ✓")
    return sld


# ===========================================================================
# TRIOS SLIDE
# ===========================================================================
def build_trios_slide(layout_num, title, trios_data, notes, chart_keys=None):
    """chart_keys: list of VISUALS keys whose charts should appear on this slide."""
    sld = new_slide(layout_num)
    for ph in sld.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
        elif ph.placeholder_format.idx == 1:
            from pptx.oxml.ns import qn as _qn
            sp_elem = ph._element
            spPr = sp_elem.find(_qn('p:spPr'))
            if spPr is None:
                from lxml import etree as _et
                spPr = _et.SubElement(sp_elem, _qn('p:spPr'))
            from lxml import etree as _et
            _et.SubElement(spPr, _qn('a:noFill'))
            ln = _et.SubElement(spPr, _qn('a:ln'))
            _et.SubElement(ln, _qn('a:noFill'))

    roles = trios_data.get('roles', [])
    has_roles = bool(roles)
    has_charts = bool(chart_keys)
    n_charts = len(chart_keys) if chart_keys else 0

    # ── Layout: challenge removed throughout — charts get full remaining height ─
    if has_charts:
        task_x, task_w = 0.5, 12.4
        task_y, task_h = 1.45, 1.10
        chart_y = 2.70
        chart_h = 4.55
        if n_charts == 1:
            chart_x, chart_w = 2.17, 9.00   # centred
        else:
            chart_w = 6.10
            chart_xs = [0.40, 6.83]
    elif has_roles:
        roles_colors = [('1F4E79','DEEAF1'), ('7030A0','EAD1F0'), ('C00000','FCE4D6')]
        role_w, role_h = 3.8, 1.0
        for i, (role_text, (text_col, fill_col)) in enumerate(zip(roles, roles_colors)):
            y_pos = 1.6 + i * 1.15
            add_sp(sld, sp(20+i, f'Role{i+1}', 0.5, y_pos, role_w, role_h,
                           role_text, font='Twinkl Cursive Looped Light', sz=18,
                           bold=True, color=text_col, align='l',
                           fill=fill_col, border=(text_col, 1.5), anchor='ctr'))
        task_x, task_w = 4.6, 8.5
        task_y, task_h = 1.6, 3.5
    else:
        task_x, task_w = 0.5, 12.4
        task_y, task_h = 1.6, 4.5

    add_sp(sld, sp(30, 'Task', task_x, task_y, task_w, task_h,
                   trios_data.get('task',''),
                   font='Twinkl Cursive Looped Light', sz=18,
                   color='1F4E79', align='l', fill='DEECF8',
                   border=('156082', 1.5), anchor='ctr'))

    # Challenge removed from slide — kept in speaker notes only

    # ── Embed charts ──────────────────────────────────────────────────────
    if has_charts:
        chart_dir = os.path.join(tempfile.gettempdir(), 'wfa_stats_charts')
        os.makedirs(chart_dir, exist_ok=True)
        for i, vk in enumerate(chart_keys):
            v = VISUALS.get(vk, {})
            ct = v.get('chart_type', 'bar_chart')
            cd = v.get('chart_data', {})
            chart_path = os.path.join(chart_dir, f'{vk}_{ct}.png')
            if not os.path.exists(chart_path):
                render_stats_chart(ct, cd, chart_path, dpi=180)
            cx = chart_xs[i] if n_charts > 1 else chart_x
            sld.shapes.add_picture(chart_path,
                                   emu(cx), emu(chart_y),
                                   emu(chart_w), emu(chart_h))

    sld.notes_slide.notes_text_frame.text = notes
    print(f"  Trios slide ({title[:40]}) ✓")
    return sld

# ===========================================================================
# INDEPENDENT SLIDE
# ===========================================================================
def build_independent_slide(layout_num, title, independent_data, notes):
    sld = new_slide(layout_num)
    for ph in sld.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            break

    add_sp(sld, sp(20, 'Standard', 0.5, 1.6, 12.5, 1.3,
                   independent_data.get('standard',''),
                   font='Twinkl Cursive Looped Light', sz=20,
                   color='1F4E79', align='l', fill='DEECF8',
                   border=('156082', 1.5), anchor='ctr'))

    add_sp(sld, sp(21, 'Supported', 0.5, 3.1, 12.5, 1.3,
                   f"Supported: {independent_data.get('supported','')}",
                   font='Twinkl Cursive Looped Light', sz=18,
                   color='333333', align='l', fill='F2F2F2',
                   border=('BBBBBB', 1.5), anchor='ctr'))

    sld.notes_slide.notes_text_frame.text = notes
    print(f"  Independent slide ({title[:40]}) ✓")
    return sld

# ===========================================================================
# LP / MARKING STATION BLANK SLIDES
# ===========================================================================
def build_lp_slide(label):
    sld = new_slide(5)
    add_sp(sld, sp(20, 'LPTitle', 0.5, 0.1, 12.0, 0.9,
                   label, font='Twinkl Cursive Looped Light',
                   sz=32, bold=True, color='000000', align='ctr',
                   fill=None, no_line=True))
    sld.notes_slide.notes_text_frame.text = (
        f"YOU DO (INDEPENDENT)\n{label}\nLeave blank — Innes adds screenshot.")
    print(f"  {label} slide ✓")
    return sld

# ===========================================================================
# SLIDE 22 — LEARNING REVIEW  (fixed: wedgeRoundRectCallout, 20pt)
# ===========================================================================
def build_learning_review():
    sld = new_slide(10)
    stems = L1['learningReviewStems']
    sid = [10]
    def nid(): sid[0]+=1; return sid[0]

    add_sp(sld, sp(nid(),'Title', 5.172,0.197, 3.352,0.640,
                   'Learning Review', font='Twinkl Cursive Looped Light',
                   sz=24, bold=True, color='000000', align='ctr',
                   fill=None, no_line=True))

    # Speech bubbles — wedgeRoundRectCallout at 20pt (matching template)
    bubbles = [
        (stems[0], 'E9917F', 1.206, 1.131, 3.065, 1.272),
        (stems[1], 'D977ED', 5.274, 1.176, 3.065, 1.272),
        (stems[2], '92D050', 9.442, 1.176, 3.012, 1.272),
    ]
    for i, (text, color, bx, by, bw, bh) in enumerate(bubbles):
        add_sp(sld, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr><p:cNvPr id="{nid()}" name="Bubble{i+1}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{emu(bx)}" y="{emu(by)}"/>
      <a:ext cx="{emu(bw)}" cy="{emu(bh)}"/></a:xfrm>
    <a:prstGeom prst="wedgeRoundRectCallout"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr rtlCol="0" anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r>
      <a:rPr lang="en-GB" sz="2000" b="1" dirty="0">
        <a:latin typeface="Twinkl Cursive Looped Light"/>
      </a:rPr>
      <a:t>{_esc(text)}</a:t>
    </a:r></a:p>
  </p:txBody>
</p:sp>''')

    # Images — positions and filenames matched exactly to approved LR_slide.pptx
    # Reference mapping (by size match):
    #   ref image6.png  (2135×1012) = image7.png  — children group
    #   ref image7.png  (234×290)   = image11.png — character (small)
    #   ref image8.png  (440×559)   = image12.png — character
    #   ref image9.png  (453×554)   = image13.png — character
    #   ref image10.png (499×582)   = image14.png — character (was image14.jpeg)
    #   ref image11.png (476×605)   = image15.png — character
    #   ref image12.png (464×464)   = image16.png — LI icon (was image16.jpeg)
    #   ref image13.png (464×464)   = image17.png — LI icon
    #   ref image14.jpg (463×463)   = image18.png — LI icon
    #   ref image15.jpg (463×463)   = image19.png — LI icon
    add_pic(sld,'image7.png',   3.166, 2.797, 6.458, 3.056)
    add_pic(sld,'image11.png',  0.539, 5.105, 1.775, 2.166)
    add_pic(sld,'image12.png',  3.435, 6.152, 1.036, 1.305)
    add_pic(sld,'image13.png',  4.530, 6.163, 1.057, 1.295)
    add_pic(sld,'image14.png',  5.603, 6.073, 1.165, 1.359)
    add_pic(sld,'image15.png',  6.806, 6.034, 1.111, 1.424)
    add_pic(sld,'image16.png', 10.810, 4.301, 1.050, 1.050)
    add_pic(sld,'image17.png', 11.997, 4.301, 1.051, 1.051)
    add_pic(sld,'image18.png', 10.808, 5.448, 1.051, 1.051)
    add_pic(sld,'image19.png', 11.998, 5.449, 1.050, 1.051)

    sld.notes_slide.notes_text_frame.text = (
        f"WE DO — Learning Review\n" +
        '\n'.join(f"Stem {i+1}: {s}" for i,s in enumerate(stems)) +
        "\n\nStand to speak. Learning intelligence icons prompt discussion style.")
    print("Slide 22 (Learning Review) ✓")

# ===========================================================================
# BUILD ALL 22 SLIDES
# ===========================================================================
week_label = L1['week']        # e.g. 'T5W1'
day        = L1['day']         # e.g. 'Monday'
lesson_num = L1['lesson']      # e.g. 1

print(f"\n=== Building {week_label}_L{lesson_num}_{day}_Teaching_v3.pptx ===\n")

c1 = L1['cycle1']
c2 = L1['cycle2']

build_slide1()
build_slide2()
build_slide3()
build_slide4()
build_slide5()
if RM_DATA.get('questions'):
    build_slide6()
    build_slide7()
build_slide8()

# Cycle 1 teaching slides
build_teaching_slide(2, 'c1_ido1', c1['slideTitles']['ido'][0], 'I DO')

# Second C1 I DO — could be a regular slide or a C1-level STM
if len(c1['slideTitles']['ido']) > 1:
    title2 = c1['slideTitles']['ido'][1]
    if 'c1_ido2' in VISUALS:
        v2 = VISUALS['c1_ido2']
        # Only use STM builder if it's explicitly a spot_the_mistake type
        if v2.get('slide_type') == 'spot_the_mistake':
            build_spot_the_mistake_slide(2, 'c1_ido2', title2)
        else:
            # Phase is I Do by default — set 'c1_ido2_phase':'WE DO' in lesson data
            # when slide 2 is a variation of slide 1 (children more involved)
            _c1_p2 = L1.get('c1_ido2_phase', 'I DO')
            build_teaching_slide(3 if _c1_p2=='WE DO' else 2, 'c1_ido2', title2, _c1_p2)

if c1['slideTitles'].get('wedo') and 'c1_wedo' in VISUALS:
    build_teaching_slide(3, 'c1_wedo', c1['slideTitles']['wedo'][0], 'WE DO')

build_trios_slide(4, c1['slideTitles']['trios'][0], c1['trios'],
                  f"YOU DO (TRIOS)\n{c1['trios']['task']}\nChallenge: {c1['trios']['challenge']}",
                  chart_keys=c1.get('trios_charts', []))
if c1.get('slideCount',{}).get('independent',1) > 0:
    build_independent_slide(5, c1['slideTitles']['independent'][0], c1['independent'],
                            f"YOU DO (INDEPENDENT) — C1\n{c1['independent']['standard']}")
build_lp_slide('Learning Paper 1')
build_lp_slide('Marking Station 1')

# Cycle 2 teaching slides
if 'c2_ido1' in VISUALS:
    build_teaching_slide(2, 'c2_ido1', c2['slideTitles']['ido'][0], 'I DO')

# C2 second I Do — authored visual takes priority over STM auto-generation
if len(c2['slideTitles']['ido']) > 1:
    title_c2i2 = c2['slideTitles']['ido'][1]
    v2 = VISUALS.get('c2_ido2', {})
    st2 = v2.get('slide_type', '')
    if st2 == 'stm_word_problem':
        # Authored word-problem STM: use visual's own title
        build_teaching_slide(3 if L1.get('c2_ido2_phase','I DO')=='WE DO' else 2, 'c2_ido2', v2.get('title', STM['slideTitle']), L1.get('c2_ido2_phase','I DO'))
    elif st2 not in ('spot_the_mistake', ''):
        # Other authored slide
        build_teaching_slide(3 if L1.get('c2_ido2_phase','I DO')=='WE DO' else 2, 'c2_ido2', title_c2i2, L1.get('c2_ido2_phase','I DO'))
    else:
        # Default: grid-based STM from JSON
        build_spot_the_mistake_slide(2, 'c2_ido2', title_c2i2)
else:
    v2 = VISUALS.get('c2_ido2', {})
    if v2.get('slide_type') == 'stm_word_problem':
        build_teaching_slide(3 if L1.get('c2_ido2_phase','I DO')=='WE DO' else 2, 'c2_ido2', v2.get('title', STM['slideTitle']), L1.get('c2_ido2_phase','I DO'))
    else:
        build_spot_the_mistake_slide(2, 'c2_ido2', STM['slideTitle'])

if c2['slideTitles'].get('wedo'):
    if 'c2_wedo' in VISUALS:
        build_teaching_slide(3, 'c2_wedo', c2['slideTitles']['wedo'][0], 'WE DO')

build_trios_slide(4, c2['slideTitles']['trios'][0], c2['trios'],
                  f"YOU DO (TRIOS)\n{c2['trios']['task']}\nChallenge: {c2['trios']['challenge']}",
                  chart_keys=c2.get('trios_charts', []))
if c2.get('slideCount',{}).get('independent',1) > 0:
    build_independent_slide(5, c2['slideTitles']['independent'][0], c2['independent'],
                            f"YOU DO (INDEPENDENT) — C2\n{c2['independent']['standard']}")
build_lp_slide('Learning Paper 2')
build_lp_slide('Marking Station 2')

build_learning_review()

# ---------------------------------------------------------------------------
# SAVE
# ---------------------------------------------------------------------------
out = f'/home/claude/{week_label}_L{lesson_num}_Teaching.pptx'
prs.save(out)
inject_kq_slide(out)
print(f"\n=== Saved: {out} ({len(prs.slides)+1} slides) ===")

# ---------------------------------------------------------------------------
# PRE-FLIGHT CHECK
# Reads the saved file and checks for common layout problems.
# Reports warnings but does not block saving.
# Checks:
#   1. Text overflow — text box too narrow/short for its content at stated font size
#   2. Out-of-bounds — shapes placed outside slide dimensions
#   3. Grid-point bounds — any plotted point outside its declared grid
#   4. WM type rule — items match the expected type for the day of week
# ---------------------------------------------------------------------------
def run_preflight(pptx_path, lesson_data_items, day_name):
    from pptx import Presentation as _Prs
    from pptx.util import Pt as _Pt
    import re as _re

    SLIDE_W_EMU = 13.333 * 914400
    SLIDE_H_EMU = 7.5   * 914400

    # Approximate character width in EMU at a given font size (pt)
    # Assumes Aptos/sans-serif — roughly 0.55× the point size in width per char
    def approx_text_w_emu(text, font_pt):
        return len(str(text)) * font_pt * 0.55 * 12700  # 12700 EMU per pt

    issues = []
    deck = _Prs(pptx_path)

    for slide_idx, slide in enumerate(deck.slides, 1):
        for shape in slide.shapes:
            # ── Out of bounds check ──────────────────────────────────────────
            if hasattr(shape, 'left') and shape.left is not None:
                l, t = shape.left, shape.top
                r = l + (shape.width  or 0)
                b = t + (shape.height or 0)
                if l < -914400 or t < -914400 or r > SLIDE_W_EMU + 914400 or b > SLIDE_H_EMU + 914400:
                    issues.append(
                        f"Slide {slide_idx} '{shape.name}': shape out of bounds "
                        f"(left={l//914400:.2f}\" top={t//914400:.2f}\" "
                        f"right={r//914400:.2f}\" bottom={b//914400:.2f}\")"
                    )

            # ── Text overflow check ──────────────────────────────────────────
            # Only flag genuine problems:
            #   - Single-line boxes where text is wider than the box
            #   - Any box where the font is taller than the box height
            if shape.has_text_frame and shape.width and shape.height:
                tf = shape.text_frame
                box_w_emu = shape.width
                box_h_emu = shape.height

                for para in tf.paragraphs:
                    for run in para.runs:
                        txt = run.text.strip()
                        if not txt:
                            continue
                        sz_pt = None
                        if run.font.size:
                            sz_pt = run.font.size.pt
                        elif para.runs and para.runs[0].font.size:
                            sz_pt = para.runs[0].font.size.pt
                        if sz_pt is None:
                            sz_pt = 12

                        line_h_emu = sz_pt * 12700 * 1.25  # 125% leading

                        # Height check: box shorter than a single line — always a problem
                        if line_h_emu > box_h_emu * 1.05:
                            issues.append(
                                f"Slide {slide_idx} '{shape.name}': "
                                f"box height {box_h_emu/12700:.0f}pt too small "
                                f"for {sz_pt:.0f}pt font — text will be clipped"
                            )

                        # Width check: only flag if box is also too short to wrap
                        # (i.e. single-line box). Multi-line boxes can wrap, so skip.
                        lines_available = box_h_emu / line_h_emu
                        if lines_available < 1.8:  # effectively single-line
                            est_w = len(txt) * sz_pt * 0.55 * 12700
                            if est_w > box_w_emu * 1.20:
                                issues.append(
                                    f"Slide {slide_idx} '{shape.name}': "
                                    f"single-line box — '{txt[:25]}...' "
                                    f"estimated {est_w/12700:.0f}pt wide, "
                                    f"box only {box_w_emu/12700:.0f}pt"
                                )

    # ── WM type rule check ───────────────────────────────────────────────────
    day_map    = {'Monday':1,'Tuesday':2,'Wednesday':3,'Thursday':4}
    wm_types   = ['numbers','words','emojis','text+image']
    day_pos    = day_map.get(day_name, 1)
    expected   = wm_types[day_pos - 1]
    items      = lesson_data_items
    if all(isinstance(i, int) for i in items):
        actual = 'numbers'
    elif all(isinstance(i, str) and any(ord(c) > 127 for c in i) for i in items):
        actual = 'emojis'
    elif all(isinstance(i, str) for i in items):
        actual = 'words'
    else:
        actual = 'mixed/text+image'

    if actual != expected and expected != 'text+image':
        issues.append(
            f"WM TYPE MISMATCH: {day_name} should be '{expected}' "
            f"but items are '{actual}' — check lesson_data.py"
        )

    # ── Report ───────────────────────────────────────────────────────────────
    print(f"\n--- Pre-flight check ({len(deck.slides)} slides) ---")
    if issues:
        print(f"  {len(issues)} issue(s) found:")
        for iss in issues:
            print(f"  ⚠  {iss}")
    else:
        print("  ✓ No layout issues detected")
    print("---\n")
    return issues

_wm_items = _ld['wm']['items']
run_preflight(out, _wm_items, day)

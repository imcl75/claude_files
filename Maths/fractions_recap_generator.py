"""
fractions_recap_generator.py — T6W5 Fractions Recap (replaces Rapid Maths)
2 slides per day × 3 days = 6 slides total.
  Slide 1: I Do (mixed→improper)  |  You Do (mixed→improper)
  Slide 2: I Do (improper→mixed)  |  You Do (improper→mixed)
You Do answers appear on click via PPTX animation XML.
"""

import os, sys, tempfile
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

OUT_FILE  = '/tmp/claude_work/T6W5_Fractions_Recap.pptx'
CHART_DIR = tempfile.mkdtemp(prefix='wfa_frac_v2_')
EMU_IN    = 914400

# ── Colours ───────────────────────────────────────────────────────────────────
C_BLUE    = RGBColor(0x17, 0x98, 0xD3)
C_ORANGE  = RGBColor(0xE5, 0x7D, 0x24)
C_GREEN   = RGBColor(0x1A, 0x5C, 0x2A)
C_GBG     = RGBColor(0xE8, 0xF5, 0xE9)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x1A, 0x1A)
C_PANEL   = RGBColor(0xDE, 0xEC, 0xF8)
C_QBG     = RGBColor(0xDE, 0xEA, 0xF1)
C_STEP    = RGBColor(0x44, 0x44, 0x44)

# ── Question data per day ─────────────────────────────────────────────────────
# mixed→improper: (whole, num, den)
# improper→mixed: (num, den)
DAYS = {
    'Monday': {
        'm2i': {'ido': (2,3,5),  'youdo': [(3,1,4), (4,2,5)]},
        'i2m': {'ido': (11,3),   'youdo': [(17,4), (13,5)]},
    },
    'Tuesday': {
        'm2i': {'ido': (2,5,7),  'youdo': [(3,3,8)]},
        'i2m': {'ido': (19,6),   'youdo': [(23,5), (27,8)]},
    },
    'Wednesday': {
        'm2i': {'ido': (5,1,3),  'youdo': [(4,3,7)]},
        'i2m': {'ido': (31,8),   'youdo': [(22,5), (16,3)]},
    },
}

# ── Fraction rendering ────────────────────────────────────────────────────────
def _frac_png(num, den, tag, colour='#1798d3', fs=32, w=0.75, h=1.05):
    path = os.path.join(CHART_DIR, f'{tag}.png')
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor('none'); ax.set_facecolor('none')
    ax.axis('off'); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.78, str(num), ha='center', va='center',
            fontsize=fs, color=colour, fontweight='bold')
    ax.plot([0.06,0.94],[0.50,0.50], color=colour, lw=2.2)
    ax.text(0.5, 0.14, str(den), ha='center', va='center',
            fontsize=fs, color=colour, fontweight='bold')
    plt.savefig(path, dpi=180, bbox_inches='tight', transparent=True)
    plt.close(fig)
    return path

def _mixed_png(W, num, den, tag, colour='#1798d3', fs=28, w=1.20, h=1.05):
    path = os.path.join(CHART_DIR, f'{tag}.png')
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor('none'); ax.set_facecolor('none')
    ax.axis('off'); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.20, 0.50, str(W), ha='center', va='center',
            fontsize=fs+4, color=colour, fontweight='bold')
    ax.text(0.63, 0.78, str(num), ha='center', va='center',
            fontsize=fs-2, color=colour, fontweight='bold')
    ax.plot([0.38,0.92],[0.50,0.50], color=colour, lw=2.0)
    ax.text(0.63, 0.14, str(den), ha='center', va='center',
            fontsize=fs-2, color=colour, fontweight='bold')
    plt.savefig(path, dpi=180, bbox_inches='tight', transparent=True)
    plt.close(fig)
    return path

# ── python-pptx helpers ───────────────────────────────────────────────────────
def I(v): return Emu(int(v * EMU_IN))

def txt(slide, text, x, y, w, h, font='Aptos', sz=14, bold=False,
        color=C_DARK, align=PP_ALIGN.LEFT, fill=None, wrap=True):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bp = tf._txBody.find(f'{{{A}}}bodyPr')
    if bp is None:
        bp = etree.SubElement(tf._txBody, f'{{{A}}}bodyPr')
    bp.set('anchor', 't')
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, I(x),I(y),I(w),I(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.fill.background()
    return s

def pic(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, I(x),I(y),I(w),I(h))

# ── Animation: appear on click ────────────────────────────────────────────────
def add_appear_on_click(slide, shape_ids):
    """Add a single click-to-appear animation group for the given shape IDs."""
    if not shape_ids:
        return
    P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    R  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    DR = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    spTree = slide.shapes._spTree
    root   = slide._element

    # Build timing XML
    timing_xml = f'''<p:timing xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>'''

    for idx, sid in enumerate(shape_ids):
        timing_xml += f'''
                            <p:par>
                              <p:cTn id="{10+idx}" presetID="1" presetClass="entr"
                                     presetSubtype="0" fill="hold" grpId="0"
                                     nodeType="clickEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr><p:cTn id="{20+idx}" dur="1" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>'''

    timing_xml += '''
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:navAttr nextAc="seek"/>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''

    timing_el = etree.fromstring(timing_xml.encode())
    root.append(timing_el)

    # Set initial visibility to hidden for each shape
    for sid in shape_ids:
        for sp in spTree.iter():
            cNvPr = sp.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr/'
                '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if cNvPr is None:
                cNvPr = sp.find(
                    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr/'
                    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if cNvPr is not None and cNvPr.get('id') == str(sid):
                spPr_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr'
                spPr = sp.find(spPr_ns)
                if spPr is None:
                    spPr_ns2 = '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr'
                    spPr = sp.find(spPr_ns2)
                if spPr is not None:
                    xfrm = spPr.find(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                    if xfrm is not None:
                        pass  # visibility handled via animation

# ── Panel drawing ─────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.333, 7.5
DIV_X  = 6.80     # dividing line between I Do and You Do
LEFT_W = DIV_X - 0.40          # I Do panel width
RGHT_X = DIV_X + 0.25
RGHT_W = SLIDE_W - RGHT_X - 0.25
CONT_Y = 1.22     # content top (below title bar)
CONT_H = SLIDE_H - CONT_Y - 0.20

def draw_title_bar(slide, day, direction_label, colour):
    rect(slide, 0, 0, SLIDE_W, 1.05, fill=colour)
    txt(slide, f'Fractions Recap — {direction_label}',
        0.30, 0.08, 9.5, 0.50,
        font='Twinkl Cursive Looped Light', sz=26, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT)
    txt(slide, day, 9.90, 0.08, 3.10, 0.50,
        font='Aptos', sz=18, bold=True,
        color=C_WHITE, align=PP_ALIGN.RIGHT)

# ── Ensure badge images are available, fetching from Shared/badges/ if needed ──
def _ensure_badges(local_dir):
    import urllib.request, os
    BASE = 'https://raw.githubusercontent.com/imcl75/claude_files/main/Shared/badges'
    TOKEN_FILE = '/mnt/skills/user/github-sync/SKILL.md'
    token = ''
    if os.path.exists(TOKEN_FILE):
        import re as _re
        m = _re.search(r'github_pat_[A-Za-z0-9_]+', open(TOKEN_FILE).read())
        if m: token = m.group()
    os.makedirs(local_dir, exist_ok=True)
    for fname in ['badge_ido.png', 'badge_wedo.png', 'badge_youdo_ind.png', 'badge_youdo_trio.png']:
        dest = os.path.join(local_dir, fname)
        if not os.path.exists(dest):
            url = f'{BASE}/{fname}'
            req = urllib.request.Request(url,
                headers={'Authorization': f'token {token}'} if token else {})
            try:
                with urllib.request.urlopen(req) as r:
                    open(dest, 'wb').write(r.read())
            except Exception as e:
                print(f'  Warning: could not fetch {fname}: {e}')

_BADGE_DIR = '/tmp/claude_work/assets'
_ensure_badges(_BADGE_DIR)

BADGE_PATHS = {
    'I Do':   f'{_BADGE_DIR}/badge_ido.png',
    'You Do': f'{_BADGE_DIR}/badge_youdo_ind.png',
    'We Do':  f'{_BADGE_DIR}/badge_wedo.png',
    'Trios':  f'{_BADGE_DIR}/badge_youdo_trio.png',
}
BADGE_W, BADGE_H = 1.20, 0.52   # display size in inches

def phase_badge(slide, label, x, y, colour=None):
    """Embed the real school phase badge image."""
    path = BADGE_PATHS.get(label)
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, I(x), I(y), I(BADGE_W), I(BADGE_H))
    else:
        # Fallback text badge
        rect(slide, x, y, 1.30, 0.42, fill=colour or C_BLUE)
        txt(slide, label, x+0.05, y+0.03, 1.20, 0.36,
            font='Aptos', sz=13, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)

def q_number_badge(slide, n, x, y, colour):
    """Render Q1/Q2 label as a matplotlib PNG — bypasses all PPTX text wrapping."""
    import matplotlib.pyplot as plt
    hex_col = f'#{colour[0]:02x}{colour[1]:02x}{colour[2]:02x}'
    path = os.path.join(CHART_DIR, f'qlabel_{n}.png')
    fig, ax = plt.subplots(figsize=(0.55, 0.34))
    fig.patch.set_facecolor(hex_col)
    ax.set_facecolor(hex_col)
    ax.axis('off')
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.text(0.50, 0.50, f'Q{n}', ha='center', va='center',
            fontsize=14, color='white', fontweight='bold')
    plt.savefig(path, dpi=180, bbox_inches='tight', facecolor=hex_col)
    plt.close(fig)
    slide.shapes.add_picture(path, I(x), I(y), I(0.55), I(0.34))

def draw_divider(slide):
    line = slide.shapes.add_connector(
        1, I(DIV_X), I(CONT_Y-0.10), I(DIV_X), I(SLIDE_H-0.15))
    line.line.color.rgb = RGBColor(0xBB,0xBB,0xBB)
    line.line.width = Pt(0.75)

# ── Slide 1: Mixed → Improper ─────────────────────────────────────────────────
def build_m2i_slide(slide, day, data):
    draw_title_bar(slide, day, 'Mixed → Improper', C_BLUE)
    draw_divider(slide)

    ido  = data['ido']   # (W, N, D)
    youdo = data['youdo']

    W_n, N_n, D_n = ido
    ans_num = W_n * D_n + N_n
    ans_den = D_n

    # ── I Do (left panel) ─────────────────────────────────────────────────────
    phase_badge(slide, 'I Do', 0.30, CONT_Y, C_BLUE)
    txt(slide, 'Watch and follow the steps.',
        BADGE_W + 0.30, CONT_Y + 0.17, 5.0, 0.38,
        font='Aptos', sz=11, color=C_STEP)

    # Question mixed number — large
    q_path = _mixed_png(W_n, N_n, D_n, f'm2i_ido_q_{day[:3]}', fs=34, w=1.40, h=1.15)
    pic(slide, q_path, 0.50, CONT_Y + 0.82, 1.40, 1.15)

    # Steps
    step_y = CONT_Y + 0.90
    step_x = 2.10
    txt(slide, f'Step 1:  Multiply whole × denominator',
        step_x, step_y, LEFT_W - step_x + 0.35, 0.36,
        font='Aptos', sz=13, bold=True, color=C_BLUE)
    txt(slide, f'{W_n} × {D_n} = {W_n * D_n}',
        step_x + 0.30, step_y + 0.38, 3.5, 0.38,
        font='Aptos', sz=16, bold=False, color=C_DARK)

    txt(slide, f'Step 2:  Add the numerator',
        step_x, step_y + 0.88, LEFT_W - step_x + 0.35, 0.36,
        font='Aptos', sz=13, bold=True, color=C_BLUE)
    txt(slide, f'{W_n * D_n} + {N_n} = {ans_num}',
        step_x + 0.30, step_y + 1.26, 3.5, 0.38,
        font='Aptos', sz=16, bold=False, color=C_DARK)

    # Answer
    txt(slide, 'Answer:', 0.50, CONT_Y + 3.00, 1.40, 0.40,
        font='Aptos', sz=13, bold=True, color=C_GREEN)
    a_path = _frac_png(ans_num, ans_den, f'm2i_ido_a_{day[:3]}',
                       colour='#1A5C2A', fs=34, w=0.80, h=1.15)
    pic(slide, a_path, 1.95, CONT_Y + 2.80, 0.80, 1.15)

    # ── You Do (right panel) ───────────────────────────────────────────────────
    phase_badge(slide, 'You Do', RGHT_X, CONT_Y, C_ORANGE)
    txt(slide, 'Now you try — convert to an improper fraction.',
        RGHT_X + BADGE_W + 0.15, CONT_Y + 0.17, RGHT_W - BADGE_W - 0.20, 0.38,
        font='Aptos', sz=11, color=C_STEP)

    answer_shape_ids = []
    for i, (Wy, Ny, Dy) in enumerate(youdo):
        ans_n = Wy * Dy + Ny
        qy = CONT_Y + 0.70 + i * 2.60

        # Question
        qp = _mixed_png(Wy, Ny, Dy, f'm2i_yd_q{i}_{day[:3]}', fs=30, w=1.25, h=1.05)
        q_number_badge(slide, i+1, RGHT_X, qy + 0.10, C_ORANGE)
        pic(slide, qp, RGHT_X + 0.48, qy - 0.05, 1.25, 1.05)
        txt(slide, '=', RGHT_X + 1.90, qy + 0.20, 0.45, 0.45,
            font='Aptos', sz=22, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

        # Answer line box
        ans_box = rect(slide, RGHT_X + 2.45, qy + 0.08, 1.40, 0.85,
                       fill=C_QBG, line=C_BLUE, lw=Pt(0.75))

        # Answer fraction (hidden, revealed on click)
        ap = _frac_png(ans_n, Dy, f'm2i_yd_a{i}_{day[:3]}',
                       colour='#1A5C2A', fs=30, w=0.75, h=1.05)
        ans_pic = pic(slide, ap, RGHT_X + 2.55, qy - 0.05, 0.75, 1.05)

        # Working note below (also hidden)
        working_tb = txt(slide, f'({Wy}×{Dy})+{Ny} = {ans_n}',
                         RGHT_X, qy + 1.10, RGHT_W, 0.35,
                         font='Aptos', sz=10, color=C_GREEN)

        # Collect shape IDs for animation
        answer_shape_ids.append(ans_pic.shape_id)
        answer_shape_ids.append(working_tb.shape_id)

    # Hide answers initially and set animation
    _set_hidden_and_animate(slide, answer_shape_ids)


# ── Slide 2: Improper → Mixed ─────────────────────────────────────────────────
def build_i2m_slide(slide, day, data):
    draw_title_bar(slide, day, 'Improper → Mixed', RGBColor(0x2B,0xAE,0x62))
    draw_divider(slide)

    ido   = data['ido']    # (N, D)
    youdo = data['youdo']

    N_n, D_n = ido
    W_ans, R_ans = divmod(N_n, D_n)

    # ── I Do ──────────────────────────────────────────────────────────────────
    phase_badge(slide, 'I Do', 0.30, CONT_Y, RGBColor(0x2B,0xAE,0x62))
    txt(slide, 'Watch and follow the steps.',
        1.70, CONT_Y + 0.04, 5.0, 0.38,
        font='Aptos', sz=11, color=C_STEP)

    q_path = _frac_png(N_n, D_n, f'i2m_ido_q_{day[:3]}', fs=34, w=0.80, h=1.15)
    pic(slide, q_path, 0.55, CONT_Y + 0.58, 0.80, 1.15)

    step_y = CONT_Y + 0.70
    step_x = 1.80
    txt(slide, f'Step 1:  Divide numerator ÷ denominator',
        step_x, step_y, LEFT_W - step_x + 0.35, 0.36,
        font='Aptos', sz=13, bold=True, color=RGBColor(0x2B,0xAE,0x62))
    txt(slide, f'{N_n} ÷ {D_n} = {W_ans}  remainder {R_ans}',
        step_x + 0.30, step_y + 0.38, 4.2, 0.38,
        font='Aptos', sz=15, bold=False, color=C_DARK)

    txt(slide, 'Step 2:  Whole = quotient, fraction = remainder / divisor',
        step_x, step_y + 0.88, LEFT_W - step_x + 0.35, 0.50,
        font='Aptos', sz=13, bold=True, color=RGBColor(0x2B,0xAE,0x62), wrap=True)

    txt(slide, 'Answer:', 0.45, CONT_Y + 3.00, 1.35, 0.40,
        font='Aptos', sz=13, bold=True, color=C_GREEN)
    if R_ans == 0:
        a_path = _frac_png(N_n, D_n, f'i2m_ido_a_{day[:3]}',
                           colour='#1A5C2A', fs=34, w=0.80, h=1.15)
        pic(slide, a_path, 1.85, CONT_Y + 2.80, 0.80, 1.15)
    else:
        a_path = _mixed_png(W_ans, R_ans, D_n, f'i2m_ido_a_{day[:3]}',
                            colour='#1A5C2A', fs=30, w=1.20, h=1.15)
        pic(slide, a_path, 1.85, CONT_Y + 2.80, 1.20, 1.15)

    # ── You Do ────────────────────────────────────────────────────────────────
    phase_badge(slide, 'You Do', RGHT_X, CONT_Y, C_ORANGE)
    txt(slide, 'Now you try — convert to a mixed number.',
        RGHT_X + BADGE_W + 0.15, CONT_Y + 0.17, RGHT_W - BADGE_W - 0.20, 0.38,
        font='Aptos', sz=11, color=C_STEP)

    answer_shape_ids = []
    for i, (Ny, Dy) in enumerate(youdo):
        Wy_ans, Ry_ans = divmod(Ny, Dy)
        qy = CONT_Y + 0.70 + i * 2.60

        qp = _frac_png(Ny, Dy, f'i2m_yd_q{i}_{day[:3]}', fs=30, w=0.80, h=1.05)
        q_number_badge(slide, i+1, RGHT_X, qy + 0.25, C_ORANGE)
        pic(slide, qp, RGHT_X + 0.48, qy - 0.05, 0.80, 1.05)
        txt(slide, '=', RGHT_X + 1.45, qy + 0.20, 0.45, 0.45,
            font='Aptos', sz=22, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

        rect(slide, RGHT_X + 2.00, qy + 0.08, 1.65, 0.85,
             fill=C_QBG, line=RGBColor(0x2B,0xAE,0x62), lw=Pt(0.75))

        if Ry_ans == 0:
            ap = _frac_png(Ny, Dy, f'i2m_yd_a{i}_{day[:3]}',
                           colour='#1A5C2A', fs=30, w=0.80, h=1.05)
            ans_pic = pic(slide, ap, RGHT_X + 2.10, qy - 0.05, 0.80, 1.05)
        else:
            ap = _mixed_png(Wy_ans, Ry_ans, Dy, f'i2m_yd_a{i}_{day[:3]}',
                            colour='#1A5C2A', fs=26, w=1.20, h=1.05)
            ans_pic = pic(slide, ap, RGHT_X + 2.10, qy - 0.05, 1.20, 1.05)

        working_tb = txt(slide, f'{Ny} ÷ {Dy} = {Wy_ans} rem {Ry_ans}',
                         RGHT_X, qy + 1.10, RGHT_W, 0.35,
                         font='Aptos', sz=10, color=C_GREEN)

        answer_shape_ids.append(ans_pic.shape_id)
        answer_shape_ids.append(working_tb.shape_id)

    _set_hidden_and_animate(slide, answer_shape_ids)


# ── Animation helper ──────────────────────────────────────────────────────────
def _set_hidden_and_animate(slide, shape_ids):
    """Hide shapes initially; each click reveals one shape."""
    if not shape_ids:
        return
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    spTree = slide.shapes._spTree

    # Mark each shape as hidden via xfrm ext (not ideal but reliable)
    # Better: use visibility attribute on cNvSpPr
    for sid in shape_ids:
        for el in spTree.iter():
            cNvPr = (el.find(f'{{{P}}}nvSpPr/{{{P}}}cNvPr') or
                     el.find(f'{{{P}}}nvPicPr/{{{P}}}cNvPr'))
            if cNvPr is not None and cNvPr.get('id') == str(sid):
                cNvPr.set('hidden', '1')

    # Build timing XML with one click group per shape
    cTn_id = [5]
    def next_id():
        cTn_id[0] += 1
        return cTn_id[0]

    click_blocks = ''
    for idx, sid in enumerate(shape_ids):
        cid1, cid2, cid3 = next_id(), next_id(), next_id()
        click_blocks += f'''
        <p:par>
          <p:cTn id="{cid1}" fill="hold">
            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
            <p:childTnLst>
              <p:par>
                <p:cTn id="{cid2}" presetID="1" presetClass="entr"
                       presetSubtype="0" fill="hold" grpId="{idx}"
                       nodeType="clickEffect">
                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                  <p:childTnLst>
                    <p:set>
                      <p:cBhvr><p:cTn id="{cid3}" dur="1" fill="hold"/>
                        <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                        <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                      </p:cBhvr>
                      <p:to><p:strVal val="visible"/></p:to>
                    </p:set>
                  </p:childTnLst>
                </p:cTn>
              </p:par>
            </p:childTnLst>
          </p:cTn>
        </p:par>'''

    timing_xml = f'''<p:timing xmlns:p="{P}" xmlns:a="{A}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{click_blocks}
              </p:childTnLst>
            </p:cTn>
            <p:navAttr nextAc="seek"/>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''

    slide._element.append(etree.fromstring(timing_xml.encode()))


# ── Build ─────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation('/tmp/claude_work/assets/rapid_maths_TEMPLATE.pptx')

    # Clear template slides
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = prs.slides._sldIdLst
    rIds = [s.get(f'{{{NS_R}}}id') for s in list(sldIdLst)]
    for rId in rIds:
        for s in list(sldIdLst):
            if s.get(f'{{{NS_R}}}id') == rId:
                sldIdLst.remove(s)
        if rId in prs.part.rels:
            prs.part.rels.pop(rId)

    blank = prs.slide_layouts[6]

    for day in ['Monday','Tuesday','Wednesday']:
        data = DAYS[day]

        s1 = prs.slides.add_slide(blank)
        s1.background.fill.solid()
        s1.background.fill.fore_color.rgb = RGBColor(0xDE,0xEC,0xF8)
        build_m2i_slide(s1, day, data['m2i'])

        s2 = prs.slides.add_slide(blank)
        s2.background.fill.solid()
        s2.background.fill.fore_color.rgb = RGBColor(0xEA,0xF7,0xEA)
        build_i2m_slide(s2, day, data['i2m'])

        print(f'  {day} ✓')

    prs.save(OUT_FILE)
    print(f'\nSaved: {OUT_FILE}')

if __name__ == '__main__':
    print('Building fractions recap...')
    build()

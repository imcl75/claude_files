#!/usr/bin/env python3
"""
inject_lp_previews.py

Post-processes a teaching PPTX to embed LP previews into the four
placeholder slides: Learning Paper 1, Marking Station 1,
Learning Paper 2, Marking Station 2.

Usage:
    python3 inject_lp_previews.py <teaching_pptx> <lp_combined_pptx>

The LP Combined PPTX has three slides:
    Slide 1 — Standard LP   (top half = LP1, bottom half = LP2)
    Slide 2 — Adapted LP    (not used here)
    Slide 3 — Marking Station (top half = MS1, bottom half = MS2)

The cut line is at exactly 50% of the LP slide height.
The script finds teaching slides by title (robust to slide count changes).
It overwrites the teaching PPTX in place.
"""

import sys
import os
import subprocess
import tempfile
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

# Teaching slide label → resolved at inject time based on LP slide count
SLIDE_LABELS = ['Learning Paper 1', 'Marking Station 1', 'Learning Paper 2', 'Marking Station 2']

def resolve_slide_map(n_lp_slides):
    """
    Map teaching slide labels to (lp_slide_index, crop_hint).
    crop_hint is 'top', 'bottom', or 'full' — overridden by per-slide
    INJECT_REPS / INJECT_REP_FRAC metadata when present.

    3-slide LP (half-page, non-arithmetic):
      slide 0 = pupil sheet (LP1 top / LP2 bottom)
      slide 1 = adapted
      slide 2 = marking station (MS1 top / MS2 bottom)

    6-slide LP (arithmetic, full-page):
      slide 0 = LP1 pupil   slide 1 = LP2 pupil
      slide 2 = LP1 adapted slide 3 = LP2 adapted
      slide 4 = LP1 MS      slide 5 = LP2 MS
    """
    if n_lp_slides >= 6:
        return {
            'Learning Paper 1':   (0, 'rep'),
            'Marking Station 1':  (4, 'rep'),
            'Learning Paper 2':   (1, 'rep'),
            'Marking Station 2':  (5, 'rep'),
        }
    else:
        return {
            'Learning Paper 1':   (0, 'top'),
            'Marking Station 1':  (2, 'top'),
            'Learning Paper 2':   (0, 'bottom'),
            'Marking Station 2':  (2, 'bottom'),
        }


def read_lp_crop_metadata(lp_pptx):
    """
    Open LP PPTX and read INJECT_REPS / INJECT_REP_FRAC from speaker notes.
    Returns dict: slide_index -> (reps, rep_frac_or_None).
    """
    import re as _re
    try:
        from pptx import Presentation as _Prs
        prs = _Prs(lp_pptx)
    except Exception:
        return {}
    meta = {}
    for i, slide in enumerate(prs.slides):
        reps, frac = 1, None
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            rm = _re.search(r'INJECT_REPS:(\d+)', notes)
            fm = _re.search(r'INJECT_REP_FRAC:([\d.]+)', notes)
            if rm:
                reps = int(rm.group(1))
            if fm:
                frac = float(fm.group(1))
        meta[i] = (reps, frac)
    return meta

# Image placement on teaching slide (inches)
IMG_TOP   = 1.05
IMG_LEFT  = 0.4
IMG_RIGHT = 12.93
IMG_BOT   = 7.30


def get_slide_title(slide):
    """Return plain text of the title shape, or ''."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            if text in SLIDE_LABELS:
                return text
    return ''


def convert_lp_to_images(lp_pptx, tmpdir, dpi=200):
    """Convert LP PPTX to per-slide JPEG images via soffice + pdftoppm."""
    pdf_out = Path(tmpdir) / (Path(lp_pptx).stem + '.pdf')

    r1 = subprocess.run(
        ['soffice', '--headless', '--convert-to', 'pdf',
         '--outdir', str(tmpdir), str(lp_pptx)],
        capture_output=True, text=True
    )
    if r1.returncode != 0 or not pdf_out.exists():
        # Try alternative soffice path
        r1 = subprocess.run(
            ['/usr/bin/soffice', '--headless', '--convert-to', 'pdf',
             '--outdir', str(tmpdir), str(lp_pptx)],
            capture_output=True, text=True
        )
    if not pdf_out.exists():
        raise RuntimeError(f"soffice failed to produce PDF.\nstdout: {r1.stdout}\nstderr: {r1.stderr}")

    prefix = str(Path(tmpdir) / 'lp_slide')
    r2 = subprocess.run(
        ['pdftoppm', '-jpeg', '-r', str(dpi), str(pdf_out), prefix],
        capture_output=True, text=True
    )
    if r2.returncode != 0:
        raise RuntimeError(f"pdftoppm failed:\n{r2.stderr}")

    images = sorted(Path(tmpdir).glob('lp_slide-*.jpg'))
    if not images:
        images = sorted(Path(tmpdir).glob('lp_slide*.jpg'))
    return images


def smart_crop(img_path, crop_hint, slide_meta):
    """
    Crop LP slide image based on crop_hint and per-slide metadata.

    crop_hint:  'top' | 'bottom' | 'full' | 'rep'
    slide_meta: (reps, rep_frac) from read_lp_crop_metadata

    'rep' mode: crop to first repetition using rep_frac from speaker notes.
         Falls back to 'full' if no metadata.
    'top' / 'bottom': traditional half-page crop (non-arithmetic LPs).
    'full': show entire slide.
    """
    img = Image.open(img_path)
    w, h = img.size

    if crop_hint == 'rep':
        reps, frac = slide_meta if slide_meta else (1, None)
        if reps > 1 and frac is not None:
            crop_h = int(h * frac)
            return img.crop((0, 0, w, crop_h))
        return img  # single rep or no metadata → full slide

    mid = h // 2
    if crop_hint == 'top':
        return img.crop((0, 0, w, mid))
    if crop_hint == 'bottom':
        return img.crop((0, mid, w, h))
    return img  # 'full'


# Keep old name for any legacy callers
def crop_half(img_path, half):
    return smart_crop(img_path, half, None)


def clear_slide_content(slide):
    """Remove all shapes except the title from a slide."""
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text in SLIDE_LABELS:
                continue  # keep the title
        to_remove.append(shape)
    sp_tree = slide.shapes._spTree
    for shape in to_remove:
        sp_tree.remove(shape._element)


def add_image_to_slide(slide, img_pil, slide_w_emu, slide_h_emu):
    """Add a PIL image to the slide, scaled to fill the usable area below the title."""
    import io
    buf = io.BytesIO()
    img_pil.save(buf, format='JPEG', quality=92)
    buf.seek(0)

    # Available area
    left_emu  = Inches(IMG_LEFT)
    top_emu   = Inches(IMG_TOP)
    right_emu = Inches(IMG_RIGHT)
    bot_emu   = Inches(IMG_BOT)
    avail_w   = right_emu - left_emu
    avail_h   = bot_emu - top_emu

    # Scale image to fit, preserving aspect ratio
    img_w, img_h = img_pil.size
    scale = min(avail_w / img_w, avail_h / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)

    # Centre within available area
    cx = left_emu + (avail_w - new_w) // 2
    cy = top_emu  + (avail_h - new_h) // 2

    slide.shapes.add_picture(buf, cx, cy, new_w, new_h)


def inject(teaching_pptx, lp_pptx):
    print(f"Teaching:  {teaching_pptx}")
    print(f"LP file:   {lp_pptx}")

    prs = Presentation(teaching_pptx)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    with tempfile.TemporaryDirectory() as tmpdir:
        print("Converting LP to images...")
        lp_images = convert_lp_to_images(lp_pptx, tmpdir)
        print(f"  Got {len(lp_images)} LP slide image(s)")

        if len(lp_images) < 3:
            raise RuntimeError(f"Expected at least 3 LP slide images, got {len(lp_images)}")

        slide_map    = resolve_slide_map(len(lp_images))
        lp_crop_meta = read_lp_crop_metadata(lp_pptx)
        injected     = 0
        for slide in prs.slides:
            title = get_slide_title(slide)
            if title not in slide_map:
                continue

            lp_idx, crop_hint = slide_map[title]
            slide_meta = lp_crop_meta.get(lp_idx)
            img = smart_crop(lp_images[lp_idx], crop_hint, slide_meta)

            clear_slide_content(slide)
            add_image_to_slide(slide, img, slide_w, slide_h)
            injected += 1
            reps = slide_meta[0] if slide_meta else 1
            desc = f"1/{reps} crop" if reps > 1 else crop_hint
            print(f"  Injected '{title}' ← LP slide {lp_idx+1} ({desc})")

    prs.save(teaching_pptx)
    print(f"\nInjected {injected} LP previews.")
    print(f"Saved: {teaching_pptx}")
    return injected


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python3 inject_lp_previews.py <teaching_pptx> <lp_combined_pptx>")
        sys.exit(1)
    inject(sys.argv[1], sys.argv[2])

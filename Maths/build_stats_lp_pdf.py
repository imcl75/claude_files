"""
build_stats_lp_pdf.py
Generates T6W5 Statistics LPs as A4 PDF using ReportLab.
Heights are calculated from actual text wrap, not guessed.
Produces 3 pages: Standard, Adapted, Marking Station.
Also creates a 3-slide wrapper PPTX for inject_lp_previews.py.
"""
import sys, os
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.platypus import Paragraph
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from PIL import Image as PILImage

# ── Register fonts ────────────────────────────────────────────────────────────
FONT_DIR = '/usr/share/fonts/truetype/liberation'
pdfmetrics.registerFont(TTFont('LibSans',     f'{FONT_DIR}/LiberationSans-Regular.ttf'))
pdfmetrics.registerFont(TTFont('LibSansBold', f'{FONT_DIR}/LiberationSans-Bold.ttf'))

# ── Page constants (points) ───────────────────────────────────────────────────
W, H    = A4                   # 595 × 841.9
M       = 18                   # page margin
CUT_Y   = H / 2               # 420.9  — cut line between LP1 and LP2
GAP     = 8                    # gap between elements

# Chart column (right side) — generous width for readability
CHART_W = 265
CHART_X = W - M - CHART_W     # 312

# Question column (left side)
Q_X     = M
Q_W     = CHART_X - Q_X - 12  # 282  (12pt gutter before chart)

# Learning Label dimensions (match build_lp_v3.js constants)
CM           = 28.35           # 1cm in points
LABEL_SCALE  = 0.72 * 0.85
LL_W         = 9.7  * CM * LABEL_SCALE   # 168pt
LL_H         = 4.24 * CM * LABEL_SCALE   # 73.5pt
LL_X         = M
ICON_W       = 18.7            # 0.26" in pt
ICON_H       = ICON_W * (103/120)
ICON_PATH    = '/tmp/claude_work/lp_assets/mathematician_icon.png'
CHART_DIR    = '/tmp/wfa_stats_charts'

# Colours (RGB 0-1)
WFA_BLUE    = (0.09, 0.60, 0.83)
Q_BG        = (0.87, 0.92, 0.95)
Q_TXT       = (0.12, 0.31, 0.47)
Q_BORDER    = (0.08, 0.38, 0.51)
GRN         = (0.10, 0.36, 0.16)
GRN_BG      = (0.91, 0.96, 0.91)
GREY        = (0.75, 0.75, 0.75)
GREY_TXT    = (0.20, 0.20, 0.20)
DARK        = (0.10, 0.10, 0.10)

# ── Paragraph styles ──────────────────────────────────────────────────────────
Q_STYLE = ParagraphStyle('Q',
    fontName='LibSansBold', fontSize=8.5, leading=12,
    textColor=(0.12, 0.31, 0.47),
    leftIndent=6, rightIndent=6, spaceBefore=0, spaceAfter=0)

INTRO_STYLE = ParagraphStyle('I',
    fontName='LibSansBold', fontSize=9, leading=13,
    textColor=DARK, leftIndent=0, rightIndent=0)

ANS_STYLE = ParagraphStyle('A',
    fontName='LibSansBold', fontSize=8.5, leading=12,
    textColor=GRN, leftIndent=6, rightIndent=6)

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_fill(c, rgb):   c.setFillColorRGB(*rgb)
def set_stroke(c, rgb): c.setStrokeColorRGB(*rgb)

def filled_rect(c, x, y_bottom, w, h, fill_rgb, stroke_rgb=None, lw=0.5):
    """Draw a rectangle. y_bottom is the BOTTOM edge (ReportLab bottom-up)."""
    set_fill(c, fill_rgb)
    if stroke_rgb:
        set_stroke(c, stroke_rgb)
        c.setLineWidth(lw)
        c.rect(x, y_bottom, w, h, fill=1, stroke=1)
    else:
        c.rect(x, y_bottom, w, h, fill=1, stroke=0)

def outline_rect(c, x, y_bottom, w, h, stroke_rgb=GREY, lw=0.5):
    """Draw an empty box."""
    set_fill(c, (1,1,1))
    set_stroke(c, stroke_rgb)
    c.setLineWidth(lw)
    c.rect(x, y_bottom, w, h, fill=1, stroke=1)

def draw_chart(c, chart_file, x, y_bottom, w, h):
    """Draw chart image, maintaining aspect ratio, centred in available space."""
    path = os.path.join(CHART_DIR, chart_file)
    if not os.path.exists(path):
        return
    img = PILImage.open(path)
    iw, ih = img.size
    aspect = iw / ih
    # Fit within (w × h) maintaining aspect
    if w / h > aspect:
        draw_h = h
        draw_w = h * aspect
    else:
        draw_w = w
        draw_h = w / aspect
    # Centre
    dx = (w - draw_w) / 2
    dy = (h - draw_h) / 2
    c.drawImage(path, x + dx, y_bottom + dy, draw_w, draw_h,
                preserveAspectRatio=True, mask='auto')

# ── Learning Label ────────────────────────────────────────────────────────────
def draw_ll(c, x, y_top, date, topic, lf, ican):
    """Draw the school LL sticker at (x, y_top). Returns actual height used."""
    PAD    = 3
    NARROW = LL_W - ICON_W - PAD*3
    FULL   = LL_W - PAD*2
    ty = y_top - PAD        # current y baseline, stepping downward

    # Icon — top right
    if os.path.exists(ICON_PATH):
        c.drawImage(ICON_PATH,
                    x + LL_W - ICON_W - PAD,
                    y_top - PAD - ICON_H,
                    ICON_W, ICON_H, mask='auto')

    # Date — 7pt
    c.setFont('LibSans', 7)
    set_fill(c, GREY_TXT)
    ty -= 7 * 0.72
    c.drawString(x + PAD, ty, date)
    ty -= 7 * 0.28 + 2

    # Topic — 9pt bold (underlined)
    c.setFont('LibSansBold', 9)
    set_fill(c, DARK)
    ty -= 9 * 0.72
    c.drawString(x + PAD, ty, topic)
    # Underline
    tw = c.stringWidth(topic, 'LibSansBold', 9)
    c.setLineWidth(0.5); set_stroke(c, DARK)
    c.line(x + PAD, ty - 1, x + PAD + tw, ty - 1)
    ty -= 9 * 0.28 + 3

    # LF — 7pt (may wrap)
    c.setFont('LibSans', 7)
    set_fill(c, GREY_TXT)
    lf_para = Paragraph(lf, ParagraphStyle('LF',
        fontName='LibSans', fontSize=7, leading=9,
        textColor=(0.2,0.2,0.2)))
    lf_w, lf_h = lf_para.wrap(FULL, 200)
    ty -= lf_h
    lf_para.drawOn(c, x + PAD, ty)
    ty -= 3

    # I can statements — 6.5pt
    for ic in ican:
        c.setFont('LibSans', 6.5)
        set_fill(c, GREY_TXT)
        ic_para = Paragraph(ic, ParagraphStyle('IC',
            fontName='LibSans', fontSize=6.5, leading=8.5,
            textColor=(0.2,0.2,0.2)))
        ic_w, ic_h = ic_para.wrap(FULL, 200)
        ty -= ic_h
        ic_para.drawOn(c, x + PAD, ty)
        ty -= 2

    return y_top - ty   # total height used

# ── Question block ────────────────────────────────────────────────────────────
def draw_questions(c, x, y_start, w, avail_h, intro, questions, is_ms):
    """
    Draw intro + questions with answer boxes.
    y_start: top of the question area (in bottom-up coords).
    avail_h: total height available.
    Returns nothing — all layout is calculated internally.
    """
    PAD = 6

    # Step 1: measure intro
    intro_para = Paragraph(intro, INTRO_STYLE)
    intro_w, intro_h = intro_para.wrap(w, 200)
    intro_h += 6   # bottom gap

    # Step 2: measure each question text
    q_paras = []
    q_heights = []
    for i, (qt, ans) in enumerate(questions):
        p = Paragraph(f'{i+1}.  {qt}', Q_STYLE)
        pw, ph = p.wrap(w - PAD*2, 200)
        box_h = ph + PAD*2   # question box height = text + padding top+bottom
        q_paras.append((p, box_h))
        q_heights.append(box_h)

    n = len(questions)
    ITEM_GAP = 6

    # Step 3: calculate answer box heights dynamically
    used = intro_h + sum(q_heights) + n * ITEM_GAP
    remaining = avail_h - used - 12   # 12pt bottom margin
    ans_h_each = max(28, remaining / n)  # at least 28pt per answer box

    # Step 4: draw everything top-down
    cy = y_start

    # Intro text
    intro_para.drawOn(c, x, cy - intro_h + 6)
    cy -= intro_h

    for i, ((p, qbox_h), (qt, ans)) in enumerate(zip(q_paras, questions)):
        # Question box
        qbox_bottom = cy - qbox_h
        filled_rect(c, x, qbox_bottom, w, qbox_h,
                    fill_rgb=Q_BG, stroke_rgb=Q_BORDER, lw=0.6)
        set_fill(c, (0,0,0))
        p.drawOn(c, x + PAD, qbox_bottom + PAD)
        cy = qbox_bottom - 3

        # Answer box
        ans_bottom = cy - ans_h_each
        if is_ms:
            # Green filled answer box
            filled_rect(c, x, ans_bottom, w, ans_h_each,
                        fill_rgb=GRN_BG, stroke_rgb=GRN, lw=0.6)
            ap = Paragraph(f'✓  {ans}', ANS_STYLE)
            ap_w, ap_h = ap.wrap(w - PAD*2, ans_h_each - PAD*2)
            ap.drawOn(c, x + PAD, ans_bottom + PAD)
        else:
            # Empty answer box (light border only)
            outline_rect(c, x, ans_bottom, w, ans_h_each, stroke_rgb=GREY, lw=0.5)

        cy = ans_bottom - ITEM_GAP

# ── Half-page builder ─────────────────────────────────────────────────────────
def draw_half(c, lp_data, region_top, region_bot, meta, show_ll, is_ms):
    """
    region_top, region_bot: y coordinates (bottom-up) bounding this LP half.
    show_ll: True for LP1 only.
    """
    PAD = 10

    # Draw chart — right column, full height of region
    chart_h = region_top - region_bot - 2*PAD
    draw_chart(c, lp_data['chart'], CHART_X, region_bot + PAD,
               CHART_W, chart_h)

    # Chart label below (centred)
    c.setFont('LibSans', 7)
    set_fill(c, GREY_TXT)
    label = lp_data.get('chart_label', '')
    tw = c.stringWidth(label, 'LibSans', 7)
    c.drawString(CHART_X + (CHART_W - tw)/2, region_bot + 2, label)

    # Learning Label (LP1 only)
    q_top = region_top - PAD
    if show_ll:
        ll_h = draw_ll(c, LL_X, region_top - PAD,
                       meta['date'], meta['topic'], meta['lf'], meta['ican'])
        q_top = region_top - PAD - ll_h - 8

    # Questions — left column
    q_bot    = region_bot + PAD
    q_avail  = q_top - q_bot
    draw_questions(c, Q_X, q_top, Q_W, q_avail,
                   lp_data['intro'], lp_data['qs'], is_ms)

# ── Cut line ──────────────────────────────────────────────────────────────────
def draw_cut_line(c):
    set_stroke(c, (0.7, 0.7, 0.7))
    c.setLineWidth(0.5)
    c.setDash([4, 4], 0)
    c.line(0, CUT_Y, W, CUT_Y)
    c.setDash()
    # Scissors symbol
    c.setFont('LibSans', 8)
    set_fill(c, (0.6, 0.6, 0.6))
    c.drawString(4, CUT_Y + 2, '✂')

# ── LP data ───────────────────────────────────────────────────────────────────
LP = {
17: {
    'date': '29/06/2026', 'topic': 'Statistics',
    'lf':   'LF: To read and interpret data from pictograms, bar charts and tables.',
    'ican': ['I can read values from a pictogram using the key.',
             'I can read values from a bar chart and a two-way table.'],
    'lp1': {
        'chart': 'c1_ido1_pictogram.png',
        'chart_label': 'Favourite sports in Year 4',
        'intro': 'Use the pictogram to answer the questions.',
        'qs': [
            ('How many pupils chose basketball?',
             '6 pupils  (3 symbols × 2)'),
            ('How many pupils chose football and volleyball altogether?',
             '12 + 8 = 20 pupils'),
            ('How many MORE pupils chose swimming than capoeira?',
             '10 − 4 = 6 more pupils'),
        ],
    },
    'lp2': {
        'chart': 'c2_ido1_table.png',
        'chart_label': 'Average daily sunshine hours',
        'intro': 'Use the table to answer the questions.',
        'qs': [
            ('How many hours of sunshine does England get in winter?',
             '2 hours'),
            ('What is the total sunshine across both countries in spring?',
             '5 + 7 = 12 hours'),
            ('Which country has more sunshine overall? By how many hours?',
             'Brazil — 31 vs 19 = 12 more hours'),
            ('In which season are England and Brazil most similar?',
             'Summer — just 1 hour apart (8 vs 9)'),
        ],
    },
},
18: {
    'date': '30/06/2026', 'topic': 'Statistics',
    'lf':   'LF: To calculate sums and differences from charts and compare two data sets.',
    'ican': ['I can add and subtract values read from a bar chart.',
             'I can compare two data sets and describe differences.'],
    'lp1': {
        'chart': 'c1_ido1_bar_chart.png',
        'chart_label': 'Animals counted at an Amazon river each day',
        'intro': 'Use the bar chart to answer the questions.',
        'qs': [
            ('How many animals were counted on Tuesday and Wednesday altogether?',
             '24 + 15 = 39 animals'),
            ('How many MORE were counted on Thursday than Monday?',
             '30 − 18 = 12 more animals'),
            ('What is the total for all five days?',
             '18 + 24 + 15 + 30 + 21 = 108 animals'),
        ],
    },
    'lp2': {
        'chart': 'c2_ido1_double_bar.png',
        'chart_label': 'Average temperature — London and Rio de Janeiro',
        'intro': 'Use the double bar chart to answer the questions.',
        'qs': [
            ('What is the temperature difference between London and Rio in winter?',
             '22 − 6 = 16°C'),
            ('In which season are the two cities closest in temperature?',
             'Summer — gap of 10°C'),
            ("What is London's temperature range across the year?",
             '18 − 6 = 12°C range'),
            ('Write a statement comparing the two cities using data from the chart.',
             'Any valid statement using values from the chart.'),
        ],
    },
},
19: {
    'date': '01/07/2026', 'topic': 'Statistics',
    'lf':   'LF: To read and interpret line graphs, including estimating between points.',
    'ican': ['I can read values from a line graph at labelled points.',
             'I can estimate values between two labelled points on a line graph.'],
    'lp1': {
        'chart': 'c1_ido1_line_graph.png',
        'chart_label': 'Temperature in Bristol on a June day',
        'intro': 'Use the line graph to answer the questions.',
        'qs': [
            ('What was the temperature at 10:00?',
             '17°C'),
            ('At what time was Bristol warmest? What was the temperature?',
             '14:00 — 24°C'),
            ('Between which two readings did the temperature rise the most?',
             '10:00 to 12:00 — rose by 4°C'),
        ],
    },
    'lp2': {
        'chart': 'c2_ido1_line_graph.png',
        'chart_label': 'Temperature in São Paulo on a June day',
        'intro': 'Use the São Paulo line graph to answer the questions.',
        'qs': [
            ('What was the temperature at 12:00?',
             '30°C'),
            ('Estimate the temperature at 09:00.',
             'About 24°C — halfway between 22°C and 26°C'),
            ('Estimate the temperature at 11:00.',
             'About 28°C — halfway between 26°C and 30°C'),
            ('At approximately what time was it 30°C on the way back down?',
             'Around 15:00 — the line falls from 32°C to 29°C'),
        ],
    },
},
}

# ── T6W6 LP data ──────────────────────────────────────────────────────────────
LP.update({
20: {
    'date': '06/07/2026', 'topic': 'Statistics',
    'lf':   'LF: To draw an accurate line graph from a given data set.',
    'ican': ['I can plot data points accurately at the correct position on given axes.',
             'I can join consecutive plotted points with straight lines and write a title.'],
    'lp1': {
        'lp_type': 'draw_graph',
        'intro': 'Use the data in the table to draw a line graph in the space on the right.',
        'table_headers': ['Month', 'Rainfall (mm)'],
        'table_rows': [['Jan', 89], ['Feb', 71], ['Mar', 58],
                       ['Apr', 48], ['May', 52], ['Jun', 45]],
        'ax': {
            'x_labels': ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'],
            'y_min': 0, 'y_max': 100, 'y_step': 10,
            'x_label': 'Month', 'y_label': 'Rainfall (mm)',
            'title_prompt': 'Give your graph a title:',
            'ms_values': [89, 71, 58, 48, 52, 45],
        },
        'instructions': [
            '1.  Check the scale — what does each grid line represent?',
            '2.  Plot each cross (×) at the correct height.',
            '3.  Join the points with straight lines — one segment at a time.',
            '4.  Write a title and label both axes.',
        ],
        'qs': [
            ('What is the wettest month shown?',  'January — 89mm'),
            ('How much less rain falls in June than in January?',
             '89 − 45 = 44mm less'),
        ],
    },
    'lp2': {
        'lp_type': 'draw_graph',
        'intro': 'Use the data below to draw a line graph.  Choose your own scale.',
        'table_headers': ['Month', 'Rainfall (mm)'],
        'table_rows': [['Jan', 130], ['Feb', 120], ['Mar', 100],
                       ['Apr', 80],  ['May', 50],  ['Jun', 40]],
        'ax': {
            'x_labels': ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'],
            'y_min': 0, 'y_max': 150, 'y_step': 25,
            'x_label': 'Month', 'y_label': 'Rainfall (mm)',
            'title_prompt': 'Give your graph a title:',
            'ms_values': [130, 120, 100, 80, 50, 40],
        },
        'instructions': [
            '1.  Find the highest value — what scale do you need?',
            '2.  Label the y-axis in equal steps.',
            '3.  Plot each point carefully and join with straight lines.',
        ],
        'qs': [
            ('Which month has the most rainfall?',     'January — 130mm'),
            ('Describe the trend from January to June.',
             'Rainfall falls every month — a clear downward trend.'),
        ],
    },
},
21: {
    'date': '08/07/2026', 'topic': 'Statistics',
    'lf':   'LF: To draw an accurate bar chart from a given data set.',
    'ican': ['I can draw bars to the correct height using a given scale.',
             'I can set up axes with a suitable scale and draw all bars from a data table.'],
    'lp1': {
        'lp_type': 'draw_bar',
        'intro': 'Three bars have been drawn for you. Add the missing bars, then answer the questions.',
        'table_headers': ['Minibeast', 'Number found'],
        'table_rows': [['Worm', 8], ['Beetle', 5], ['Snail', 4], ['Spider', 12], ['Ant', 19]],
        'ax': {
            'categories': ['Worm', 'Beetle', 'Snail', 'Spider', 'Ant'],
            'y_min': 0, 'y_max': 20, 'y_step': 2,
            'x_label': 'Minibeast', 'y_label': 'Number found',
            'title_prompt': 'Give your chart a title:',
            'partial_values': [None, 5, 4, 12, None],
            'ms_values':      [8,    5, 4, 12, 19],
        },
        'qs': [
            ('Which minibeast was found most often?',  'Ants — 19'),
            ('How many more spiders than snails were found?', '12 \u2212 4 = 8 more'),
        ],
    },
    'lp2': {
        'lp_type': 'draw_bar',
        'intro': 'Use the data table to draw the bar chart from scratch. Then answer the questions.',
        'table_headers': ['How they travel', 'Number of pupils'],
        'table_rows': [['Walk', 10], ['Cycle', 4], ['Car', 8], ['Bus', 3], ['Scooter', 1]],
        'ax': {
            'categories': ['Walk', 'Cycle', 'Car', 'Bus', 'Scooter'],
            'y_min': 0, 'y_max': 12, 'y_step': 2,
            'x_label': 'How they travel', 'y_label': 'Number of pupils',
            'title_prompt': 'Give your chart a title:',
            'partial_values': [None, None, None, None, None],
            'ms_values':      [10,   4,    8,   3,    1],
        },
        'qs': [
            ('How many more pupils walk than travel by car?', '10 \u2212 8 = 2 more'),
            ('How many pupils are in the class altogether?',  '10+4+8+3+1 = 26'),
        ],
    },
},
22: {
    'date': '09/07/2026', 'topic': 'Statistics',
    'lf':   'LF: To evaluate claims about data using evidence from charts and tables.',
    'ican': ['I can decide whether a claim is true or false using values from a chart.',
             'I can write a verdict with evidence: "I [agree/disagree] because…"'],
    'lp1': {
        'chart': 'c1_ido1_double_bar.png',
        'chart_label': 'Average monthly rainfall (mm) — Bristol vs Manaus',
        'intro': 'For each claim: write Agree, Disagree or Not enough information. Then show your evidence.',
        'qs': [
            ('Claim: "Manaus gets more rain than Bristol in every month shown."\nAgree / Disagree / Not enough information?',
             'Agree — every Manaus bar is taller. e.g. Jan: 260 > 89.'),
            ('Claim: "Bristol and Manaus get a similar amount of rain in June."\nAgree / Disagree / Not enough information?',
             'Disagree — Bristol 45mm, Manaus 100mm. 100 \u2212 45 = 55mm difference.'),
            ('Claim: "Manaus gets more than 200mm of rain in every month from Jan to Jun."\nCheck using the chart.',
             'Disagree — June is only 100mm, which is less than 200mm.'),
        ],
    },
    'lp2': {
        'chart': 'c2_ido1_table.png',
        'chart_label': 'Hours of daylight — Bristol and S\u00e3o Paulo, by season',
        'intro': 'For each claim: write Agree, Disagree or Not enough information. Show your evidence.',
        'qs': [
            ('Claim: "Bristol always has fewer hours of daylight than S\u00e3o Paulo."\nAgree / Disagree / Not enough information?',
             'Disagree — in summer Bristol has 17h, S\u00e3o Paulo 11h. Bristol has MORE.'),
            ('Claim: "The total daylight hours are the same for both cities."\nCalculate to check.',
             'Disagree — Bristol: 8+12+17+12 = 49h. S\u00e3o Paulo: 11+12+11+12 = 46h.'),
            ('Claim: "S\u00e3o Paulo has the same hours of daylight in spring and autumn."\nAgree / Disagree?',
             'Agree — S\u00e3o Paulo: Spring 12h, Autumn 12h. They are the same.'),
        ],
    },
},
23: {
    'date': '11/07/2026', 'topic': 'Statistics',
    'lf':   'LF: To complete the full data cycle \u2014 tally, represent, draw and analyse.',
    'ican': ['I can tally a list of data and choose a suitable representation.',
             'I can draw a bar chart and answer questions from my own chart.'],
    'lp1': {
        'lp_type': 'tally_draw',
        'intro': 'Tally the responses below, then draw a bar chart and answer the questions.',
        'raw_data': (
            'Maths, Science, Maths, Art, Maths, English, Science, Maths, Art, English, '
            'Maths, Science, Maths, Art, English, Maths, Science, English, Maths, Maths, '
            'Science, Art, English, Maths, Science, Art'
        ),
        'tally_categories': ['Maths', 'Science', 'English', 'Art', 'PE'],
        'ms_tallies': [8, 6, 5, 4, 3],
        'ax': {
            'categories': ['Maths', 'Science', 'English', 'Art', 'PE'],
            'y_min': 0, 'y_max': 10, 'y_step': 2,
            'x_label': 'Favourite subject', 'y_label': 'Number of pupils',
            'title_prompt': 'Give your chart a title:',
            'ms_values': [8, 6, 5, 4, 3],
        },
        'qs': [
            ('Which subject did the most pupils choose?',
             'Maths \u2014 8 pupils'),
            ('How many more pupils chose Maths than Art?',
             '8 \u2212 4 = 4 more'),
            ('Why is a bar chart a better choice than a line graph for this data?',
             'Subjects are separate categories, not continuous over time.'),
        ],
    },
    'lp2': {
        'lp_type': 'tally_draw',
        'intro': 'Tally the responses below, then draw a bar chart and answer the questions.',
        'raw_data': (
            'Football, Gymnastics, Football, Swimming, Football, Athletics, Gymnastics, '
            'Football, Athletics, Swimming, Football, Gymnastics, Football, Athletics, '
            'Swimming, Football, Gymnastics, Swimming, Football, Football, Gymnastics, '
            'Athletics, Swimming, Football, Gymnastics, Cycling'
        ),
        'tally_categories': ['Football', 'Gymnastics', 'Swimming', 'Athletics', 'Cycling'],
        'ms_tallies': [10, 6, 5, 4, 1],
        'ax': {
            'categories': ['Football', 'Gymnastics', 'Swimming', 'Athletics', 'Cycling'],
            'y_min': 0, 'y_max': 12, 'y_step': 2,
            'x_label': 'Favourite sport', 'y_label': 'Number of pupils',
            'title_prompt': 'Give your chart a title:',
            'ms_values': [10, 6, 5, 4, 1],
        },
        'qs': [
            ('Which sport did most pupils choose?',
             'Football \u2014 10 pupils'),
            ('How many pupils were surveyed in total?',
             '10+6+5+4+1 = 26'),
            ('If you used a pictogram with key = 2, how many symbols would Football need?',
             '5 whole symbols  (10 \u00f7 2 = 5)'),
        ],
    },
},
})

DAYS = {17: 'Monday',  18: 'Tuesday',  19: 'Wednesday',
        20: 'Monday',  21: 'Wednesday', 22: 'Thursday',  23: 'Friday'}

WEEK = {17: 'T6W5', 18: 'T6W5', 19: 'T6W5',
        20: 'T6W6', 21: 'T6W6', 22: 'T6W6', 23: 'T6W6'}

# ── Blank-axes LP (for drawing lessons) ───────────────────────────────────────
def draw_blank_axes_area(c, ax, x, y_bottom, w, h, is_ms):
    """
    Draw blank graph paper with pre-labelled axes.
    ax: dict with x_labels, y_min, y_max, y_step, x_label, y_label, ms_values.
    x, y_bottom: bottom-left of the drawing area (bottom-up coords).
    """
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    import numpy as np
    import io

    x_labels = ax['x_labels']
    y_min    = ax.get('y_min', 0)
    y_max    = ax['y_max']
    y_step   = ax['y_step']
    x_label  = ax.get('x_label', '')
    y_label  = ax.get('y_label', '')
    ms_vals  = ax.get('ms_values', [])

    n_x = len(x_labels)
    n_y = int((y_max - y_min) / y_step)

    fig, ax_plt = plt.subplots(figsize=(3.2, 3.0), dpi=180)
    ax_plt.set_facecolor('white')
    fig.patch.set_facecolor('white')

    ax_plt.set_xlim(-0.5, n_x - 0.5)
    ax_plt.set_ylim(y_min, y_max)
    ax_plt.set_xticks(range(n_x))
    ax_plt.set_xticklabels(x_labels, fontsize=7)
    ax_plt.set_yticks(range(y_min, y_max + 1, y_step))
    ax_plt.tick_params(axis='both', labelsize=7)
    ax_plt.set_xlabel(x_label, fontsize=7, labelpad=2)
    ax_plt.set_ylabel(y_label, fontsize=7, labelpad=2)
    ax_plt.grid(True, linewidth=0.4, color='#AAAAAA', alpha=0.6)
    ax_plt.spines['top'].set_visible(False)
    ax_plt.spines['right'].set_visible(False)
    ax_plt.spines['left'].set_linewidth(0.8)
    ax_plt.spines['bottom'].set_linewidth(0.8)

    if is_ms and ms_vals and len(ms_vals) == n_x:
        ax_plt.plot(range(n_x), ms_vals,
                    color='#C83030', linewidth=1.2, marker='x',
                    markersize=5, markeredgewidth=1.5, zorder=3)

    buf = io.BytesIO()
    plt.tight_layout(pad=0.3)
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor='white')
    plt.close(fig)
    buf.seek(0)

    import tempfile
    tmp = tempfile.mktemp(suffix='.png')
    with open(tmp, 'wb') as f:
        f.write(buf.read())

    img = PILImage.open(tmp)
    iw, ih = img.size
    aspect = iw / ih
    draw_h = h
    draw_w = h * aspect
    if draw_w > w:
        draw_w = w
        draw_h = w / aspect
    dx = (w - draw_w) / 2
    dy = (h - draw_h) / 2
    c.drawImage(tmp, x + dx, y_bottom + dy, draw_w, draw_h,
                preserveAspectRatio=True, mask='auto')
    os.unlink(tmp)


def draw_blank_bar_chart_area(c, ax, x, y_bottom, w, h, is_ms):
    """Blank (or completed) bar chart axes for pupil drawing."""
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import numpy as np
    import io, tempfile

    cats     = ax['categories']
    y_min    = ax.get('y_min', 0)
    y_max    = ax['y_max']
    y_step   = ax['y_step']
    x_label  = ax.get('x_label', '')
    y_label  = ax.get('y_label', '')
    partial  = ax.get('partial_values', [None]*len(cats))
    ms_vals  = ax.get('ms_values', [])
    n = len(cats)

    fig, ax_plt = plt.subplots(figsize=(3.2, 3.0), dpi=180)
    ax_plt.set_facecolor('white')
    fig.patch.set_facecolor('white')
    ax_plt.set_ylim(y_min, y_max)
    ax_plt.set_xticks(range(n))
    ax_plt.set_xticklabels(cats, fontsize=6, rotation=15, ha='right')
    ax_plt.set_yticks(range(y_min, y_max + 1, y_step))
    ax_plt.tick_params(axis='both', labelsize=7)
    ax_plt.set_xlabel(x_label, fontsize=7, labelpad=2)
    ax_plt.set_ylabel(y_label, fontsize=7, labelpad=2)
    ax_plt.grid(axis='y', linewidth=0.4, color='#AAAAAA', alpha=0.6)
    ax_plt.spines['top'].set_visible(False)
    ax_plt.spines['right'].set_visible(False)
    ax_plt.spines['left'].set_linewidth(0.8)
    ax_plt.spines['bottom'].set_linewidth(0.8)

    if is_ms and ms_vals:
        ax_plt.bar(range(n), ms_vals, color='#1798d3', alpha=0.85,
                   edgecolor='#0f6fa0', linewidth=0.6, width=0.6)
    else:
        for i, v in enumerate(partial):
            if v is not None:
                ax_plt.bar(i, v, color='#1798d3', alpha=0.85,
                           edgecolor='#0f6fa0', linewidth=0.6, width=0.6)

    buf = io.BytesIO()
    plt.tight_layout(pad=0.3)
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    buf.seek(0)
    tmp = tempfile.mktemp(suffix='.png')
    with open(tmp, 'wb') as f:
        f.write(buf.read())

    img = PILImage.open(tmp)
    iw, ih = img.size
    aspect = iw / ih
    draw_h = h
    draw_w = h * aspect
    if draw_w > w:
        draw_w = w
        draw_h = w / aspect
    dx = (w - draw_w) / 2
    dy = (h - draw_h) / 2
    c.drawImage(tmp, x + dx, y_bottom + dy, draw_w, draw_h,
                preserveAspectRatio=True, mask='auto')
    os.unlink(tmp)


def draw_half_bar(c, lp_data, region_top, region_bot, meta, show_ll, is_ms):
    """Half-page LP for bar chart drawing lessons."""
    PAD = 10
    q_top = region_top - PAD
    if show_ll:
        ll_h = draw_ll(c, LL_X, region_top - PAD,
                       meta['date'], meta['topic'], meta['lf'], meta['ican'])
        q_top = region_top - PAD - ll_h - 8

    intro_para = Paragraph(lp_data['intro'], INTRO_STYLE)
    _, intro_h = intro_para.wrap(Q_W, 200)
    intro_h += 4
    cy = q_top
    intro_para.drawOn(c, Q_X, cy - intro_h)
    cy -= intro_h + 4

    # Data table
    hdrs = lp_data['table_headers']
    rows = lp_data['table_rows']
    ROW_H = 13
    COL_W = [Q_W * 0.55, Q_W * 0.35]
    TBL_W = sum(COL_W)
    filled_rect(c, Q_X, cy - ROW_H, TBL_W, ROW_H,
                fill_rgb=(0.09, 0.60, 0.83), stroke_rgb=None)
    for ci, hdr in enumerate(hdrs):
        x_off = Q_X + sum(COL_W[:ci])
        c.setFont('LibSansBold', 7.5); set_fill(c, (1, 1, 1))
        c.drawString(x_off + 3, cy - ROW_H + 3, hdr)
    cy -= ROW_H
    for ri, row in enumerate(rows):
        fill = (0.93, 0.97, 1.0) if ri % 2 == 0 else (1, 1, 1)
        filled_rect(c, Q_X, cy - ROW_H, TBL_W, ROW_H,
                    fill_rgb=fill, stroke_rgb=(0.75, 0.75, 0.75), lw=0.4)
        for ci, val in enumerate(row):
            x_off = Q_X + sum(COL_W[:ci])
            c.setFont('LibSans', 7.5); set_fill(c, (0.1, 0.1, 0.1))
            c.drawString(x_off + 3, cy - ROW_H + 3, str(val))
        cy -= ROW_H
    cy -= 6

    for qi, (qt, ans) in enumerate(lp_data.get('qs', [])):
        qp = Paragraph(f'{qi+1}.  {qt}', Q_STYLE)
        _, qph = qp.wrap(Q_W - 6, 200)
        qbox_h = qph + 8
        filled_rect(c, Q_X, cy - qbox_h, Q_W, qbox_h,
                    fill_rgb=Q_BG, stroke_rgb=Q_BORDER, lw=0.6)
        set_fill(c, (0, 0, 0)); qp.drawOn(c, Q_X + 3, cy - qbox_h + 4)
        cy -= qbox_h + 3
        ANS_H = 18
        if is_ms:
            filled_rect(c, Q_X, cy - ANS_H, Q_W, ANS_H,
                        fill_rgb=GRN_BG, stroke_rgb=GRN, lw=0.6)
            ap = Paragraph(f'\u2713  {ans}', ANS_STYLE)
            ap.wrap(Q_W - 6, ANS_H - 4); ap.drawOn(c, Q_X + 3, cy - ANS_H + 3)
        else:
            outline_rect(c, Q_X, cy - ANS_H, Q_W, ANS_H, stroke_rgb=GREY, lw=0.5)
        cy -= ANS_H + 4

    prompt = lp_data['ax'].get('title_prompt', 'Title:')
    c.setFont('LibSans', 7); set_fill(c, (0.3, 0.3, 0.3))
    if cy > region_bot + 20:
        c.drawString(Q_X, cy - 8, prompt)
        outline_rect(c, Q_X, cy - 22, Q_W, 14, stroke_rgb=GREY, lw=0.5)

    graph_h = region_top - region_bot - 2 * PAD
    draw_blank_bar_chart_area(c, lp_data['ax'],
                              CHART_X, region_bot + PAD,
                              CHART_W, graph_h, is_ms)


def draw_half_tally(c, lp_data, region_top, region_bot, meta, show_ll, is_ms):
    """Half-page LP for tally+draw lessons."""
    PAD = 10
    q_top = region_top - PAD
    if show_ll:
        ll_h = draw_ll(c, LL_X, region_top - PAD,
                       meta['date'], meta['topic'], meta['lf'], meta['ican'])
        q_top = region_top - PAD - ll_h - 8

    # Intro
    intro_para = Paragraph(lp_data['intro'], INTRO_STYLE)
    _, ih = intro_para.wrap(Q_W, 200)
    intro_para.drawOn(c, Q_X, q_top - ih)
    cy = q_top - ih - 4

    # Raw data display
    raw_p = Paragraph(
        '<font size="7" color="#333333">' + lp_data['raw_data'] + '</font>',
        ParagraphStyle('RD', fontName='LibSans', fontSize=7, leading=10,
                       backColor=(0.95, 0.97, 1.0), borderPad=4,
                       textColor=(0.1, 0.1, 0.1))
    )
    _, raw_h = raw_p.wrap(Q_W, 200)
    raw_h = min(raw_h + 8, 60)
    filled_rect(c, Q_X, cy - raw_h, Q_W, raw_h,
                fill_rgb=(0.95, 0.97, 1.0), stroke_rgb=(0.8, 0.8, 0.8), lw=0.5)
    raw_p.drawOn(c, Q_X + 3, cy - raw_h + 3)
    cy -= raw_h + 5

    # Tally table
    cats  = lp_data['tally_categories']
    ms_t  = lp_data.get('ms_tallies', [0]*len(cats))
    ROW_H = 14
    C1, C2, C3 = Q_W * 0.38, Q_W * 0.38, Q_W * 0.22
    filled_rect(c, Q_X, cy - ROW_H, Q_W, ROW_H,
                fill_rgb=(0.09, 0.60, 0.83), stroke_rgb=None)
    for txt, x_off in [('Category', Q_X + 3),
                       ('Tally', Q_X + C1 + 3),
                       ('Total', Q_X + C1 + C2 + 3)]:
        c.setFont('LibSansBold', 7); set_fill(c, (1, 1, 1))
        c.drawString(x_off, cy - ROW_H + 3, txt)
    cy -= ROW_H
    for ri, (cat, mt) in enumerate(zip(cats, ms_t)):
        fill = (0.93, 0.97, 1.0) if ri % 2 == 0 else (1, 1, 1)
        filled_rect(c, Q_X, cy - ROW_H, Q_W, ROW_H,
                    fill_rgb=fill, stroke_rgb=(0.75, 0.75, 0.75), lw=0.4)
        c.setFont('LibSans', 7.5); set_fill(c, (0.1, 0.1, 0.1))
        c.drawString(Q_X + 3, cy - ROW_H + 3, cat)
        if is_ms:
            tally_str = '\u2225 ' * (mt // 5) + '| ' * (mt % 5)
            c.setFont('LibSans', 7)
            c.drawString(Q_X + C1 + 3, cy - ROW_H + 3, tally_str.strip())
            c.drawString(Q_X + C1 + C2 + 3, cy - ROW_H + 3, str(mt))
        cy -= ROW_H
    cy -= 5

    # Questions
    for qi, (qt, ans) in enumerate(lp_data.get('qs', [])):
        qp = Paragraph(f'{qi+1}.  {qt}', Q_STYLE)
        _, qph = qp.wrap(Q_W - 6, 200)
        qbox_h = qph + 8
        filled_rect(c, Q_X, cy - qbox_h, Q_W, qbox_h,
                    fill_rgb=Q_BG, stroke_rgb=Q_BORDER, lw=0.6)
        set_fill(c, (0, 0, 0)); qp.drawOn(c, Q_X + 3, cy - qbox_h + 4)
        cy -= qbox_h + 3
        ANS_H = 18
        if is_ms:
            filled_rect(c, Q_X, cy - ANS_H, Q_W, ANS_H,
                        fill_rgb=GRN_BG, stroke_rgb=GRN, lw=0.6)
            ap = Paragraph(f'\u2713  {ans}', ANS_STYLE)
            ap.wrap(Q_W - 6, ANS_H - 4); ap.drawOn(c, Q_X + 3, cy - ANS_H + 3)
        else:
            outline_rect(c, Q_X, cy - ANS_H, Q_W, ANS_H, stroke_rgb=GREY, lw=0.5)
        cy -= ANS_H + 4

    prompt = lp_data['ax'].get('title_prompt', 'Title:')
    c.setFont('LibSans', 7); set_fill(c, (0.3, 0.3, 0.3))
    if cy > region_bot + 20:
        c.drawString(Q_X, cy - 8, prompt)
        outline_rect(c, Q_X, cy - 22, Q_W, 14, stroke_rgb=GREY, lw=0.5)

    graph_h = region_top - region_bot - 2 * PAD
    draw_blank_bar_chart_area(c, lp_data['ax'],
                              CHART_X, region_bot + PAD,
                              CHART_W, graph_h, is_ms)


def draw_half_draw(c, lp_data, region_top, region_bot, meta, show_ll, is_ms):
    """Half-page LP for draw-graph (line graph) lessons."""
    PAD = 10
    q_top = region_top - PAD
    if show_ll:
        ll_h = draw_ll(c, LL_X, region_top - PAD,
                       meta['date'], meta['topic'], meta['lf'], meta['ican'])
        q_top = region_top - PAD - ll_h - 8

    intro_para = Paragraph(lp_data['intro'], INTRO_STYLE)
    _, intro_h = intro_para.wrap(Q_W, 200)
    intro_h += 4
    cy = q_top
    intro_para.drawOn(c, Q_X, cy - intro_h)
    cy -= intro_h + 4

    hdrs = lp_data['table_headers']
    rows = lp_data['table_rows']
    ROW_H = 13
    COL_W = [Q_W * 0.45, Q_W * 0.45]
    TBL_W = sum(COL_W)
    filled_rect(c, Q_X, cy - ROW_H, TBL_W, ROW_H,
                fill_rgb=(0.09, 0.60, 0.83), stroke_rgb=None)
    for ci, hdr in enumerate(hdrs):
        x_off = Q_X + sum(COL_W[:ci])
        c.setFont('LibSansBold', 7.5); set_fill(c, (1, 1, 1))
        c.drawString(x_off + 3, cy - ROW_H + 3, hdr)
    cy -= ROW_H
    for ri, row in enumerate(rows):
        fill = (0.93, 0.97, 1.0) if ri % 2 == 0 else (1, 1, 1)
        filled_rect(c, Q_X, cy - ROW_H, TBL_W, ROW_H,
                    fill_rgb=fill, stroke_rgb=(0.75, 0.75, 0.75), lw=0.4)
        for ci, val in enumerate(row):
            x_off = Q_X + sum(COL_W[:ci])
            c.setFont('LibSans', 7.5); set_fill(c, (0.1, 0.1, 0.1))
            c.drawString(x_off + 3, cy - ROW_H + 3, str(val))
        cy -= ROW_H
    cy -= 6

    for step in lp_data.get('instructions', []):
        p = Paragraph(step, ParagraphStyle('ST', fontName='LibSans', fontSize=7,
                      leading=10, textColor=(0.2, 0.2, 0.2)))
        _, ph = p.wrap(Q_W, 200)
        p.drawOn(c, Q_X, cy - ph)
        cy -= ph + 2
    cy -= 4

    for qi, (qt, ans) in enumerate(lp_data.get('qs', [])):
        qp = Paragraph(f'{qi+1}.  {qt}', Q_STYLE)
        _, qph = qp.wrap(Q_W - 6, 200)
        qbox_h = qph + 8
        filled_rect(c, Q_X, cy - qbox_h, Q_W, qbox_h,
                    fill_rgb=Q_BG, stroke_rgb=Q_BORDER, lw=0.6)
        set_fill(c, (0, 0, 0)); qp.drawOn(c, Q_X + 3, cy - qbox_h + 4)
        cy -= qbox_h + 3
        ANS_H = 18
        if is_ms:
            filled_rect(c, Q_X, cy - ANS_H, Q_W, ANS_H,
                        fill_rgb=GRN_BG, stroke_rgb=GRN, lw=0.6)
            ap = Paragraph(f'\u2713  {ans}', ANS_STYLE)
            ap.wrap(Q_W - 6, ANS_H - 4); ap.drawOn(c, Q_X + 3, cy - ANS_H + 3)
        else:
            outline_rect(c, Q_X, cy - ANS_H, Q_W, ANS_H, stroke_rgb=GREY, lw=0.5)
        cy -= ANS_H + 4

    prompt = lp_data['ax'].get('title_prompt', 'Title:')
    c.setFont('LibSans', 7); set_fill(c, (0.3, 0.3, 0.3))
    if cy > region_bot + 20:
        c.drawString(Q_X, cy - 8, prompt)
        outline_rect(c, Q_X, cy - 22, Q_W, 14, stroke_rgb=GREY, lw=0.5)

    graph_h = region_top - region_bot - 2 * PAD
    draw_blank_axes_area(c, lp_data['ax'],
                         CHART_X, region_bot + PAD,
                         CHART_W, graph_h, is_ms)


# ── Build ─────────────────────────────────────────────────────────────────────
def build(lesson_num):
    d   = LP[lesson_num]
    day = DAYS[lesson_num]
    wk  = WEEK.get(lesson_num, 'T6W5')
    _dn = {'Monday':1,'Tuesday':2,'Wednesday':3,'Thursday':4,'Friday':5}
    out = f'/tmp/claude_work/{wk} - {_dn.get(day,1)} - {day} - MathsLP.pdf'

    c = canvas.Canvas(out, pagesize=A4)
    c.setTitle(f'{wk} {day} L{lesson_num} — Statistics LP')

    is_draw = d['lp1'].get('lp_type') == 'draw_graph'
    is_bar   = d['lp1'].get('lp_type') == 'draw_bar'
    is_tally = d['lp1'].get('lp_type') == 'tally_draw'

    for page_type in ['standard', 'adapted', 'marking']:
        is_ms = (page_type == 'marking')
        c.setPageSize(A4)

        # White background
        c.setFillColorRGB(1,1,1)
        c.rect(0, 0, W, H, fill=1, stroke=0)

        if is_draw:
            draw_half_draw(c, d['lp1'], H - M, CUT_Y + M, d,
                           show_ll=True, is_ms=is_ms)
            draw_half_draw(c, d['lp2'], CUT_Y - M, M, d,
                           show_ll=False, is_ms=is_ms)
        elif is_bar:
            draw_half_bar(c, d['lp1'], H - M, CUT_Y + M, d,
                          show_ll=True, is_ms=is_ms)
            draw_half_bar(c, d['lp2'], CUT_Y - M, M, d,
                          show_ll=False, is_ms=is_ms)
        elif is_tally:
            draw_half_tally(c, d['lp1'], H - M, CUT_Y + M, d,
                            show_ll=True, is_ms=is_ms)
            draw_half_tally(c, d['lp2'], CUT_Y - M, M, d,
                            show_ll=False, is_ms=is_ms)
        else:
            draw_half(c, d['lp1'], H - M, CUT_Y + M, d,
                      show_ll=True, is_ms=is_ms)
            draw_half(c, d['lp2'], CUT_Y - M, M, d,
                      show_ll=False, is_ms=is_ms)

        draw_cut_line(c)
        c.showPage()

    c.save()
    print(f'Saved: {out}')
    return out


# ── PPTX wrapper for inject_lp_previews.py ────────────────────────────────────
def make_pptx_wrapper(pdf_path, pptx_path):
    """
    Converts each PDF page to an image and creates a 3-slide A4 PPTX.
    inject_lp_previews.py will then crop top/bottom halves for the preview.
    """
    import subprocess, tempfile, glob
    from pptx import Presentation
    from pptx.util import Inches, Emu

    SLIDE_W, SLIDE_H = 7.5, 10.833

    # Rasterise PDF pages to images
    tmp = tempfile.mkdtemp()
    subprocess.run(['pdftoppm', '-jpeg', '-r', '150', pdf_path,
                    os.path.join(tmp, 'page')], capture_output=True)
    pages = sorted(glob.glob(os.path.join(tmp, 'page-*.jpg')))
    if not pages:
        pages = sorted(glob.glob(os.path.join(tmp, 'page*.jpg')))

    prs = Presentation()
    prs.slide_width  = Emu(int(SLIDE_W * 914400))
    prs.slide_height = Emu(int(SLIDE_H * 914400))
    blank = prs.slide_layouts[6]

    for page_img in pages[:3]:
        sld = prs.slides.add_slide(blank)
        sld.background.fill.solid()
        sld.background.fill.fore_color.rgb = __import__(
            'pptx.dml.color', fromlist=['RGBColor']).RGBColor(0xFF,0xFF,0xFF)
        sld.shapes.add_picture(page_img,
                               Emu(0), Emu(0),
                               Emu(int(SLIDE_W*914400)),
                               Emu(int(SLIDE_H*914400)))

    prs.save(pptx_path)
    print(f'PPTX wrapper: {pptx_path}  ({len(pages)} slides)')


if __name__ == '__main__':
    n = int(sys.argv[1]) if len(sys.argv) > 1 else 17
    pdf  = build(n)
    day  = DAYS[n]
    wk   = WEEK.get(n, 'T6W5')
    _dn  = {'Monday':1,'Tuesday':2,'Wednesday':3,'Thursday':4,'Friday':5}
    pptx = f'/tmp/claude_work/{wk} - {_dn.get(day,1)} - {day} - MathsLP.pptx'
    make_pptx_wrapper(pdf, pptx)

"""
build_stats_lp.py
Generates A4 portrait LP PPTXs for T6W5 Statistics lessons (L17, L18, L19).
Each file has ONE slide (printed and cut in half):
  LP1 (top half): chart image on right, 3 questions with answer lines on left
  LP2 (bottom half): 3-4 further questions + marking station answers on right

Run: python3 build_stats_lp.py 17   (or 18 or 19)
"""

import sys, os, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Constants ─────────────────────────────────────────────────────────────────
W       = Inches(7.5)
H       = Inches(10.833)      # A4 portrait
MID_Y   = Inches(5.4165)      # cut line
LABEL_H = Inches(1.021)       # Learning Label height
MARGIN  = Inches(0.28)
GAP     = Inches(0.10)

WFA_BLUE   = RGBColor(0x17, 0x98, 0xD3)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
GREY_LINE  = RGBColor(0xCC, 0xCC, 0xCC)
MARK_BG    = RGBColor(0xE8, 0xF5, 0xE9)
MARK_TEXT  = RGBColor(0x1A, 0x5C, 0x2A)
Q_BG       = RGBColor(0xDE, 0xEA, 0xF1)

CHART_DIR  = '/tmp/wfa_stats_charts'
EMU_IN     = 914400

def emu(inches): return Emu(int(inches * EMU_IN))

# ── LP data ───────────────────────────────────────────────────────────────────
LP_DATA = {

    17: {
        'li':   'I can read and interpret pictograms, bar charts and tables.',
        'lp1': {
            'chart': 'c1_ido1_pictogram.png',
            'chart_title': 'Favourite sports in Year 4',
            'intro': 'Use the pictogram to answer these questions.',
            'questions': [
                ('How many pupils chose basketball?', '6 pupils  (3 symbols × 2)'),
                ('How many pupils chose football and volleyball altogether?',
                 '12 + 8 = 20 pupils'),
                ('How many MORE pupils chose swimming than capoeira?',
                 '10 − 4 = 6 more pupils'),
            ],
        },
        'lp2': {
            'chart': 'c2_ido1_table.png',
            'chart_title': 'Average daily sunshine hours',
            'intro': 'Use the table to answer these questions.',
            'questions': [
                ('How many hours of sunshine does England get in winter?',
                 '2 hours'),
                ('What is the total sunshine across both countries in spring?',
                 '5 + 7 = 12 hours'),
                ('Which country has more sunshine overall? By how many hours?',
                 'Brazil — 31 vs 19, so 12 more hours'),
                ('In which season are England and Brazil most similar?',
                 'Summer — just 1 hour apart (8 vs 9)'),
            ],
        },
    },

    18: {
        'li':   'I can calculate the sum and difference from charts and compare two data sets.',
        'lp1': {
            'chart': 'c1_ido1_bar_chart.png',
            'chart_title': 'Animals counted at an Amazon river each day',
            'intro': 'Use the bar chart to answer these questions.',
            'questions': [
                ('How many animals were counted on Tuesday and Wednesday altogether?',
                 '24 + 15 = 39 animals'),
                ('How many MORE were counted on Thursday than Monday?',
                 '30 − 18 = 12 more animals'),
                ('What is the total for all five days?',
                 '18 + 24 + 15 + 30 + 21 = 108 animals'),
            ],
        },
        'lp2': {
            'chart': 'c2_ido1_double_bar.png',
            'chart_title': 'Average temperature (°C) — London and Rio de Janeiro',
            'intro': 'Use the double bar chart to answer these questions.',
            'questions': [
                ('What is the temperature difference between London and Rio in winter?',
                 '22 − 6 = 16°C'),
                ('In which season are the two cities closest in temperature?',
                 'Summer — 28°C vs 18°C, a gap of 10°C'
                 ' (smallest gap across all seasons)'),
                ('What is London\'s temperature range across the year?',
                 '18 − 6 = 12°C range'),
                ('Write a statement comparing the two cities using data.',
                 'Any reasonable statement using values from the chart.'),
            ],
        },
    },

    19: {
        'li':   'I can read line graphs and estimate values between labelled points.',
        'lp1': {
            'chart': 'c1_ido1_line_graph.png',
            'chart_title': 'Temperature in Bristol on a June day',
            'intro': 'Use the line graph to answer these questions.',
            'questions': [
                ('What was the temperature at 10:00?', '17°C'),
                ('At what time was it warmest? What was the temperature?',
                 '14:00 — 24°C'),
                ('Between which two readings did temperature rise most?',
                 '10:00 to 12:00 — rose by 4°C'),
            ],
        },
        'lp2': {
            'chart': 'c2_ido1_line_graph.png',
            'chart_title': 'Temperature in São Paulo on a June day',
            'intro': 'Use the line graph to answer these questions.',
            'questions': [
                ('What was the temperature in São Paulo at 12:00?', '30°C'),
                ('Estimate the temperature at 09:00.',
                 'About 24°C — halfway between 22 and 26'),
                ('Estimate the temperature at 11:00.',
                 'About 28°C — halfway between 26 and 30'),
                ('At approximately what time was it 30°C on the way back down?',
                 'Around 15:00 — the line falls from 32 to 29, passing 30 near the start'),
            ],
        },
    },
}

LESSON_META = {
    17: {'day': 'Monday',    'week': 'T6W5', 'topic': 'Statistics'},
    18: {'day': 'Tuesday',   'week': 'T6W5', 'topic': 'Statistics'},
    19: {'day': 'Wednesday', 'week': 'T6W5', 'topic': 'Statistics'},
}

# ── Drawing helpers ───────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        emu(x), emu(y), emu(w), emu(h)
    )
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_name='Aptos', size=10,
             bold=False, color=DARK, align=PP_ALIGN.LEFT,
             fill=None, wrap=True, anchor='top'):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    from pptx.oxml.ns import qn
    from lxml import etree as _et
    _A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    _bp = tf._txBody.find(f'{{{_A}}}bodyPr')
    if _bp is None:
        _bp = _et.SubElement(tf._txBody, f'{{{_A}}}bodyPr')
    _bp.set('anchor', 't' if anchor == 'top' else 'ctr')
    if fill:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def draw_learning_label(slide, y_top, week, day, topic, li, lp_num):
    """Draw the Learning Label (LL) header band for one LP half."""
    iw = 7.5  # slide width in inches

    # Blue band
    add_rect(slide, 0, y_top, iw, 1.021, fill=WFA_BLUE)

    # LP number badge (small white box, top-left)
    add_text(slide, f'LP{lp_num}', 0.08, y_top + 0.07, 0.38, 0.30,
             font_name='Aptos', size=9, bold=True, color=WHITE,
             fill=None, wrap=False)

    # Week / day / topic line
    add_text(slide, f'{week}  ·  {day}  ·  {topic}',
             0.50, y_top + 0.04, 6.60, 0.28,
             font_name='Aptos', size=8, bold=False, color=WHITE, wrap=False)

    # LI (Learning Intention)
    add_text(slide, li,
             0.50, y_top + 0.35, 6.60, 0.60,
             font_name='Aptos', size=9, bold=True, color=WHITE, wrap=True)


def draw_lp_half(slide, data, y_offset, week, day, topic, li, lp_num):
    """Draw one LP half (LP1 or LP2) within the slide."""
    chart_file = os.path.join(CHART_DIR, data['chart'])
    intro      = data['intro']
    questions  = data['questions']

    # Learning Label
    draw_learning_label(slide, y_offset, week, day, topic, li, lp_num)

    content_top = y_offset + 1.021 + 0.10  # below LL
    half_h      = 5.4165                    # each half
    content_bot = y_offset + half_h - 0.15
    content_h   = content_bot - content_top

    # ── Left side: questions (55% width) ──────────────────────────────────────
    Q_X  = 0.28
    Q_W  = 3.90
    MS_X = 4.30  # marking station / chart left edge
    MS_W = 2.95  # chart / marking station width

    # Intro line
    add_text(slide, intro, Q_X, content_top, Q_W, 0.28,
             font_name='Aptos', size=8.5, bold=True, color=DARK)

    n_q  = len(questions)
    q_h  = (content_h - 0.35) / n_q   # height per question slot
    Q_TEXT_H  = 0.38
    ANS_LINES = max(1, int((q_h - Q_TEXT_H - 0.08) / 0.28))

    for i, (q_text, answer) in enumerate(questions):
        qy = content_top + 0.33 + i * q_h

        # Question text box (light blue)
        add_text(slide, f'{i+1}.  {q_text}',
                 Q_X, qy, Q_W, Q_TEXT_H,
                 font_name='Aptos', size=8.5, bold=False,
                 color=RGBColor(0x1F, 0x4E, 0x79),
                 fill=Q_BG, wrap=True)

        # Answer lines
        for ln in range(ANS_LINES):
            line_y = qy + Q_TEXT_H + 0.06 + ln * 0.28
            add_rect(slide, Q_X, line_y, Q_W, 0.22,
                     fill=None, line_color=GREY_LINE, line_w=Pt(0.5))

    # ── Right side: chart image ────────────────────────────────────────────────
    chart_top = content_top
    chart_h   = content_h * 0.58
    if os.path.exists(chart_file):
        slide.shapes.add_picture(
            chart_file,
            emu(MS_X), emu(chart_top), emu(MS_W), emu(chart_h)
        )

    # Chart title below image
    add_text(slide, data.get('chart_title', ''),
             MS_X, chart_top + chart_h + 0.04, MS_W, 0.25,
             font_name='Aptos', size=7.5, bold=False,
             color=RGBColor(0x44, 0x44, 0x44),
             align=PP_ALIGN.CENTER, wrap=False)

    # ── Marking station (below chart on right) ─────────────────────────────────
    ms_top = chart_top + chart_h + 0.33
    ms_h   = content_bot - ms_top - 0.05
    if ms_h > 0.4:
        add_rect(slide, MS_X, ms_top, MS_W, ms_h,
                 fill=MARK_BG, line_color=MARK_TEXT, line_w=Pt(0.75))
        add_text(slide, 'Answers', MS_X + 0.08, ms_top + 0.06,
                 MS_W - 0.16, 0.22,
                 font_name='Aptos', size=8, bold=True, color=MARK_TEXT)
        ans_y = ms_top + 0.30
        for i, (_, answer) in enumerate(questions):
            if ans_y + 0.22 > content_bot:
                break
            add_text(slide, f'{i+1}. {answer}',
                     MS_X + 0.08, ans_y, MS_W - 0.16, 0.30,
                     font_name='Aptos', size=7, bold=False, color=MARK_TEXT,
                     wrap=True)
            # Estimate lines needed
            lines = max(1, len(textwrap.wrap(answer, width=38)))
            ans_y += 0.18 + (lines - 1) * 0.16

    # ── Cut line ──────────────────────────────────────────────────────────────
    if lp_num == 1:
        cut_y = y_offset + 5.4165
        line = slide.shapes.add_connector(1, emu(0), emu(cut_y), emu(7.5), emu(cut_y))
        line.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        line.line.width = Pt(0.5)
        line.line.dash_style = 4   # DASH


# ── Main build ────────────────────────────────────────────────────────────────

def build_stats_lp(lesson_num):
    data   = LP_DATA[lesson_num]
    meta   = LESSON_META[lesson_num]
    week   = meta['week']
    day    = meta['day']
    topic  = meta['topic']
    li     = data['li']

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(blank_layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    draw_lp_half(slide, data['lp1'], 0,       week, day, topic, li, lp_num=1)
    draw_lp_half(slide, data['lp2'], 5.4165,  week, day, topic, li, lp_num=2)

    out = f'T6W5_{day}_L{lesson_num}_LP.pptx'
    prs.save(out)
    print(f'Saved: {out}')
    return out


if __name__ == '__main__':
    ln = int(sys.argv[1]) if len(sys.argv) > 1 else 17
    build_stats_lp(ln)

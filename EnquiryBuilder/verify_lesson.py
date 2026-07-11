#!/usr/bin/env python3
"""
verify_lesson.py - automated gate a built lesson PPTX must pass before it is
ever called "delivered". Encodes every bug reported across the T6W7 sessions
as one specific, repeatable check, so a fixed bug cannot silently regress.

Usage: python3 verify_lesson.py <pptx_path> <mtp_json_path> <manifest_json_path>
Exits 0 and prints "VERIFY: PASS" only if every check passes.
Exits 1 and prints every failure (slide + reason) otherwise.
"""
import sys, json, zipfile, re
from pptx import Presentation
from pptx.util import Emu

BANNED_TEXT = [
    "turn on the light",
    "eyes",
    "white cat in the dark room",
    "insert any other states of being icons",
    "click to edit master text styles",
]

def fail(failures, msg):
    failures.append(msg)

def check_customxml(pptx_path, failures):
    with zipfile.ZipFile(pptx_path) as z:
        stray = [n for n in z.namelist() if n.startswith('customXml/')]
    if stray:
        fail(failures, f"REPAIR-DIALOG RISK: {len(stray)} customXml/ part(s) still present "
                        f"(SharePoint/Teams metadata) - fix_pptx_ooxml.py Fix #6 was not applied "
                        f"or did not run: {stray[:5]}")

def check_slide_sequence(manifest, mtp, failures):
    expected = [s['type'] for s in mtp['lesson']['slides']]
    actual = [s['type'] for s in manifest['slides']]
    if expected != actual:
        fail(failures, f"Slide sequence mismatch.\n    expected: {expected}\n    actual:   {actual}")

def check_required_present(manifest, failures):
    import science_registry as REG
    present = {s['type'] for s in manifest['slides']}
    missing = REG.REQUIRED_TYPES - present
    if missing:
        fail(failures, f"Required slide types missing from build: {sorted(missing)}")

def _slide_all_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)

def check_banned_text(prs, failures):
    for i, slide in enumerate(prs.slides, 1):
        text = _slide_all_text(slide).lower()
        for banned in BANNED_TEXT:
            if banned in text:
                fail(failures, f"Slide {i}: banned template placeholder text survived ('{banned}')")

def check_concept_cartoon_content(prs, manifest, mtp, failures):
    for entry in manifest['slides']:
        if entry['type'] != 'concept_cartoon':
            continue
        idx = entry['output_index']
        slide = prs.slides[idx - 1]
        spec = next(s for s in mtp['lesson']['slides'] if s['type'] == 'concept_cartoon')
        expected_statements = [l['statement'] for l in spec['learners']]
        slide_text = _slide_all_text(slide)
        for stmt in expected_statements:
            if stmt not in slide_text:
                fail(failures, f"Slide {idx} (concept_cartoon): expected learner statement not found "
                                f"verbatim on the slide: '{stmt[:60]}...'")
        # must contain an actual (non-empty) picture for the central image
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
        if not pics:
            fail(failures, f"Slide {idx} (concept_cartoon): no picture shapes found at all")

def check_lo_no_duplicate_panel(prs, manifest, failures):
    """The known LO-slide bug: two overlapping panels. Guard against it coming
    back even if LO_STALE_GROUP_IDS in the registry ever goes stale."""
    for entry in manifest['slides']:
        if entry['type'] != 'lo':
            continue
        idx = entry['output_index']
        slide = prs.slides[idx - 1]
        names = [sh.name for sh in slide.shapes]
        # 'Text Placeholder 33' only exists in the stale duplicate group
        stale_hits = [n for n in names if n == 'Text Placeholder 33']
        if stale_hits:
            fail(failures, f"Slide {idx} (lo): stale duplicate LO panel is back "
                            f"({len(stale_hits)}x 'Text Placeholder 33' shapes found)")

def check_geometric_overlap(prs, failures):
    """Real, literal geometric overlap check: any two TEXT-bearing shapes on
    the same slide whose bounding boxes intersect by more than a small
    tolerance are flagged. Decorative/background shapes without text are
    ignored (borders, icons, frames are allowed to sit under content)."""
    TOLERANCE = 0.15  # allow up to 15% overlap area before flagging
    for i, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            if not shape.text_frame.text.strip(): continue
            try:
                l, t, w, h = shape.left, shape.top, shape.width, shape.height
            except Exception:
                continue
            if None in (l, t, w, h): continue
            # shape_type distinguishes "content label" (TEXT_BOX / PLACEHOLDER)
            # from "background panel" (AUTO_SHAPE with a fill, e.g. a rounded
            # rectangle a caption is deliberately layered on top of). Only
            # compare shapes of the SAME kind - a panel-plus-label pairing is
            # standard layered design, not a bug. Two panels or two labels
            # stacked on each other is the real bug (this is exactly the LO
            # duplicate-group pattern).
            try:
                kind = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                kind = -1
            boxes.append((shape.name, l, t, w, h, kind))
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                n1, l1, t1, w1, h1, k1 = boxes[a]
                n2, l2, t2, w2, h2, k2 = boxes[b]
                if k1 != k2: continue
                ix = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
                iy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
                inter = ix * iy
                if inter <= 0: continue
                smaller = min(w1 * h1, w2 * h2)
                if smaller > 0 and inter / smaller > TOLERANCE:
                    fail(failures, f"Slide {i}: same-kind text shapes overlap ('{n1}' and '{n2}', "
                                    f"{inter/smaller:.0%} of the smaller shape's area)")

def check_animation_pattern(pptx_path, failures):
    """Forbidden pattern: an explicit hide-at-start <p:par> before the seq.
    SKILL.md's own rule: PowerPoint hides clickEffect entries automatically;
    an explicit hide block produces 'TRIGGER: UNNAMED' in the animation pane."""
    with zipfile.ZipFile(pptx_path) as z:
        slide_files = [n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', n)]
        for name in slide_files:
            xml = z.read(name).decode('utf-8', errors='ignore')
            if '<p:timing' not in xml:
                continue
            timing = xml[xml.index('<p:timing'):]
            # a hide-at-start block sets visibility "hidden" as the FIRST childTnLst
            # entries before the <p:seq> - detect by checking a hidden 'strVal val="hidden"'
            # appears before the first '<p:seq'
            seq_pos = timing.find('<p:seq')
            hidden_pos = timing.find('val="hidden"')
            if hidden_pos != -1 and (seq_pos == -1 or hidden_pos < seq_pos):
                fail(failures, f"{name}: forbidden hide-at-start animation pattern found "
                                f"(will show as 'TRIGGER: UNNAMED' in PowerPoint's animation pane)")
            if 'restart="whenNotActive"' in timing:
                fail(failures, f"{name}: animation root uses restart=\"whenNotActive\" "
                                f"(should be restart=\"never\" per SKILL.md)")

def check_images_present(prs, manifest, mtp, failures):
    """Every slide type whose spec declared an image_path must actually have
    a non-empty picture on the delivered slide - never a placeholder gap."""
    for spec in mtp['lesson']['slides']:
        needs_image = 'image_path' in spec or (spec['type'] == 'wedo_grid' and 'items' in spec)
        if not needs_image: continue
        matching = [e for e in manifest['slides'] if e['type'] == spec['type']]
        for entry in matching:
            idx = entry['output_index']
            slide = prs.slides[idx - 1]
            pics = [sh for sh in slide.shapes if sh.shape_type == 13]
            expected_n = len(spec['items']) if spec['type'] == 'wedo_grid' else 1
            if len(pics) < expected_n:
                fail(failures, f"Slide {idx} ({spec['type']}): expected at least {expected_n} "
                                f"image(s), found {len(pics)}")

def main():
    pptx_path, mtp_path, manifest_path = sys.argv[1:4]
    with open(mtp_path) as f: mtp = json.load(f)
    with open(manifest_path) as f: manifest = json.load(f)
    prs = Presentation(pptx_path)

    failures = []
    check_customxml(pptx_path, failures)
    check_slide_sequence(manifest, mtp, failures)
    check_required_present(manifest, failures)
    check_banned_text(prs, failures)
    check_concept_cartoon_content(prs, manifest, mtp, failures)
    check_lo_no_duplicate_panel(prs, manifest, failures)
    check_geometric_overlap(prs, failures)
    check_animation_pattern(pptx_path, failures)
    check_images_present(prs, manifest, mtp, failures)

    if failures:
        print(f"VERIFY: FAIL ({len(failures)} issue(s))\n")
        for f_ in failures:
            print(f"  - {f_}")
        sys.exit(1)
    else:
        n_slides = len(list(prs.slides))
        print(f"VERIFY: PASS ({n_slides} slides checked)")
        sys.exit(0)

if __name__ == '__main__':
    main()

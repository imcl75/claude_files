#!/usr/bin/env python3
"""
build_lp.py — Generic enquiry Learning Paper builder.
v2: supports both legacy flat-sections schema and new three-level schema.

Entry points:
    build_lp(lesson_data, out_path, resource_base, level='standard')
        Build one level's LP PPTX.

    build_lp_all_levels(lesson_data, out_dir, resource_base, base_name)
        Build all three levels, return dict {level: path}.

Schema detection:
    LEGACY: lp = {"sections": [...], "date": ..., "lf": ..., ...}
    NEW:    lp = {"standard": {"elements": [...]}, "adapted": {...}, "further_adapted": {...}}

New element types (v2):
    answer_lines, cloze, matching, diagram
    (plus all legacy types: heading, instruction, write_lines, word_bank,
     reference_image, row_boxes, pair_boxes, table, graph_template,
     sentence_starter, spacer, sort_table, sort_table_answers,
     marking_station, answer_text)

Called automatically from build_science_lesson.py after VERIFY: PASS.
"""

import os, sys, json

_DEFAULT_RESOURCE_BASE = '/tmp/t6w7'
if _DEFAULT_RESOURCE_BASE not in sys.path:
    sys.path.insert(0, _DEFAULT_RESOURCE_BASE)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from PIL import Image as _PILImage
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# ── Slide dimensions (A4 portrait) ─────────────────────────────────────────
SW_IN, SH_IN = 7.5, 10.833
SW, SH = int(SW_IN * 914400), int(SH_IN * 914400)
MARGIN = 0.25
LBL_Y  = MARGIN
CONT_X = MARGIN
CONT_W = SW_IN - 2 * MARGIN

BLUE   = RGBColor(0x17, 0x98, 0xD3)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREEN  = RGBColor(0x4F, 0xAD, 0x5B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xCC, 0xCC, 0xCC)
EXBOOK = RGBColor(0xBD, 0xD7, 0xEE)
CREAM  = RGBColor(0xFE, 0xF9, 0xE7)
PINK   = RGBColor(0xC0, 0x15, 0x7B)

FONT_BODY = 'Twinkl Cursive Looped'

# ── Text-height model ────────────────────────────────────────────────────────
# All text rendering uses these constants so behaviour is consistent everywhere.
#
# _TCF: chars per (pt × inch) of usable width — calibrated for Twinkl Cursive
#       Looped, which is wider than a typical body font. 0.46 (the old value)
#       allowed 99 chars/line at 11pt on 7 in and classified a 92-char label
#       as fitting in one line; it does not. 0.58 gives 78 chars/line and
#       correctly wraps the same label to two lines.
# _LHF: line-height as a multiple of point size (1.45 × pt / 72 = inches/line).
# _TPAD: fixed padding added to both textbox height and y-advance. Gives the
#        text frame headroom against rounding error without burning much space.
_TCF  = 0.58    # chars per (pt × inch) — Twinkl Cursive calibration
_LHF  = 1.45    # line-height factor (× size_pt / 72 = inches per line)
_TPAD = 0.06    # textbox padding added to both height and y-advance (inches)

# answer_lines sub-layout gaps (all in inches):
_AL_LABEL_LINE_GAP   = 0.05   # between bottom of label text and first ruled line
_AL_STARTER_LINE_GAP = 0.04   # between bottom of sentence starter and first ruled line
_AL_LINE_PITCH       = 0.315  # ruled-line pitch (≈ 8 mm) — matches exercise-book ruling
_AL_POST_GAP         = 0.13   # breathing room after last ruled line before next element

# ── Minimum font sizes ───────────────────────────────────────────────────────
# Every renderer runs its size_pt through _safe_pt() before use.
# These floors ensure nothing on a children's LP becomes unreadable in print.
#
#   BODY     – text children must read independently: instructions, question labels
#   SUBLABEL – secondary labels, sentence starters, column headers in tables
#   CELL     – pre-filled cell content (e.g. material names); printed, brief
#   MARKING  – marking-station answers; teacher-read, smaller is acceptable
#
MIN_PT_BODY     = 12   # pt
MIN_PT_SUBLABEL = 11   # pt
MIN_PT_CELL     = 10   # pt
MIN_PT_MARKING  = 10   # pt

# ── Write-box geometry ───────────────────────────────────────────────────────
# Any box, cell, or row where a child is expected to write uses _write_box_h().
# Height = n_lines × _AL_LINE_PITCH + 2 × _WRITE_BOX_PAD.
# At 1 line: 0.315 + 0.14 = 0.455 in (≈ 11.6 mm).
# The old sort-table default of 0.36 in left only 0.023 in above and below the
# writing — less than 1 mm — and is replaced throughout.
_WRITE_BOX_PAD = 0.07  # top + bottom internal padding, each side (inches)

# Differentiation level colours and labels
LEVEL_COLOURS = {
    'standard':        BLUE,
    'adapted':         ORANGE,
    'further_adapted': PINK,
}
LEVEL_LABELS = {
    'standard':        'Standard',
    'adapted':         'Adapted',
    'further_adapted': 'Further Adapted',
}

# Sort table column setup (legacy — kept for reuse)
_SORT_COL_LABELS  = ['Material', 'State?', 'Reason (use the particle model)']
_SORT_COL_RATIO   = [0.26, 0.14, 0.60]


# ══════════════════════════════════════════════════════════════════
# Low-level PPTX helpers
# ══════════════════════════════════════════════════════════════════

def _new_prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SW), Emu(SH)
    blank = prs.slide_layouts[6]
    prs.slides.add_slide(blank)
    prs.slides.add_slide(blank)
    return prs


def _clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            sp_tree.remove(child)


def _i(v):
    return Inches(v)


def _add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))


def _set_para(para, text, bold=False, italic=False, size_pt=10,
              color=None, font=FONT_BODY, align=PP_ALIGN.LEFT, underline=False):
    para.alignment = align
    run = para.add_run()
    run.text = text
    rf = run.font
    rf.name = font
    rf.size = Pt(size_pt)
    rf.bold = bold
    rf.italic = italic
    rf.underline = underline
    if color:
        rf.color.rgb = color
    return run


def _add_line(slide, x, y, w, color=None, width_pt=1.0):
    c = color if color is not None else EXBOOK
    hex_col = f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'
    lw = int(Pt(width_pt))
    sp_tree = slide.shapes._spTree
    uid = max((int(el.get('id', 0)) for el in sp_tree.iter()
               if el.get('id') and el.get('id').isdigit()), default=100) + 1
    xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{uid}" name="Line{uid}"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{int(Inches(x))}" y="{int(Inches(y))}"/>'
        f'<a:ext cx="{int(Inches(w))}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom><a:noFill/>'
        f'<a:ln w="{lw}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:ln>'
        f'<a:effectLst/></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    sp_tree.append(etree.fromstring(xml))


def _add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0.75):
    shape = slide.shapes.add_shape(1, _i(x), _i(y), _i(w), _i(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _wrap_lines(text, width_in, size_pt):
    """
    Return wrapped-line count for text in a box of width_in inches at size_pt.
    Uses _TCF = 0.58 (conservative calibration for Twinkl Cursive Looped).
    At 11 pt on 7 in: 78 chars/line. Old value 0.46 gave 99 chars/line and
    silently classified 92-char labels as 1-line — they are not.
    """
    if not text:
        return 1
    chars_per_line = max(1, int(width_in * 72 / (size_pt * _TCF)))
    words = str(text).split()
    lines, cur = 1, ''
    for word in words:
        test = (cur + ' ' + word).strip()
        if len(test) <= chars_per_line:
            cur = test
        else:
            lines += 1
            cur = word
    return max(1, lines)


def _text_h(text, width_in, size_pt):
    """
    Height in inches for a text block (excluding _TPAD).
    = wrapped-line count × (size_pt × _LHF / 72).
    At 11 pt: 0.222 in per line.
    Callers add _TPAD when setting both textbox height and y-advance.
    """
    return _wrap_lines(text, width_in, size_pt) * (size_pt * _LHF / 72)


# Alias — legacy renderers that call _wrap_line_count still work unchanged.
def _wrap_line_count(text, width_in, size_pt):
    return _wrap_lines(text, width_in, size_pt)


def _safe_pt(size_pt, minimum=None):
    """
    Clamp size_pt to a minimum.  Default minimum is MIN_PT_BODY (12).
    Pass the appropriate MIN_PT_* constant for the context:
        _safe_pt(pt)                  → body text / labels (≥ 12)
        _safe_pt(pt, MIN_PT_SUBLABEL) → sub-labels, starters, table headers (≥ 11)
        _safe_pt(pt, MIN_PT_CELL)     → pre-filled cell content (≥ 10)
        _safe_pt(pt, MIN_PT_MARKING)  → marking-station answers (≥ 10)
    """
    if minimum is None:
        minimum = MIN_PT_BODY
    return max(float(size_pt), float(minimum))


def _write_box_h(n_lines=1):
    """
    Minimum height in inches for a box where n_lines of child handwriting fits.
    = n_lines × _AL_LINE_PITCH + 2 × _WRITE_BOX_PAD
    Callers may pass this as row_height or use max(specified, _write_box_h(n)).
    """
    return n_lines * _AL_LINE_PITCH + 2 * _WRITE_BOX_PAD


def _resolve_path(path, resource_base):
    if os.path.isabs(path) and os.path.exists(path):
        return path
    candidate = os.path.join(resource_base, path)
    if os.path.exists(candidate):
        return candidate
    raise FileNotFoundError(
        f"LP resource not found: '{path}' (also tried '{candidate}')")


# ══════════════════════════════════════════════════════════════════
# Legacy section renderers
# ══════════════════════════════════════════════════════════════════

def _render_heading(slide, y, sec):
    text = sec['text']
    size_pt = _safe_pt(sec.get('size_pt', 12))
    w = sec.get('w', CONT_W)
    h = _text_h(text, w, size_pt) + _TPAD
    tb = _add_textbox(slide, CONT_X, y, w, h)
    tb.text_frame.word_wrap = True
    _set_para(tb.text_frame.paragraphs[0], text, bold=True, size_pt=size_pt, color=BLUE)
    return y + h + 0.06


def _render_instruction(slide, y, sec):
    text = sec['text']
    size_pt = _safe_pt(sec.get('size_pt', 12))
    w = sec.get('w', CONT_W)
    h = _text_h(text, w, size_pt) + _TPAD
    tb = _add_textbox(slide, CONT_X, y, w, h)
    tb.text_frame.word_wrap = True
    _set_para(tb.text_frame.paragraphs[0], text, size_pt=size_pt, color=DARK)
    return y + h + 0.08


def _render_write_lines(slide, y, sec):
    n = sec.get('n', 3)
    gap = sec.get('gap_in', 0.315)
    w = sec.get('w', CONT_W)
    x = sec.get('x', CONT_X)
    for _ in range(n):
        _add_line(slide, x, y, w)
        y += gap
    return y + 0.10


def _render_word_bank(slide, y, sec):
    words = sec['words']
    # words can be a list or a string
    if isinstance(words, list):
        words = '  •  '.join(words)
    size_pt = _safe_pt(sec.get('size_pt', 12))
    w = sec.get('w', CONT_W)
    inner_w = w - 0.10
    full_text = 'Word bank: ' + words
    box_h = _text_h(full_text, inner_w, size_pt) + _TPAD + 0.10  # +0.10 for border padding
    _add_rect(slide, CONT_X, y, w, box_h, fill_rgb=CREAM, line_rgb=ORANGE, line_pt=1.0)
    tb = _add_textbox(slide, CONT_X + 0.05, y + 0.05, inner_w, box_h - 0.08)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'Word bank: '
    r1.font.name = FONT_BODY
    r1.font.bold = True
    r1.font.size = Pt(size_pt)
    r1.font.color.rgb = ORANGE
    r2 = p.add_run()
    r2.text = words
    r2.font.name = FONT_BODY
    r2.font.size = Pt(size_pt)
    r2.font.color.rgb = DARK
    return y + box_h + 0.06


def _render_reference_image(slide, y, sec, resource_base):
    path = _resolve_path(sec['path'], resource_base)
    max_h = sec.get('max_h', 1.35)
    max_w = sec.get('max_w', CONT_W)
    iw, ih = _PILImage.open(path).size
    ratio = iw / ih
    h = max_h
    w = h * ratio
    if w > max_w:
        w = max_w
        h = w / ratio
    x = CONT_X + (CONT_W - w) / 2
    slide.shapes.add_picture(path, _i(x), _i(y), width=_i(w), height=_i(h))
    return y + h + 0.10


def _render_row_boxes(slide, y, sec):
    items   = sec['items']
    size_pt = _safe_pt(sec.get('size_pt', 10), MIN_PT_CELL)
    n       = len(items)
    gap     = 0.10
    box_w   = (CONT_W - gap * (n - 1)) / n
    inner_w = box_w - 0.14
    x = CONT_X
    for label in items:
        # Height: enough for the wrapped text, never below a write-box minimum
        text_h  = _text_h(label, inner_w, size_pt) + _TPAD
        box_h   = max(_write_box_h(1), sec.get('height', text_h), text_h)
        _add_rect(slide, x, y, box_w, box_h, line_rgb=BLUE, line_pt=1.0)
        tb = _add_textbox(slide, x + 0.05, y + _WRITE_BOX_PAD,
                           inner_w, box_h - 2 * _WRITE_BOX_PAD)
        tb.text_frame.word_wrap = True
        _set_para(tb.text_frame.paragraphs[0], label, size_pt=size_pt, color=DARK)
        x += box_w + gap
    # Use the last computed box_h (all boxes share the tallest height)
    return y + box_h + 0.08


def _render_pair_boxes(slide, y, sec):
    items   = sec['items']
    size_pt = _safe_pt(sec.get('size_pt', 11), MIN_PT_SUBLABEL)
    half_w  = CONT_W / 2 - 0.1
    inner_w = half_w - 0.12
    x = CONT_X
    y_next = y
    for label in items[:2]:
        text_h = _text_h(label, inner_w, size_pt) + _TPAD
        box_h  = max(_write_box_h(1), sec.get('height', text_h), text_h)
        _add_rect(slide, x, y, half_w, box_h, line_rgb=BLUE, line_pt=1.0)
        tb = _add_textbox(slide, x + 0.06, y + _WRITE_BOX_PAD,
                           inner_w, box_h - 2 * _WRITE_BOX_PAD)
        tb.text_frame.word_wrap = True
        _set_para(tb.text_frame.paragraphs[0], label, size_pt=size_pt, color=DARK)
        x += half_w + 0.20
        y_next = y + box_h + 0.08
    return y_next


def _render_table(slide, y, sec):
    headers   = sec['headers']
    rows      = sec.get('rows', [])
    # write_rows: set of row indices where children write (affects min row height)
    write_rows = set(sec.get('write_rows', range(len(rows))))
    hdr_col   = sec.get('header_color', BLUE)
    body_size = _safe_pt(sec.get('body_size_pt', 11), MIN_PT_CELL)
    hdr_size  = _safe_pt(sec.get('header_size_pt', 11), MIN_PT_SUBLABEL)
    # Minimum row height: write rows use _write_box_h, display rows use text height
    default_row_h = _write_box_h(1)
    row_h = max(sec.get('row_height', default_row_h), default_row_h)

    col_ws = [CONT_W * h['frac'] for h in headers]
    x = CONT_X
    # Header height: tallest wrapped header across all columns
    hdr_h = max(
        _text_h(h['text'], cw - 0.08, hdr_size) + _TPAD
        for h, cw in zip(headers, col_ws)
    )
    hdr_h = max(hdr_h, 0.30)
    for hdr, cw in zip(headers, col_ws):
        al = PP_ALIGN.CENTER if hdr.get('align') == 'center' else PP_ALIGN.LEFT
        _add_rect(slide, x, y, cw - 0.02, hdr_h, fill_rgb=hdr_col)
        tb = _add_textbox(slide, x + 0.04, y + 0.04, cw - 0.10, hdr_h - 0.06)
        tb.text_frame.word_wrap = True
        _set_para(tb.text_frame.paragraphs[0], hdr['text'],
                  bold=True, size_pt=hdr_size, color=WHITE, align=al)
        x += cw
    y += hdr_h

    for i, row in enumerate(rows):
        fill = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
        x = CONT_X
        # Write rows get minimum write-box height; display rows just need text height
        min_h = _write_box_h(1) if i in write_rows else 0.0
        # Row height is the max of: explicit row_h, write-box minimum, tallest cell text
        this_h = max(
            row_h,
            min_h,
            *(_text_h(cell or '', cw - 0.08, body_size) + _TPAD
              for cell, cw in zip(row, col_ws))
        )
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            al = PP_ALIGN.CENTER if headers[j].get('align') == 'center' else PP_ALIGN.LEFT
            _add_rect(slide, x, y, cw - 0.02, this_h,
                      fill_rgb=fill, line_rgb=LGREY, line_pt=0.5)
            if cell:
                tb = _add_textbox(slide, x + 0.04, y + 0.06, cw - 0.10, this_h - 0.08)
                tb.text_frame.word_wrap = True
                _set_para(tb.text_frame.paragraphs[0], cell,
                          size_pt=body_size, color=DARK, align=al)
            x += cw
        y += this_h
    return y


def _render_graph_template(slide, y, sec, resource_base):
    filename = sec['filename']
    if os.path.isabs(filename):
        graph_path = filename
    else:
        graph_path = os.path.join(resource_base, 'images', filename)

    if not os.path.exists(graph_path):
        os.makedirs(os.path.dirname(graph_path), exist_ok=True)
        fig, ax = plt.subplots(figsize=(6, 3.5))
        ax.set_xlim(*sec.get('xlim', [0, 10]))
        ax.set_ylim(*sec.get('ylim', [0, 100]))
        ax.set_xlabel(sec.get('xlabel', ''), fontsize=11)
        ax.set_ylabel(sec.get('ylabel', ''), fontsize=11)
        ax.set_title(sec.get('title', ''), fontsize=12, fontweight='bold')
        if 'xticks' in sec:
            ax.set_xticks(sec['xticks'])
        if 'yticks' in sec:
            ax.set_yticks(sec['yticks'])
        ax.grid(True, alpha=0.4, linestyle='--')
        for spine in ax.spines.values():
            spine.set_color('#1A3A5C')
        plt.tight_layout()
        plt.savefig(graph_path, dpi=120, bbox_inches='tight', facecolor='white')
        plt.close()

    max_h = sec.get('max_h', 2.0)
    iw, ih = _PILImage.open(graph_path).size
    ratio = iw / ih
    h = max_h
    w = min(h * ratio, CONT_W)
    if w == CONT_W:
        h = w / ratio
    x = CONT_X + (CONT_W - w) / 2
    slide.shapes.add_picture(graph_path, _i(x), _i(y), width=_i(w), height=_i(h))
    return y + h + 0.10


def _render_sentence_starter(slide, y, sec):
    text    = sec['text']
    size_pt = _safe_pt(sec.get('size_pt', 12))
    inner_w = CONT_W - 0.16
    # Height driven by text wrap; never below a write-box minimum
    text_h  = _text_h(text, inner_w, size_pt) + _TPAD
    box_h   = max(_write_box_h(1), text_h + 0.14)  # +0.14 for top/bottom border padding
    _add_rect(slide, CONT_X, y, CONT_W, box_h,
              fill_rgb=RGBColor(0xDE, 0xEC, 0xF8), line_rgb=BLUE, line_pt=1.5)
    tb = _add_textbox(slide, CONT_X + 0.08, y + _WRITE_BOX_PAD,
                       inner_w, box_h - 2 * _WRITE_BOX_PAD)
    tb.text_frame.word_wrap = True
    _set_para(tb.text_frame.paragraphs[0], text, italic=True,
              size_pt=size_pt, color=DARK)
    return y + box_h + 0.08


def _render_sort_table(slide, y, sec):
    materials  = sec['materials']
    size_pt    = _safe_pt(sec.get('size_pt', 11), MIN_PT_SUBLABEL)
    col_labels = sec.get('col_labels', _SORT_COL_LABELS)
    col_ratio  = sec.get('col_ratio', _SORT_COL_RATIO)
    col_ws     = [CONT_W * r for r in col_ratio]
    # Write cells: all columns except the first (pre-filled material name)
    # must be at least _write_box_h(1) tall.
    min_row_h = max(sec.get('row_height', _write_box_h(1)), _write_box_h(1))

    # Header: height from tallest wrapped header label
    hdr_h = max(
        _text_h(lbl, cw - 0.08, size_pt) + _TPAD
        for lbl, cw in zip(col_labels, col_ws)
    )
    hdr_h = max(hdr_h, 0.32)
    x = CONT_X
    for lbl, cw in zip(col_labels, col_ws):
        _add_rect(slide, x, y, cw, hdr_h, fill_rgb=BLUE)
        tb = _add_textbox(slide, x + 0.04, y + 0.04, cw - 0.08, hdr_h - 0.06)
        tb.text_frame.word_wrap = True
        _set_para(tb.text_frame.paragraphs[0], lbl, bold=True, size_pt=size_pt, color=WHITE)
        x += cw
    y += hdr_h

    for i, mat in enumerate(materials):
        # Row height driven by material-name wrap, but never below _write_box_h(1)
        mat_h = _text_h(mat, col_ws[0] - 0.08, size_pt) + _TPAD
        this_h = max(min_row_h, mat_h)
        fill = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
        x = CONT_X
        for j, cw in enumerate(col_ws):
            _add_rect(slide, x, y, cw, this_h, fill_rgb=fill, line_rgb=LGREY, line_pt=0.5)
            if j == 0:
                tb = _add_textbox(slide, x + 0.04, y + _WRITE_BOX_PAD,
                                   cw - 0.08, this_h - 2 * _WRITE_BOX_PAD)
                tb.text_frame.word_wrap = True
                _set_para(tb.text_frame.paragraphs[0], mat, size_pt=size_pt, color=DARK)
            x += cw
        y += this_h
    return y


def _render_sort_table_answers(slide, y, sec):
    """Marking-station version: all cells filled in green. Size ≥ MIN_PT_MARKING."""
    answers  = sec['answers']
    size_pt  = _safe_pt(sec.get('size_pt', 10), MIN_PT_MARKING)
    col_ratio = sec.get('col_ratio', _SORT_COL_RATIO)
    col_ws   = [CONT_W * r for r in col_ratio]

    for row in answers:
        # Row height: tallest cell after text wrap, never below marking min
        this_h = max(
            _text_h(cell, cw - 0.08, size_pt) + _TPAD
            for cell, cw in zip(row, col_ws)
        )
        this_h = max(this_h, 0.28)
        x = CONT_X
        for cell, cw in zip(row, col_ws):
            tb = _add_textbox(slide, x + 0.04, y + 0.04, cw - 0.08, this_h - 0.06)
            tb.text_frame.word_wrap = True
            _set_para(tb.text_frame.paragraphs[0], cell, size_pt=size_pt, color=GREEN)
            x += cw
        y += this_h
    return y


def _render_marking_station(slide, y, sec):
    title = sec.get('title', 'Marking Station')
    tb = _add_textbox(slide, CONT_X, y, CONT_W, 0.40)
    _set_para(tb.text_frame.paragraphs[0], title, bold=True, size_pt=16, color=GREEN)
    return y + 0.45


def _render_answer_text(slide, y, sec):
    text = sec['text']
    bold = sec.get('bold', False)
    size_pt = _safe_pt(sec.get('size_pt', 10), MIN_PT_MARKING)
    w = sec.get('w', CONT_W)
    h = _text_h(text, w, size_pt) + _TPAD
    tb = _add_textbox(slide, CONT_X, y, w, h)
    tb.text_frame.word_wrap = True
    _set_para(tb.text_frame.paragraphs[0], text, bold=bold, size_pt=size_pt, color=GREEN)
    return y + h + 0.06


# ══════════════════════════════════════════════════════════════════
# v2 element renderers
# ══════════════════════════════════════════════════════════════════

def _render_answer_lines(slide, y, sec):
    """
    Labelled ruled write-lines — the primary pupil-response element.

    sec keys:
        label            – question text, rendered bold above the lines (optional)
        n / count        – number of ruled lines (default 3)
        gap_in           – line pitch in inches (default _AL_LINE_PITCH = 0.315 ≈ 8 mm)
        w                – content width in inches (default CONT_W)
        x                – left edge in inches (default CONT_X)
        size_pt          – font size for label (default 11)
        sentence_starter – italic one-line prompt printed above the first line (optional)

    Y-advance breakdown (all constants documented at top of file):
        [label present]   _text_h(label) + _TPAD   +  _AL_LABEL_LINE_GAP
        [starter present] _text_h(starter) + _TPAD  +  _AL_STARTER_LINE_GAP
        [lines]           n × gap_in
        [post-block]      _AL_POST_GAP
    """
    label   = sec.get('label', '')
    n       = sec.get('n', sec.get('count', 3))
    gap     = sec.get('gap_in', _AL_LINE_PITCH)
    w       = sec.get('w', CONT_W)
    x       = sec.get('x', CONT_X)
    size_pt = _safe_pt(sec.get('size_pt', 12))       # label: body minimum
    starter = sec.get('sentence_starter', '')

    # ── 1. Label ──────────────────────────────────────────────────────────────
    if label:
        h = _text_h(label, w, size_pt) + _TPAD
        tb = _add_textbox(slide, CONT_X, y, w, h)
        tb.text_frame.word_wrap = True
        _set_para(tb.text_frame.paragraphs[0], label, bold=True,
                  size_pt=size_pt, color=DARK)
        y += h + _AL_LABEL_LINE_GAP

    # ── 2. Sentence starter ───────────────────────────────────────────────────
    if starter:
        st_pt = _safe_pt(size_pt - 0.5, MIN_PT_SUBLABEL)   # never below sub-label min
        h = _text_h(starter, w, st_pt) + _TPAD
        tb = _add_textbox(slide, CONT_X, y, w, h)
        tb.text_frame.word_wrap = True
        _set_para(tb.text_frame.paragraphs[0], starter, italic=True,
                  size_pt=st_pt, color=BLUE)
        y += h + _AL_STARTER_LINE_GAP

    # ── 3. Ruled lines ────────────────────────────────────────────────────────
    for _ in range(n):
        _add_line(slide, x, y, w)
        y += gap

    # ── 4. Post-block breathing room ─────────────────────────────────────────
    return y + _AL_POST_GAP


def _render_cloze(slide, y, sec, resource_base=None):
    """
    Cloze passage with ___ blanks.
    sec = {"type": "cloze",
           "text": "The ___ began with the ___.",
           "blanks": ["universe", "Big Bang"],
           "word_bank": "Big Bang  •  universe  •  Milky Way",
           "show_answers": false}
    On the task slide: renders text as-is with blanks shown as ___.
    On the marking slide: caller sets show_answers=True and blanks replace
    ___ in green.
    """
    text     = sec['text']
    blanks   = sec.get('blanks', [])
    size_pt  = _safe_pt(sec.get('size_pt', 12))
    w        = sec.get('w', CONT_W)
    show_ans = sec.get('show_answers', False)

    # Height from text-wrap model; never below MIN_PT_BODY
    h = _text_h(text, w, size_pt) + _TPAD

    if show_ans:
        parts = text.split('___')
        tb = _add_textbox(slide, CONT_X, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        for i, part in enumerate(parts):
            r = para.add_run()
            r.text = part
            r.font.name = FONT_BODY
            r.font.size = Pt(size_pt)
            r.font.color.rgb = DARK
            if i < len(blanks):
                rb = para.add_run()
                rb.text = blanks[i]
                rb.font.name = FONT_BODY
                rb.font.size = Pt(size_pt)
                rb.font.bold = True
                rb.font.color.rgb = GREEN
        y += h + 0.06
    else:
        tb = _add_textbox(slide, CONT_X, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        _set_para(tf.paragraphs[0], text, size_pt=size_pt, color=DARK)
        y += h + 0.06

    wb = sec.get('word_bank', '')
    if wb and not show_ans:
        y = _render_word_bank(slide, y, {'words': wb, 'size_pt': size_pt - 1})

    return y


def _render_matching(slide, y, sec, resource_base=None):
    """
    Two-column matching activity.
    sec = {"type": "matching",
           "left": ["universe", "galaxy", "Big Bang"],
           "right": ["everything that exists", "huge collection of stars",
                     "theory of how universe began"],
           "answer_pairs": [[0,0],[1,1],[2,2]]}
    Renders left column (bold term) and right column (definition).
    Pupil draws arrows in the gap. On marking station, show_answers=True
    adds answer_pairs as coloured bullet connections.
    """
    left  = sec['left']
    right = sec['right']
    size_pt = _safe_pt(sec.get('size_pt', 11), MIN_PT_SUBLABEL)
    # Minimum row height: always at least one write-box height so children can draw arrows
    min_row_h = max(sec.get('row_height', _write_box_h(1)), _write_box_h(1))

    col_left_w  = CONT_W * 0.32
    col_mid_w   = CONT_W * 0.08
    col_right_w = CONT_W * 0.60

    # Column headers
    hdr_h = max(
        _text_h('Word',    col_left_w  - 0.08, size_pt) + _TPAD,
        _text_h('Meaning', col_right_w - 0.08, size_pt) + _TPAD,
        0.28,
    )
    tb_l = _add_textbox(slide, CONT_X, y, col_left_w, hdr_h)
    _set_para(tb_l.text_frame.paragraphs[0], 'Word', bold=True, size_pt=size_pt, color=BLUE)
    tb_r = _add_textbox(slide, CONT_X + col_left_w + col_mid_w, y, col_right_w, hdr_h)
    _set_para(tb_r.text_frame.paragraphs[0], 'Meaning', bold=True, size_pt=size_pt, color=BLUE)
    y += hdr_h + 0.04

    n = max(len(left), len(right))
    for i in range(n):
        l_text = left[i]  if i < len(left)  else ''
        r_text = right[i] if i < len(right) else ''

        lh = (_text_h(l_text, col_left_w  - 0.08, size_pt) + _TPAD) if l_text else 0
        rh = (_text_h(r_text, col_right_w - 0.08, size_pt) + _TPAD) if r_text else 0
        this_h = max(min_row_h, lh, rh)
        fill = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE

        # Left cell
        _add_rect(slide, CONT_X, y, col_left_w, this_h,
                  fill_rgb=fill, line_rgb=LGREY, line_pt=0.5)
        if l_text:
            tb = _add_textbox(slide, CONT_X + 0.05, y + 0.05,
                               col_left_w - 0.10, this_h - 0.08)
            tb.text_frame.word_wrap = True
            _set_para(tb.text_frame.paragraphs[0], l_text,
                      size_pt=size_pt, color=DARK, bold=True)

        # Right cell
        rx = CONT_X + col_left_w + col_mid_w
        _add_rect(slide, rx, y, col_right_w, this_h,
                  fill_rgb=fill, line_rgb=LGREY, line_pt=0.5)
        if r_text:
            tb = _add_textbox(slide, rx + 0.05, y + 0.05,
                               col_right_w - 0.10, this_h - 0.08)
            tb.text_frame.word_wrap = True
            _set_para(tb.text_frame.paragraphs[0], r_text,
                      size_pt=size_pt, color=DARK)

        # Arrow hint dot in gap
        mid_x = CONT_X + col_left_w + 0.02
        _add_rect(slide, mid_x, y + this_h / 2 - 0.02, col_mid_w - 0.04, 0.04,
                  fill_rgb=LGREY)

        y += this_h + 0.04

    return y + 0.06


def _render_diagram(slide, y, sec, resource_base):
    """
    Image or placeholder. Gracefully handles missing files.
    sec = {"type": "diagram", "path": "images/foo.png",
           "description": "diagram of the solar system", "max_h": 2.0}
    """
    path = sec.get('path', '')
    if path:
        try:
            return _render_reference_image(slide, y, sec, resource_base)
        except FileNotFoundError:
            pass  # fall through to placeholder

    desc  = sec.get('description', 'Diagram')
    box_h = sec.get('max_h', sec.get('max_h', 1.5))
    _add_rect(slide, CONT_X, y, CONT_W, box_h, line_rgb=LGREY, line_pt=1.0)
    tb = _add_textbox(slide, CONT_X + 0.2, y + box_h / 2 - 0.15,
                       CONT_W - 0.4, 0.30)
    _set_para(tb.text_frame.paragraphs[0],
              f'[{desc}]', italic=True, size_pt=11, color=LGREY,
              align=PP_ALIGN.CENTER)
    return y + box_h + 0.10


def _render_image(slide, y, sec, resource_base):
    """Alias for diagram (element type 'image')."""
    return _render_diagram(slide, y, sec, resource_base)


# ── dispatch ────────────────────────────────────────────────────────────────
_SECTION_RENDERERS = {
    # legacy types
    'heading':            lambda sl, y, sec, rb: _render_heading(sl, y, sec),
    'instruction':        lambda sl, y, sec, rb: _render_instruction(sl, y, sec),
    'write_lines':        lambda sl, y, sec, rb: _render_write_lines(sl, y, sec),
    'word_bank':          lambda sl, y, sec, rb: _render_word_bank(sl, y, sec),
    'reference_image':    _render_reference_image,
    'row_boxes':          lambda sl, y, sec, rb: _render_row_boxes(sl, y, sec),
    'pair_boxes':         lambda sl, y, sec, rb: _render_pair_boxes(sl, y, sec),
    'table':              lambda sl, y, sec, rb: _render_table(sl, y, sec),
    'graph_template':     _render_graph_template,
    'sentence_starter':   lambda sl, y, sec, rb: _render_sentence_starter(sl, y, sec),
    'spacer':             lambda sl, y, sec, rb: y + sec.get('h', 0.1),
    'sort_table':         lambda sl, y, sec, rb: _render_sort_table(sl, y, sec),
    'sort_table_answers': lambda sl, y, sec, rb: _render_sort_table_answers(sl, y, sec),
    'marking_station':    lambda sl, y, sec, rb: _render_marking_station(sl, y, sec),
    'answer_text':        lambda sl, y, sec, rb: _render_answer_text(sl, y, sec),
    # v2 types
    'answer_lines':       lambda sl, y, sec, rb: _render_answer_lines(sl, y, sec),
    'cloze':              _render_cloze,
    'matching':           _render_matching,
    'diagram':            _render_diagram,
    'image':              _render_image,
    # v2 aliases
    'timeline_diagram':   _render_diagram,
    'graph_axes':         _render_graph_template,
    'sorting_record':     lambda sl, y, sec, rb: _render_sort_table(sl, y, sec),
}


def _render_section(slide, sec, y, resource_base):
    t = sec['type']
    renderer = _SECTION_RENDERERS.get(t)
    if renderer is None:
        raise ValueError(f"Unknown LP section type: '{t}'")
    return renderer(slide, y, sec, resource_base)


# ══════════════════════════════════════════════════════════════════
# Label builder
# ══════════════════════════════════════════════════════════════════

def _add_label(slide, lf, ican1, ican2, date_str, key_question,
               resource_base, png_dest, year='Y4'):
    """Place the WFA enquiry label and return the y position below it."""
    try:
        if resource_base not in sys.path:
            sys.path.insert(0, resource_base)
        import build_enquiry_label as _bel
        _bel.ASSETS = os.path.join(resource_base, 'll_assets')
        from label_builder import build_enquiry_label, LL_W

        label_h = build_enquiry_label(
            slide,
            SW_IN - LL_W - MARGIN,
            LBL_Y,
            date_str=date_str,
            key_q=key_question,
            lf=lf,
            ican1=ican1,
            ican2=ican2,
            icon_path=None,
            subject='scientist',
            year=year,
            png_dest=png_dest,
        )
        LABEL_SCALE = 0.707
        label_pic = slide.shapes[-1]
        new_w = LL_W * LABEL_SCALE
        new_h = label_h * LABEL_SCALE
        label_pic.left   = _i(SW_IN - new_w - MARGIN)
        label_pic.top    = _i(LBL_Y)
        label_pic.width  = _i(new_w)
        label_pic.height = _i(new_h)
        return LBL_Y + new_h + 0.15
    except (ImportError, ModuleNotFoundError, FileNotFoundError):
        # label_builder not available in this environment — skip gracefully
        return LBL_Y + 0.10


# ══════════════════════════════════════════════════════════════════
# Level badge (v2 only) — small coloured tag top-right of slide 1
# ══════════════════════════════════════════════════════════════════

def _add_level_badge(slide, level):
    colour = LEVEL_COLOURS.get(level, BLUE)
    label  = LEVEL_LABELS.get(level, level.replace('_', ' ').title())
    badge_w, badge_h = 1.5, 0.28
    bx = SW_IN - badge_w - MARGIN
    by = SH_IN - badge_h - MARGIN
    _add_rect(slide, bx, by, badge_w, badge_h, fill_rgb=colour)
    tb = _add_textbox(slide, bx + 0.06, by + 0.04, badge_w - 0.12, badge_h - 0.08)
    _set_para(tb.text_frame.paragraphs[0], label,
              bold=True, size_pt=10, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# Schema detection
# ══════════════════════════════════════════════════════════════════

def _is_new_schema(lp_spec):
    """Return True if lp_spec uses the three-level standard/adapted/further_adapted schema."""
    return 'standard' in lp_spec or 'adapted' in lp_spec or 'further_adapted' in lp_spec


# ══════════════════════════════════════════════════════════════════
# Core builder (internal)
# ══════════════════════════════════════════════════════════════════

def _build_one_level(lesson_data, level_spec, level, lp_top, out_path,
                     resource_base, lesson_meta, year='Y5'):
    """
    Build a single-level LP PPTX from a level_spec dict
    (the standard/adapted/further_adapted sub-dict of lp).

    lesson_meta: dict with optional keys: key_question, date, lf, ican1, ican2
    """
    prs = _new_prs()
    s1, s2 = prs.slides[0], prs.slides[1]
    _clear_slide(s1)
    _clear_slide(s2)

    # Resolve label fields: level_spec first, then lesson_meta fallback
    lf     = level_spec.get('lf')     or lesson_meta.get('lf', '')
    ican1  = level_spec.get('ican1')  or lesson_meta.get('ican1', '')
    ican2  = level_spec.get('ican2')  or lesson_meta.get('ican2', '')
    date_s = lp_top.get('date', '')   or lesson_meta.get('date', '')
    key_q  = lesson_meta.get('key_question', '')

    _tmp_dir = '/tmp'
    png_dest = os.path.join(_tmp_dir, f'_lp_label_{os.getpid()}_{level}.png')
    y1 = _add_label(s1, lf, ican1, ican2, date_s, key_q,
                    resource_base, png_dest, year=year)

    # Level badge on slide 1
    _add_level_badge(s1, level)

    # Slide 1: elements up to first 'marking_station'
    elements = level_spec.get('elements', [])
    try:
        split_at = next(i for i, e in enumerate(elements)
                        if e['type'] == 'marking_station')
    except StopIteration:
        split_at = len(elements)

    s1_elements = elements[:split_at]
    s2_elements_from_list = elements[split_at:]

    for el in s1_elements:
        y1 = _render_section(s1, el, y1, resource_base)

    # Slide 2: explicit answers list (preferred) or remainder of elements list
    answers = level_spec.get('answers', [])
    y2 = 0.30
    # Auto marking station header
    y2 = _render_marking_station(s2, y2, {'title': 'Marking Station'})
    _add_level_badge(s2, level)

    if answers:
        for el in answers:
            y2 = _render_section(s2, el, y2, resource_base)
    else:
        for el in s2_elements_from_list:
            y2 = _render_section(s2, el, y2, resource_base)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)

    try:
        if os.path.exists(png_dest):
            os.remove(png_dest)
    except OSError:
        pass

    size_kb = os.path.getsize(out_path) // 1024
    print(f"  LP [{level}] -> {out_path} ({size_kb} KB)")
    return out_path


def _build_legacy(lesson_data, lp_spec, out_path, resource_base):
    """Build using the original flat-sections schema (unchanged behaviour)."""
    lesson = lesson_data['lesson']

    if resource_base not in sys.path:
        sys.path.insert(0, resource_base)
    try:
        import build_enquiry_label as _bel
        _bel.ASSETS = os.path.join(resource_base, 'll_assets')
    except ImportError:
        pass

    prs = _new_prs()
    s1, s2 = prs.slides[0], prs.slides[1]
    _clear_slide(s1)
    _clear_slide(s2)

    _tmp_dir = '/tmp'
    png_dest = os.path.join(_tmp_dir, f'_lp_label_{os.getpid()}.png')

    lf    = lp_spec.get('lf') or lesson.get('lo', '')
    ican1 = lp_spec.get('ican1', '')
    ican2 = lp_spec.get('ican2', '')
    date_s = lp_spec.get('date', '')
    key_q  = lesson.get('key_question', '')

    y1 = _add_label(s1, lf, ican1, ican2, date_s, key_q,
                    resource_base, png_dest, year='Y4')

    all_sections = lp_spec.get('sections', [])
    try:
        split_at = next(i for i, s in enumerate(all_sections)
                        if s['type'] == 'marking_station')
    except StopIteration:
        split_at = len(all_sections)

    for sec in all_sections[:split_at]:
        y1 = _render_section(s1, sec, y1, resource_base)

    y2 = 0.30
    for sec in all_sections[split_at:]:
        y2 = _render_section(s2, sec, y2, resource_base)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)

    try:
        if os.path.exists(png_dest):
            os.remove(png_dest)
    except OSError:
        pass

    size_kb = os.path.getsize(out_path) // 1024
    print(f"  LP -> {out_path} ({size_kb} KB)")
    return out_path


# ══════════════════════════════════════════════════════════════════
# Public entry points
# ══════════════════════════════════════════════════════════════════

def build_lp(lesson_json_or_path, out_path, resource_base=None, level='standard'):
    """
    Build the LP PPTX for a single lesson (one differentiation level).

    Args:
        lesson_json_or_path: path to lesson JSON or already-loaded dict.
                             For the new MTP schema, pass the full MTP dict
                             with 'lessons' list, or a single lesson dict
                             that has a 'lp' key.
        out_path:            where to write the output PPTX.
        resource_base:       directory for images, ll_assets, etc.
        level:               'standard' | 'adapted' | 'further_adapted'
                             (ignored for legacy flat-sections schema)

    Returns:
        out_path on success.
    """
    # Load JSON
    if isinstance(lesson_json_or_path, (str, os.PathLike)):
        json_path = str(lesson_json_or_path)
        with open(json_path) as f:
            lesson_json = json.load(f)
        if resource_base is None:
            resource_base = os.path.dirname(os.path.abspath(json_path))
    else:
        lesson_json = lesson_json_or_path
        if resource_base is None:
            resource_base = _DEFAULT_RESOURCE_BASE

    # Support both wrapped {"lesson": {...}} and bare lesson dicts
    if 'lesson' in lesson_json:
        lesson_data = lesson_json
        lp_spec = lesson_data['lesson'].get('lp')
    elif 'lp' in lesson_json:
        lesson_data = {'lesson': lesson_json}
        lp_spec = lesson_json['lp']
    else:
        raise ValueError("Cannot locate 'lp' key in provided data")

    if lp_spec is None:
        raise ValueError("Lesson has no 'lp' key — add an lp spec before calling build_lp()")

    if _is_new_schema(lp_spec):
        level_spec = lp_spec.get(level)
        if level_spec is None:
            raise ValueError(
                f"Level '{level}' not found in lp spec. "
                f"Available: {[k for k in lp_spec if k in ('standard','adapted','further_adapted')]}")
        lesson = lesson_data['lesson']
        year = lesson.get('year_group', 'Y5')
        lesson_meta = {
            'key_question': lesson.get('key_question', lesson.get('building_block_text', '')),
            'date': lp_spec.get('date', ''),
            'lf':   lesson.get('what', lesson.get('lo', '')),
            'ican1': '',
            'ican2': '',
        }
        # Split success criteria into two I can statements if possible
        success = lesson.get('success', '')
        if success:
            parts = success.split(',')
            lesson_meta['ican1'] = parts[0].strip()
            lesson_meta['ican2'] = ', '.join(parts[1:]).strip() if len(parts) > 1 else ''
        return _build_one_level(lesson_data, level_spec, level, lp_spec,
                                out_path, resource_base, lesson_meta, year=year)
    else:
        return _build_legacy(lesson_data, lp_spec, out_path, resource_base)


def build_lp_all_levels(lesson_json_or_path, out_dir, resource_base=None,
                         base_name=None):
    """
    Build all three LP levels from a lesson with a new three-level lp spec.

    Args:
        lesson_json_or_path: path or dict (same as build_lp).
        out_dir:             directory to write the three PPTX files.
        resource_base:       images / assets directory.
        base_name:           filename prefix, e.g. 'L1_Universe_LP'.
                             Defaults to 'LP'.

    Returns:
        dict: {'standard': path, 'adapted': path, 'further_adapted': path}
              Only levels present in the lp spec are included.
    """
    if isinstance(lesson_json_or_path, (str, os.PathLike)):
        json_path = str(lesson_json_or_path)
        with open(json_path) as f:
            lesson_json = json.load(f)
        if resource_base is None:
            resource_base = os.path.dirname(os.path.abspath(json_path))
    else:
        lesson_json = lesson_json_or_path
        if resource_base is None:
            resource_base = _DEFAULT_RESOURCE_BASE

    if 'lesson' in lesson_json:
        lesson_data = lesson_json
    elif 'lp' in lesson_json:
        lesson_data = {'lesson': lesson_json}
    else:
        raise ValueError("Cannot locate lesson data in provided dict")

    lp_spec = lesson_data['lesson'].get('lp')
    if lp_spec is None:
        raise ValueError("Lesson has no 'lp' key")
    if not _is_new_schema(lp_spec):
        raise ValueError("build_lp_all_levels() requires new three-level schema. "
                         "For legacy schema use build_lp().")

    os.makedirs(out_dir, exist_ok=True)
    prefix = base_name or 'LP'
    level_suffixes = {
        'standard':        '_S',
        'adapted':         '_A',
        'further_adapted': '_FA',
    }

    results = {}
    for level in ('standard', 'adapted', 'further_adapted'):
        if level not in lp_spec:
            continue
        out_path = os.path.join(out_dir, f'{prefix}{level_suffixes[level]}.pptx')
        build_lp(lesson_data, out_path, resource_base=resource_base, level=level)
        results[level] = out_path

    return results


if __name__ == '__main__':
    import argparse
    ap = argparse.ArgumentParser(description='Build an enquiry LP from a lesson JSON.')
    ap.add_argument('lesson_json')
    ap.add_argument('out_pptx')
    ap.add_argument('--resource-base', default=None)
    ap.add_argument('--level', default='standard',
                    choices=['standard', 'adapted', 'further_adapted'],
                    help='Differentiation level (new schema only)')
    ap.add_argument('--all-levels', action='store_true',
                    help='Build all three levels into the same directory as out_pptx')
    args = ap.parse_args()

    if args.all_levels:
        out_dir = os.path.dirname(os.path.abspath(args.out_pptx))
        base = os.path.splitext(os.path.basename(args.out_pptx))[0]
        paths = build_lp_all_levels(args.lesson_json, out_dir,
                                     resource_base=args.resource_base,
                                     base_name=base)
        for lvl, p in paths.items():
            print(f"  {lvl}: {p}")
    else:
        build_lp(args.lesson_json, args.out_pptx,
                 resource_base=args.resource_base, level=args.level)

#!/usr/bin/env python3
"""
Build T6W4 LP4–LP6 (standard + adapted) as PPTX.
Correct WFA Set 1 Enquiry learning label: top-right, 2.75" x 1.20", no borders, no colour.
Globe icon + Aptos plain text.
Slide: 7.5" x 10.833" (A4 portrait matching school files).
"""
import os, sys, base64
sys.path.insert(0, '/home/claude')
from label_builder import build_enquiry_label, LL_W, LL_H
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

ICON_PATH = '/home/claude/ll_icons/geographer.png'

# ── Slide dimensions ─────────────────────────────────────────────
SW_IN  = 7.5
SH_IN  = 10.833
SW     = int(SW_IN * 914400)
SH     = int(SH_IN * 914400)
CM     = 1 / 2.54
MARGIN = 0.25

# ── Learning label position ───────────────────────────────────────
LBL_X  = SW_IN - LL_W - MARGIN
LBL_Y  = MARGIN

# Content area
CONT_X = MARGIN
CONT_Y = LBL_Y + LL_H + 0.15
CONT_W = SW_IN - 2 * MARGIN

# ── Colours ──────────────────────────────────────────────────────
BLUE   = RGBColor(0x17, 0x98, 0xD3)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREEN  = RGBColor(0x4F, 0xAD, 0x5B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xCC, 0xCC, 0xCC)
EXBOOK = RGBColor(0xBD, 0xD7, 0xEE)  # pale blue exercise-book lines
LTBLUE = RGBColor(0xEA, 0xF5, 0xFB)
CREAM  = RGBColor(0xFE, 0xF9, 0xE7)

# ── Fonts ────────────────────────────────────────────────────────
FONT_BODY  = 'Twinkl Cursive Looped'
FONT_LABEL = 'Aptos'

import sys; sys.path.insert(0, '/home/claude')
from label_builder import build_enquiry_label, LL_W, LL_H

ICON_PATH = '/home/claude/ll_icons/geographer.png'
OUT = '/mnt/user-data/outputs'
os.makedirs(OUT, exist_ok=True)


# ══════════════════════════════════════════════════════════════════
# PPTX helpers
# ══════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)
    blank = prs.slide_layouts[6]
    prs.slides.add_slide(blank)
    prs.slides.add_slide(blank)
    return prs

def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp','pic','graphicFrame','grpSp','cxnSp'):
            sp_tree.remove(child)

def _i(v): return Inches(v)

def add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))

def set_para(para, text, bold=False, italic=False, size_pt=10,
             color=None, font=FONT_BODY, align=PP_ALIGN.LEFT, underline=False):
    para.alignment = align
    run = para.add_run()
    run.text = text
    rf = run.font
    rf.name = font
    rf.size = Pt(size_pt)
    rf.bold = bold
    rf.italic = italic
    rf.underline = underline
    if color:
        rf.color.rgb = color
    return run

def add_line(slide, x, y, w, color=None, width_pt=1.0):
    """Plain horizontal line injected as raw XML — no theme style, no shadow ever."""
    from pptx.util import Inches as _In, Pt as _Pt
    from lxml import etree as _et

    c = color if color is not None else EXBOOK
    hex_col = f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'
    lw = int(_Pt(width_pt))

    sp_tree = slide.shapes._spTree
    uid = max((int(el.get('id', 0))
               for el in sp_tree.iter() if el.get('id') and el.get('id').isdigit()),
              default=100) + 1

    xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="{uid}" name="Line{uid}"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'<p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{int(_In(x))}" y="{int(_In(y))}"/>'
        f'<a:ext cx="{int(_In(w))}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{lw}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:ln>'
        f'<a:effectLst/>'
        f'</p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        f'</p:sp>'
    )
    sp_tree.append(_et.fromstring(xml))

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, _i(x), _i(y), _i(w), _i(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════════════
# LEARNING LABEL (Set 1 Enquiry — top-right, plain Aptos text)
# Layout within label box (2.75" × 1.20"):
#   Globe icon shape on far right (0.30" wide)
#   Date text (small, top-left of label)
#   "Key Question" bold underline
#   Question text bold underline
#   LF: ... plain
#   I can 1 plain
#   I can 2 plain
# ══════════════════════════════════════════════════════════════════
LABEL_SZ      = 6.5    # pt — Aptos label font size
LABEL_SZ_DATE = 6.0    # pt


# ══════════════════════════════════════════════════════════════════
# CONTENT HELPERS
# ══════════════════════════════════════════════════════════════════

def heading(slide, y, text, w=None):
    """Blue bold heading."""
    w = w or CONT_W
    tb = add_textbox(slide, CONT_X, y, w, 0.30)
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf.paragraphs[0], text, bold=True, size_pt=12,
             color=BLUE, font=FONT_BODY)
    return y + 0.33

def instruction(slide, y, text, w=None):
    """Instruction text."""
    w = w or CONT_W
    tb = add_textbox(slide, CONT_X, y, w, 0.25)
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf.paragraphs[0], text, size_pt=9.5, color=DARK, font=FONT_BODY)
    return y + 0.30

def body_text(slide, y, text, bold=False, sz=10, col=None, x=None, w=None):
    col = col or DARK
    x = x or CONT_X
    w = w or CONT_W
    tb = add_textbox(slide, x, y, w, 0.25)
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf.paragraphs[0], text, bold=bold, size_pt=sz, color=col, font=FONT_BODY)
    return y + 0.28

def write_lines(slide, y, n=3, w=None, x=None):
    """n writing lines, exactly 8mm (0.315") apart."""
    w = w or CONT_W
    x = x if x is not None else CONT_X
    for _ in range(n):
        add_line(slide, x, y, w)
        y += 0.315
    return y + 0.10

def word_bank(slide, y, words, w=None):
    """Orange word bank box."""
    w = w or CONT_W
    box = add_rect(slide, CONT_X, y, w, 0.37,
                   fill_rgb=CREAM, line_rgb=ORANGE, line_pt=1.0)
    tb = add_textbox(slide, CONT_X + 0.05, y + 0.04, w - 0.1, 0.29)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'Word bank: '
    r1.font.name = FONT_BODY
    r1.font.bold = True
    r1.font.size = Pt(9)
    r1.font.color.rgb = ORANGE
    r2 = p.add_run()
    r2.text = words
    r2.font.name = FONT_BODY
    r2.font.size = Pt(9)
    r2.font.color.rgb = DARK
    return y + 0.41

def marking_heading(slide, y, text='Marking Station'):
    tb = add_textbox(slide, CONT_X, y, CONT_W, 0.40)
    tf = tb.text_frame
    set_para(tf.paragraphs[0], text, bold=True, size_pt=16, color=GREEN, font=FONT_BODY)
    return y + 0.45

def answer_text(slide, y, text, bold=False):
    tb = add_textbox(slide, CONT_X, y, CONT_W, 0.25)
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf.paragraphs[0], text, bold=bold, size_pt=9.5, color=GREEN, font=FONT_BODY)
    return y + 0.27


# ── Table helpers ──────────────────────────────────────────────────
COL_WIDTHS_IN = [
    CONT_W * (2800000/5960000),   # Land use column
    CONT_W * (1400000/5960000),   # E/B/Both column
    CONT_W * (1760000/5960000),   # Type column
]

def table_header(slide, y):
    labels = ['Land use', 'England / Brazil / Both', 'Type']
    x = CONT_X
    for lbl, cw in zip(labels, COL_WIDTHS_IN):
        box = add_rect(slide, x, y, cw, 0.26, fill_rgb=BLUE)
        tb = add_textbox(slide, x + 0.04, y + 0.03, cw - 0.08, 0.20)
        tf = tb.text_frame
        set_para(tf.paragraphs[0], lbl, bold=True, size_pt=8.5,
                 color=WHITE, font=FONT_BODY)
        x += cw
    return y + 0.26

def table_row(slide, y, row_text, even=True):
    label_fill = RGBColor(0xF5, 0xF5, 0xF5) if even else WHITE
    x = CONT_X
    for i, (txt, cw) in enumerate(zip([row_text, '', ''], COL_WIDTHS_IN)):
        # Land use column keeps alternating fill; answer columns always white
        cell_fill = label_fill if i == 0 else WHITE
        add_rect(slide, x, y, cw, 0.22,
                 fill_rgb=cell_fill, line_rgb=LGREY, line_pt=0.5)
        if txt:
            tb = add_textbox(slide, x + 0.04, y + 0.03, cw - 0.08, 0.16)
            tf = tb.text_frame
            set_para(tf.paragraphs[0], txt, size_pt=8.0, color=DARK, font=FONT_BODY)
        x += cw
    return y + 0.22

def save_pptx(prs, filename):
    path = os.path.join(OUT, filename)
    prs.save(path)
    print(f'  {filename}: {os.path.getsize(path)//1024}KB')
    return path


# ══════════════════════════════════════════════════════════════════
# LP4 STANDARD
# ══════════════════════════════════════════════════════════════════
def build_lp4_standard():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    build_enquiry_label(s1, LBL_X, LBL_Y, '06/07/2026',
        'Are England and Brazil different?',
        'describe and compare land use in England and Brazil',
        'describe land use in Brazil',
        'compare land use using geographical vocabulary',
        ICON_PATH)

    y = CONT_Y
    y = heading(s1, y, 'Part A   Land use sort')
    y = instruction(s1, y, 'For each land use below, write E (England), B (Brazil) or Both. Then write the type of land use.')

    y = table_header(s1, y)
    rows = [
        'Coffee plantation','Coal mine / quarry','Cattle ranch / farmland',
        'Offshore wind farm','Iron ore mine','Arable crop field (wheat, barley)',
        'Oil rig or power station','Terraced housing / urban suburb',
        'Hydro-electric dam','Shopping centre / commercial area',
        'Port / container terminal','Moorland / national park',
    ]
    for ri, row in enumerate(rows):
        y = table_row(s1, y, row, ri % 2 == 0)

    y += 0.10
    y = heading(s1, y, 'Part B   Comparison sentences')
    y = instruction(s1, y, 'Write one sentence about what England and Brazil have in common and one about how they differ. Use at least one word from the vocabulary bank.')
    y = body_text(s1, y, 'England and Brazil are similar because…', bold=True)
    y = write_lines(s1, y, 3)
    y = body_text(s1, y, 'However, they are different because…', bold=True)
    y = write_lines(s1, y, 3)
    y += 0.05
    word_bank(s1, y, 'land use  •  natural resource  •  trade  •  economic activity  •  agricultural  •  industrial')

    # Marking station slide 2
    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part A   Answers')
    answers = [
        ('Coffee plantation', 'B', 'Agricultural'),
        ('Coal mine / quarry', 'E (or Both)', 'Extractive / industrial'),
        ('Cattle ranch / farmland', 'Both', 'Agricultural'),
        ('Offshore wind farm', 'E', 'Energy / industrial'),
        ('Iron ore mine', 'B (mainly)', 'Extractive'),
        ('Arable crop field', 'Both', 'Agricultural'),
        ('Oil rig / power station', 'Both', 'Energy'),
        ('Terraced housing', 'E', 'Residential / urban'),
        ('Hydro-electric dam', 'B', 'Energy'),
        ('Shopping centre', 'Both', 'Commercial / service'),
        ('Port / terminal', 'Both', 'Industrial / transport'),
        ('Moorland / national park', 'E', 'Conservation / recreational'),
    ]
    for row_t, eb, typ in answers:
        y = answer_text(s2, y, f'{row_t}  →  {eb}  |  {typ}')

    save_pptx(prs, 'T6W4_LP4_Geographers_Human_Geography.pptx')


# ══════════════════════════════════════════════════════════════════
# LP4 ADAPTED
# ══════════════════════════════════════════════════════════════════
def build_lp4_adapted():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    build_enquiry_label(s1, LBL_X, LBL_Y, '06/07/2026',
        'Are England and Brazil different?',
        'describe land use in England and Brazil',
        'describe land use in Brazil',
        'compare land use in both countries',
        ICON_PATH)

    y = CONT_Y
    y = heading(s1, y, 'Part A   Land use match')
    y = instruction(s1, y, 'Draw a line to match each land use to the correct country.')

    col_w = CONT_W / 3
    labels = ['England', 'Both', 'Brazil']
    x_starts = [CONT_X + i * col_w for i in range(3)]
    for lbl, x in zip(labels, x_starts):
        add_rect(s1, x, y, col_w, 0.26, fill_rgb=BLUE)
        tb = add_textbox(s1, x + 0.04, y + 0.03, col_w - 0.08, 0.20)
        set_para(tb.text_frame.paragraphs[0], lbl, bold=True, size_pt=9,
                 color=WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER)
    y += 0.26

    items_e    = ['Offshore wind farm', 'Terraced housing', 'Moorland / national park']
    items_both = ['Cattle ranch / farmland', 'Port / container terminal', 'Arable crop field']
    items_b    = ['Coffee plantation', 'Iron ore mine', 'Hydro-electric dam']
    for ri in range(3):
        fill = RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else WHITE
        for items, x in zip([items_e, items_both, items_b], x_starts):
            add_rect(s1, x, y, col_w, 0.23, fill_rgb=fill, line_rgb=LGREY, line_pt=0.5)
            tb = add_textbox(s1, x + 0.04, y + 0.04, col_w - 0.08, 0.19)
            set_para(tb.text_frame.paragraphs[0], items[ri], size_pt=8, color=DARK, font=FONT_BODY)
        y += 0.23

    y += 0.12
    y = heading(s1, y, 'Part B   Cloze sentences')
    y = instruction(s1, y, 'Fill in the missing words using the word bank.')
    for line in [
        'Brazil and England both use land for _________________________.',
        'England uses land for _________________________ (e.g. wind farms and quarries).',
        'Brazil mainly uses land for _________________________ (e.g. coffee plantations).',
        'One difference is that Brazil has more _________________________.',
    ]:
        y = body_text(s1, y, line, sz=9.5)
    y += 0.05
    word_bank(s1, y, 'farming  •  industry  •  agriculture  •  mining  •  energy  •  natural resources')

    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part A   Match answers')
    for country, items in [('England', items_e), ('Both', items_both), ('Brazil', items_b)]:
        y = answer_text(s2, y, f'{country}: {", ".join(items)}', bold=True)
    y += 0.1
    y = heading(s2, y, 'Part B   Model answers')
    for a in ['farming / agriculture', 'industry (or energy)', 'agriculture', 'natural resources / mining']:
        y = answer_text(s2, y, f'→ {a}')

    save_pptx(prs, 'T6W4_LP4_Geographers_Human_Geography_adapted.pptx')


# ══════════════════════════════════════════════════════════════════
# LP5 STANDARD
# ══════════════════════════════════════════════════════════════════
def build_lp5_standard():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    build_enquiry_label(s1, LBL_X, LBL_Y, '07/07/2026',
        'Are England and Brazil different?',
        'use maps to describe places',
        'read a grid reference',
        'describe what a map shows about a place',
        ICON_PATH)

    y = CONT_Y
    y = heading(s1, y, 'Part A   Grid reference questions')
    y = instruction(s1, y, 'Use the Westhaven map on the board to answer these questions.')

    for qn, opts in [
        ('1. What is found at grid reference 3448?',
         ['A. A road', 'B. The woodland', 'C. The beach', 'D. A quarry']),
        ('2. What is the grid reference for the school?',
         ['A. 3346', 'B. 3248', 'C. 3150', 'D. 3447']),
        ('3. Why might grid reference 3150 be hard to build on?',
         ['A. The land is flat', 'B. The land is wet', 'C. The land is steep', 'D. There is a river']),
    ]:
        y = body_text(s1, y, qn, bold=True, sz=9.5)
        for opt in opts:
            y = body_text(s1, y, opt, sz=9, x=CONT_X + 0.2, w=CONT_W - 0.2)
        y += 0.05

    y += 0.05
    y = heading(s1, y, 'Part B   Comparison sentence')
    y = instruction(s1, y, 'Fill in the missing words using the word bank.')
    for line in [
        'The OS map of Westhaven shows ___________________ and ___________________.',
        'This is different from the satellite image of Brazil because ___________________________',
        '___________________________________________.',
    ]:
        y = body_text(s1, y, line, sz=9.5)
    y += 0.05
    word_bank(s1, y, 'woodland  •  roads  •  hills  •  settlement  •  rainforest  •  flat  •  land use')

    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part A   Answers')
    for a in ['1 → B. The woodland', '2 → B. 3248', '3 → C. The land is steep']:
        y = answer_text(s2, y, a)
    y += 0.1
    y = heading(s2, y, 'Part B   Model sentence')
    for line in [
        'The OS map of Westhaven shows woodland and roads.',
        'This is different from the satellite image of Brazil because Brazil shows',
        'flat farmland and rainforest with very different land use patterns.',
    ]:
        y = answer_text(s2, y, line)

    save_pptx(prs, 'T6W4_LP5_Geographers_Map_Skills.pptx')


# ══════════════════════════════════════════════════════════════════
# LP5 ADAPTED
# ══════════════════════════════════════════════════════════════════
def build_lp5_adapted():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    build_enquiry_label(s1, LBL_X, LBL_Y, '07/07/2026',
        'Are England and Brazil different?',
        'use a map to find information',
        'read a grid reference',
        'say what the map shows me',
        ICON_PATH)

    y = CONT_Y
    y = heading(s1, y, 'Part A   Grid reference questions')
    y = instruction(s1, y, 'Use the Westhaven map on the board. Circle A, B, C or D.')

    for qn, opts in [
        ('1. What is found at grid reference 3448?',
         ['A. A road', 'B. The woodland', 'C. The beach', 'D. A quarry']),
        ('2. What is the grid reference for the school?',
         ['A. 3346', 'B. 3248', 'C. 3150', 'D. 3447']),
    ]:
        y = body_text(s1, y, qn, bold=True, sz=9.5)
        for opt in opts:
            y = body_text(s1, y, opt, sz=9, x=CONT_X + 0.2, w=CONT_W - 0.2)
        y += 0.05

    y += 0.05
    y = heading(s1, y, 'Part B   Comparison sentence (cloze)')
    y = instruction(s1, y, 'Fill in the missing words.')
    for line in [
        'The OS map of Westhaven shows ___________________ and ___________________.',
        'This is different from the satellite image of Brazil because ___________________________',
        '____________________________________________________________.',
    ]:
        y = body_text(s1, y, line, sz=9.5)
    y += 0.05
    word_bank(s1, y, 'woodland  •  roads  •  hills  •  settlement  •  rainforest  •  flat  •  land use')

    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part A   Answers')
    for a in ['1 → B. The woodland', '2 → B. 3248']:
        y = answer_text(s2, y, a)
    y += 0.1
    y = heading(s2, y, 'Part B   Model sentence')
    for line in [
        'The OS map of Westhaven shows woodland and roads.',
        'This is different from the satellite image of Brazil because Brazil shows flat farmland and rainforest.',
    ]:
        y = answer_text(s2, y, line)

    save_pptx(prs, 'T6W4_LP5_Geographers_Map_Skills_adapted.pptx')


# ══════════════════════════════════════════════════════════════════
# LP6 STANDARD
# ══════════════════════════════════════════════════════════════════
def build_lp6_standard():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    build_enquiry_label(s1, LBL_X, LBL_Y, '08/07/2026',
        'Are England and Brazil different?',
        'explain how humans affect the environment',
        'describe one human impact on Brazil',
        'write a comparison using geographical vocabulary',
        ICON_PATH)

    y = CONT_Y
    y = heading(s1, y, 'Part A   Before and after')
    y = instruction(s1, y, 'Look at the images on the board. For each pair, record what you observe.')

    for pair in ['Image pair 1: Amazon rainforest', 'Image pair 2: English landscape']:
        # Title
        tb = add_textbox(s1, CONT_X, y, CONT_W, 0.18)
        set_para(tb.text_frame.paragraphs[0], pair, size_pt=9, color=DARK, font=FONT_BODY)
        y += 0.26

        for q_text, lw in [('What changed?', 1.3),
                            ('What caused it?', 1.35),
                            ('What might the geographical impact be?', 3.1)]:
            tb = add_textbox(s1, CONT_X, y, lw, 0.18)
            set_para(tb.text_frame.paragraphs[0], q_text, size_pt=8.5, color=DARK, font=FONT_BODY)
            add_line(s1, CONT_X + lw + 0.05, y + 0.14, CONT_W - lw - 0.05, BLUE)
            y += 0.315

        # Extra line below for additional writing space
        add_line(s1, CONT_X, y + 0.10, CONT_W, BLUE)
        y += 0.50

    y += 0.10
    y = heading(s1, y, 'Part B   Geographical comparison')
    y = instruction(s1, y, 'Write your comparison. Tick each word in the vocabulary checklist when you use it.')

    VC_TERMS = ['hemisphere', 'biome', 'climate zone', 'topography', 'land use',
                'natural resource', 'trade', 'deforestation', 'urbanisation', 'temperate', 'tropical']
    VC_W = 1.50
    VC_X = CONT_X + CONT_W - VC_W
    WW   = CONT_W - VC_W - 0.1

    # Vocab checklist box
    vc_h = 0.22 + len(VC_TERMS) * 0.175
    box = add_rect(s1, VC_X, y, VC_W, vc_h, fill_rgb=CREAM, line_rgb=ORANGE, line_pt=0.75)
    vc_tb = add_textbox(s1, VC_X + 0.05, y + 0.05, VC_W - 0.1, vc_h - 0.1)
    vc_tf = vc_tb.text_frame
    vc_tf.word_wrap = True
    set_para(vc_tf.paragraphs[0], 'Vocabulary checklist', bold=True, size_pt=7, color=ORANGE, font=FONT_BODY)
    for t in VC_TERMS:
        p = vc_tf.add_paragraph()
        set_para(p, f'\u25a1 {t}', size_pt=6.5, color=DARK, font=FONT_BODY)

    # Writing prompts and lines
    wy = y
    for prompt, nl in [
        ('Physical geography — how the two countries compare:', 3),
        ('Human geography — how land use compares:', 3),
        ('Environmental impact — how humans are affecting each place:', 3),
    ]:
        body_text(s1, wy, prompt, bold=True, sz=9, w=WW)
        wy += 0.28
        for _ in range(nl):
            add_line(s1, CONT_X, wy + 0.16, WW, EXBOOK, 0.75)
            wy += 0.20
        wy += 0.06

    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part A   Key points')
    for lbl_t, note in [
        ('Amazon rainforest', 'Deforestation: cleared for cattle, soya, mining. Impact: loss of biome, species, carbon storage.'),
        ('English landscape', 'Urban growth: farmland covered by housing and roads. Quarrying changes highland landscapes.'),
    ]:
        y = answer_text(s2, y, lbl_t, bold=True)
        y = answer_text(s2, y, note)
    y += 0.1
    y = heading(s2, y, 'Part B   Model comparison (extract)')
    for line in [
        'Physically, England has a temperate maritime climate with deciduous woodland, while Brazil has a tropical climate.',
        'In terms of human geography, Brazil\u2019s main land uses are agriculture and mining,',
        'while England focuses more on arable farming and services.',
        'Humans are having a greater impact in Brazil: around 20% of the Amazon has been deforested.',
        'In England, urban growth has covered farmland around cities like Bristol.',
    ]:
        y = answer_text(s2, y, line)

    save_pptx(prs, 'T6W4_LP6_Geographers_Environmental_Impact.pptx')


# ══════════════════════════════════════════════════════════════════
# LP6 ADAPTED
# ══════════════════════════════════════════════════════════════════
def build_lp6_adapted():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    build_enquiry_label(s1, LBL_X, LBL_Y, '08/07/2026',
        'Are England and Brazil different?',
        'describe how humans affect the environment',
        'name one cause of Amazon deforestation',
        'complete comparison sentences',
        ICON_PATH)

    y = CONT_Y
    y = heading(s1, y, 'Part A   What has changed?')
    y = instruction(s1, y, 'Look at the images on the board. Tick the boxes that apply to each image pair.')

    tick_items = ['Trees / vegetation removed', 'Buildings added', 'Farmland expanded', 'Roads or infrastructure built']
    for pair in ['Image pair 1: Amazon rainforest', 'Image pair 2: English landscape']:
        box = add_rect(s1, CONT_X, y, CONT_W, 0.72, fill_rgb=WHITE, line_rgb=BLUE, line_pt=0.75)
        tb = add_textbox(s1, CONT_X + 0.06, y + 0.05, CONT_W - 0.12, 0.62)
        tf = tb.text_frame
        tf.word_wrap = True
        set_para(tf.paragraphs[0], pair, bold=True, size_pt=9, color=BLUE, font=FONT_BODY)
        for t in tick_items:
            p = tf.add_paragraph()
            set_para(p, f'\u25a1 {t}', size_pt=8.5, color=DARK, font=FONT_BODY)
        y += 0.77

    y += 0.10
    y = heading(s1, y, 'Part B   Cloze comparison')
    y = instruction(s1, y, 'Fill in the missing words using the word bank.')
    for line in [
        'Physically, England has a _______________ climate, while Brazil has a _______________ climate.',
        'England\u2019s main biome is _______________ forest.',
        '',
        'For human geography, Brazil uses land mainly for _______________ such as coffee and soya,',
        'while England uses land more for _______________ and services.',
        '',
        'Humans are affecting Brazil by _______________ the Amazon.',
        'In England, _______________ growth has covered farmland around cities.',
    ]:
        if line == '':
            y += 0.07
        else:
            y = body_text(s1, y, line, sz=9.5)
    y += 0.05
    word_bank(s1, y, 'temperate  •  tropical  •  deciduous  •  agriculture  •  arable farming  •  deforesting  •  urban')

    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part B   Cloze answers')
    for a, lbl_t in [
        ('temperate', 'Climate — England'), ('tropical', 'Climate — Brazil'),
        ('deciduous', 'England biome'), ('agriculture', 'Brazil land use'),
        ('arable farming', 'England land use'), ('deforesting', 'Human impact'),
        ('urban', 'England impact type'),
    ]:
        y = answer_text(s2, y, f'\u2192 {a}   ({lbl_t})')

    save_pptx(prs, 'T6W4_LP6_Geographers_Environmental_Impact_adapted.pptx')


# ══════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    print('Building LP4...')
    build_lp4_standard()
    build_lp4_adapted()
    print('Building LP5...')
    build_lp5_standard()
    build_lp5_adapted()
    print('Building LP6...')
    build_lp6_standard()
    build_lp6_adapted()
    print('All done.')

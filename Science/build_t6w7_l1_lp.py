#!/usr/bin/env python3
"""
Build the T6W7 L1 (States of Matter) Learning Paper as PPTX.
Follows the pattern established in Geography/build_geo_lps_pptx_v3.py:
  slide 1 = pupil page, slide 2 = marking station.
Uses the real WFA label builder (label_builder.py -> generate_label_png.py
-> build_enquiry_label.py), not a hand-rolled label.
"""
import os, sys

WORK = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, WORK)

# Patch the hardcoded /home/claude asset path (not writable in this sandbox)
# to this session's copy before importing anything that reads it.
import build_enquiry_label as _bel
_bel.ASSETS = os.path.join(WORK, 'll_assets')

from label_builder import build_enquiry_label, LL_W
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

ICON_PATH = os.path.join(WORK, 'll_icons', 'scientist.png')  # legacy param, ignored by builder
PARTICLE_MODEL_IMG = os.path.join(WORK, 'particle_model.png')

# ── Slide dimensions (A4 portrait, matches school LP files) ────────────────
SW_IN, SH_IN = 7.5, 10.833
SW, SH = int(SW_IN * 914400), int(SH_IN * 914400)
MARGIN = 0.25

LBL_X = SW_IN - LL_W - MARGIN
LBL_Y = MARGIN
CONT_X = MARGIN
CONT_W = SW_IN - 2 * MARGIN

BLUE   = RGBColor(0x17, 0x98, 0xD3)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREEN  = RGBColor(0x4F, 0xAD, 0x5B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xCC, 0xCC, 0xCC)
EXBOOK = RGBColor(0xBD, 0xD7, 0xEE)
CREAM  = RGBColor(0xFE, 0xF9, 0xE7)

FONT_BODY = 'Twinkl Cursive Looped'

OUT = os.path.join(WORK, 'out')
os.makedirs(OUT, exist_ok=True)


# ══════════════════════════════════════════════════════════════════
# PPTX helpers (adapted from Geography/build_geo_lps_pptx_v3.py)
# ══════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SW), Emu(SH)
    blank = prs.slide_layouts[6]
    prs.slides.add_slide(blank)
    prs.slides.add_slide(blank)
    return prs

def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            sp_tree.remove(child)

def _i(v): return Inches(v)

def add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))

def set_para(para, text, bold=False, italic=False, size_pt=10,
             color=None, font=FONT_BODY, align=PP_ALIGN.LEFT, underline=False):
    para.alignment = align
    run = para.add_run()
    run.text = text
    rf = run.font
    rf.name = font
    rf.size = Pt(size_pt)
    rf.bold = bold
    rf.italic = italic
    rf.underline = underline
    if color:
        rf.color.rgb = color
    return run

def add_line(slide, x, y, w, color=None, width_pt=1.0):
    c = color if color is not None else EXBOOK
    hex_col = f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'
    lw = int(Pt(width_pt))
    sp_tree = slide.shapes._spTree
    uid = max((int(el.get('id', 0)) for el in sp_tree.iter()
               if el.get('id') and el.get('id').isdigit()), default=100) + 1
    xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="{uid}" name="Line{uid}"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{int(Inches(x))}" y="{int(Inches(y))}"/>'
        f'<a:ext cx="{int(Inches(w))}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom><a:noFill/>'
        f'<a:ln w="{lw}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:ln>'
        f'<a:effectLst/></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    sp_tree.append(etree.fromstring(xml))

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0.75):
    shape = slide.shapes.add_shape(1, _i(x), _i(y), _i(w), _i(h))
    if fill_rgb:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb; shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def _wrap_line_count(text, width_in, size_pt):
    """Estimate wrapped line count for Twinkl Cursive Looped.

    Round 12 (11 Jul 2026) correction: this used the 0.52 chars-per-pt-width
    ratio from EnquiryBuilder/lib_ooxml.py's force_shrink_to_fit, on the
    assumption it was a safe, already-proven value. It wasn't safe here -
    Twinkl Cursive Looped isn't installed anywhere in this sandbox (no font
    file in the repo, not on the system), so every QA render in this file
    used a LibreOffice substitute font to check wrapping, and that
    substitute is wider than the real thing. Innes confirmed directly from
    his own PowerPoint: "Balloon (filled with air)" (25 chars) fits on ONE
    line in the real font at 10pt in a 1.74" column - this function's old
    ratio estimated 24 chars/line there, one short, and wrongly forced a
    2-line row. Corrected to 0.46 chars-per-pt-width ratio (chars_per_line
    26 for that same case - a small margin past the confirmed-fitting 25,
    not a razor's-edge match to a single data point).

    This is still a heuristic, not a measurement - there's no real Twinkl
    Cursive Looped font file anywhere to measure against directly, and
    0.46 is calibrated from exactly one confirmed real-PowerPoint example.
    Do not assume it's precise; if a future LP shows text overflowing a
    box that this function said would fit, or wrapping unnecessarily
    when Innes's real PowerPoint shows it fitting, that's another real
    data point - adjust the ratio again rather than re-guessing from
    scratch, and note the correction here.

    NOT shared with lib_ooxml.py's force_shrink_to_fit - that function's
    0.52 ratio is calibrated for OOXML slide runs (confirmed working
    there, e.g. the concept cartoon speech bubbles) and is left alone;
    this is a separate, LP-table-cell-specific calibration."""
    usable_w_pt = width_in * 72
    chars_per_line = max(1, int(usable_w_pt / (size_pt * 0.46)))
    words = text.split()
    lines, cur = 1, ''
    for w in words:
        test = (cur + ' ' + w).strip()
        if len(test) <= chars_per_line:
            cur = test
        else:
            lines += 1
            cur = w
    return max(1, lines)

def heading(slide, y, text, w=None, size_pt=12):
    w = w or CONT_W
    n = _wrap_line_count(text, w, size_pt)
    h = (size_pt / 12) * 0.24 * n
    tb = add_textbox(slide, CONT_X, y, w, h)
    tb.text_frame.word_wrap = True
    set_para(tb.text_frame.paragraphs[0], text, bold=True, size_pt=size_pt, color=BLUE)
    return y + h + 0.06

def instruction(slide, y, text, w=None, size_pt=12):
    # Innes: 12pt is the comfortable reading size on a pupil-facing LP -
    # default to it. Only go smaller for content that genuinely needs the
    # space and isn't extended prose (short table labels etc, handled by
    # their own functions below, not this one).
    w = w or CONT_W
    n = _wrap_line_count(text, w, size_pt)
    h = (size_pt / 9.5) * 0.19 * n
    tb = add_textbox(slide, CONT_X, y, w, h)
    tb.text_frame.word_wrap = True
    set_para(tb.text_frame.paragraphs[0], text, size_pt=size_pt, color=DARK)
    return y + h + 0.08

def body_text(slide, y, text, bold=False, sz=10, col=None, x=None, w=None):
    col = col or DARK; x = x if x is not None else CONT_X; w = w or CONT_W
    tb = add_textbox(slide, x, y, w, 0.25)
    tb.text_frame.word_wrap = True
    set_para(tb.text_frame.paragraphs[0], text, bold=bold, size_pt=sz, color=col)
    return y + 0.28

def write_lines(slide, y, n=3, w=None, x=None):
    w = w or CONT_W; x = x if x is not None else CONT_X
    for _ in range(n):
        add_line(slide, x, y, w); y += 0.315
    return y + 0.10

def word_bank(slide, y, words, w=None, size_pt=9):
    w = w or CONT_W
    n = _wrap_line_count('Word bank: ' + words, w - 0.1, size_pt)
    box_h = (size_pt / 9) * 0.20 * n + 0.14
    add_rect(slide, CONT_X, y, w, box_h, fill_rgb=CREAM, line_rgb=ORANGE, line_pt=1.0)
    tb = add_textbox(slide, CONT_X + 0.05, y + 0.05, w - 0.1, box_h - 0.08)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'Word bank: '
    r1.font.name = FONT_BODY; r1.font.bold = True; r1.font.size = Pt(size_pt); r1.font.color.rgb = ORANGE
    r2 = p.add_run(); r2.text = words
    r2.font.name = FONT_BODY; r2.font.size = Pt(size_pt); r2.font.color.rgb = DARK
    return y + box_h + 0.06

def add_reference_image(slide, y, img_path, max_w=None, max_h=1.35):
    """Place a reference image at native aspect ratio, capped by max_h,
    centred within CONT_W. If a task tells pupils to 'use X to help you',
    X should be on the page, not just named - this exists for that."""
    from PIL import Image as _PILImage
    max_w = max_w or CONT_W
    iw, ih = _PILImage.open(img_path).size
    ratio = iw / ih
    h = max_h
    w = h * ratio
    if w > max_w:
        w = max_w
        h = w / ratio
    x = CONT_X + (CONT_W - w) / 2
    slide.shapes.add_picture(img_path, _i(x), _i(y), width=_i(w), height=_i(h))
    return y + h + 0.10

def marking_heading(slide, y, text='Marking Station'):
    tb = add_textbox(slide, CONT_X, y, CONT_W, 0.40)
    set_para(tb.text_frame.paragraphs[0], text, bold=True, size_pt=16, color=GREEN)
    return y + 0.45

def answer_text(slide, y, text, bold=False, w=None):
    w = w or CONT_W
    n = _wrap_line_count(text, w, 9.5)
    h = 0.19 * n
    tb = add_textbox(slide, CONT_X, y, w, h)
    tb.text_frame.word_wrap = True
    set_para(tb.text_frame.paragraphs[0], text, bold=bold, size_pt=9.5, color=GREEN)
    return y + h + 0.08

def save_pptx(prs, filename):
    path = os.path.join(OUT, filename)
    prs.save(path)
    print(f'  {filename}: {os.path.getsize(path)//1024}KB')
    return path


# ── Materials sorting table (Material | State | Reason) ────────────────────
COL_LABELS = ['Material', 'State?', 'Reason (use the particle model)']
COL_W_RATIO = [0.26, 0.14, 0.60]
COL_WIDTHS_IN = [CONT_W * r for r in COL_W_RATIO]

def sort_table_header(slide, y, size_pt=8):
    n = max(_wrap_line_count(lbl, cw - 0.08, size_pt) for lbl, cw in zip(COL_LABELS, COL_WIDTHS_IN))
    row_h = (size_pt / 8) * 0.18 * n + 0.10
    x = CONT_X
    for lbl, cw in zip(COL_LABELS, COL_WIDTHS_IN):
        add_rect(slide, x, y, cw, row_h, fill_rgb=BLUE)
        tb = add_textbox(slide, x + 0.04, y + 0.03, cw - 0.08, row_h - 0.06)
        tb.text_frame.word_wrap = True
        set_para(tb.text_frame.paragraphs[0], lbl, bold=True, size_pt=size_pt, color=WHITE)
        x += cw
    return y + row_h

def sort_table_row(slide, y, material, even=True, row_h=0.42, size_pt=8.5):
    # row_h is a floor, not a fixed value - a genuinely wrapping material
    # name needs a taller row or its second line gets clipped against the
    # row below. Only grow past the floor when the corrected wrap estimate
    # (see _wrap_line_count) actually says it needs more than one line -
    # a single-line estimate must not push the row taller than the floor,
    # which is already sized generously for handwriting room in the blank
    # State/Reason cells.
    n = _wrap_line_count(material, COL_WIDTHS_IN[0] - 0.08, size_pt)
    if n > 1:
        needed_h = 0.08 + n * (size_pt * 1.3 / 72)
        row_h = max(row_h, needed_h)
    fill = RGBColor(0xF5, 0xF5, 0xF5) if even else WHITE
    x = CONT_X
    for i, cw in enumerate(COL_WIDTHS_IN):
        add_rect(slide, x, y, cw, row_h, fill_rgb=fill, line_rgb=LGREY, line_pt=0.5)
        if i == 0:
            tb = add_textbox(slide, x + 0.04, y + 0.04, cw - 0.08, row_h - 0.08)
            tb.text_frame.word_wrap = True
            set_para(tb.text_frame.paragraphs[0], material, size_pt=size_pt, color=DARK)
        x += cw
    return y + row_h

def sort_table_answer_row(slide, y, material, state, reason):
    cells = [material, state, reason]
    n = max(_wrap_line_count(c, cw - 0.08, 8) for c, cw in zip(cells, COL_WIDTHS_IN))
    row_h = 0.16 * n + 0.10
    x = CONT_X
    for cell, cw in zip(cells, COL_WIDTHS_IN):
        tb = add_textbox(slide, x + 0.04, y + 0.02, cw - 0.08, row_h - 0.04)
        tb.text_frame.word_wrap = True
        set_para(tb.text_frame.paragraphs[0], cell, size_pt=8, color=GREEN)
        x += cw
    return y + row_h


# ══════════════════════════════════════════════════════════════════
# T6W7 L1 — States of Matter — sorting LP
# ══════════════════════════════════════════════════════════════════
MATERIALS = ['Ice', 'Water', 'Steam (water vapour)', 'Wood',
             'Sand', 'Milk', 'Balloon (filled with air)', 'Honey']

ANSWERS = [
    ('Ice', 'Solid', 'Particles packed tightly in a fixed arrangement, so ice keeps its shape.'),
    ('Water', 'Liquid', 'Particles are close but can slide past each other, so water flows and takes the shape of its container.'),
    ('Steam (water vapour)', 'Gas', 'Particles move quickly and are spread far apart, so steam fills the space available.'),
    ('Wood', 'Solid', 'Particles are fixed in place, so wood keeps its shape and does not flow.'),
    ('Sand', 'Solid (tricky!)', 'Each grain is a solid with its own fixed shape, but lots of grains together can pour like a liquid.'),
    ('Milk', 'Liquid', 'Particles slide past each other, so milk flows and takes the shape of its container.'),
    ('Balloon (filled with air)', 'Gas', 'The air particles inside move quickly and spread out to fill the balloon.'),
    ('Honey', 'Liquid (tricky!)', 'Honey flows very slowly, but it is still a liquid because it takes the shape of its container.'),
]

def build_lp():
    prs = new_prs()
    s1 = prs.slides[0]
    clear_slide(s1)

    label_h = build_enquiry_label(
        s1, LBL_X, LBL_Y,
        date_str='13/07/2026',
        key_q='Can materials change their state?',
        lf='compare and group materials',
        ican1='sort materials based on their properties',
        ican2='explain my sorting using the particle model',
        icon_path=ICON_PATH,
        subject='scientist',
        year='Y4',
        png_dest=os.path.join(WORK, 'lp_label.png'),
    )
    # Innes: label was too big - shrink it (matches his manual resize,
    # ~70.7% of natural render size, aspect ratio preserved) rather than
    # embedding at the label pipeline's fixed natural width.
    LABEL_SCALE = 0.707
    label_pic = s1.shapes[-1]
    new_w, new_h = LL_W * LABEL_SCALE, label_h * LABEL_SCALE
    label_pic.left = _i(SW_IN - new_w - MARGIN)
    label_pic.top = _i(LBL_Y)
    label_pic.width, label_pic.height = _i(new_w), _i(new_h)
    label_h = new_h
    CONT_Y = LBL_Y + label_h + 0.15

    # Innes: pupil-page text was too small for children to read - sizes
    # below (16/12/10/10/14/10.5/12) match his own edit exactly. Marking
    # station (slide 2, below) is unchanged - it's for the teacher, not
    # read from across the room, so stays compact.
    y = CONT_Y
    y = heading(s1, y, 'Part A   Sort the materials', size_pt=16)
    y = instruction(s1, y, 'Sort each material into solid, liquid or gas. Write one reason for '
                            'each, using the particle model below to help you.', size_pt=12)
    y = add_reference_image(s1, y, PARTICLE_MODEL_IMG, max_h=1.37)

    y = sort_table_header(s1, y, size_pt=10)
    for i, m in enumerate(MATERIALS):
        y = sort_table_row(s1, y, m, even=(i % 2 == 0), row_h=0.36, size_pt=10)

    y += 0.12
    y = heading(s1, y, 'Part B   Challenge', size_pt=14)
    y = instruction(s1, y, 'Can you find one material that is difficult to classify? Why?', size_pt=12)
    # Innes: gap before the first write-line must be at least the same
    # 0.8cm (0.315") used between the lines themselves, so there's
    # actually room to write on that first line - instruction() only
    # leaves 0.08" by default, top up to the full 0.315".
    y += 0.315 - 0.08
    y = write_lines(s1, y, 3)
    y += 0.05
    word_bank(s1, y, 'solid  •  liquid  •  gas  •  particles  •  '
                      'closely packed  •  slide past  •  spread out  •  fixed shape  •  flows',
              size_pt=12)

    # ── Marking station ──
    s2 = prs.slides[1]
    clear_slide(s2)
    y = 0.3
    y = marking_heading(s2, y)
    y = heading(s2, y, 'Part A   Answers')
    y = sort_table_header(s2, y)
    for i, (m, state, reason) in enumerate(ANSWERS):
        y = sort_table_answer_row(s2, y, m, state, reason)

    y += 0.15
    y = heading(s2, y, 'Part B   Challenge - model answer')
    for line in [
        'Sand and honey are the best examples.',
        'Sand: a single grain is a solid with its own fixed shape, but a big pile of grains can be poured '
        'and takes the shape of its container - a bit like a liquid.',
        'Honey: it is a liquid because it takes the shape of its container, but it flows so slowly it can '
        'look almost solid.',
    ]:
        y = answer_text(s2, y, line)

    save_pptx(prs, 'T6W7_L1_LP1_Scientist_States_of_Matter.pptx')


if __name__ == '__main__':
    build_lp()
    print('Done.')
